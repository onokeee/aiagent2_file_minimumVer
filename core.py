"""core.py — 業務ロジック一式。画面(web.py)と設定(config.py)と統計の計算(advanced.py)以外は全部ここ（元: 35モジュール）。

元は以下のファイルに分かれていた。中身は変えずに1つにまとめている:
  db.py
  filecheck.py
  chats.py
  history.py
  catalog_history.py
  prefs.py
  models.py
  catalog.py
  verify.py
  sqlusage.py
  usage.py
  exports.py
  excel.py
  charts.py
  figures.py
  docx_report.py
  pptx_report.py
  business.py
  importer.py
  jobs.py
  scheduler.py
  cleanup.py
  mailer.py
  custom_tools.py
  tools/results.py
  tools/common.py
  tools/schemas.py
  tools/business.py
  tools/files.py
  tools/mail.py
  tools/query.py
  tools/reports.py
  tools/stats.py
  tools/usage.py
  tools/__init__.py
  llm.py

統合前と同じく `import db` / `db.run_select(...)` と書けるよう、
このファイルを元のモジュール名でも参照できるよう登録している。
そのため呼び出し側のコードは統合前のまま動く。
"""
from __future__ import annotations

import sys as _sys

# 元のモジュール名でも import できるようにする（呼び出し側を変えないため）
for _alias in ("db", "filecheck", "chats", "history", "catalog_history", "prefs", "models", "catalog", "verify", "sqlusage", "usage", "exports", "excel", "charts", "figures", "docx_report", "pptx_report", "business", "importer", "jobs", "scheduler", "cleanup", "mailer", "custom_tools", "llm", "tools"):
    _sys.modules[_alias] = _sys.modules[__name__]
del _alias

import advanced  # noqa: F401  （統合前のモジュール名を登録させる）

results = _sys.modules[__name__]  # 統合前の書き方をそのまま使えるようにする
_results = _sys.modules[__name__]  # 統合前の書き方をそのまま使えるようにする


# ==========================================================================
# ===== 元 db.py
# SQLite アクセス層（複数DB対応）。
#
# data/ フォルダの .db ファイルを列挙し、選択されたDB群を読み取り専用で
# ATTACH した1つの接続に対して SELECT を実行する。複数DBを選択した場合は
# `エイリアス.テーブル名` でファイルをまたいだ JOIN が可能。
#
# 最重要: ユーザー(LLM)が生成したSQLは SELECT のみ 実行を許可する。
# 多層防御で守る:
#   1. 構文チェック   : 単一ステートメント / SELECT・WITH で始まる / 書込キーワード禁止
#   2. 読み取り専用接続: mode=ro で ATTACH するのでそもそも書込不可
#   3. オーソライザ    : SQLite の authorizer で SELECT/READ 以外を DENY
#   4. タイムアウト    : progress handler で暴走クエリを中断
# ==========================================================================
import re
import sqlite3
import time
from pathlib import Path

import config

# --- data/ フォルダのDBファイル ----------------------------------------------

def list_db_files() -> list[Path]:
    """data/ 直下の .db ファイル一覧（名前順）。"""
    if not config.DATA_DIR.exists():
        return []
    return sorted(p for p in config.DATA_DIR.glob("*.db") if p.is_file())


def path_for(name) -> Path:
    """画面から渡されたDB名を data/ の実ファイルに解決する。

    名前を data/ に連結するのではなく、列挙済みの一覧から名前が一致するものを
    探す。こうしておくと "../" のような指定が入っても data/ の外には出ない。
    """
    target = Path(str(name or "")).name          # ディレクトリ部分は捨てる
    for p in list_db_files():
        if p.name == target:
            return p
    raise FileNotFoundError(f"DBが見つかりません: {name}")


# 記号と空白だけを潰す。日本語などのマルチバイト文字はそのまま残す
# （SQLiteは非ASCIIの識別子をクオート無しで扱えるので、"売上.db" は 売上.受注 と書ける）。
# ここでASCIIだけに絞ると「店舗マスタ」が "_____" になり、
# 複数の日本語DBを選んだときに区別できなくなる。
_ALIAS_BAD = re.compile(r"[^\w]", re.UNICODE)
_RESERVED_ALIASES = {"main", "temp"}


def alias_for(path: Path) -> str:
    """ファイル名から SQL で使うエイリアス名（英数字と_のみ）を作る。"""
    a = _ALIAS_BAD.sub("_", Path(path).stem)
    if not a or a[0].isdigit():
        a = "db_" + a
    if a.lower() in _RESERVED_ALIASES:
        a += "_db"
    return a


def aliases_for(paths: list[Path]) -> list[str]:
    """複数ファイルに一意なエイリアスを割り当てる（衝突時は連番を付ける）。"""
    result: list[str] = []
    used: set[str] = set()
    for p in paths:
        a = alias_for(p)
        base, n = a, 2
        while a.lower() in used:
            a = f"{base}_{n}"
            n += 1
        used.add(a.lower())
        result.append(a)
    return result


# --- 接続ヘルパ ---------------------------------------------------------------

def _ro_uri(path) -> str:
    return Path(path).resolve().as_uri() + "?mode=ro"


def connect_ro(path) -> sqlite3.Connection:
    """単一DBへの読み取り専用接続（プロファイリング用）。"""
    return sqlite3.connect(_ro_uri(path), uri=True)


# SQLiteが同時にATTACHできる数の上限（既定10）。main を1つ使うので実質これだけ。
MAX_ATTACHED = 10


def connect_scope(paths_aliases: list[tuple]) -> sqlite3.Connection:
    """空の :memory: を main とし、各DBを読み取り専用で ATTACH した接続を作る。

    paths_aliases: [(path, alias), ...]
    """
    if len(paths_aliases) > MAX_ATTACHED:
        raise ValueError(
            f"1つのSQLで扱えるDBは{MAX_ATTACHED}個までです"
            f"（この問い合わせは{len(paths_aliases)}個を必要としています）。SQLiteの制限です。"
            "テーブル名を『DB名.テーブル名』の形で書けば、実際に使うDBだけを繋ぐので"
            "多くの場合はこの制限に当たりません。")
    conn = sqlite3.connect("file::memory:", uri=True)
    for path, alias in paths_aliases:
        # alias は英数字と_のみに正規化済みなので識別子として安全
        conn.execute(f'ATTACH DATABASE ? AS "{alias}"', (_ro_uri(path),))
    return conn


def dbs_named_in(sql: str) -> list[str]:
    """SQLが「エイリアス.」の形で名指ししているDBファイル名。"""
    out = []
    for p in list_db_files():
        a = alias_for(p)
        if a and re.search(r'(?<![\w."])' + re.escape(a) + r'\s*\.', sql, re.IGNORECASE):
            out.append(p.name)
    return out


def widen_scope(sql: str, scope: list[dict]) -> list[dict]:
    """SQLが必要とするDBを、選ばれていなくても繋ぐ。

    ユーザー定義ツールは作った人がDBを意識せずに書くので、SQLが別DBに入ることがある。
    選択中のDBだけを繋ぐと、正しいツールが "no such table" で落ちる。
    読むだけであり、DBの選択はもともと「見る範囲を絞る」ためのもので
    アクセス制御ではない（README参照）ため、必要なものは繋いでよい。

    ATTACH の上限があるので、そこで打ち止める（超えた分は元のエラーで気づける）。
    """
    out = list(scope or [])
    have = {str(s.get("alias") or "").lower() for s in out}
    for p in list_db_files():
        if len(out) >= MAX_ATTACHED:
            break
        a = alias_for(p)
        if a.lower() in have:
            continue
        if re.search(r'(?<![\w."])' + re.escape(a) + r'\s*\.', sql, re.IGNORECASE):
            out.append({"path": str(p), "alias": a, "name": p.name, "tables": None})
            have.add(a.lower())
    return out


def narrow_scope(sql: str, scope: list[dict]) -> list[dict]:
    """そのSQLに関係するDBだけに絞る。

    選択中のDBを全部つなぐ必要はない。SQLiteは一度に10個までしかATTACHできないので、
    11個以上選んでいると、2つのテーブルを見るだけの問い合わせも実行できなくなっていた。
    「エイリアス.テーブル」で名指しされたDBと、修飾なしのテーブル名が一致するDBだけを残す。

    どちらでも判断できないときは、今までどおり全部を返す（勝手に減らして
    「no such table」にするより、元の分かりやすいエラーの方がよい）。
    """
    if len(scope) <= 1:
        return scope

    picked, seen = [], set()

    def add(s):
        key = str(s.get("path"))
        if key not in seen:
            seen.add(key)
            picked.append(s)

    # 名指しされているDB（複数DBを選んでいるときは必ずこの形で書かせている）
    for s in scope:
        alias = str(s.get("alias") or "")
        if alias and re.search(r'(?<![\w."])' + re.escape(alias) + r'\s*\.',
                               sql, re.IGNORECASE):
            add(s)
    # 修飾なしのテーブル名で参照されているDB。上と混在したSQLでも取りこぼさない
    for s in scope:
        for t in (s.get("tables") or []):
            if re.search(r'(?<![\w."])' + re.escape(str(t)) + r'(?![\w"])',
                         sql, re.IGNORECASE):
                add(s)
                break

    return picked[:MAX_ATTACHED] if picked else scope


# --- SELECT 専用ガード --------------------------------------------------------

# replace はここに入れない。SQLite の replace(X,Y,Z) は文字列関数で、
# 「株式会社」を落とすといった用途でごく普通に使う。書き込みになるのは
# REPLACE INTO の形だけなので、それは下で別に見る。
_FORBIDDEN = re.compile(
    r"\b(insert|update|delete|drop|alter|create|truncate|attach|detach|"
    r"reindex|vacuum|pragma|grant|revoke|begin|commit|rollback|savepoint|merge)\b",
    re.IGNORECASE,
)
_REPLACE_INTO = re.compile(r"\breplace\s+into\b", re.IGNORECASE)

#: 文字列リテラルと引用符付き識別子。'' や "" のエスケープも1つの塊として食う。
_QUOTED = re.compile(r"'(?:[^']|'')*'|\"(?:[^\"]|\"\")*\"|`(?:[^`]|``)*`|\[[^\]]*\]")


def _strip_sql_comments(sql: str) -> str:
    sql = re.sub(r"--[^\n]*", " ", sql)                     # 行コメント
    sql = re.sub(r"/\*.*?\*/", " ", sql, flags=re.DOTALL)   # ブロックコメント
    return sql


def _blank_quoted(sql: str) -> str:
    """引用符で囲まれた中身を空にした、検査用のコピーを作る。

    キーワードや ';' をそのまま探すと、値の中の文字まで拾ってしまう。
    WHERE status = 'delete' や WHERE note = ';' が「危険なSQL」として
    弾かれていた。実行するのは元のSQLで、これは検査にしか使わない。
    """
    return _QUOTED.sub(lambda m: m.group(0)[0] + m.group(0)[-1], sql)


def validate_select(sql: str) -> str:
    """SELECT文として安全か検証し、整形済みSQLを返す。問題があれば ValueError。"""
    if not sql or not sql.strip():
        raise ValueError("SQLが空です。")
    cleaned = _strip_sql_comments(sql).strip().rstrip(";").strip()
    if not cleaned:
        raise ValueError("実行可能なSQLがありません。")
    probe = _blank_quoted(cleaned)          # 検査は中身を抜いたコピーに対して行う
    if ";" in probe:
        raise ValueError("複数ステートメントは実行できません(SELECT文を1つだけ指定してください)。")
    low = probe.lower()
    if not (low.startswith("select") or low.startswith("with")):
        raise ValueError("SELECT文(または WITH ... SELECT)のみ実行できます。")
    m = _FORBIDDEN.search(probe) or _REPLACE_INTO.search(probe)
    if m:
        raise ValueError(f"書き込み・DDL系のキーワード '{m.group(0)}' は使用できません。読み取り専用です。")
    return cleaned


# SQLite オーソライザで許可するアクション
_ALLOWED_ACTIONS = {sqlite3.SQLITE_SELECT, sqlite3.SQLITE_READ, sqlite3.SQLITE_FUNCTION}
for _name in ("SQLITE_RECURSIVE",):  # 環境によって存在しない場合がある
    if hasattr(sqlite3, _name):
        _ALLOWED_ACTIONS.add(getattr(sqlite3, _name))


def _authorizer(action, arg1, arg2, db_name, trigger):
    if action in _ALLOWED_ACTIONS:
        return sqlite3.SQLITE_OK
    return sqlite3.SQLITE_DENY


# SQLiteに無い関数 → 代わりに使うもの。
# エラーメッセージにこれを添えないと、LLMは同じ関数で何度も書き直す。
_MISSING_FUNC_HINTS = {
    "stddev": "analyze_stats(method='describe')", "stdev": "analyze_stats(method='describe')",
    "stddev_samp": "analyze_stats(method='describe')",
    "variance": "analyze_stats(method='describe')", "var_samp": "analyze_stats(method='describe')",
    "median": "analyze_stats(method='describe')",
    "percentile": "analyze_stats(method='describe')",
    "percentile_cont": "analyze_stats(method='describe')",
    "percentile_disc": "analyze_stats(method='describe')",
    "corr": "analyze_stats(method='correlation')",
    "regr_slope": "regression", "stddev_pop": "analyze_stats(method='describe')",
    "sqrt": "analyze_stats か advanced 系のツール",
    "power": "掛け算で書き換える（x*x など）",
    "date_trunc": "strftime('%Y-%m', 列) など strftime を使う",
    "now": "date('now') / datetime('now')",
    "year": "strftime('%Y', 列)", "month": "strftime('%m', 列)",
    "day": "strftime('%d', 列)", "concat": "|| で連結する",
    "ifnull_": "IFNULL は使える", "listagg": "group_concat",
    "string_agg": "group_concat", "top": "LIMIT",
}


def explain_error(e: Exception) -> str:
    """SQLの失敗を、次に何をすればよいかまで書いた文にする。"""
    msg = str(e)
    m = re.search(r"no such function:\s*([A-Za-z_0-9]+)", msg)
    if m:
        fn = m.group(1)
        hint = _MISSING_FUNC_HINTS.get(fn.lower())
        if hint:
            return (f"{msg} … SQLite には {fn}() がありません。"
                    f"SQLで書き直そうとせず、{hint} を使ってください。")
        return (f"{msg} … SQLite には {fn}() がありません。"
                "標準のSQLite関数だけで書き直すか、専用の分析ツールを使ってください。")
    m = re.search(r"no such column:\s*(\S+)", msg)
    if m:
        return (f"{msg} … 列名が違います。describe_table でテーブルの列を確認してから"
                "書き直してください（推測で列名を作らないこと）。")
    if "syntax error" in msg:
        return (f"{msg} … SQLite で解釈できない書き方です。"
                "ウィンドウ関数の一部・WITHIN GROUP・PIVOT などは使えません。"
                "集計や統計は専用ツール（pivot_table / analyze_stats）に任せてください。")
    return msg


def run_select(sql: str, scope: list[dict], max_rows: int | None = None,
               timeout_s: int | None = None, params: dict | None = None):
    """検証済みSELECTを、選択スコープのDB群に対して実行する。

    scope:  [{"path": str, "alias": str, ...}, ...]（tables キーは無視。
            安全性はDB単位のATTACH + SELECT専用ガードで担保する）
    params: バインド変数（:name）に渡す値。値はSQL文字列に埋め込まれず
            プレースホルダ経由で渡るため、SQLインジェクションは起こらない。
    戻り値: (columns, rows, truncated)
    """
    if not scope:
        raise ValueError("対象のDBがありません。data/ にDBを置いてください。")
    safe_sql = validate_select(sql)
    max_rows = max_rows or config.MAX_RESULT_ROWS
    timeout_s = timeout_s or config.QUERY_TIMEOUT_SEC

    # 選択中のDBを全部つながない。このSQLが要るものだけを繋ぐ（ATTACHは10個まで）
    use = narrow_scope(safe_sql, scope)
    conn = connect_scope([(s["path"], s["alias"]) for s in use])
    try:
        conn.set_authorizer(_authorizer)
        start = time.time()
        conn.set_progress_handler(lambda: 1 if (time.time() - start) > timeout_s else 0, 10000)
        try:
            cur = conn.execute(safe_sql, params or {})  # 単一ステートメントのみ実行可能
        except sqlite3.Error as e:
            # 「何が悪いか」だけでなく「代わりに何を使うか」まで返す
            raise sqlite3.OperationalError(explain_error(e)) from e
        columns = [d[0] for d in cur.description] if cur.description else []
        rows = cur.fetchmany(max_rows + 1)
        truncated = len(rows) > max_rows
        rows = [tuple(r) for r in rows[:max_rows]]
        return columns, rows, truncated
    finally:
        conn.close()


if __name__ == "__main__":
    # SELECT専用ガード + ATTACH横断クエリのセルフテスト
    import tempfile, os

    ok_cases = [
        "SELECT 1",
        "select a.x from t a join u b on a.id=b.id",
        "WITH t AS (SELECT 1 AS a) SELECT a FROM t",
        # 文字列リテラルの中のキーワードや記号で弾かないこと
        "SELECT replace(name,'株式会社','') FROM t",
        "SELECT * FROM t WHERE x='delete me'",
        "SELECT * FROM t WHERE note=';'",
        "SELECT * FROM t WHERE s='don''t drop it'",
        'SELECT "delete" FROM t',
    ]
    ng_cases = [
        "DELETE FROM t",
        "DROP TABLE t",
        "UPDATE t SET x=1",
        "SELECT 1; DELETE FROM t",
        "INSERT INTO t VALUES(1)",
        "PRAGMA table_info(t)",
        "ATTACH DATABASE 'x.db' AS z",
        "REPLACE INTO t VALUES(1)",
        "SELECT * FROM t WHERE x='a'; DROP TABLE t",
    ]
    for s in ok_cases:
        validate_select(s)
        print("OK   ", s)
    for s in ng_cases:
        try:
            validate_select(s)
            print("!! ガードすり抜け:", s)
        except ValueError as e:
            print("BLOCK", s, "=>", e)

    # ATTACH 横断クエリ
    d = tempfile.mkdtemp()
    p1, p2 = os.path.join(d, "a.db"), os.path.join(d, "b.db")
    c = sqlite3.connect(p1); c.execute("CREATE TABLE t(id INTEGER, v TEXT)")
    c.execute("INSERT INTO t VALUES(1,'x'),(2,'y')"); c.commit(); c.close()
    c = sqlite3.connect(p2); c.execute("CREATE TABLE u(id INTEGER, w TEXT)")
    c.execute("INSERT INTO u VALUES(1,'A'),(2,'B')"); c.commit(); c.close()
    scope = [{"path": p1, "alias": "a"}, {"path": p2, "alias": "b"}]
    cols, rows, tr = run_select("SELECT t.v, u.w FROM a.t t JOIN b.u u ON t.id=u.id", scope)
    print("CROSS-DB JOIN:", cols, rows)
    # 書込は物理的に拒否されるか
    try:
        run_select("SELECT 1", scope)  # ガード通過の確認
        conn = connect_scope([(p1, "a")])
        conn.execute("INSERT INTO a.t VALUES(9,'z')")
        print("!! 読み取り専用が効いていない")
    except sqlite3.OperationalError as e:
        print("RO-GUARD:", e)


# ==========================================================================
# ===== 元 filecheck.py
# そのファイルが「そのまま取り込める表」かを判定する。
#
# 取り込みは 1行=1レコード / 1列=1項目 の素直な表を前提にしている。
# ところが現場のExcelは、見出しがセル結合されていたり、月が横に並んでいたり、
# 合計行が混ざっていたりする。そのまま取り込むと、列名が「Unnamed: 3」になったり、
# 合計が二重に数えられたりして、後の集計が静かに狂う。
#
# ここでは中身を読む前に形を見て、次のどれかを返す。
#   そのまま取り込める / 手直しが要る / 取り込みに向かない / 対応していない形式
#
# 判定は「取り込みボタンを押す前に気づけるようにする」ためのもので、
# 最終的に決めるのは人。だから理由と直し方を必ず添える。
# ==========================================================================
import csv
import re
from pathlib import Path

import config

#: 形を見るために読む最大行数。これ以上は見なくても判断できる。
MAX_SCAN_ROWS = 200
#: 見出し行を探す範囲（先頭から何行目まで）。
HEADER_SEARCH_ROWS = 12
#: 結合セルの調査は通常読み込みが要る（メモリを食う）ので、この大きさまで。
MERGE_CHECK_MAX_MB = 20

VERDICTS = ("そのまま取り込める", "手直しが要る", "取り込みに向かない", "対応していない形式")

#: 合計・小計の行に出やすい言葉。混ざったまま取り込むと二重計上になる。
_TOTAL_WORDS = ("合計", "総計", "小計", "計", "累計", "total", "subtotal", "sum")
#: 見出しが日付・期間になっている＝横持ち（クロス表）の目印
_PERIOD_RE = re.compile(
    r"^\s*(?:"
    r"(?:19|20)\d{2}[-/年.]?(?:0?[1-9]|1[0-2])?[月]?"      # 2026-04 / 2026年4月
    r"|(?:0?[1-9]|1[0-2])月"                                # 4月
    r"|[QＱ][1-4]|第[1-4一二三四]四半期"                     # Q1 / 第1四半期
    r"|上期|下期|上半期|下半期"
    r")\s*$")


def _blank(v) -> bool:
    return v is None or str(v).strip() == ""


def _issue(level: str, text: str, fix: str = "") -> dict:
    return {"level": level, "text": text, "fix": fix}


# =============================================================================
# ファイルを「素の格子」として読む（見出しがどこかは、まだ決めつけない）
# =============================================================================

def _grid_excel(path: Path, sheet: str | None) -> tuple:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        names = list(wb.sheetnames)
        target = sheet if (sheet and sheet in names) else names[0]
        ws = wb[target]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_SCAN_ROWS:
                break
            rows.append(list(row))
        return rows, names, target
    finally:
        wb.close()


def _merged_ranges(path: Path, sheet: str) -> list | None:
    """結合セルの範囲。読み取り専用モードでは取れないので通常読み込みする。

    大きいファイルで開くと重いので、その場合は調べずに None を返す
    （「分からなかった」と「無かった」を混同しないため）。
    """
    try:
        if path.stat().st_size > MERGE_CHECK_MAX_MB * 1024 * 1024:
            return None
    except OSError:
        return None
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=False, data_only=True)
        try:
            ws = wb[sheet] if sheet in wb.sheetnames else wb[wb.sheetnames[0]]
            return [(r.min_row, r.min_col, r.max_row, r.max_col)
                    for r in ws.merged_cells.ranges]
        finally:
            wb.close()
    except Exception:
        return None


def _grid_text(path: Path) -> tuple:
    """CSV/TSV/TXT。区切り文字と文字コードもここで見当をつける。"""
    from importer import CSV_ENCODINGS

    raw = None
    for enc in CSV_ENCODINGS:
        try:
            raw = path.read_text(encoding=enc)
            used = enc
            break
        except (UnicodeDecodeError, OSError):
            continue
    if raw is None:
        raise ValueError("文字コードを判定できませんでした（UTF-8 か Shift_JIS で保存し直してください）。")

    head = "\n".join(raw.splitlines()[:MAX_SCAN_ROWS])
    if path.suffix.lower() == ".tsv":
        delim = "\t"
    else:
        try:
            delim = csv.Sniffer().sniff(head[:4000], delimiters=",\t;|").delimiter
        except csv.Error:
            delim = "\t" if head.count("\t") > head.count(",") else ","
    rows = [r for r in csv.reader(head.splitlines(), delimiter=delim)]
    return rows, used, delim


# =============================================================================
# 形を見る
# =============================================================================

def _guess_header(rows: list) -> int:
    """見出しの行番号（0始まり）を当てる。

    「文字が並んでいて、その下に中身が続いている行」を見出しとみなす。
    タイトル行（1セルだけ埋まっている）や空行は飛ばす。
    """
    best, best_score = 0, -1.0
    for i, row in enumerate(rows[:HEADER_SEARCH_ROWS]):
        filled = [v for v in row if not _blank(v)]
        if len(filled) < 2:
            continue                       # タイトル行や空行
        below = rows[i + 1] if i + 1 < len(rows) else []
        if not [v for v in below if not _blank(v)]:
            continue                       # 下に中身が無いなら見出しではない
        texty = sum(1 for v in filled if not _looks_number(v))
        # 埋まり具合＋文字らしさ。上の行ほど見出しらしいので少し優遇する
        score = (len(filled) / max(len(row), 1)) + (texty / len(filled)) - i * 0.06
        if score > best_score:
            best, best_score = i, score
    return best


def _looks_number(v) -> bool:
    s = str(v).strip().replace(",", "")
    if not s:
        return False
    try:
        float(s)
        return True
    except ValueError:
        return False


def _is_total_row(row: list) -> bool:
    head = " ".join(str(v) for v in row[:2] if not _blank(v)).strip().lower()
    return bool(head) and any(w in head for w in _TOTAL_WORDS)


def _analyze(rows: list, header_row: int) -> dict:
    """見出し行を決めた上で、中身の形を調べる。"""
    header = rows[header_row] if header_row < len(rows) else []
    body = rows[header_row + 1:]
    width = max((len(r) for r in rows), default=0)

    names = [("" if _blank(v) else str(v).strip()) for v in header]
    names += [""] * (width - len(names))

    empty_names = [i for i, n in enumerate(names) if not n]
    dup = sorted({n for n in names if n and names.count(n) > 1})
    period_cols = [n for n in names if n and _PERIOD_RE.match(n)]
    numeric_names = [n for n in names if n and _looks_number(n)]

    blank_rows = sum(1 for r in body if all(_blank(v) for v in r))
    total_rows = [i for i, r in enumerate(body) if _is_total_row(r)]
    ragged = sum(1 for r in body if len([v for v in r if not _blank(v)]) > len(names))
    multiline = any(isinstance(v, str) and "\n" in v for r in body[:50] for v in r)

    # 全部空の列（見出しだけあって中身が無い／見出しも中身も無い）
    empty_cols = []
    for c in range(width):
        col = [r[c] for r in body if c < len(r)]
        if col and all(_blank(v) for v in col):
            empty_cols.append(names[c] or f"{c + 1}列目")

    # 数字と文字が混ざる列（"-" や "N/A" が入ると、数値として取り込めない）
    mixed = []
    for c in range(width):
        col = [r[c] for r in body if c < len(r) and not _blank(r[c])]
        if len(col) < 4:
            continue
        nums = sum(1 for v in col if _looks_number(v))
        if 0.6 <= nums / len(col) < 1.0:
            odd = [str(v) for v in col if not _looks_number(v)][:3]
            mixed.append((names[c] or f"{c + 1}列目", odd))

    return {
        "names": names, "width": width, "body_rows": len(body),
        "empty_names": empty_names, "dup_names": dup,
        "period_cols": period_cols, "numeric_names": numeric_names,
        "blank_rows": blank_rows, "total_rows": total_rows,
        "ragged": ragged, "multiline": multiline,
        "empty_cols": empty_cols, "mixed": mixed,
    }


def _blocks(rows: list, header_row: int) -> int:
    """1シートに表がいくつ入っていそうか（空行で切れて、また見出しが始まる）。"""
    blocks, in_block, gap = 1, True, 0
    for r in rows[header_row + 1:]:
        if all(_blank(v) for v in r):
            gap += 1
            in_block = False
        else:
            if not in_block and gap >= 2:
                blocks += 1
            in_block, gap = True, 0
    return blocks


# =============================================================================
# 判定
# =============================================================================

def inspect(path, sheet: str | None = None) -> dict:
    """1ファイル（Excelは1シート）の形を見て、取り込めるかを判定する。"""
    p = Path(path)
    ext = p.suffix.lower()
    out = {"file": p.name, "sheet": sheet, "sheets": [], "header_row": 0,
           "verdict": "", "issues": [], "shape": {}}

    if ext not in config.IMPORT_EXTENSIONS:
        out["verdict"] = "対応していない形式"
        out["issues"] = [_issue(
            "高", f"{ext or '拡張子なし'} は取り込みに対応していません。",
            f"扱えるのは {'、'.join(config.IMPORT_EXTENSIONS)} です。"
            "元のシステムからCSVで出し直すか、Excelで開いて「名前を付けて保存」で"
            ".xlsx か .csv にしてください。")]
        return out

    try:
        if ext in (".xlsx", ".xlsm"):
            rows, names, target = _grid_excel(p, sheet)
            out["sheets"], out["sheet"] = names, target
            merged = _merged_ranges(p, target)
        else:
            rows, enc, delim = _grid_text(p)
            merged = []
            out["encoding"] = enc
            out["delimiter"] = {"\t": "タブ", ",": "カンマ", ";": "セミコロン",
                                "|": "パイプ"}.get(delim, delim)
    except Exception as e:
        out["verdict"] = "取り込みに向かない"
        out["issues"] = [_issue("高", f"ファイルを開けませんでした: {e}",
                                "壊れているか、パスワードが掛かっている可能性があります。")]
        return out

    rows = [r for r in rows if r is not None]
    if not any(any(not _blank(v) for v in r) for r in rows):
        out["verdict"] = "取り込みに向かない"
        out["issues"] = [_issue("高", "中身が空です。", "データの入ったファイルを指定してください。")]
        return out

    header_row = _guess_header(rows)
    info = _analyze(rows, header_row)
    blocks = _blocks(rows, header_row)
    out["header_row"] = header_row
    out["shape"] = {"列数": info["width"], "読んだ行数": info["body_rows"],
                    "見出し行": header_row + 1}
    out["columns"] = info["names"]

    issues: list[dict] = []

    # --- そのままでは取り込めないもの ---------------------------------------
    if merged:
        hrow = header_row + 1                    # 1始まりの行番号にそろえる
        # 見出しの行と、そのすぐ上をまたぐ横方向の結合＝多段見出し。
        # いちばん上のタイトル行（1セルだけの飾り）は、これに含めない。
        in_header = [m for m in merged
                     if m[1] != m[3] and m[0] <= hrow and m[2] >= hrow - 1
                     and not (m[0] == m[2] == 1 and hrow > 2)]
        in_body = [m for m in merged if m[0] > hrow]
        if in_header:
            issues.append(_issue(
                "高", f"見出しがセル結合されています（{len(in_header)}箇所）。"
                      "多段の見出しは1行の列名にできません。",
                "結合を解除し、見出しを1行にまとめてください"
                "（例:「上期／4月」→「上期_4月」）。"))
        if in_body:
            issues.append(_issue(
                "高", f"データ部分にセル結合があります（{len(in_body)}箇所）。"
                      "結合されたセルは先頭以外が空になり、行が正しく揃いません。",
                "結合を解除し、空いたセルに同じ値を埋めてください。"))

    if len(info["period_cols"]) >= 3:
        issues.append(_issue(
            "高", f"月や期間が横に並んでいます（{'、'.join(info['period_cols'][:5])}…）。"
                  "いわゆるクロス表で、1行=1レコードになっていません。",
            "「年月」「値」の2列に縦持ちへ直してください"
            "（Excelなら [データ]→[パワークエリ]→[列のピボット解除]）。"))
    elif len(info["numeric_names"]) >= 3:
        issues.append(_issue(
            "高", f"見出しが数字になっています（{'、'.join(info['numeric_names'][:5])}…）。"
                  "見出し行の位置が違うか、横持ちの表の可能性があります。",
            "1行目に列名が来るようにしてください。"))

    if blocks > 1:
        issues.append(_issue(
            "高", f"1つのシートに表が{blocks}個あるように見えます（間に空行があります）。",
            "表ごとにシートを分けてください。取り込みは1シート=1テーブルです。"))

    # --- 直せば取り込めるもの ------------------------------------------------
    if header_row > 0:
        issues.append(_issue(
            "中", f"{header_row + 1}行目が見出しに見えます（1行目ではありません）。"
                  "上にタイトルや空行が入っています。",
            f"取り込み画面の「見出しの行」に {header_row + 1} を指定するか、"
            "上の行を削除してください。"))

    if info["total_rows"]:
        issues.append(_issue(
            "中", f"合計・小計らしい行が {len(info['total_rows'])} 行あります。"
                  "そのまま取り込むと二重に数えられます。",
            "合計行を削除してから取り込んでください（集計はアプリ側でできます）。"))

    if info["empty_names"]:
        issues.append(_issue(
            "中", f"列名が空の列が {len(info['empty_names'])} 個あります。",
            "列名を付けてください（空のままだと自動で仮の名前が付きます）。"))

    if info["dup_names"]:
        issues.append(_issue(
            "中", f"同じ列名が複数あります: {'、'.join(info['dup_names'][:5])}",
            "区別できる名前に変えてください（取り込み時は連番が付きます）。"))

    if info["ragged"]:
        issues.append(_issue(
            "中", f"見出しより列が多い行が {info['ragged']} 行あります。"
                  "区切り文字がデータの中に入っている可能性があります。",
            "その列を引用符で囲むか、区切り文字を変えて出し直してください。"))

    # --- 気に留めておく程度 --------------------------------------------------
    for name, odd in info["mixed"][:3]:
        issues.append(_issue(
            "低", f"「{name}」は数字の列に見えますが、文字が混ざっています"
                  f"（{'、'.join(odd)}）。",
            "空欄や「-」「N/A」は空にしておくと、数値として取り込めます。"))
    if info["empty_cols"]:
        issues.append(_issue(
            "低", f"中身が空の列があります: {'、'.join(info['empty_cols'][:5])}",
            "取り込む列の選択から外せます。"))
    if info["blank_rows"]:
        issues.append(_issue(
            "低", f"途中に空行が {info['blank_rows']} 行あります。", "空行は取り込み時に残ります。"))
    if info["multiline"]:
        issues.append(_issue(
            "低", "セルの中で改行しているところがあります。",
            "表示は崩れませんが、検索や集計がしにくくなります。"))
    if merged is None:
        issues.append(_issue(
            "低", "ファイルが大きいため、セル結合までは調べていません。", ""))

    levels = {i["level"] for i in issues}
    out["issues"] = issues
    out["verdict"] = ("取り込みに向かない" if "高" in levels else
                      "手直しが要る" if "中" in levels else
                      "そのまま取り込める")
    return out


def summary_line(res: dict) -> str:
    """一覧に出す一言。"""
    high = [i for i in res["issues"] if i["level"] == "高"]
    if res["verdict"] == "そのまま取り込める":
        return "そのまま取り込める"
    if high:
        return f"{res['verdict']}（{high[0]['text'][:40]}）"
    mid = [i for i in res["issues"] if i["level"] == "中"]
    return f"{res['verdict']}（{mid[0]['text'][:40]}）" if mid else res["verdict"]


# ==========================================================================
# ===== 元 chats.py
# ユーザーごとのチャット履歴。
#
#   data/users/<ユーザー>/chats/index.json … 一覧（タイトルと日時だけ）
#   data/users/<ユーザー>/chats/<ID>.json  … 会話の中身
#
# 会話1件に保存するのは次の2つ。
#
#   messages    … LLMに送るメッセージ列。これが無いと「続きから」会話できない
#   render_log  … 画面に描くアイテム（テキスト・SQL・表・グラフ・作成ファイル）
#
# 一覧を別ファイルにしているのは、サイドバーを描くたびに全会話を読み込まないため。
# index.json が壊れた/消えた場合は、置いてある会話ファイルから作り直す。
#
# 古い会話は2つの条件で自動的に消える。
#   本数   … CHAT_HISTORY_LIMIT を超えたぶん（古い順）
#   保存期間 … 最後に使った日から CHAT_HISTORY_DAYS を過ぎたもの（既定90日）
# 掃除は一覧を読むついでに行う。常駐の掃除役を置かずに済ませるため。
#
# Excel等の作成ファイルはバイト列なのでJSONに入らない。上限までは base64 で埋め込み、
# 超えるものは本体を捨てる（過去の会話を開いても再ダウンロードはできない）。
# ==========================================================================
import base64
import json
import re
import uuid
from datetime import datetime
from pathlib import Path

import config

_INDEX_NAME = "index.json"
_ID_RE = re.compile(r"[^0-9a-zA-Z_-]")
_TITLE_MAX = 40


# --- 場所 ---------------------------------------------------------------------

def chats_dir(user) -> Path:
    key = getattr(user, "safe_key", None) or str(user)
    return config.USER_META_DIR / key / "chats"


def _safe_id(chat_id: str) -> str:
    """ファイル名に使う前に無害化する（.. や / を混ぜられないように）。"""
    return _ID_RE.sub("", str(chat_id))[:64]


def _chat_file(user, chat_id: str) -> Path:
    return chats_dir(user) / f"{_safe_id(chat_id)}.json"


def new_id() -> str:
    # 先頭に日時を置いて、ファイル名を見ただけで新しい順に並ぶようにする
    return datetime.now().strftime("%Y%m%d-%H%M%S-") + uuid.uuid4().hex[:6]


def _now() -> str:
    return datetime.now().isoformat(timespec="seconds")


def now() -> str:
    """表示物に打つ時刻。会話を積む側（chat_bp）から使う。"""
    return _now()


# --- ファイル入出力 -------------------------------------------------------------

def _read_json(p: Path):
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"[chats] 読めませんでした: {p} ({e})")
        return None


def _write_json(p: Path, data) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(data, ensure_ascii=False, default=str), encoding="utf-8")


# --- 一覧 ---------------------------------------------------------------------

def _index_path(user) -> Path:
    return chats_dir(user) / _INDEX_NAME


def _rebuild_index(user) -> list[dict]:
    """会話ファイルから一覧を作り直す（index.json を失った場合の保険）。"""
    items = []
    for p in chats_dir(user).glob("*.json"):
        if p.name == _INDEX_NAME:
            continue
        data = _read_json(p)
        if not isinstance(data, dict) or not data.get("id"):
            continue
        items.append(_summary(data))
    items.sort(key=lambda c: c.get("updated_at") or "", reverse=True)
    items = _drop_expired(user, items)
    if items:
        _write_json(_index_path(user), {"chats": items})
    return items


def _expired(summary: dict) -> bool:
    """保存期間を過ぎた会話か。

    数えるのは「最後に使った日」から。開いて続きを話した会話は寿命が延びる。
    日付が読めないものは、消して困る方が大きいので残す。
    """
    if config.CHAT_HISTORY_DAYS <= 0:
        return False
    stamp = summary.get("updated_at") or summary.get("created_at") or ""
    try:
        used = datetime.fromisoformat(str(stamp))
    except ValueError:
        return False
    return (datetime.now() - used).days > config.CHAT_HISTORY_DAYS


def _drop_expired(user, items: list[dict]) -> list[dict]:
    """期限切れの会話を実体ごと消して、残ったぶんを返す。"""
    keep, gone = [], 0
    for c in items:
        if _expired(c):
            f = _chat_file(user, c.get("id", ""))
            if f.exists():
                f.unlink()
            gone += 1
        else:
            keep.append(c)
    if gone:
        print(f"[chats] 保存期間({config.CHAT_HISTORY_DAYS}日)を過ぎた会話を"
              f"{gone}件削除しました（{_key_label(user)}）")
    return keep


def _key_label(user) -> str:
    return getattr(user, "username", None) or str(user)


def list_chats(user) -> list[dict]:
    """会話の一覧（新しい順）。中身は読まない。"""
    if user is None or not chats_dir(user).exists():
        return []
    data = _read_json(_index_path(user)) if _index_path(user).exists() else None
    items = data.get("chats") if isinstance(data, dict) else data
    if not isinstance(items, list):
        items = None
    if not items:
        return _rebuild_index(user)
    # 実体が消えているものは一覧からも落とす
    items = [c for c in items
             if isinstance(c, dict) and c.get("id") and _chat_file(user, c["id"]).exists()]
    items.sort(key=lambda c: c.get("updated_at") or "", reverse=True)
    # 期限切れは、一覧を出すついでに片付ける（掃除専用の仕組みを持たない）
    kept = _drop_expired(user, items)
    if len(kept) != len(items):
        _save_index(user, kept)
    return kept


def _summary(chat: dict) -> dict:
    return {
        "id": chat.get("id"),
        "title": chat.get("title") or "（無題）",
        "created_at": chat.get("created_at") or "",
        "updated_at": chat.get("updated_at") or "",
        "db_names": list(chat.get("db_names") or []),
        "n_turns": sum(1 for m in (chat.get("messages") or []) if m.get("role") == "user"),
    }


def _save_index(user, items: list[dict]) -> None:
    _write_json(_index_path(user), {"chats": items})


def _upsert_index(user, summary: dict) -> list[dict]:
    items = [c for c in list_chats(user) if c.get("id") != summary["id"]]
    items.insert(0, summary)
    items.sort(key=lambda c: c.get("updated_at") or "", reverse=True)

    # 上限を超えた古い会話は実体ごと削除
    for old in items[config.CHAT_HISTORY_LIMIT:]:
        f = _chat_file(user, old.get("id", ""))
        if f.exists():
            f.unlink()
    items = items[:config.CHAT_HISTORY_LIMIT]
    _save_index(user, items)
    return items


# --- 作成ファイル(bytes)の出し入れ ------------------------------------------------

def _encode_item(item: dict) -> dict:
    out = dict(item)
    # 表示物1つずつに時刻を持たせる。会話の created_at だけでは
    # 「その日に始めた」までしか分からず、いつ何を聞いたのかを追えない。
    # 質問には積んだ時刻が入っているので、ここで入るのは応答側の完了時刻になる。
    # 差を取れば、その質問にどれだけ待たされたかも分かる。
    out.setdefault("at", _now())
    data = out.get("data")
    if isinstance(data, (bytes, bytearray)):
        if len(data) <= config.CHAT_EMBED_FILE_MAX_BYTES:
            out["data"] = base64.b64encode(bytes(data)).decode("ascii")
            out["_b64"] = True
        else:                              # 大きすぎるので中身は保存しない
            out.pop("data", None)
            out["_no_data"] = True
    return out


def _decode_item(item: dict) -> dict:
    out = dict(item)
    if out.pop("_b64", False):
        try:
            out["data"] = base64.b64decode(out.get("data") or "")
        except Exception:
            out.pop("data", None)
            out["_no_data"] = True
    return out


# --- 読み書き -------------------------------------------------------------------

def make_title(messages: list[dict]) -> str:
    """最初のユーザー発言をタイトルにする。"""
    for m in messages or []:
        if m.get("role") != "user" or not m.get("content"):
            continue
        content = m["content"]
        if isinstance(content, list):
            # 画像つきの発言は content が配列。文字の部分だけ拾う。
            content = "".join(p.get("text", "") for p in content
                              if isinstance(p, dict) and p.get("type") == "text")
        t = " ".join(str(content).split())
        if t:
            return t[:_TITLE_MAX] + ("…" if len(t) > _TITLE_MAX else "")
    return "（無題）"


def save_chat(user, chat_id: str, messages: list[dict], render_log: list[dict],
              db_names=None, tables=None, title: str = "", created_at: str = "") -> dict:
    """会話を保存し、一覧用の要約を返す。"""
    chat = {
        "id": chat_id,
        "title": title or make_title(messages),
        "created_at": created_at or _now(),
        "updated_at": _now(),
        "db_names": list(db_names or []),
        "tables": dict(tables or {}),
        # system prompt は開くたびに作り直すので保存しない（カタログ変更に追従させる）
        "messages": [m for m in (messages or []) if m.get("role") != "system"],
        "render_log": [_encode_item(i) for i in (render_log or [])],
    }
    _write_json(_chat_file(user, chat_id), chat)
    summary = _summary(chat)
    _upsert_index(user, summary)
    return summary


def load_chat(user, chat_id: str) -> dict | None:
    p = _chat_file(user, chat_id)
    if not p.exists():
        return None
    data = _read_json(p)
    if not isinstance(data, dict):
        return None
    data["render_log"] = [_decode_item(i) for i in (data.get("render_log") or [])]
    data.setdefault("messages", [])
    return data


def rename_chat(user, chat_id: str, title: str) -> bool:
    chat = load_chat(user, chat_id)
    if chat is None:
        return False
    chat["title"] = (title or "").strip()[:_TITLE_MAX] or make_title(chat.get("messages"))
    chat["render_log"] = [_encode_item(i) for i in (chat.get("render_log") or [])]
    _write_json(_chat_file(user, chat_id), chat)
    _upsert_index(user, _summary(chat))
    return True


def delete_chat(user, chat_id: str) -> bool:
    p = _chat_file(user, chat_id)
    existed = p.exists()
    if existed:
        p.unlink()
    _save_index(user, [c for c in list_chats(user) if c.get("id") != chat_id])
    return existed


def label(summary: dict) -> str:
    """サイドバーのプルダウンに出す1行。"""
    stamp = (summary.get("updated_at") or "")[5:16].replace("T", " ")   # MM-DD HH:MM
    title = summary.get("title") or "（無題）"
    return f"{title}　（{stamp}）" if stamp else title


# ==========================================================================
# ===== 元 history.py
# 取り込みの更新履歴。
#
# ジョブ定義（import_jobs.yaml）が持っているのは直前1回ぶんの結果だけなので、
# 「先週の火曜は何行入ったのか」「いつから失敗し続けているのか」を追えない。
# そこで、1回の取り込みにつき1件をここに追記していく。
#
# 置き場所は data/import_history.jsonl（1行1件のJSON）。
# YAML ではなく追記型にしているのは、実行のたびに全件を書き直したくないため。
# 手動の取り込みも定期取り込みも同じ形で残し、kind で区別する。
# ==========================================================================
import json
import threading
from datetime import datetime
from pathlib import Path

import config

_history_lock = threading.Lock()
# 追記のたびに全件を数え直さないよう、行数はプロセス内で覚えておく。
# 別プロセス（refresh.py など）が書くとずれるが、間引きは後追いで効けばよい。
_count: int | None = None

IMPORT_RECORD_KINDS = {"manual": "手動", "auto": "定期", "job": "定期（手動実行）"}


def _history_path() -> Path:
    return config.IMPORT_HISTORY_FILE


def add_import_record(db_file: str, table: str, ok: bool, message: str, *,
        kind: str = "manual", mode: str = "replace", rows: int = 0,
        removed: int = 0, kept=None, keep=None, source: str = "",
        sheet: str | None = None, job_id: str | None = None,
        job_name: str | None = None, user: str | None = None,
        started: datetime | None = None) -> dict:
    """1回ぶんの結果を残す。記録に失敗しても取り込み自体は止めない。"""
    now = datetime.now()
    rec = {
        "at": (started or now).isoformat(timespec="seconds"),
        "db_file": db_file, "table": table,
        "ok": bool(ok), "kind": kind, "mode": mode,
        "rows": int(rows or 0), "removed": int(removed or 0),
        "kept": kept, "keep": keep,
        "source": str(source or ""), "sheet": sheet,
        "job_id": job_id, "job_name": job_name, "user": user,
        "message": message,
        "seconds": round((now - started).total_seconds(), 1) if started else None,
    }
    global _count
    try:
        with _history_lock:
            p = _history_path()
            p.parent.mkdir(parents=True, exist_ok=True)
            if _count is None:
                _count = _line_count(p)
            with p.open("a", encoding="utf-8") as f:
                f.write(json.dumps(rec, ensure_ascii=False) + "\n")
            _count += 1
            _trim_if_needed(p)
    except Exception as e:                       # 履歴が書けなくても取り込みは成功扱い
        print(f"[history] 記録できませんでした: {e}")
    return rec


def _line_count(p: Path) -> int:
    if not p.exists():
        return 0
    with p.open("rb") as f:
        return sum(1 for line in f if line.strip())


def _trim_if_needed(p: Path) -> None:
    """行数が上限を超えたら、新しい方から上限ぶんだけ残す。

    毎回書き直すと重いので、1割ぶん超えてからまとめて間引く。
    """
    global _count
    limit = max(1, config.IMPORT_HISTORY_MAX)
    if (_count or 0) <= limit * 1.1:
        return
    lines = [x for x in p.read_text(encoding="utf-8").splitlines() if x.strip()]
    keep = lines[-limit:]
    p.write_text("\n".join(keep) + "\n", encoding="utf-8")
    _count = len(keep)


def _read_all() -> list[dict]:
    p = _history_path()
    if not p.exists():
        return []
    out = []
    try:
        for line in p.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
            except json.JSONDecodeError:
                continue                          # 壊れた行は飛ばす
            if isinstance(rec, dict):
                out.append(rec)
    except Exception as e:
        print(f"[history] 読めませんでした: {p} ({e})")
    return out


def _newest_first(items: list[dict]) -> list[dict]:
    """新しい順に並べる。

    at は秒までしか持たないので、同じ秒の中は「後に書いた方が新しい」で決める。
    先に並びを逆にしてから安定ソートすると、同着がその順で残る。
    """
    items = list(reversed(items))
    items.sort(key=lambda r: r.get("at") or "", reverse=True)
    return items


def for_table(db_file: str, table: str, limit: int = 30) -> list[dict]:
    """あるテーブルの履歴を新しい順で。"""
    hit = _newest_first([r for r in _read_all()
                         if r.get("db_file") == db_file and r.get("table") == table])
    return hit[:limit] if limit else hit


def recent_import_records(limit: int = 100) -> list[dict]:
    """テーブルを問わず、新しい順に。取り込み全体の傾向を見るとき用。"""
    hit = _newest_first(_read_all())
    return hit[:limit] if limit else hit


def counts() -> dict[tuple, int]:
    """(DB, テーブル) ごとの件数。"""
    out: dict[tuple, int] = {}
    for r in _read_all():
        key = (r.get("db_file"), r.get("table"))
        out[key] = out.get(key, 0) + 1
    return out


def latest_by_source() -> dict[str, dict]:
    """取り込み元ファイルごとの、いちばん新しい記録。

    「このファイルはもう取り込んだのか」「いつ・どのテーブルに入ったのか」を
    ファイルの一覧と突き合わせるために使う。キーはファイルパス。
    """
    out: dict[str, dict] = {}
    for r in _newest_first(_read_all()):
        src = str(r.get("source") or "")
        if src and src not in out:
            out[src] = r
    return out


# ==========================================================================
# ===== 元 catalog_history.py
# 用語集・例文の変更履歴。誰が・いつ・何を・どう変えたかを残す。
#
# カタログは全員共通の土台で、チャットからは一般ユーザーも書けるようにした。
# 書けるようにした以上、「いつの間にか定義が変わっていた」が起きるので、
# 変更のたびに1件を追記して、後から辿れるようにする。
#
# 置き場所は data/catalog_history.jsonl（1行1件のJSON・追記型）。
# import_history と同じ考え方で、YAMLに混ぜない（メタ情報は「現在の定義」だけを
# 持ち、履歴で膨らませない。normalize が知らないキーを消す作りとも衝突しない）。
# ==========================================================================
import json
import threading
from datetime import datetime
from pathlib import Path

import config

_catalog_history_lock = threading.Lock()

#: 表示用のラベル
CATALOG_CHANGE_KINDS = {"glossary": "用語", "example": "例文"}
OPS = {"add": "新規", "update": "変更", "remove": "削除"}


def _catalog_history_path() -> Path:
    return config.CATALOG_HISTORY_FILE


def add_catalog_change(kind: str, op: str, db_file: str, name: str, *,
        user: str | None = None, table: str | None = None,
        before=None, after=None, source: str = "chat") -> None:
    """1件追記する。失敗しても本体の保存は止めない（履歴は本体より弱い）。"""
    rec = {
        "at": datetime.now().isoformat(timespec="seconds"),
        "kind": kind, "op": op, "db": db_file, "table": table or "",
        "name": name, "user": user or "不明", "source": source,
        "before": before, "after": after,
    }
    try:
        with _catalog_history_lock:
            p = _catalog_history_path()
            p.parent.mkdir(parents=True, exist_ok=True)
            with p.open("a", encoding="utf-8") as f:
                f.write(json.dumps(rec, ensure_ascii=False, default=str) + "\n")
            _trim(p)
    except Exception as e:
        print(f"[catalog_history] 書けませんでした: {e}")


def _trim(p: Path) -> None:
    """上限を超えたら古い行から捨てる（毎回数えず、たまに間引く）。"""
    lines = p.read_text(encoding="utf-8").splitlines()
    if len(lines) > config.CATALOG_HISTORY_MAX * 1.2:
        keep = lines[-config.CATALOG_HISTORY_MAX:]
        p.write_text("\n".join(keep) + "\n", encoding="utf-8")


def recent_catalog_changes(limit: int = 50) -> list[dict]:
    """新しい順に。カタログ画面の「変更履歴」に出す。"""
    p = _catalog_history_path()
    if not p.exists():
        return []
    out = []
    for line in p.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            rec = json.loads(line)
        except json.JSONDecodeError:
            continue
        if isinstance(rec, dict):
            out.append(rec)
    out.reverse()
    return out[:limit] if limit else out


def summarize_catalog_changes(rec: dict) -> str:
    """1件を画面向けの短い日本語に。"""
    where = f"{rec.get('db', '')}" + (f" の {rec['table']}" if rec.get("table") else "")
    label = f"{CATALOG_CHANGE_KINDS.get(rec.get('kind'), rec.get('kind'))}「{rec.get('name', '')}」"
    op = OPS.get(rec.get("op"), rec.get("op"))
    return f"{where}: {label} を{op}"


# ==========================================================================
# ===== 元 prefs.py
# ログインユーザーごとの画面の状態。
#
#   data/users/<ユーザー>/prefs.yaml
#
# 覚えておくのは次の2つ。
#
#   selection … 対象データの選択（{DBファイル名: [テーブル名, ...]}）
#   model     … 使うモデル
#
# セッション（Cookie）に置くとログアウトやブラウザを閉じたときに消えてしまう。
# 毎回選び直すのは手間なので、そのユーザーのフォルダにファイルとして残す。
# カタログやチャット履歴と同じ場所に置くので、退職者のデータを消すときは
# そのユーザーのフォルダごと消せばよい。
# ==========================================================================
import threading

import yaml

import config

_prefs_lock = threading.Lock()

# ここに挙げたキーだけを読み書きする（余計なものが混ざっても無視する）
KEYS = ("selection", "model")


def _key(user) -> str:
    """保存先のフォルダ名。catalog / chats と同じ決め方にする。"""
    return getattr(user, "safe_key", None) or str(user)


def _prefs_path(user):
    return config.USER_META_DIR / _key(user) / "prefs.yaml"


def load(user) -> dict:
    if user is None:
        return {}
    p = _prefs_path(user)
    if not p.exists():
        return {}
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}
    except Exception as e:
        print(f"[prefs] 読めませんでした: {p} ({e})")
        return {}
    return {k: v for k, v in data.items() if k in KEYS} if isinstance(data, dict) else {}


def _save(user, data: dict) -> None:
    p = _prefs_path(user)
    p.parent.mkdir(parents=True, exist_ok=True)
    with _prefs_lock:
        p.write_text(yaml.safe_dump({k: data[k] for k in KEYS if k in data},
                                    allow_unicode=True, sort_keys=False),
                     encoding="utf-8")


def set_value(user, key: str, value) -> None:
    """1項目だけ更新する。他の項目は触らない。"""
    if user is None or key not in KEYS:
        return
    data = load(user)
    data[key] = value
    _save(user, data)


# --- 対象データの選択 -------------------------------------------------------------

def get_selection(user) -> dict:
    sel = load(user).get("selection")
    if not isinstance(sel, dict):
        return {}
    # 保存後にDBやテーブルが消えている場合もあるが、そこは build_scope 側で弾かれる
    return {str(k): [str(t) for t in (v or [])] for k, v in sel.items()}


def set_selection(user, selection: dict) -> dict:
    sel = {str(k): [str(t) for t in (v or [])] for k, v in (selection or {}).items()}
    set_value(user, "selection", sel)
    return sel


# 対象データの選択（selection）は選択UIの廃止で書き込まれなくなったが、
# 既存ユーザーのファイルに残った選択を cleanup.py が掃除するため、読み書きは残す。

# --- モデルの選択 ----------------------------------------------------------------

def get_model(user) -> str:
    return str(load(user).get("model") or "").strip()


def set_model(user, model: str) -> None:
    set_value(user, "model", str(model or "").strip())


# ==========================================================================
# ===== 元 models.py
# 使うモデルの選択。
#
# 2段構えになっている。
#   管理者 … 「モデル設定」画面で、選ばせる候補・既定・画像対応の判定を決める
#   利用者 … チャット画面のプルダウンで、その候補から自分の1つを選ぶ
#
# 利用者が選んだモデルは prefs.py（ユーザーごとのファイル）に残るので、
# ログアウトしても次に入ったときは同じモデルのまま。
#
# 候補の決まり方は 管理者の設定 > env の OPENAI_MODELS > APIの /models の順。
# 「そのモデルは画像を送れるか」も、ここで一元的に判断する。
# ==========================================================================
import threading
import time

import yaml

import config
import prefs

_models_lock = threading.Lock()
_models_cache: dict = {"at": 0.0, "models": []}
_CACHE_SEC = 300


# =============================================================================
# 管理者が決める設定（data/model_settings.yaml）
#
# env を初期値として、このファイルの内容で上書きする。
# メール設定と同じ考え方で、env は「まだ画面で決めていないときの値」。
# =============================================================================

ADMIN_KEYS = ("models", "default", "vision", "context_overrides")

#: カタログのインライン上限として認める範囲。
#: 下限は「1DBぶんの詳細（実測で平均5.3K字）が入る」ことを目安にした。
#: 上限は、いちばん広いモデルでも文脈を食い尽くさないところで止める。
INLINE_LIMIT_MIN = 4_000
INLINE_LIMIT_MAX = 400_000


def _read_admin() -> dict:
    p = config.MODEL_SETTINGS_FILE
    if not p.exists():
        return {}
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}
    except Exception as e:
        print(f"[models] 設定を読めませんでした: {p} ({e})")
        return {}
    return {k: v for k, v in data.items() if k in ADMIN_KEYS} \
        if isinstance(data, dict) else {}


def _write_admin(data: dict) -> None:
    p = config.MODEL_SETTINGS_FILE
    p.parent.mkdir(parents=True, exist_ok=True)
    with _models_lock:
        p.write_text(yaml.safe_dump({k: data[k] for k in ADMIN_KEYS if k in data},
                                    allow_unicode=True, sort_keys=False),
                     encoding="utf-8")


def _vision_keys() -> list[str]:
    ov = _read_admin().get("vision")
    keys = ov if isinstance(ov, list) else None
    return [str(k).strip().lower() for k in (keys or config.OPENAI_VISION_MODELS)
            if str(k).strip()]


def default_model() -> str:
    """未選択のユーザーが使うモデル。"""
    return str(_read_admin().get("default") or config.OPENAI_MODEL or "").strip()


def is_vision(model: str) -> bool:
    """画像を送れるモデルか。名前に手がかりが含まれるかで判断する。

    モデル名は環境によって違うので、完全一致ではなく部分一致にしている
    （「モデル設定」画面、または env の OPENAI_VISION_MODELS で調整できる）。
    """
    low = str(model or "").lower()
    return any(key in low for key in _vision_keys())


def prompt_inline_limit() -> int:
    """カタログをそのまま入れる量の天井（文字）。

    実効値はモデルの文脈から自動で決まる（inline_limit_for）。この天井は
    「文脈が100万トークンあるモデルでも、カタログに割く量はここまで」という
    安全弁で、画面から変えるものではない。
    """
    return INLINE_LIMIT_MAX


#: モデルの文脈のうち、カタログに使ってよい割合。
#: 残り半分はツール定義・会話の履歴・SQL結果・回答のために空けておく。
CATALOG_CONTEXT_RATIO = 0.5


def context_overrides() -> dict:
    """管理者が「モデル設定」画面で登録した文脈量 {モデル名(小文字): トークン}。

    公式の表（config.MODEL_CONTEXT_WINDOWS）に無いモデル — ゲートウェイ独自の名前や
    他社モデル — の実力を教えるための口。表より優先する。
    """
    raw = _read_admin().get("context_overrides") or {}
    out = {}
    for k, v in (raw.items() if isinstance(raw, dict) else []):
        try:
            n = int(v)
        except (TypeError, ValueError):
            continue
        if n > 0 and str(k).strip():
            out[str(k).strip().lower()] = n
    return out


def inline_limit_for(model: str | None = None) -> int:
    """カタログをそのまま入れる上限（文字）。選択中モデルの文脈量から自動で決める。

        上限 = 文脈(トークン) × 0.5 ÷ 0.55(日本語1文字あたりの概算トークン)

    文脈の半分をカタログに、残りをツール定義・履歴・回答に使う配分。
    文脈量は 管理者の登録 > 公式の表 > 既定値 の順で決まる（context_window）。
    天井（INLINE_LIMIT_MAX）と床（INLINE_LIMIT_MIN）で丸める。
    """
    if not model:
        return prompt_inline_limit()
    import llm                     # 循環importを避ける（llm側もmodelsを遅延importしている）
    context, _ = context_window(model)
    capacity = int(context * CATALOG_CONTEXT_RATIO / llm.TOKENS_PER_CHAR_TEXT)
    return max(INLINE_LIMIT_MIN, min(prompt_inline_limit(), capacity))


def catalog_total_chars() -> int:
    """全DBの詳細カタログの合計文字数（キャッシュ済みテキストを測るだけ）。"""
    import catalog as catalog_mod
    import db as db_mod
    return catalog_mod.inline_length(
        [{"path": str(f), "alias": db_mod.alias_for(f), "tables": None}
         for f in db_mod.list_db_files()])


def context_window(model: str) -> tuple:
    """そのモデルが一度に読める量（トークン）と、それが確かな値かどうか。

    戻り値: (トークン数, 分かっているモデルか)
    名前は環境によって違うので、前方一致の長い方から当てる。
    """
    low = str(model or "").lower()
    # 管理者の登録（完全一致 → 部分一致）> 公式の表（長い名前から部分一致）> 既定値（推定）
    ov = context_overrides()
    if low in ov:
        return ov[low], True
    for key in sorted(ov, key=len, reverse=True):
        if key and key in low:
            return ov[key], True
    for key in sorted(config.MODEL_CONTEXT_WINDOWS, key=len, reverse=True):
        if key in low:
            return config.MODEL_CONTEXT_WINDOWS[key], True
    return config.MODEL_CONTEXT_DEFAULT, False


def _from_api() -> list[str]:
    """APIに聞ける環境なら、使えるモデルの一覧を取ってくる。"""
    import llm
    if not llm.is_configured():
        return []
    now = time.time()
    with _models_lock:
        if _models_cache["models"] and now - _models_cache["at"] < _CACHE_SEC:
            return list(_models_cache["models"])
    try:
        got = sorted(m.id for m in llm.client().models.list().data)
    except Exception as e:
        print(f"[models] 一覧を取得できませんでした: {e}")
        got = []
    with _models_lock:
        _models_cache["at"], _models_cache["models"] = now, got
    return list(got)


def source() -> str:
    """候補がどこから来ているか。"admin" | "env" | "default"

    画面に出す文言を、実態とずれないようにするためのもの。
    """
    if [str(m).strip() for m in (_read_admin().get("models") or []) if str(m).strip()]:
        return "admin"
    return "env" if config.OPENAI_MODELS else "default"


def available(refresh: bool = False) -> list[str]:
    """チャット画面のプルダウンに出す候補。

    決まり方は 管理者の設定 > env の OPENAI_MODELS > 既定＋利用中のモデル。

    APIが返す一覧はここでは使わない。以前は最後の手段として使っていたが、
    それだと何も設定していないときに babbage-002 のような使えないモデルまで
    100件以上並び、「モデル設定」で絞ったつもりが効いていないように見えた。

    何も決めていないときは既定だけにしたいところだが、それだと以前の一覧から
    選んでいた人が黙って別のモデルに変わってしまう。決まるまでの間は、
    すでに誰かが選んでいるモデルも残す（画面で候補を決めれば、そちらが優先）。
    """
    if refresh:
        with _models_lock:
            _models_cache["at"] = 0.0
    admin = [str(m).strip() for m in (_read_admin().get("models") or []) if str(m).strip()]
    names = admin or list(config.OPENAI_MODELS) or sorted(users_by_model())
    # 既定のモデルは必ず候補に入れる（一覧に出てこないAPIもあるため）
    d = default_model()
    if d and d not in names:
        names.insert(0, d)
    return names


def users_by_model() -> dict:
    """いま誰がどのモデルを選んでいるか。{モデル名: [ユーザー名, ...]}

    候補から外すとその人は既定に戻る。外す前に影響が見えるようにするため、
    利用者のフォルダを読んで集める（読むだけで、何も書き換えない）。
    """
    out: dict = {}
    root = config.USER_META_DIR
    if not root.exists():
        return out
    for d in sorted(root.iterdir()):
        p = d / "prefs.yaml"
        if not p.is_file():
            continue
        try:
            data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}
        except Exception:
            continue
        m = str((data or {}).get("model") or "").strip()
        if m:
            out.setdefault(m, []).append(d.name)
    return out


def model_catalog(refresh: bool = False) -> list[str]:
    """管理者が候補を選ぶときに見せる「選べる全部」。

    available() は絞り込んだ後の一覧なので、管理画面ではこちらを使う。
    """
    if refresh:
        with _models_lock:
            _models_cache["at"] = 0.0
    return sorted(set(list(config.OPENAI_MODELS) + _from_api()))


def current(user=None) -> str:
    """そのユーザーがいま使うモデル。

    選んでいても、管理者が候補から外していれば既定に戻す。
    外したモデルを使い続けられると、絞り込んだ意味がなくなるため。
    """
    if user:
        chosen = prefs.get_model(user)
        if chosen and chosen in available():
            return chosen
    return default_model()


def choose(user, model: str) -> str:
    """モデルを選ぶ。管理者が決めた候補の中からだけ。"""
    model = str(model or "").strip()
    if not model:
        raise ValueError("モデル名が空です。")
    allowed = available()
    if model not in allowed:
        raise ValueError(f"{model} は選べません。"
                         f"選べるのは {'、'.join(allowed) or '（候補なし）'} です。")
    prefs.set_model(user, model)
    who = getattr(user, "username", None) or user
    print(f"[models] {who} のモデルを {model} にしました")
    return model


def _scope_note(total: int, limit: int) -> str:
    """カタログがそのモデルに収まらないときの、画面向けの説明文。"""
    head = (f"データカタログ全体（約{total:,}字）が、このモデルで一度に読める量"
            f"（約{limit:,}字）を超えています。")
    if config.SCOPE_MODE == "all":
        return (head + "詳細（列名・コード値）が渡らない要約モードになります。"
                "文脈の大きいモデルを選ぶか、env の SCOPE_MODE を見直してください。")
    return (head + "質問ごとに関係するDBだけへ自動で絞って、詳細を保ちます"
            "（答えに必要なDBが絞られるだけで、使えるデータは変わりません）。")


def status(user: str | None = None, refresh: bool = False) -> dict:
    cur = current(user)
    names = available(refresh)
    total = catalog_total_chars()
    limit = inline_limit_for(cur)
    fits = total <= limit
    return {
        "current": cur,
        "models": [{"id": m, "vision": is_vision(m),
                    "catalog_fits": total <= inline_limit_for(m)} for m in names],
        "vision": is_vision(cur),
        "from_env": bool(config.OPENAI_MODELS),
        "image_max_mb": config.IMAGE_MAX_MB,
        "image_max_count": config.IMAGE_MAX_COUNT,
        # 選択中モデルにカタログ全体が収まるか（チャット画面の警告表示に使う）
        "scope": {"mode": config.SCOPE_MODE, "catalog_chars": total,
                  "limit_chars": limit, "fits": fits,
                  "note": "" if fits else _scope_note(total, limit)},
    }


# --- 管理画面向け -----------------------------------------------------------

def admin_status(refresh: bool = False, scope: list[dict] | None = None) -> dict:
    """「モデル設定」画面に渡す内容。

    scope を渡すと、そのデータ範囲で「文脈をどれだけ使うか」も一緒に返す
    （上限を決めるのに、いまの実測値が要るため）。
    """
    ov = _read_admin()
    chosen = [str(m).strip() for m in (ov.get("models") or []) if str(m).strip()]
    names = chosen or list(config.OPENAI_MODELS)
    out = {
        "models": names,
        "default": default_model(),
        "vision": _vision_keys(),
        "catalog": model_catalog(refresh),
        "source": source(),
        "effective": available(),
        "in_use": users_by_model(),
        "from_env": not chosen,
        "env_models": list(config.OPENAI_MODELS),
        "env_default": config.OPENAI_MODEL,
        "settings_file": str(config.MODEL_SETTINGS_FILE),
        "llm_ready": _llm_ready(),
        "context_overrides": context_overrides(),
        "env_context_default": config.MODEL_CONTEXT_DEFAULT,
        "catalog_chars": catalog_total_chars(),
        # 候補ごとの文脈量とカタログ上限。出所も返す（登録 / 公式の表 / 推定）
        "contexts": [_context_row(m) for m in names],
    }
    if scope is not None:
        import llm
        base = llm.budget(scope, model=default_model(), admin=True)
        out["budget"] = base
        # 候補それぞれで、いまのカタログがどれだけ文脈を食うか
        out["per_model"] = []
        for m in (names or [default_model()]):
            ctx, known = context_window(m)
            out["per_model"].append({
                "id": m, "context": ctx, "context_known": known,
                "now_pct": round(base["now_tokens"] / ctx * 100, 1) if ctx else 0.0,
                "at_limit_pct": round(base["at_limit_tokens"] / ctx * 100, 1) if ctx else 0.0,
            })
    return out


def _context_row(model: str) -> dict:
    """モデル設定画面の1行ぶん。文脈量がどこから来た値かも添える。"""
    low = str(model or "").lower()
    ov = context_overrides()
    if low in ov or any(k in low for k in ov):
        source = "override"
    elif any(k in low for k in config.MODEL_CONTEXT_WINDOWS):
        source = "table"
    else:
        source = "default"
    ctx, _ = context_window(model)
    limit = inline_limit_for(model)
    return {"id": model, "context": ctx, "source": source, "limit_chars": limit,
            "fits": catalog_total_chars() <= limit}


def _llm_ready() -> bool:
    import llm
    return llm.is_configured()


def save_admin(data: dict, user: str | None = None) -> dict:
    """「モデル設定」画面からの保存。"""
    models = [str(m).strip() for m in (data.get("models") or []) if str(m).strip()]
    if not models:
        raise ValueError("選択できるモデルを1つ以上残してください。")
    if len(models) != len(set(models)):
        raise ValueError("同じモデルが重複しています。")
    for m in models:
        if len(m) > 120:
            raise ValueError(f"モデル名が長すぎます: {m[:40]}…")

    default = str(data.get("default") or "").strip() or models[0]
    if default not in models:
        raise ValueError(f"既定のモデル {default} が候補に入っていません。")

    vision = [str(v).strip().lower() for v in (data.get("vision") or []) if str(v).strip()]

    # 文脈量の登録 {モデル名: トークン}。表に無いモデルの実力を教える口
    overrides = {}
    for k, v in (data.get("context_overrides") or {}).items():
        name = str(k).strip().lower()
        if not name:
            continue
        try:
            n = int(str(v).replace(",", "").replace("_", ""))
        except (TypeError, ValueError):
            raise ValueError(f"「{k}」の文脈量は数字（トークン数）で指定してください。") from None
        if n < 1_000 or n > 10_000_000:
            raise ValueError(f"「{k}」の文脈量 {n:,} は範囲外です（1,000〜10,000,000）。")
        overrides[name] = n

    _write_admin({"models": models, "default": default, "vision": vision,
                  "context_overrides": overrides})
    print(f"[models] モデル設定を更新しました（{user or '不明'}）: "
          f"候補{len(models)}件 / 既定={default} / 画像判定={len(vision)}件 / "
          f"文脈量の登録={len(overrides)}件")
    return admin_status()


# ==========================================================================
# ===== 元 catalog.py
# データカタログ層。
#
# 「自動プロファイル（機械の知識）」と「サイドカーYAML（人間の知識）」を統合し、
# UI表示・ER図・LLM用 system prompt を **同じ情報源** から生成する。
#
# ファイル配置:
#   data/sales.db                                 … DB本体（読み取り専用で扱う）
#   data/sales.db.meta.yaml                       … メタ情報（全員で1つ。編集は管理者のみ）
#   data/.profile_cache/sales.db.profile.json     … 自動プロファイル（mtime+sizeで自動再生成）
#
# メタ情報(YAML)の構造:
#   title: 受注管理DB
#   description: |            # 何のデータか＋AIが知らないと間違える前提（※で始める行）
#     受注と請求。金額は明細側にしかない。
#     ※ 退職者も employees に残る。現役だけなら active_flag = 1 で絞る。
#   tables:
#     orders:
#       description: 受注明細。1行 = 1受注明細行。
#       ai_draft: true          # AI下書きのまま人間が未確認ならtrue
#       columns:
#         status: { description: 受注状態, values: { "1": 受付, "2": 出荷済 } }
#       glossary:               # そのテーブル固有の業務用語
#         有効な受注:
#           description: キャンセル以外の、実際に売上になる受注   # 自然言語だけでもよい
#           sql: status != '9'                                  # あればAIはこの式をそのまま使う
#   relationships:
#     - { from: orders.customer_id, to: customers.id, cardinality: "N:1" }
#       # to には "他DBエイリアス.テーブル.列" の3要素形式も書ける
#   glossary:                   # テーブルをまたぐ業務用語だけをここに書く
#     稼働率: { description: 実働時間÷所定時間 }
#   examples:
#     - q: 今月の売上は？
#       description: 締め日は月末。キャンセルは除く   # 任意。この例の読み方をAIに伝える
#       sql: SELECT ...
# ==========================================================================
import json
import sqlite3
import time
from datetime import datetime
from pathlib import Path

import yaml

import config
import db

# =============================================================================
# メタ情報（サイドカーYAML）
# =============================================================================

_META_KEYS = ("title", "description", "tables", "relationships", "glossary",
              "examples", "checks", "er_layout", "er_external", "tools", "builtin_tools")


# カタログは全員で1つ。DBの中身が何かは人によって変わらないので、
# 定義を分けると「同じ質問なのに人によって答えが違う」ことになる。
#   data/<DB>.db.meta.yaml … 唯一のカタログ。書き換えるのは管理者だけ
#                            （画面側は web/catalog_bp.py が admin_required で守る）

def meta_path(db_path) -> Path:
    """カタログの置き場所。DBファイルの隣に同じ名前で置く。"""
    return Path(str(db_path) + ".meta.yaml")


def _read_yaml(p: Path) -> dict:
    if not p.exists():
        return {}
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else {}
    except Exception as e:
        print(f"[catalog] メタ情報を読めませんでした: {p} ({e})")
        return {}


def load_meta(db_path) -> dict:
    """カタログを読む（全員が同じものを見る）。"""
    return _read_yaml(meta_path(db_path))


def merge_caveats(description, caveats) -> str:
    """説明と、かつて別欄だった注意書き(caveats)を1つの文章にする。

    以前は「説明」と「注意書き（1行に1つ）」の2欄だったが、AIへの渡り方は
    同じ場所に続けて書かれた文章で、分ける意味が薄かった。いまは「説明」1欄で、
    注意したい事実は行頭に ※ を付けて書く。古いYAMLの caveats はここで合流させる。
    """
    lines = [str(description or "").strip()]
    for c in (caveats or []):
        c = str(c or "").strip()
        if c:
            lines.append(c if c.startswith(("※", "⚠")) else f"※ {c}")
    return "\n".join(l for l in lines if l)


def db_description(meta: dict) -> str:
    """DBの説明（古い caveats があれば ※ 行として末尾に合流させた1本の文章）。"""
    return merge_caveats(meta.get("description"), meta.get("caveats"))


def save_meta(db_path, meta: dict) -> None:
    """カタログを保存する（内容をまるごと書く）。呼べるのは管理者の画面だけ。"""
    target = meta_path(db_path)
    # 説明は1欄。古い caveats が残っていれば説明に合流させてから書く
    if meta.get("caveats"):
        meta["description"] = merge_caveats(meta.get("description"), meta.get("caveats"))
        meta.pop("caveats", None)
    cleaned = {}
    for k in _META_KEYS:
        v = meta.get(k)
        if v in (None, "", [], {}):
            continue
        cleaned[k] = v

    target.parent.mkdir(parents=True, exist_ok=True)
    if not cleaned:
        target.write_text("", encoding="utf-8")
        return
    target.write_text(
        yaml.dump(cleaned, Dumper=_MetaDumper, allow_unicode=True, sort_keys=False,
                  default_flow_style=False),
        encoding="utf-8",
    )


class _MetaDumper(yaml.SafeDumper):
    """複数行の文字列（説明・SQL）は '...' の折り返しではなく | ブロックで書く。
    手で開いて読める・直せるファイルにしておくため。"""


def _repr_str(dumper, data: str):
    if "\n" in data:
        # 行末の空白があると PyYAML は | を使えず引用符に落ちるので、先に落とす
        data = "\n".join(l.rstrip() for l in data.splitlines())
        return dumper.represent_scalar("tag:yaml.org,2002:str", data, style="|")
    return dumper.represent_scalar("tag:yaml.org,2002:str", data)


_MetaDumper.add_representer(str, _repr_str)


# =============================================================================
# 業務用語（用語集）
# =============================================================================
#
# 用語は「テーブル固有」と「テーブルをまたぐもの」の2種類あるので、置き場所も2つ。
#   meta["tables"][テーブル名]["glossary"]  … そのテーブルの用語（基本はこちら）
#   meta["glossary"]                        … 複数テーブルにまたがる用語
# テーブル側に置くと、そのテーブルが選択されているときだけプロンプトに載る。
#
# 1つの用語は次の2つを持つ。どちらか一方だけでもよい。
#   description … 自然言語の説明（AIはこれを読んで自分でSQLを組み立てる）
#   sql         … SQLの条件式や計算式（あればAIはこの式をそのまま使う）

def normalize_glossary(gl) -> dict:
    """用語集を {用語: {"description":…, "sql":…}} の形に揃える。

    値は必ず辞書。手でYAMLを書いて文字列になっていた場合は「説明」として扱う。
    以前はSQL式として扱っていたが、説明文が書かれていると
    「この式をそのまま使う」とAIに渡してしまい、構文エラーのSQLを作らせていた。
    説明として扱えば、間違っていてもAIが列情報から組み立て直せる。
    """
    out = {}
    for term, val in (gl or {}).items():
        term = str(term).strip()
        if not term:
            continue
        if isinstance(val, dict):
            desc = str(val.get("description") or "").strip()
            sql = str(val.get("sql") or "").strip()
        else:
            print(f"[catalog] 用語 '{term}' が古い書き方です。説明として扱います。")
            desc, sql = str(val or "").strip(), ""
        if desc or sql:
            out[term] = {"description": desc, "sql": sql}
    return out


def table_glossary(meta: dict, tname: str) -> dict:
    """テーブル固有の用語。"""
    return normalize_glossary(((meta.get("tables") or {}).get(tname) or {}).get("glossary"))


def db_glossary(meta: dict) -> dict:
    """テーブルをまたぐ用語。"""
    return normalize_glossary(meta.get("glossary"))


def set_table_glossary(meta: dict, tname: str, gl: dict) -> None:
    """テーブル固有の用語を書き戻す（空なら削除）。"""
    tm = meta.setdefault("tables", {}).setdefault(tname, {})
    if gl:
        tm["glossary"] = gl
    else:
        tm.pop("glossary", None)
        if not tm:
            meta["tables"].pop(tname, None)


def glossary_lines(gl: dict) -> list[str]:
    """プロンプトに載せる用語の行。"""
    lines = []
    for term, e in gl.items():
        desc, sql = e.get("description") or "", e.get("sql") or ""
        lines.append(f"- {term}: {desc}" if desc else f"- {term}:")
        if sql:
            lines.append(f"    SQL式: {sql}   ← この式をそのまま使う")
        else:
            lines.append("    （SQL式は未登録。上の列情報をもとに自分で組み立てる）")
    return lines


def glossary_count(meta: dict) -> int:
    """DB全体＋全テーブルの用語数。"""
    n = len(db_glossary(meta))
    for tname in (meta.get("tables") or {}):
        n += len(table_glossary(meta, tname))
    return n


# =============================================================================
# 自動プロファイル
# =============================================================================

def _qi(name: str) -> str:
    """SQLite識別子のクオート。"""
    return '"' + str(name).replace('"', '""') + '"'


def _cache_path(db_path) -> Path:
    return config.PROFILE_CACHE_DIR / (Path(db_path).name + ".profile.json")


def _make_timeout(conn: sqlite3.Connection, seconds: float):
    """接続にタイムアウトを仕掛け、クエリごとに呼ぶ reset 関数を返す。"""
    box = {"t": time.time()}
    conn.set_progress_handler(lambda: 1 if (time.time() - box["t"]) > seconds else 0, 100000)

    def reset():
        box["t"] = time.time()
    return reset


def _profile_table(conn: sqlite3.Connection, name: str, reset) -> dict:
    t = _qi(name)
    info: dict = {"columns": [], "fks": [], "row_count": None,
                  "sample_columns": [], "sample_rows": [], "col_stats": {}}

    reset()
    # PRAGMA table_info の pk は 0=非キー / 1以上=複合主キー内の順番。
    # 複合キーの構成順は「1行が何を表すか」の手がかりになるので pk_seq に残す。
    for cid, cname, ctype, notnull, dflt, pk in conn.execute(f"PRAGMA table_info({t})"):
        info["columns"].append({"name": cname, "type": ctype or "", "notnull": bool(notnull),
                                "pk": bool(pk), "pk_seq": int(pk or 0)})

    reset()
    try:
        for row in conn.execute(f"PRAGMA foreign_key_list({t})"):
            # (id, seq, table, from, to, on_update, on_delete, match)
            info["fks"].append({"from": row[3], "table": row[2], "to": row[4] or "id"})
    except sqlite3.Error:
        pass

    reset()
    try:
        info["row_count"] = conn.execute(f"SELECT COUNT(*) FROM {t}").fetchone()[0]
    except sqlite3.Error:
        pass  # タイムアウト等 → 行数不明として続行

    reset()
    try:
        cur = conn.execute(f"SELECT * FROM {t} LIMIT {config.PROFILE_SAMPLE_ROWS}")
        info["sample_columns"] = [d[0] for d in cur.description] if cur.description else []
        info["sample_rows"] = [[_jsonable(v) for v in r] for r in cur.fetchall()]
    except sqlite3.Error:
        pass

    # 列統計（巨大テーブルはスキップ）
    rc = info["row_count"]
    if rc is not None and rc <= config.PROFILE_STATS_MAX_ROWS and rc > 0:
        limit = config.PROFILE_LOW_CARDINALITY
        for col in info["columns"]:
            c = _qi(col["name"])
            stat: dict = {}
            reset()
            try:
                vals = conn.execute(
                    f"SELECT {c} AS v, COUNT(*) AS n FROM {t} GROUP BY 1 ORDER BY n DESC LIMIT {limit + 1}"
                ).fetchall()
                if len(vals) <= limit:
                    stat["values"] = [[_jsonable(v), n] for v, n in vals]
                else:
                    reset()
                    mn, mx = conn.execute(f"SELECT MIN({c}), MAX({c}) FROM {t}").fetchone()
                    stat["min"], stat["max"] = _jsonable(mn), _jsonable(mx)
            except sqlite3.Error:
                pass
            if stat:
                info["col_stats"][col["name"]] = stat
    return info


def _jsonable(v):
    if isinstance(v, bytes):
        return f"<BLOB {len(v)} bytes>"
    return v


def profile_db(db_path, force: bool = False) -> dict:
    """DBを読み取り専用でプロファイリング。mtime+sizeが一致するキャッシュがあれば再利用。"""
    db_path = Path(db_path)
    st = db_path.stat()
    # v はプロファイルの構造バージョン。上げると古いキャッシュが無効になる。
    key = {"v": 2, "mtime": st.st_mtime, "size": st.st_size}

    cache = _cache_path(db_path)
    if not force and cache.exists():
        try:
            data = json.loads(cache.read_text(encoding="utf-8"))
            if data.get("key") == key:
                return data
        except Exception:
            pass

    conn = db.connect_ro(db_path)
    try:
        reset = _make_timeout(conn, config.PROFILE_TIMEOUT_SEC)
        tables: dict = {}
        reset()
        rows = conn.execute(
            "SELECT name, type FROM sqlite_master "
            "WHERE type IN ('table','view') AND name NOT LIKE 'sqlite_%' ORDER BY name"
        ).fetchall()
        for name, typ in rows:
            try:
                t = _profile_table(conn, name, reset)
                t["type"] = typ
                tables[name] = t
            except sqlite3.Error as e:
                tables[name] = {"type": typ, "error": str(e), "columns": [], "fks": [],
                                "row_count": None, "sample_columns": [], "sample_rows": [],
                                "col_stats": {}}
    finally:
        conn.close()

    profile = {
        "file": db_path.name,
        "key": key,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "tables": tables,
    }
    config.PROFILE_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    cache.write_text(json.dumps(profile, ensure_ascii=False, default=str), encoding="utf-8")
    return profile


# =============================================================================
# 乖離検知・結合候補・カバレッジ
# =============================================================================

def drift_warnings(profile: dict, meta: dict) -> list[str]:
    """メタ情報がスキーマの実体からズレている箇所を警告として返す。"""
    warns = []
    ptables = profile.get("tables", {})
    for tname, tmeta in (meta.get("tables") or {}).items():
        if tname not in ptables:
            warns.append(f"メタ情報のテーブル '{tname}' はDBに存在しません（改名/削除された可能性）。")
            continue
        pcols = {c["name"] for c in ptables[tname]["columns"]}
        for cname in ((tmeta or {}).get("columns") or {}):
            if cname not in pcols:
                warns.append(f"メタ情報の列 '{tname}.{cname}' はDBに存在しません。")
        for cname in ((tmeta or {}).get("primary_key") or []):
            if cname not in pcols:
                warns.append(f"指定された主キーの列 '{tname}.{cname}' はDBに存在しません。")
    for rel in (meta.get("relationships") or []):
        for end in (rel.get("from", ""), rel.get("to", "")):
            parts = str(end).split(".")
            if len(parts) == 2:  # table.col（同一DB内）のみ検証。db付き3要素は他DBなので対象外
                tname, cname = parts
                if tname in ptables:
                    if cname not in {c["name"] for c in ptables[tname]["columns"]}:
                        warns.append(f"結合定義の '{end}' に対応する列がありません。")
                else:
                    warns.append(f"結合定義の '{end}' に対応するテーブルがありません。")
    return warns


def join_suggestions(profile: dict, meta: dict) -> list[dict]:
    """列名ヒューリスティックによる結合候補（FK宣言済み・登録済みは除く）。"""
    ptables = profile.get("tables", {})
    existing = set()
    for rel in (meta.get("relationships") or []):
        existing.add((str(rel.get("from", "")).lower(), str(rel.get("to", "")).lower()))
    for tname, t in ptables.items():
        for fk in t.get("fks", []):
            existing.add((f"{tname}.{fk['from']}".lower(), f"{fk['table']}.{fk['to']}".lower()))

    sugs = []
    for tname, t in ptables.items():
        for col in t.get("columns", []):
            cname = col["name"]
            low = cname.lower()
            if not low.endswith("_id") and not low.endswith("id"):
                continue
            base = low[:-3] if low.endswith("_id") else None
            if not base:
                continue
            # 候補テーブル名: base / base+"s" / base+"es"
            for cand in (base, base + "s", base + "es"):
                target = next((n for n in ptables if n.lower() == cand), None)
                if not target or target == tname:
                    continue
                tcols = ptables[target]["columns"]
                pk = next((c["name"] for c in tcols if c["pk"]), None)
                # 複合主キーの相手に1列だけで結合する候補は誤りになるので出さない
                if len([c for c in tcols if c["pk"]]) > 1:
                    continue
                to_col = pk or next((c["name"] for c in tcols if c["name"].lower() in ("id", low)), None)
                if not to_col:
                    continue
                frm, to = f"{tname}.{cname}", f"{target}.{to_col}"
                if (frm.lower(), to.lower()) in existing:
                    continue
                sugs.append({"from": frm, "to": to, "cardinality": "N:1",
                             "reason": f"列名 '{cname}' → テーブル '{target}' の推測"})
                break
    return sugs


def coverage(profile: dict, meta: dict) -> dict:
    """メタ情報の充実度。カタログページの案内表示に使う。"""
    ptables = profile.get("tables", {})
    mtables = meta.get("tables") or {}
    n_tables = len(ptables)
    n_tdesc = sum(1 for t in ptables if (mtables.get(t) or {}).get("description"))
    n_cols = sum(len(t["columns"]) for t in ptables.values())
    n_cdesc = 0
    for tname, t in ptables.items():
        mcols = (mtables.get(tname) or {}).get("columns") or {}
        for c in t["columns"]:
            cm = mcols.get(c["name"]) or {}
            if cm.get("description") or cm.get("values"):
                n_cdesc += 1
    return {
        "tables": (n_tdesc, n_tables),
        "columns": (n_cdesc, n_cols),
        "relationships": len(meta.get("relationships") or []),
        "glossary": glossary_count(meta),
        "examples": len(meta.get("examples") or []),
    }


# =============================================================================
# 結合の端点表記（"table.col" / "alias.table.col" の相互変換）
# =============================================================================

def parse_endpoint(end: str, default_alias: str):
    """'table.col' または 'alias.table.col' を (alias, table, column) に解く。"""
    parts = [p.strip() for p in str(end).split(".")]
    if len(parts) == 3:
        return parts[0], parts[1], parts[2]
    if len(parts) == 2:
        return default_alias, parts[0], parts[1]
    return None


def node_id(alias: str, table: str) -> str:
    """テーブル（親ノード）のID。"""
    return f"{alias}.{table}"


# 列ノード（親テーブルの中に並ぶ子ノード）のID。
# テーブルIDが "alias.table" なので、列との区切りには "::" を使う。
COL_SEP = "::"


def col_node_id(alias: str, table: str, column: str) -> str:
    return f"{alias}.{table}{COL_SEP}{column}"


def edge_label(cardinality: str | None) -> str:
    """IPA表記の関連ラベル。線は列ノード同士を結ぶので、列名はラベルに出さず
    多重度だけを示す（始点側 ─ 終点側）。例: "* ─ 1"
    """
    tail, head = _CARD_ENDS.get(cardinality or "N:1", ("*", "1"))
    return f"{tail} ─ {head}"


def collect_edges(entries: list[dict]) -> list[dict]:
    """キャンバス/ER図に描く結合を集める。

    entries: [{"alias": str, "profile": dict, "meta": dict}, ...]
    戻り値の各要素:
      {"id", "source", "target", "label", "kind": "fk"|"meta", "owner", "index"}
      kind="fk"   … DBに宣言されたFOREIGN KEY（削除不可）
      kind="meta" … .meta.yaml の relationships（編集・削除可。index は配列位置）
    """
    nodes = {node_id(e["alias"], t) for e in entries for t in e["profile"].get("tables", {})}

    # メタ側の端点集合（FKと重複したら FK 側を出さない）
    meta_pairs = set()
    for e in entries:
        for rel in (e["meta"].get("relationships") or []):
            a = parse_endpoint(rel.get("from", ""), e["alias"])
            b = parse_endpoint(rel.get("to", ""), e["alias"])
            if a and b:
                meta_pairs.add((a, b))

    def valid(p):
        """端点(alias, table, column)が実在し、キャンバス上にあるか。"""
        if node_id(p[0], p[1]) not in nodes:
            return False
        e = next((x for x in entries if x["alias"] == p[0]), None)
        cols = {c["name"] for c in (e["profile"]["tables"].get(p[1]) or {}).get("columns", [])}
        return p[2] in cols

    edges: list[dict] = []
    for e in entries:
        alias = e["alias"]
        for tname, t in e["profile"].get("tables", {}).items():
            for fk in t.get("fks", []):
                a = (alias, tname, fk["from"])
                b = (alias, fk["table"], fk["to"])
                if not valid(a) or not valid(b) or (a, b) in meta_pairs:
                    continue
                edges.append({
                    "id": f"fk||{a[0]}.{a[1]}.{a[2]}||{b[0]}.{b[1]}.{b[2]}",
                    "source": col_node_id(*a), "target": col_node_id(*b),
                    "from": a, "to": b,
                    "label": edge_label("N:1"), "cardinality": "N:1",
                    "kind": "fk", "owner": alias, "index": None,
                })
        for i, rel in enumerate(e["meta"].get("relationships") or []):
            a = parse_endpoint(rel.get("from", ""), alias)
            b = parse_endpoint(rel.get("to", ""), alias)
            if not a or not b or not valid(a) or not valid(b):
                continue
            card = rel.get("cardinality") or "N:1"
            edges.append({
                "id": f"rel||{alias}||{i}",
                "source": col_node_id(*a), "target": col_node_id(*b),
                "from": a, "to": b,
                "label": edge_label(card), "cardinality": card,
                "kind": "meta", "owner": alias, "index": i,
            })
    return edges


def declared_pk(profile: dict, tname: str) -> list[str]:
    """DBが宣言している主キー（複合キーは構成順）。宣言が無ければ空リスト。"""
    t = profile.get("tables", {}).get(tname) or {}
    cols = [c for c in t.get("columns", []) if c.get("pk")]
    cols.sort(key=lambda c: c.get("pk_seq") or 0)
    return [c["name"] for c in cols]


def effective_pk(profile: dict, meta: dict, tname: str):
    """実際に主キーとして扱う列と、その出所を返す。

    戻り値: (列名リスト, "override" | "declared" | "none")
    メタの tables.<name>.primary_key があれば、DB宣言より優先する。
    主キーが宣言されていないテーブル（CSV取込など）に人が指定できるようにするため。
    """
    valid = [c["name"] for c in (profile.get("tables", {}).get(tname) or {}).get("columns", [])]
    ov = ((meta.get("tables") or {}).get(tname) or {}).get("primary_key")
    if ov:
        cols = [c for c in ov if c in valid]
        if cols:
            return cols, "override"
    d = declared_pk(profile, tname)
    return (d, "declared") if d else ([], "none")


def fk_columns(entries: list[dict], alias: str, tname: str) -> set:
    """外部キーとして扱う列（FK宣言 + メタの relationships の from 側）。"""
    out = set()
    for e in entries:
        if e["alias"] == alias:
            t = e["profile"].get("tables", {}).get(tname) or {}
            for fk in t.get("fks", []):
                out.add(fk["from"])
        for rel in (e["meta"].get("relationships") or []):
            p = parse_endpoint(rel.get("from", ""), e["alias"])
            if p and p[0] == alias and p[1] == tname:
                out.add(p[2])
    return out


# --- IPA表記のノードラベル -------------------------------------------------------
#: 1DBあたりに持つ例文の上限。例文は毎回 system prompt に載るので、
#: 増えるほどテーブル定義の説明が押し出される。多くても効果は上がらない。
EXAMPLES_MAX = 20


def _norm_sql(sql: str) -> str:
    """比べるためだけの正規化。空白の入れ方と大小の違いを無視する。"""
    return " ".join(str(sql or "").split()).lower()


def dedupe_examples(examples: list[dict]) -> list[dict]:
    """例文から重複を落とす。

    同じSQLが複数あると、毎回のプロンプトが太るうえ、AIがその型を
    過剰に当てはめるようになる。SQLが同じものは最初の1件だけ残す。
    質問文が同じものも、後から入れた方（確認し直した方）を残す。
    説明は任意なので、書かれているものだけを残す。
    """
    by_sql: dict = {}
    for ex in examples or []:
        q = str(ex.get("q") or "").strip()
        sql = str(ex.get("sql") or "").strip()
        if not q or not sql:
            continue
        key = _norm_sql(sql)
        if key in by_sql:
            continue                      # 同じSQLは1件でよい
        desc = str(ex.get("description") or "").strip()
        by_sql[key] = {"q": q, **({"description": desc} if desc else {}), "sql": sql}

    # 質問文の重複は後勝ち（同じ問いに対する新しいSQLを正とする）
    by_q: dict = {}
    for ex in by_sql.values():
        by_q[ex["q"]] = ex
    return list(by_q.values())[:EXAMPLES_MAX]


def find_example(examples: list[dict], sql: str) -> dict | None:
    """同じSQLの例文が既にあれば返す。"""
    key = _norm_sql(sql)
    for ex in examples or []:
        if _norm_sql(ex.get("sql")) == key:
            return ex
    return None


def load_layout(meta: dict) -> dict:
    """メタからノード座標を読む。{'alias.table': (x, y)}"""
    raw = meta.get("er_layout") or {}
    if not isinstance(raw, dict):
        return {}
    out = {}
    for k, v in raw.items():
        # 期待する形は [x, y]。辞書や文字列など形の違う項目は読み飛ばす
        # （1つ壊れているだけでER図全体が出なくなるのを避ける）
        if not isinstance(v, (list, tuple)) or len(v) < 2:
            continue
        try:
            out[str(k)] = (float(v[0]), float(v[1]))
        except (TypeError, ValueError):
            continue
    return out


def _dbmod():
    """db モジュール。循環importを避けるため使うときに読む。"""
    import db
    return db


def cross_db_tables(alias: str, meta: dict, extra: list | None = None) -> dict:
    """このキャンバスに引き込む他DBのテーブル。{alias: {table: 由来}}

    由来は3つ。
      "out" … 自DBの関連が指している先（demo_sales → demo_master.customers）
      "in"  … 他DBの関連が、こちらを指しているもの（相手が持っている関連）
      "pin" … 画面から明示的に追加されたもの（関連ゼロのDBともつなげるように）
    必要なぶんだけ借りる。相手DBの全テーブルは出さない。

    由来を分けるのは、"in" が多くなりがちだから。マスタ系のDBは
    あらゆるDBから参照されるので、全部出すと自分のテーブルが埋もれる。
    既定では隠し、画面のボタンで出す。
    """
    want: dict = {}

    def add(ep, why):
        if not ep or ep[0] == alias:
            return
        cur = want.setdefault(ep[0], {})
        # 自分が張った関連・明示追加のほうが強い（"in" で上書きしない）
        if cur.get(ep[1]) != "in" and ep[1] in cur:
            return
        cur[ep[1]] = why

    for rel in (meta.get("relationships") or []):
        add(parse_endpoint(rel.get("from", ""), alias), "out")
        add(parse_endpoint(rel.get("to", ""), alias), "out")

    for f in _dbmod().list_db_files():
        other = _dbmod().alias_for(f)
        if other == alias:
            continue
        for rel in (load_meta(f).get("relationships") or []):
            eps = [parse_endpoint(rel.get("from", ""), other),
                   parse_endpoint(rel.get("to", ""), other)]
            if any(ep and ep[0] == alias for ep in eps):
                for ep in eps:
                    add(ep, "in")

    for item in (extra or []):
        parts = str(item).split(".")
        if len(parts) == 2 and parts[0] != alias:
            want.setdefault(parts[0], {})[parts[1]] = "pin"
    return want


def er_payload(path, profile: dict | None = None,
               meta: dict | None = None) -> dict:
    """自前キャンバス用に、テーブル・列・関連を素のJSONで渡す。

    DBまたぎの関連も描けるよう、関連が指している他DBのテーブルを
    「借りたノード」として一緒に返す（external=true）。
    """
    profile = profile if profile is not None else profile_db(path)
    meta = meta if meta is not None else load_meta(path)
    alias = _dbmod().alias_for(path)
    entries = [{"alias": alias, "profile": profile, "meta": meta,
                "path": path, "editable": True}]

    # 他DBを引き込む。相手のメタも読むのは、向こうが持っている関連
    # （こちらを指しているもの）を線として描くため
    borrowed = cross_db_tables(alias, meta, meta.get("er_external") or [])
    others: dict = {}
    for f in _dbmod().list_db_files():
        a = _dbmod().alias_for(f)
        if a in borrowed:
            others[a] = {"profile": profile_db(f), "meta": load_meta(f)}
            entries.append({"alias": a, "profile": others[a]["profile"],
                            "meta": others[a]["meta"], "path": f, "editable": False})

    layout = load_layout(meta)
    nodes = []
    for i, (tname, t) in enumerate(profile["tables"].items()):
        fks = fk_columns(entries, alias, tname)
        pk = set(effective_pk(profile, meta, tname)[0])
        nid = node_id(alias, tname)
        pos = layout.get(nid) or [40 + (i % 4) * 300, 40 + (i // 4) * 320]
        nodes.append({
            "id": nid, "alias": alias, "table": tname, "external": False,
            "x": pos[0], "y": pos[1], "rows": t.get("row_count"),
            "columns": [{"name": c["name"], "type": c["type"],
                         "pk": c["name"] in pk, "fk": c["name"] in fks}
                        for c in t["columns"]],
        })

    # 借りたノード。位置はこのDBの er_layout に覚える（DBごとに自分の絵を持てる）
    k, base_row = 0, len(nodes) // 4 + 1
    for a, info in others.items():
        for tname, why in sorted(borrowed[a].items()):
            t = (info["profile"].get("tables") or {}).get(tname)
            if t is None:
                continue
            pk = set(effective_pk(info["profile"], info["meta"], tname)[0])
            nid = node_id(a, tname)
            pos = layout.get(nid) or [40 + (k % 4) * 300, 40 + (base_row + k // 4) * 320]
            k += 1
            nodes.append({
                "id": nid, "alias": a, "table": tname, "external": True,
                "incoming": why == "in", "pinned": why == "pin",
                "x": pos[0], "y": pos[1], "rows": t.get("row_count"),
                "columns": [{"name": c["name"], "type": c["type"],
                             "pk": c["name"] in pk, "fk": False}
                            for c in t["columns"]],
            })

    # 借りたDBは全テーブルぶんのプロフィールを持っているので、
    # 描いていないテーブルに向かう線が混じらないよう、置いたノードで絞る
    placed = {n["id"] for n in nodes}
    edges = []
    for e in collect_edges(entries):
        if (node_id(*e["from"][:2]) not in placed
                or node_id(*e["to"][:2]) not in placed):
            continue
        # 他DBのメタに書かれた関連は、ここからは直せない（持ち主が別）
        owner = e.get("owner")
        editable = e["kind"] == "meta" and owner == alias
        edges.append({"id": e["id"], "kind": e["kind"], "label": e["label"],
                      "cardinality": e["cardinality"], "index": e.get("index"),
                      "from": list(e["from"]), "to": list(e["to"]),
                      "owner": owner, "editable": editable,
                      "cross": e["from"][0] != e["to"][0]})
    return {"nodes": nodes, "edges": edges, "alias": alias,
            "extra": sorted(f"{a}.{t}" for a, ts in borrowed.items() for t in ts)}


# =============================================================================
# 関連の向きと多重度（ER図とデータ検査が共有する規則）
# =============================================================================

# IPA表記の多重度ラベル: 線の両端に "1" と "*" を置く
_CARD_ENDS = {
    "N:1": ("*", "1"),   # from(多側) ─ to(1側)
    "1:N": ("1", "*"),
    "1:1": ("1", "1"),
    "N:M": ("*", "*"),
}

#: 向きを入れ替えたときの多重度。1:1 と N:M は入れ替えても同じ。
_CARD_FLIP = {"N:1": "1:N", "1:N": "N:1", "1:1": "1:1", "N:M": "N:M"}


def _is_sole_pk(profile: dict, meta: dict, table: str, column: str) -> bool:
    """その列が、そのテーブルの主キー全体か（単独主キーか）。"""
    pk, _ = effective_pk(profile, meta, table)
    return len(pk) == 1 and pk[0] == column


def normalize_direction(a: tuple, b: tuple, cardinality: str, lookup) -> tuple:
    """関連の向きを「子（外部キー側）→ 親（主キー側）」に揃える。

    ER図はIPA表記なので矢印を描かない。見た目に向きが無いぶん、人は
    好きな方向にドラッグする。ところが from/to は単なる描画順ではなく、
    「どちらが参照している側か」を表しており、参照整合性の検査
    （親に居ない子を数える）はこの向きに依存する。逆向きに登録されると
    「入金の無い請求」を異常として数えるような、意味の反転が起きる。

    lookup(alias) は (profile, meta) を返す関数。判断できないときは触らない。

    戻り値: (from, to, cardinality)
    """
    card = cardinality or "N:1"
    try:
        pa, ma = lookup(a[0])
        pb, mb = lookup(b[0])
    except Exception:
        return a, b, card
    if not (pa and pb):
        return a, b, card
    a_is_pk = _is_sole_pk(pa, ma, a[1], a[2])
    b_is_pk = _is_sole_pk(pb, mb, b[1], b[2])
    # 片方だけが主キーなら、そちらを親（to）にする
    if a_is_pk and not b_is_pk:
        return b, a, _CARD_FLIP.get(card, card)
    return a, b, card


def _sample_values(profile: dict, table: str, column: str) -> list:
    """プロファイルに残っている実値（あれば）。警告文に例として添える用。"""
    st = ((profile or {}).get("tables", {}).get(table) or {}).get("col_stats", {}).get(column) or {}
    vals = st.get("values") or []
    return [v[0] if isinstance(v, (list, tuple)) else v for v in vals]


def link_check(child: tuple, parent: tuple, lookup, path_of) -> dict:
    """この2列を関連として結んでよいかを、実データを見て判定する。

    child / parent は (alias, table, column)。normalize_direction を通した後の向き。
    lookup(alias) は (profile, meta)、path_of(alias) は DBファイルのパスを返す。

    戻り値: {"level": "ok" | "warn" | "block", "issues": [{level, title, detail}]}
      block … 結んではいけない（値が全く重ならない等）。保存しない
      warn  … 結べるが、意味を確かめてほしい（型が違う・親が一意でない等）。確認して保存
      ok    … 問題なし

    なぜ止めるかを人が読める形で必ず添える。ER図の線は「この列で JOIN してよい」という
    AIへの指示なので、実データで JOIN が成立しない線を引くと、AIが自信を持って
    間違った結合を書くようになる。
    """
    issues: list[dict] = []
    ca, ct, cc = child
    pa, pt, pc = parent

    def add(level, title, detail):
        issues.append({"level": level, "title": title, "detail": detail})

    # --- カタログ上の情報（プロファイル）で分かること ---------------------------------
    prof_c, meta_c = lookup(ca)
    prof_p, meta_p = lookup(pa)
    col_c = next((c for c in (prof_c["tables"].get(ct) or {}).get("columns", [])
                  if c["name"] == cc), {}) if prof_c else {}
    col_p = next((c for c in (prof_p["tables"].get(pt) or {}).get("columns", [])
                  if c["name"] == pc), {}) if prof_p else {}
    type_c = str(col_c.get("type") or "").upper()
    type_p = str(col_p.get("type") or "").upper()

    def kind(t):
        if any(k in t for k in ("INT",)):
            return "整数"
        if any(k in t for k in ("REAL", "FLOA", "DOUB", "NUM", "DEC")):
            return "小数"
        if any(k in t for k in ("CHAR", "TEXT", "CLOB")):
            return "文字"
        if "DATE" in t or "TIME" in t:
            return "日時"
        return t or "不明"

    if type_c and type_p and kind(type_c) != kind(type_p):
        add("warn", "型が違います",
            f"{ct}.{cc} は {type_c}（{kind(type_c)}）、{pt}.{pc} は {type_p}（{kind(type_p)}）です。"
            "SQLite は型が違っても比較できてしまいますが、たいてい別の意味の列です"
            "（例: 数値のIDと文字のコード）。本当に同じものを指すか確かめてください。")

    pk_c = set(effective_pk(prof_c, meta_c, ct)[0]) if prof_c else set()
    pk_p = set(effective_pk(prof_p, meta_p, pt)[0]) if prof_p else set()
    if cc in pk_c and pk_c == {cc} and pc in pk_p and pk_p == {pc} and (ct != pt or ca != pa):
        add("warn", "主キー同士を結んでいます",
            f"{ct}.{cc} も {pt}.{pc} もそれぞれのテーブルの主キーです。"
            "1対1の関連（同じIDを持つ2つのテーブル）なら正しいですが、"
            "「たまたま両方IDという名前」なら結ぶべきではありません。")

    # --- 実データで分かること（読み取り専用で数える） ---------------------------------
    try:
        pc_path, pp_path = path_of(ca), path_of(pa)
        conn = db.connect_scope([(pc_path, "c"), (pp_path, "p")] if pc_path != pp_path
                                else [(pc_path, "c")])
        pal = "c" if pc_path == pp_path else "p"
        q = lambda s: '"' + str(s).replace('"', '""') + '"'
        C = f'"c".{q(ct)}', q(cc)
        P = f'"{pal}".{q(pt)}', q(pc)

        n_child = conn.execute(f"SELECT COUNT(*) FROM {C[0]} WHERE {C[1]} IS NOT NULL").fetchone()[0]
        n_parent = conn.execute(f"SELECT COUNT(*) FROM {P[0]} WHERE {P[1]} IS NOT NULL").fetchone()[0]
        n_parent_distinct = conn.execute(
            f"SELECT COUNT(DISTINCT {P[1]}) FROM {P[0]} WHERE {P[1]} IS NOT NULL").fetchone()[0]
        # 子の値のうち親に存在するもの / しないもの
        matched = conn.execute(
            f"SELECT COUNT(*) FROM {C[0]} c0 WHERE c0.{C[1]} IS NOT NULL "
            f"AND EXISTS (SELECT 1 FROM {P[0]} p0 WHERE p0.{P[1]} = c0.{C[1]})").fetchone()[0]
        # 子の「異なる値」の数と、親の値のうち子から参照されている数（親側のカバー率）。
        # 「status(1,2,3,9) → product_id(1〜40)」のような偶然の一致は、子の値は全部
        # 親に見つかるのに、親の値はほとんど参照されない。本物の外部キーなら親の多くが
        # 参照される。値の一致だけでは見抜けないので、この角度を足す。
        n_child_distinct = conn.execute(
            f"SELECT COUNT(DISTINCT {C[1]}) FROM {C[0]} WHERE {C[1]} IS NOT NULL").fetchone()[0]
        parent_hit = conn.execute(
            f"SELECT COUNT(DISTINCT p0.{P[1]}) FROM {P[0]} p0 "
            f"WHERE EXISTS (SELECT 1 FROM {C[0]} c0 WHERE c0.{C[1]} = p0.{P[1]})").fetchone()[0]
        conn.close()

        if n_child and n_parent and matched == 0:
            add("block", "値が1件も一致しません",
                f"{ct}.{cc} の {n_child:,} 件は、{pt}.{pc} の {n_parent:,} 件のどれとも一致しません。"
                "この2列で JOIN しても結果は必ず0行になります。別の意味の列です。")
        elif n_child and matched:
            miss = n_child - matched
            rate = miss / n_child * 100
            if rate >= 30:
                add("warn", "一致しない値が多すぎます",
                    f"{ct}.{cc} の {n_child:,} 件のうち {miss:,} 件（{rate:.0f}%）が {pt}.{pc} に存在しません。"
                    "外部キーなら親に無い値はごく少数のはずです。列の取り違えの可能性があります。")
            elif miss:
                add("info", "親に無い値があります",
                    f"{ct}.{cc} の {miss:,} 件（{rate:.1f}%）が {pt}.{pc} に存在しません"
                    "（未登録・削除済みの参照。数が少なければ通常の範囲です）。")
            # 子の値の種類が極端に少なく、親のごく一部にしか当たらない → 区分値とIDの偶然の一致
            if n_parent_distinct >= 10 and n_child_distinct <= 10                     and parent_hit / n_parent_distinct < 0.5:
                add("warn", "区分値とIDを結んでいる可能性があります",
                    f"{ct}.{cc} は値の種類が {n_child_distinct} 種類しかなく"
                    f"（{', '.join(str(v) for v in _sample_values(prof_c, ct, cc)[:6])} など）、"
                    f"{pt}.{pc} の {n_parent_distinct:,} 種類のうち {parent_hit} 種類にしか当たりません。"
                    "ステータスや区分のような「コード値」の列を、番号がたまたま重なるIDの列に"
                    "結ぼうとしていませんか。")
        if n_parent and n_parent_distinct < n_parent:
            dup = n_parent - n_parent_distinct
            add("warn", "参照先（1側）の値が一意ではありません",
                f"{pt}.{pc} は {n_parent:,} 件中 {dup:,} 件が重複しています。"
                "「1側」は本来ユニークです。重複したまま JOIN すると行が増えて集計が膨らみます。"
                "多重度を N:M にするか、参照先を主キー列に変えてください。")
    except Exception as e:
        add("info", "実データでの確認ができませんでした", str(e)[:120])

    level = "ok"
    if any(i["level"] == "block" for i in issues):
        level = "block"
    elif any(i["level"] == "warn" for i in issues):
        level = "warn"
    return {"level": level, "issues": issues}


def child_parent(entries: list[dict], edge: dict) -> tuple:
    """この関連の (子, 親)。参照整合性の検査はこの向きでしか意味を持たない。

    保存済みの from/to を鵜呑みにせず、主キーがどちら側にあるかで決め直す。
    手で書いた .meta.yaml が逆向きでも、検査は正しい向きで走る。
    """
    frm, to = tuple(edge["from"]), tuple(edge["to"])
    by_alias = {e["alias"]: e for e in entries}

    def sole_pk(ep):
        e = by_alias.get(ep[0])
        if not e:
            return None
        return _is_sole_pk(e["profile"], e.get("meta") or {}, ep[1], ep[2])

    f_pk, t_pk = sole_pk(frm), sole_pk(to)
    if f_pk and not t_pk:
        return to, frm            # from が親だった。入れ替える
    return frm, to


# =============================================================================
# LLM用テキスト生成（プロンプト＝カタログの直列化）
# =============================================================================

def _fmt_value_list(values, col_meta_values: dict) -> str:
    """実値一覧を '1=受付(120), 2=出荷済(300)' 形式で。メタのコード値辞書で意味を補完。"""
    parts = []
    for v, n in values:
        key = "" if v is None else str(v)
        label = (col_meta_values or {}).get(key)
        disp = "NULL" if v is None else str(v)
        if label:
            disp += f"={label}"
        parts.append(f"{disp}({n})")
    return ", ".join(parts)


def table_text(alias: str, tname: str, profile: dict, meta: dict, full: bool) -> str:
    """1テーブル分の説明テキスト。full=False なら1行要約のみ。"""
    t = profile["tables"].get(tname)
    if t is None:
        return f"- {alias}.{tname} : (プロファイル未取得)"
    tmeta = ((meta.get("tables") or {}).get(tname)) or {}
    desc = (tmeta.get("description") or "").strip()
    draft = "（AI推測・未確認）" if tmeta.get("ai_draft") else ""
    rc = t.get("row_count")
    rc_s = f"{rc:,}行" if rc is not None else "行数不明"
    head = f"{alias}.{tname}（{rc_s}）"
    if not full:
        line = f"- {head}" + (f" : {desc}{draft}" if desc else "")
        terms = list(table_glossary(meta, tname))
        # 用語があることだけ知らせる。定義は describe_table で取りに行かせる
        return line + (f" / 業務用語: {', '.join(terms)}" if terms else "")

    lines = [f"### {head}"]
    if desc:
        lines.append(f"{desc}{draft}")

    # 主キーは「1行が何を表すか（粒度）」の手がかりなので、複合キーは構成順で明示する
    pk_cols, pk_src = effective_pk(profile, meta, tname)
    note = {"override": "（人が指定）", "declared": "", "none": ""}[pk_src]
    if len(pk_cols) > 1:
        lines.append(f"主キー{note}: ({', '.join(pk_cols)}) の複合キー → この組み合わせで1行が一意。"
                     f"結合するときは{len(pk_cols)}列すべてを条件にする。")
    elif len(pk_cols) == 1:
        lines.append(f"主キー{note}: {pk_cols[0]}")
    else:
        lines.append("主キー: なし（DBに宣言が無く、指定もされていない）。"
                     "重複行があり得るので COUNT(DISTINCT ...) の要否に注意する。")

    mcols = tmeta.get("columns") or {}
    lines.append("列:")
    for c in t["columns"]:
        cm = (mcols.get(c["name"])) or {}
        parts = [f"- {c['name']} {c['type']}".rstrip()]
        if c["pk"]:
            parts.append("PK")
        if cm.get("description"):
            parts.append(f": {cm['description']}")
        stat = t.get("col_stats", {}).get(c["name"]) or {}
        if "values" in stat:
            parts.append(f"/ 値: {_fmt_value_list(stat['values'], cm.get('values'))}")
        elif cm.get("values"):
            vv = ", ".join(f"{k}={v}" for k, v in cm["values"].items())
            parts.append(f"/ コード値: {vv}")
        elif "min" in stat:
            parts.append(f"/ 範囲: {stat['min']} 〜 {stat['max']}")
        lines.append(" ".join(parts))
    if t.get("sample_rows"):
        lines.append(f"サンプル行 {t['sample_columns']}:")
        for r in t["sample_rows"][:3]:
            lines.append(f"  {r}")

    tgl = table_glossary(meta, tname)
    if tgl:
        lines.append(f"{tname} の業務用語（質問にこの言葉が出たら必ずこの定義に従う）:")
        lines.extend(glossary_lines(tgl))
    return "\n".join(lines)


def db_text(alias: str, db_path, tables: list[str] | None, full: bool) -> str:
    """1DB分の説明テキスト（プロファイル＋メタの合成）。"""
    profile = profile_db(db_path)
    meta = load_meta(db_path)
    names = tables or list(profile["tables"].keys())

    lines = []
    title = meta.get("title") or ""
    lines.append(f"## DB: {alias}" + (f"（{title}）" if title else "") + f" — ファイル: {Path(db_path).name}")
    desc = db_description(meta)
    if desc:
        lines.append(desc)
    lines.append("")
    if full:
        for tname in names:
            lines.append(table_text(alias, tname, profile, meta, full=True))
            lines.append("")
    else:
        lines.append("テーブル一覧:")
        for tname in names:
            lines.append(table_text(alias, tname, profile, meta, full=False))
        lines.append("")

    rels = [r for r in (meta.get("relationships") or [])]
    fk_lines = []
    for tname in names:
        for fk in profile["tables"].get(tname, {}).get("fks", []):
            fk_lines.append(f"- {alias}.{tname}.{fk['from']} = {alias}.{fk['table']}.{fk['to']} (FK宣言)")
    if rels or fk_lines:
        lines.append("結合キー（JOINにはこれを使う）:")
        lines.extend(fk_lines)
        for r in rels:
            card = f" ({r['cardinality']})" if r.get("cardinality") else ""
            lines.append(f"- {r.get('from')} = {r.get('to')}{card}")
        lines.append("")

    gl = db_glossary(meta)
    if gl:
        lines.append("テーブルをまたぐ業務用語（質問にこの言葉が出たら必ずこの定義に従う）:")
        lines.extend(glossary_lines(gl))
        lines.append("")

    exs = meta.get("examples") or []
    if exs:
        lines.append("正しいと確認済みの質問とSQLの例:")
        for ex in exs:
            lines.append(f"Q: {ex.get('q')}")
            # 説明は「この例をどう読むか」の注意書き。人が書いたときだけ載せる
            if ex.get("description"):
                lines.append(f"補足: {ex['description']}")
            lines.append(f"SQL: {ex.get('sql')}")
        lines.append("")
    return "\n".join(lines)


#: 組み立て済みのカタログ本文。DBが多いと1回あたり数十msかかり、
#: 質問のたび・対象を選び直すたびに作り直すのは無駄なので覚えておく。
_TEXT_CACHE: dict = {}


def _text_key(alias: str, path, tables, full: bool):
    """中身が変わったら別物になるキー。DB・メタ・プロファイルの更新時刻を見る。"""
    p = Path(path)

    def stamp(f: Path) -> int:
        try:
            return f.stat().st_mtime_ns
        except OSError:
            return 0

    return (alias, str(p), tuple(tables or ()), full,
            stamp(p), stamp(meta_path(p)), stamp(_cache_path(p)))


def forget(db_path) -> None:
    """そのDBについて覚えているものを捨てる。DBを消したときに呼ぶ。

    本文のキャッシュは更新時刻で自動的に切り替わるが、プロファイルの
    キャッシュはファイルとして残る。DBが無くなったあとも残っていると、
    同じ名前で作り直したときに古い中身が出てくる。
    """
    p = Path(db_path)
    try:
        _cache_path(p).unlink(missing_ok=True)
    except OSError as e:
        print(f"[catalog] キャッシュを消せませんでした: {e}")
    for key in [k for k in _TEXT_CACHE if k[1] == str(p)]:
        _TEXT_CACHE.pop(key, None)


def db_text_cached(alias: str, path, tables=None, full: bool = True) -> str:
    """db_text の結果を使い回す版。カタログを直せば自動で作り直される。"""
    try:
        key = _text_key(alias, path, tables, full)
    except Exception:
        return db_text(alias, path, tables, full=full)
    hit = _TEXT_CACHE.get(key)
    if hit is None:
        hit = db_text(alias, path, tables, full=full)
        if len(_TEXT_CACHE) > 64:            # 古い世代が溜まりすぎないように
            _TEXT_CACHE.clear()
        _TEXT_CACHE[key] = hit
    return hit


def inline_length(scope: list[dict]) -> int:
    """詳細版カタログの文字数（列名までAIに渡せるかの判断に使う）。"""
    return sum(len(db_text_cached(s["alias"], s["path"], s.get("tables"), full=True))
               for s in (scope or []))


def inline_limit() -> int:
    """カタログを全文のまま入れる上限の天井（モデルを知らない呼び出し向け）。

    実効値は選択中モデルの文脈量から自動で決まる（models.inline_limit_for）ので、
    通常は prompt_for_scope に limit を渡す。ここはモデルが分からない場面の
    フォールバック（天井値）。読めなければ env の初期値に落とす。
    """
    try:
        import models
        return models.prompt_inline_limit()
    except Exception:
        return config.PROMPT_INLINE_LIMIT_CHARS


def prompt_for_scope(scope: list[dict], limit: int | None = None) -> str:
    """選択スコープ全体のカタログテキスト。

    全文が上限（limit。省略時は管理者設定/env）以下なら詳細をインライン、
    超えるなら要約のみ（詳細は describe_table ツールで取得させる）。
    limit は「選択中のモデルが一度に読める量」から呼び出し側が渡せる
    （models.inline_limit_for 参照。固定値だと小さいモデルで溢れるため）。
    """
    if not scope:
        return "（対象にできるDBがありません。「データ取り込み」でDBを作るよう案内してください。）"
    full = "\n".join(db_text_cached(s["alias"], s["path"], s.get("tables"), full=True)
                     for s in scope)
    if len(full) <= (limit if limit is not None else inline_limit()):
        return full
    compact = "\n".join(db_text_cached(s["alias"], s["path"], s.get("tables"), full=False)
                        for s in scope)
    # ここに載っているのはテーブル単位の説明までで、列名は入っていない。
    # それを言わずに渡すと「その列は無い」と早合点して、できることまで断ってしまう。
    return (compact + "\n"
            "【重要】対象のDBが多いため、上には各テーブルの説明までしか載せていません。"
            "**列名は1つも載っていません。**\n"
            "そのため、上に見当たらないという理由で「その列は無い」「そのテーブルは無い」と"
            "判断してはいけません。必要な列があるかどうかは、必ず describe_table を呼んで"
            "確かめること。名前から中身が推測できるテーブル（商品なら products、"
            "社員なら employees など）は、まず describe_table で列を見てから答えること。\n"
            "ユーザーに「その情報は無い」と答えてよいのは、関係しそうなテーブルを"
            "describe_table で実際に確認した後だけです。")


def describe_table_text(scope: list[dict], db_alias: str, tname: str) -> str:
    """describe_table ツールの実体。alias と テーブル名から詳細テキストを返す。"""
    # "alias.table" 形式で渡された場合に対応
    if "." in tname and not db_alias:
        db_alias, tname = tname.split(".", 1)
    entry = next((s for s in scope if s["alias"].lower() == str(db_alias).lower()), None)
    if entry is None:
        aliases = ", ".join(s["alias"] for s in scope)
        return f"エラー: DBエイリアス '{db_alias}' は選択されていません。選択中: {aliases}"
    profile = profile_db(entry["path"])
    meta = load_meta(entry["path"])
    if tname not in profile["tables"]:
        cand = ", ".join(profile["tables"].keys())
        return f"エラー: テーブル '{tname}' は {db_alias} にありません。存在するテーブル: {cand}"
    return table_text(entry["alias"], tname, profile, meta, full=True)


# ==========================================================================
# ===== 元 verify.py
# 相互検証（検算）。同じ数字を独立した2つの経路で計算して突き合わせる。
#
# text-to-SQL の最大のリスクは「もっともらしいが間違っているSQL」ではなく、
# 「正しいSQLなのに、業務的には別の数字を指している」ことにある。
# 実際、demo_sales の「売上」は明細から数えると1.23億、請求から数えると0.85億で、
# どちらのSQLも正しい。差の3,866万円は未請求の受注339件だった。
# どのSQLを書くかで答えが1.5倍変わるのに、聞いた人にはそれが見えない。
#
# そこで、カタログに「一致するはずの2つの式」を検算ルールとして登録しておき、
# AIがそのテーブルに触れるSQLを実行するたびに突き合わせる。
#
#   data/<DB>.db.meta.yaml:
#     checks:
#       - name: 入金と請求（入金済）の一致
#         left:  {label: 入金の合計,          sql: SELECT SUM(amount) FROM demo_sales.payments}
#         right: {label: 請求のうち入金済み,   sql: SELECT SUM(amount) FROM demo_sales.invoices WHERE paid_flag = 1}
#         tolerance_pct: 0.1        # 許容差（%）。これ以内なら一致とみなす
#         drilldown: SELECT ...     # 不一致のとき、差の実体を見せるSQL（任意）
#         enabled: true
#
# 設計上の約束:
#   * 左右のSQLは「1行1列のスカラ」を返すこと（SUM や COUNT）。
#   * 検算は質問のたびに走るが、結果はデータの版（DBファイルのmtime）で
#     キャッシュするので、実際に実行されるのはデータが変わった後の最初の1回だけ。
#   * 壊れた検算ルール（SQLエラー）は黙って飛ばす。質問への回答を止めないため。
#     ルール自体の点検は、カタログ画面の「検算」から人が行う。
# ==========================================================================
import re
from pathlib import Path

import db

#: 許容差の既定（%）。丸め誤差を拾わない程度に小さく。
DEFAULT_TOLERANCE_PCT = 0.5
#: 不一致時に内訳SQLで見せる行数。
DRILL_ROWS = 8
#: 検算結果のキャッシュ。キーは（ルールの中身, データの版）。
_verify_cache: dict = {}
_CACHE_MAX = 300


# =============================================================================
# ルールの読み出し
# =============================================================================

def normalize(raw) -> list[dict]:
    """meta の checks をあるべき形に揃える。壊れた項目は落とす。"""
    out = []
    for c in (raw or []):
        if not isinstance(c, dict):
            continue
        left, right = c.get("left") or {}, c.get("right") or {}
        lsql = str(left.get("sql") or "").strip()
        rsql = str(right.get("sql") or "").strip()
        if not lsql or not rsql:
            continue
        try:
            tol = float(c.get("tolerance_pct", DEFAULT_TOLERANCE_PCT))
        except (TypeError, ValueError):
            tol = DEFAULT_TOLERANCE_PCT
        out.append({
            "name": str(c.get("name") or "検算").strip(),
            "left": {"label": str(left.get("label") or "左"), "sql": lsql},
            "right": {"label": str(right.get("label") or "右"), "sql": rsql},
            "tolerance_pct": max(0.0, tol),
            "drilldown": str(c.get("drilldown") or "").strip(),
            "enabled": c.get("enabled", True) is not False,
        })
    return out


def checks_for(scope: list[dict]) -> list[dict]:
    """選択中のDB群に登録されている検算ルール（有効なものだけ）。"""
    import catalog

    out = []
    for s in scope or []:
        meta = s.get("meta") or catalog.load_meta(s["path"])
        for c in normalize(meta.get("checks")):
            if c["enabled"]:
                out.append({**c, "owner": s.get("alias") or ""})
    return out


# =============================================================================
# 「このSQLはどのテーブルに触れているか」
# =============================================================================

def tables_in(sql: str, scope: list[dict]) -> set:
    """SQLが触れている (alias, table) の集合。名前の照合だけで判定する。"""
    found = set()
    for s in scope or []:
        alias = str(s.get("alias") or "")
        for t in (s.get("tables") or []):
            name = str(t)
            qualified = alias and re.search(
                r'(?<![\w."])' + re.escape(alias) + r'\s*\.\s*"?' + re.escape(name) + r'"?(?![\w])',
                sql, re.IGNORECASE)
            bare = re.search(r'(?<![\w."])"?' + re.escape(name) + r'"?(?![\w])',
                             sql, re.IGNORECASE)
            if qualified or bare:
                found.add((alias, name))
    return found


# =============================================================================
# 実行
# =============================================================================

def _scalar(sql: str, scope: list[dict]):
    """1行1列のSELECTを実行して数値を返す。数値でなければ ValueError。"""
    columns, rows, _ = db.run_select(sql, scope, max_rows=1)
    v = rows[0][0] if rows else None
    if v is None:
        return 0.0                      # SUMが空のときのNULLは0として扱う
    return float(v)


def _fingerprint(check: dict) -> tuple:
    return (check["name"], check["left"]["sql"], check["right"]["sql"],
            check["tolerance_pct"], check["drilldown"])


def _data_version(check: dict, scope: list[dict]) -> tuple:
    """検算が読むDBファイルの版。これが変わったら計算し直す。"""
    text = " ".join([check["left"]["sql"], check["right"]["sql"], check["drilldown"]])
    stamps = []
    for s in db.narrow_scope(text, scope):
        try:
            stamps.append((str(s["path"]), Path(s["path"]).stat().st_mtime_ns))
        except OSError:
            stamps.append((str(s["path"]), 0))
    return tuple(sorted(stamps))


def run_check(check: dict, scope: list[dict], use_cache: bool = True) -> dict:
    """検算を1本実行する。

    戻り値:
      {"ok_run": bool, "match": bool, "left": float, "right": float,
       "diff": float, "pct": float|None, "version": str,
       "drill": {"columns", "rows", "truncated"} | None, "error": str|None}
    """
    version = _data_version(check, scope)
    key = (_fingerprint(check), version)
    if use_cache and key in _verify_cache:
        return _verify_cache[key]

    res: dict = {"ok_run": False, "match": True, "left": None, "right": None,
                 "diff": None, "pct": None, "drill": None, "error": None,
                 "version": str(hash(version))}
    try:
        lv = _scalar(check["left"]["sql"], scope)
        rv = _scalar(check["right"]["sql"], scope)
    except Exception as e:
        res["error"] = str(e).splitlines()[0][:200]
        _remember(key, res)
        return res

    diff = lv - rv
    base = max(abs(lv), abs(rv))
    pct = (abs(diff) / base * 100) if base else 0.0
    match = pct <= check["tolerance_pct"]
    res.update({"ok_run": True, "match": match, "left": lv, "right": rv,
                "diff": diff, "pct": round(pct, 2)})

    if not match and check["drilldown"]:
        try:
            columns, rows, truncated = db.run_select(
                check["drilldown"], scope, max_rows=DRILL_ROWS)
            res["drill"] = {"columns": columns,
                            "rows": [list(r) for r in rows],
                            "truncated": truncated}
        except Exception as e:
            res["drill"] = {"error": str(e).splitlines()[0][:160]}

    _remember(key, res)
    return res


def _remember(key, res) -> None:
    if len(_verify_cache) > _CACHE_MAX:
        _verify_cache.clear()
    _verify_cache[key] = res


def clear_cache() -> None:
    """テスト用。"""
    _verify_cache.clear()


# =============================================================================
# 質問への割り込み（ツール実行後に呼ばれる）
# =============================================================================

def alerts_for(sql_texts: list[str], scope: list[dict]) -> list[dict]:
    """実行されたSQL群に関係する検算を走らせ、不一致だけを返す。

    一致した検算・実行できなかった検算は何も言わない
    （毎回「問題ありません」と言われても読まれなくなるだけ）。
    """
    texts = [t for t in (sql_texts or []) if t and t.strip()]
    if not texts or not scope:
        return []
    try:
        checks = checks_for(scope)
    except Exception:
        return []
    if not checks:
        return []

    touched = set()
    for t in texts:
        touched |= tables_in(t, scope)
    if not touched:
        return []

    alerts = []
    for check in checks:
        involved = tables_in(check["left"]["sql"] + " " + check["right"]["sql"], scope)
        if not (involved & touched):
            continue
        res = run_check(check, scope)
        if not res["ok_run"] or res["match"]:
            continue
        alerts.append({
            "key": f"verify||{check['owner']}||{check['name']}||{res['version']}",
            "name": check["name"],
            "left_label": check["left"]["label"], "left": res["left"],
            "right_label": check["right"]["label"], "right": res["right"],
            "diff": res["diff"], "pct": res["pct"],
            "tolerance_pct": check["tolerance_pct"],
            "drill": res["drill"],
        })
    return alerts


def _fmt(v) -> str:
    if v is None:
        return "—"
    return f"{v:,.4f}".rstrip("0").rstrip(".") if v % 1 else f"{int(v):,}"


def llm_note(alert: dict) -> dict:
    """LLMのツール結果に混ぜる、検算の注意書き。"""
    note = {
        "check": alert["name"],
        alert["left_label"]: alert["left"],
        alert["right_label"]: alert["right"],
        "difference": alert["diff"],
        "difference_pct": alert["pct"],
        "instruction": ("この2つの数字は一致するはずですが食い違っています。"
                        "回答では、どちらの数字を使ったのかと、この差異があることを"
                        "必ず注記してください。差異の理由を推測で断定しないこと。"),
    }
    drill = alert.get("drill") or {}
    if drill.get("rows"):
        note["difference_detail_sample"] = {
            "columns": drill["columns"], "rows": drill["rows"][:3]}
    return note


def render_item(alert: dict) -> dict:
    """画面に出す検算カード。分析結果と同じ report の形で描ける。"""
    tables = [{
        "name": "2つの経路の比較",
        "columns": ["経路", "値"],
        "rows": [(alert["left_label"], alert["left"]),
                 (alert["right_label"], alert["right"]),
                 ("差", alert["diff"]),
                 ("差の割合", f"{alert['pct']}%（許容 {alert['tolerance_pct']}%）")],
    }]
    notes = [f"「{alert['left_label']}」と「{alert['right_label']}」は一致するはず"
             f"ですが、{_fmt(abs(alert['diff']))}（{alert['pct']}%）食い違っています。"]
    drill = alert.get("drill") or {}
    if drill.get("rows"):
        tables.append({"name": "差の内訳（先頭のみ）",
                       "columns": drill["columns"], "rows": drill["rows"]})
        notes.append("内訳の表は差の実体の一部です。全体はデータカタログの"
                     "「検算」で確認できます。")
    elif drill.get("error"):
        notes.append(f"内訳SQLは実行できませんでした: {drill['error']}")
    notes.append("この検算ルールはデータカタログの「用語集・例文 → 検算」で"
                 "管理されています。差が正しい業務状態なら、許容差を広げるか"
                 "ルールを無効にしてください。")
    return {"role": "assistant", "kind": "report",
            "title": f"⚠ 検算: {alert['name']}",
            "tables": tables, "notes": notes,
            "verify_key": alert["key"]}


# ==========================================================================
# ===== 元 sqlusage.py
# 過去の分析で「実際に使われた」結合を数える。
#
# ER図が描いているのは宣言された関連で、実際に通った道ではない。
# チャット履歴には実行されたSQLが全部残っているので、そこからJOINを取り出して
# 数えると、次の3つが見えるようになる。
#
#   よく通る道       … 太く描く。分析の主要動線
#   誰も通らない道   … 灰色にする。検算されていない経路でもある
#                      （実際、demo_sales で3,866万円の未請求が見つかったのは
#                        一度も使われていなかった invoices への経路の上だった）
#   登録の無い道     … AIが実際に結合しているのにカタログに無い。
#                      「関連の候補」に実績つきで出す。登録すべきか、
#                      AIが誤った結合をしているかのどちらかで、どちらでも知る価値がある
#
# 解析は正規表現＋カタログのプロファイル（テーブル・列の一覧）で行う。
# SQLパーサは入れない。ここでの用途は「多い・少ない・ゼロ」が分かればよく、
# 多少の取りこぼしで結論が変わらないため。
# ==========================================================================
import json
import re
from pathlib import Path

import catalog
import config
import db

#: エイリアスとして解釈してはいけない語。
_sqlusage_RESERVED = {"on", "using", "where", "group", "order", "left", "right", "inner",
             "outer", "cross", "natural", "join", "as", "select", "from", "limit",
             "having", "union", "all", "and", "or", "not", "set", "by"}

_FROM_RE = re.compile(
    r'\b(?:from|join)\s+("?[\w一-龠ぁ-んァ-ヶ．.]+"?)(?:\s+(?:as\s+)?([A-Za-z_]\w*))?',
    re.IGNORECASE)
_USING_RE = re.compile(
    r'\bjoin\s+("?[\w一-龠ぁ-んァ-ヶ．.]+"?)(?:\s+(?:as\s+)?([A-Za-z_]\w*))?'
    r'\s+using\s*\(\s*"?(\w+)"?\s*\)', re.IGNORECASE)
_ON_RE = re.compile(
    r'\bon\s+("?[\w一-龠ぁ-んァ-ヶ．.]+"?)\s*=\s*("?[\w一-龠ぁ-んァ-ヶ．.]+"?)',
    re.IGNORECASE)


# =============================================================================
# 履歴からSQLを集める
# =============================================================================

def _walk_sql(node, acc: list) -> None:
    if isinstance(node, dict):
        for k, v in node.items():
            if k == "sql" and isinstance(v, str) and v.strip():
                acc.append(v)
            else:
                _walk_sql(v, acc)
    elif isinstance(node, list):
        for v in node:
            _walk_sql(v, acc)


def collect_sqls() -> tuple:
    """全ユーザーのチャット履歴から実行SQLを集める。

    同じSQLが画面用の写しとツール呼び出しの両方に残っているので、
    会話単位で重複を除く。戻り値: (SQLのリスト, 会話数)
    """
    sqls: list[str] = []
    users = Path(config.USER_META_DIR)
    chats = 0
    if not users.exists():
        return sqls, 0
    for f in users.glob("*/chats/*.json"):
        if f.name == "index.json":
            continue
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            continue
        chats += 1
        seen: set = set()
        for item in data.get("render_log") or []:
            if item.get("kind") == "sql" and item.get("sql"):
                seen.add(str(item["sql"]).strip())
        buf: list = []
        for m in data.get("messages") or []:
            for tc in (m.get("tool_calls") or []):
                try:
                    args = json.loads((tc.get("function") or {}).get("arguments") or "{}")
                except Exception:
                    continue
                _walk_sql(args, buf)
        seen |= {s.strip() for s in buf}
        sqls.extend(seen)
    return sqls, chats


# =============================================================================
# JOIN の解決
# =============================================================================

def _entries() -> list[dict]:
    """全DBのプロファイル（テーブル・列）。名前解決の台帳になる。"""
    out = []
    for p in db.list_db_files():
        try:
            out.append({"alias": db.alias_for(p), "path": p,
                        "profile": catalog.profile_db(p)})
        except Exception:
            continue
    return out


def _resolve_table(raw: str, entries: list[dict], hint_aliases: set):
    """'demo_sales.orders' や 'orders' を (DBエイリアス, テーブル) にする。"""
    name = raw.strip().strip('"')
    if "." in name:
        prefix, _, rest = name.partition(".")
        rest = rest.strip('"')
        for e in entries:
            if e["alias"].lower() == prefix.lower() and rest in e["profile"]["tables"]:
                return (e["alias"], rest)
        return None
    hits = [e["alias"] for e in entries if name in e["profile"]["tables"]]
    if len(hits) == 1:
        return (hits[0], name)
    if hits:
        # 同名テーブルが複数DBにある。同じSQLに出てきたDBを優先する
        for a in hits:
            if a in hint_aliases:
                return (a, name)
    return None


def _columns_of(entries: list[dict], alias: str, table: str) -> set:
    for e in entries:
        if e["alias"] == alias:
            t = e["profile"]["tables"].get(table) or {}
            return {c["name"] for c in t.get("columns", [])}
    return set()


def _edge_key(a: tuple, b: tuple) -> str:
    x, y = ".".join(a), ".".join(b)
    return f"{x}||{y}" if x <= y else f"{y}||{x}"


def joins_in(sql: str, entries: list[dict]) -> list[tuple]:
    """1本のSQLから、(端点, 端点) のリストを取り出す。端点 = (alias, table, column)。"""
    flat = " ".join(sql.split())

    # 出現順のテーブルと、エイリアス→テーブルの対応
    order: list[tuple] = []
    alias_map: dict = {}
    hint = {e["alias"] for e in entries
            if re.search(r'(?<![\w."])' + re.escape(e["alias"]) + r'\s*\.',
                         flat, re.IGNORECASE)}
    for m in _FROM_RE.finditer(flat):
        raw, al = m.group(1), (m.group(2) or "")
        resolved = _resolve_table(raw, entries, hint)
        if resolved is None:
            continue
        order.append(resolved)
        alias_map[resolved[1].lower()] = resolved          # テーブル名でも引ける
        if al and al.lower() not in _sqlusage_RESERVED:
            alias_map[al.lower()] = resolved

    out: list[tuple] = []

    # JOIN ... USING(col): 相手は「それより前に出た、同じ列を持つテーブル」
    for m in _USING_RE.finditer(flat):
        raw, col = m.group(1), m.group(3)
        right = _resolve_table(raw, entries, hint)
        if right is None:
            continue
        try:
            pos = order.index(right)
        except ValueError:
            continue
        left = next((t for t in reversed(order[:pos])
                     if col in _columns_of(entries, *t)), None)
        if left and col in _columns_of(entries, *right):
            out.append(((*left, col), (*right, col)))

    # JOIN ... ON a.x = b.y
    for m in _ON_RE.finditer(flat):
        ends = []
        for side in (m.group(1), m.group(2)):
            side = side.strip().strip('"')
            if "." not in side:
                break
            qual, _, col = side.rpartition(".")
            qual = qual.strip('"')
            t = None
            if "." in qual:                       # demo_sales.orders.customer_id
                t = _resolve_table(qual, entries, hint)
            else:                                 # o.customer_id / orders.customer_id
                t = alias_map.get(qual.lower())
            if t is None or col not in _columns_of(entries, *t):
                break
            ends.append((*t, col))
        if len(ends) == 2 and ends[0][:2] != ends[1][:2]:
            out.append((ends[0], ends[1]))
    return out


# =============================================================================
# 集計とAPI
# =============================================================================

def usage_counts() -> dict:
    """全履歴のJOINを数える。{edge_key: {"from","to","count"}}"""
    entries = _entries()
    sqls, chats = collect_sqls()
    edges: dict = {}
    tables: dict = {}
    for sql in sqls:
        for a, b in joins_in(sql, entries):
            key = _edge_key(a, b)
            hit = edges.setdefault(key, {"from": list(a), "to": list(b), "count": 0})
            hit["count"] += 1
        low = sql.lower()
        for e in entries:
            for t in e["profile"]["tables"]:
                if re.search(r'(?<![\w."])' + re.escape(t.lower()) + r'(?![\w])', low):
                    tables[f"{e['alias']}.{t}"] = tables.get(f"{e['alias']}.{t}", 0) + 1
    return {"edges": edges, "tables": tables,
            "scanned": {"chats": chats, "sqls": len(sqls)}}


def _declared_pairs(entries: list[dict]) -> set:
    """カタログ/FKに登録済みの結合（端点の組）。"""
    cat_entries = [{"alias": e["alias"], "profile": e["profile"],
                    "meta": catalog.load_meta(e["path"])} for e in entries]
    out = set()
    for edge in catalog.collect_edges(cat_entries):
        out.add(_edge_key(edge["from"], edge["to"]))
    return out


def usage_for(alias: str) -> dict:
    """ER図に重ねるためのデータ（このDBのキャンバス向け）。"""
    data = usage_counts()
    return {
        "edges": {k: v["count"] for k, v in data["edges"].items()},
        "tables": {k: n for k, n in data["tables"].items()
                   if k.startswith(alias + ".")},
        "scanned": data["scanned"],
    }


def suggestions_for(alias: str, profile: dict, meta: dict) -> list[dict]:
    """実際に使われているのにカタログに無い結合を「関連の候補」に出す。

    形は catalog.join_suggestions と同じ。from 側は必ずこのDBのテーブルにする
    （関連はそのDBの .meta.yaml に書かれるため）。
    """
    entries = _entries()
    if not entries:
        return []
    declared = _declared_pairs(entries)
    data = usage_counts()

    out = []
    for key, e in sorted(data["edges"].items(), key=lambda kv: -kv[1]["count"]):
        if key in declared:
            continue
        a, b = tuple(e["from"]), tuple(e["to"])
        # from 側をこのDBに揃える。どちらもこのDBでなければ、この画面では出さない
        if a[0] != alias and b[0] != alias:
            continue
        if a[0] != alias:
            a, b = b, a
        frm = f"{a[1]}.{a[2]}"
        to = f"{b[1]}.{b[2]}" if b[0] == alias else f"{b[0]}.{b[1]}.{b[2]}"
        out.append({"from": frm, "to": to, "cardinality": "N:1",
                    "reason": f"過去の分析で{e['count']}回使われています（未登録）"})
    return out[:8]


# ==========================================================================
# ===== 元 usage.py
# このアプリ自身が、誰にどう使われているかを数える。
#
# sqlusage.py が「どの結合が通ったか」だけを見るのに対し、こちらは利用そのものを見る。
# 見たいのは利用者数ではなく、次の3つ。
#
#   伸びているか   … 使われ続けているのか、最初の週だけだったのか
#   何に使われるか … よく呼ばれる機能と、まったく呼ばれない機能
#   どこで転ぶか   … 失敗した質問。これがカタログを直す入口になる
#
# 失敗の中身は分けて数える。「列が無い」はカタログ不足で人間が直せるが、
# 「LLM呼び出しに失敗」は設定や回線の問題で、カタログをいくら直しても減らない。
# 混ぜて「エラー率5%」と出すと、直せないものを直そうとして時間を溶かす。
#
# 材料は data/users/<ユーザー>/chats/*.json（会話の実体）と、
# data/import_history.jsonl（取り込みの記録）。どちらも読むだけで書き換えない。
#
# 戻り値の形は advanced.py / business.py と同じ {"title", "tables", "notes", "meta"}。
# 画面もLLMも同じ入れ物で受け取れる。
#
# 時系列は発言ごとの時刻（表示物の at）で数える。この仕組みを入れる前の
# 古い会話には at が無いので、会話の開始時刻で代用し、その旨を所見に明示する。
# ==========================================================================
import json
import re
from collections import Counter
from datetime import datetime, timedelta
from pathlib import Path

import config
import history

#: 質問文・エラー文をそのまま並べるときの上限。多すぎると読まれない。
MAX_LIST = 40

METHODS = {
    "summary": "全体像（期間・利用者・会話数・失敗率）",
    "users": "利用者ごとの利用量",
    "trend": "日ごと・曜日・時間帯の推移",
    "tools": "呼ばれた機能の回数",
    "databases": "実際に使われたDB",
    "errors": "失敗の内訳と、直し方の当たり",
    "questions": "実際に聞かれた質問",
}

#: 失敗の分類。上から順に当てる（先に当たったものを採る）。
#: 「誰が直せるか」で分ける。カタログ担当・管理者・利用者では打ち手が違う。
_ERROR_KINDS = (
    ("カタログ不足（列・テーブルの取り違え）",
     r"no such (table|column)|列名が違います|テーブルが見つかりません",
     "describe_table で確認できる情報が足りていない。"
     "カタログの列説明・コード値・結合定義を書き足すと減る。"),
    ("SQLの誤り（構文・集計）",
     r"SQL実行エラー|syntax error|ambiguous|misuse of aggregate",
     "例文（Q&SQL）を足すと、AIが型を真似るので減る。"),
    ("分析に足りるデータが無い",
     r"行しかなく|0行でした|データが0行|足りません",
     "抽出条件が狭すぎる。期間を広げるか、分析の指定（説明変数など）を減らす。"),
    ("ツールの引数不足",
     r"(には|は).{0,30}(必要|指定してください)|引数|列が結果にありません|指定列",
     "ツールの説明文を具体的にすると、AIの指定ミスが減る。"),
    ("LLM・API側の問題",
     r"LLM呼び出しに失敗|Error code:|timeout|接続",
     "カタログでは直らない。モデル設定・APIキー・回線を確認する。"),
    ("実行時間切れ",
     r"時間がかかりすぎ|タイムアウト|interrupted",
     "対象データを絞るか、集計済みのユーザー定義ツールを用意する。"),
)

_WEEKDAYS = ("月", "火", "水", "木", "金", "土", "日")


def _usage_dt(value) -> datetime | None:
    try:
        return datetime.fromisoformat(str(value))
    except (TypeError, ValueError):
        return None


def _usage_table(name: str, columns: list, rows: list) -> dict:
    return {"name": name, "columns": columns, "rows": [tuple(r) for r in rows]}


def _usage_out(title: str, tables: list, notes: list, meta: dict | None = None) -> dict:
    return {"title": title, "tables": tables, "notes": notes, "meta": meta or {}}


def classify_error(message: str) -> tuple[str, str]:
    """エラー文を「誰が直せるか」で分類する。戻り値: (分類, 打ち手)"""
    for name, pattern, fix in _ERROR_KINDS:
        if re.search(pattern, message, re.IGNORECASE):
            return name, fix
    return "その他", "内容を読んで個別に判断する。"


# =============================================================================
# 材料集め
# =============================================================================

def _asked(log: list, fallback: datetime | None) -> list[dict]:
    """1発言ぶんの記録。質問の時刻と、答え終わった時刻を組にする。

    表示物には at が入っている（質問は積んだ時刻、応答は保存した時刻）。
    この2つの差が、その質問で待たされた時間になる。
    at を持たない古い会話は、会話の開始時刻で代用する。
    """
    out: list[dict] = []
    for item in log:
        at = _usage_dt(item.get("at")) or fallback
        if item.get("role") == "user" and item.get("kind") == "text":
            out.append({"at": at, "text": str(item.get("content") or "").strip(),
                        "done": None, "failed": False, "errors": [],
                        "exact": bool(item.get("at"))})
        elif out:
            if at and out[-1]["at"] and at >= out[-1]["at"]:
                out[-1]["done"] = at
            if item.get("kind") == "error":
                out[-1]["failed"] = True
                out[-1]["errors"].append(str(item.get("message") or ""))
    return out


def collect(days: int | None = None, user: str | None = None) -> list[dict]:
    """会話ファイルを1会話1レコードに畳む。

    days を指定すると、その日数より前に始まった会話は捨てる。
    捨てるのは開始時刻で判定するので、古い会話を今日まで続けていた場合も対象外。
    """
    root = Path(config.USER_META_DIR)
    if not root.exists():
        return []
    limit = datetime.now() - timedelta(days=days) if days else None

    out = []
    for f in root.glob("*/chats/*.json"):
        if f.name == "index.json":
            continue
        who = f.parent.parent.name
        if user and who.lower() != user.lower():
            continue
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            continue                      # 壊れた1本のために全体を止めない
        created = _usage_dt(data.get("created_at"))
        if limit and created and created < limit:
            continue

        log = data.get("render_log") or []
        asked = _asked(log, created)
        # 最初の質問より前に出た失敗は、どの質問のものとも言えない
        first_q = next((i for i, x in enumerate(log)
                        if x.get("role") == "user" and x.get("kind") == "text"), len(log))
        orphans = [str(x.get("message") or "") for x in log[:first_q]
                   if x.get("kind") == "error"]
        kinds = Counter(str(i.get("kind") or "") for i in log)
        tools = [(tc.get("function") or {}).get("name")
                 for m in (data.get("messages") or [])
                 for tc in (m.get("tool_calls") or [])]
        out.append({
            "user": who,
            "id": data.get("id") or f.stem,
            "title": str(data.get("title") or "").strip(),
            "created": created,
            "updated": _usage_dt(data.get("updated_at")),
            "asked": asked,
            "questions": [a["text"] for a in asked if a["text"]],
            "tools": [t for t in tools if t],
            "errors": [str(i.get("message") or "") for i in log
                       if i.get("kind") == "error"],
            "orphan_errors": orphans,
            "dbs": list(data.get("db_names") or []),
            "sqls": kinds.get("sql", 0),
            "charts": kinds.get("chart", 0),
            "tables": kinds.get("table", 0),
            "files": kinds.get("file", 0),
            "reports": kinds.get("report", 0),
        })
    out.sort(key=lambda r: r["created"] or datetime.min)
    return out


def _period_note(records: list[dict]) -> str:
    days = [r["created"] for r in records if r["created"]]
    if not days:
        return "期間: 不明"
    return f"期間: {min(days):%Y-%m-%d} 〜 {max(days):%Y-%m-%d}"


def _empty(days: int | None) -> dict:
    span = f"直近{days}日には" if days else ""
    return _usage_out("利用状況", [], [f"{span}会話の記録がありませんでした。"
                                "まだ誰も使っていないか、data/users/ が空です。"])


# =============================================================================
# 分析
# =============================================================================

def summary(records: list[dict], days: int | None = None) -> dict:
    """全体像。まずこれを見て、気になった軸を他のメソッドで掘る。"""
    if not records:
        return _empty(days)

    n_chats = len(records)
    n_turns = sum(len(r["questions"]) for r in records)
    users = {r["user"] for r in records}
    errs = [e for r in records for e in r["errors"]]
    chats_with_err = sum(1 for r in records if r["errors"])
    active_days = {r["created"].date() for r in records if r["created"]}

    rows = [
        ("会話", f"{n_chats:,} 件"),
        ("質問", f"{n_turns:,} 回"),
        ("利用者", f"{len(users)} 人"),
        ("使われた日", f"{len(active_days)} 日"),
        ("1会話あたりの質問", f"{n_turns / n_chats:.1f} 回"),
        ("失敗を含む会話", f"{chats_with_err} 件（{chats_with_err / n_chats * 100:.0f}%）"),
        ("作った表・グラフ", f"表 {sum(r['tables'] for r in records):,} / "
                            f"グラフ {sum(r['charts'] for r in records):,}"),
        ("出したファイル", f"{sum(r['files'] for r in records):,} 件"),
    ]

    notes = [_period_note(records)]
    # 1回で終わった会話が多いなら、続けて聞ける場になっていない可能性がある
    one_shot = sum(1 for r in records if len(r["questions"]) <= 1)
    if n_chats >= 5:
        notes.append(
            f"1問だけで終わった会話が {one_shot}/{n_chats} 件"
            f"（{one_shot / n_chats * 100:.0f}%）。"
            + ("会話を続けて掘り下げる使い方が根づいています。"
               if one_shot / n_chats < 0.5 else
               "多くが単発です。最初の答えで満足したか、続きを諦めたかのどちらかなので、"
               "失敗の内訳（errors）も合わせて見てください。"))
    if errs:
        kinds = Counter(classify_error(e)[0] for e in errs)
        top, n = kinds.most_common(1)[0]
        notes.append(f"失敗 {len(errs)} 件のうち最も多いのは「{top}」{n} 件。"
                     "内訳は errors で確認できます。")
    else:
        notes.append("記録された失敗はありません。")

    heavy = Counter(r["user"] for r in records).most_common(1)
    if heavy and len(users) > 1:
        who, cnt = heavy[0]
        notes.append(f"最も使っているのは {who}（{cnt} 件 / 全体の "
                     f"{cnt / n_chats * 100:.0f}%）。")

    return _usage_out("利用状況の全体像", [_usage_table("全体", ["項目", "値"], rows)], notes,
                {"chats": n_chats, "turns": n_turns, "users": len(users),
                 "errors": len(errs), "active_days": len(active_days)})


def by_user(records: list[dict]) -> dict:
    """利用者ごと。誰が使っていて、誰が離れたかを見る。"""
    if not records:
        return _empty(None)

    per: dict[str, list] = {}
    for r in records:
        per.setdefault(r["user"], []).append(r)

    rows = []
    for who, rs in sorted(per.items(), key=lambda kv: -len(kv[1])):
        last = max((r["updated"] or r["created"] for r in rs
                    if (r["updated"] or r["created"])), default=None)
        turns = sum(len(r["questions"]) for r in rs)
        errs = sum(len(r["errors"]) for r in rs)
        rows.append((who, len(rs), turns, round(turns / len(rs), 1), errs,
                     f"{last:%Y-%m-%d}" if last else "—"))
    notes = [_period_note(records),
             "「最終利用」が古い人は、使えなかったのか、必要が無かったのかを直接聞くのが早いです。"]

    today = datetime.now()
    stale = [r[0] for r in rows
             if r[5] != "—" and (today - datetime.fromisoformat(r[5])).days >= 14]
    if stale:
        notes.append(f"2週間以上使っていないのは {len(stale)} 人（{'、'.join(stale[:5])}"
                     f"{' ほか' if len(stale) > 5 else ''}）。")
    return _usage_out("利用者ごとの利用量",
                [_usage_table("利用者別", ["利用者", "会話", "質問", "1会話あたり", "失敗", "最終利用"],
                        rows)],
                notes, {"users": len(rows)})


def trend(records: list[dict]) -> dict:
    """日ごと・曜日・時間帯。定着したのか、一度きりだったのかを見る。

    数えるのは会話ではなく質問1件ずつ。1本の会話で何度も聞いていれば、
    その回数だけ数える（実際にどれだけ使われたかは、そちらの方が近い）。
    """
    if not records:
        return _empty(None)

    asked = [a for r in records for a in r["asked"] if a["at"]]
    if not asked:
        return _empty(None)

    daily = Counter(a["at"].date() for a in asked)
    dow = Counter(a["at"].weekday() for a in asked)
    hour = Counter(a["at"].hour for a in asked)

    day_rows = [(str(d), n) for d, n in sorted(daily.items())]
    dow_rows = [(_WEEKDAYS[i], dow.get(i, 0)) for i in range(7)]
    hour_rows = [(f"{h:02d}時", hour.get(h, 0)) for h in range(24) if hour.get(h)]

    notes = [_period_note(records), f"質問 {len(asked)} 件を、聞かれた時刻で数えています。"]
    rough = sum(1 for a in asked if not a["exact"])
    if rough:
        notes.append(f"うち {rough} 件は時刻を持たない古い会話で、"
                     "会話の開始時刻で代用しています。")

    # 応答にかかった時間。待たされているなら、対象データの絞り込みや
    # ユーザー定義ツールの用意で短くできる。
    waits = [(a["done"] - a["at"]).total_seconds() for a in asked
             if a["done"] and a["at"] and (a["done"] - a["at"]).total_seconds() >= 0]
    if waits:
        waits.sort()
        mid = waits[len(waits) // 2]
        slow = sum(1 for w in waits if w >= 30)
        notes.append(f"回答までの時間は中央値 {mid:.0f} 秒、最長 {waits[-1]:.0f} 秒。"
                     + (f"30秒以上待った質問が {slow} 件あります。" if slow else ""))
    if len(daily) >= 2:
        days_sorted = sorted(daily)
        span = (days_sorted[-1] - days_sorted[0]).days + 1
        notes.append(f"{span} 日のうち {len(daily)} 日に利用がありました"
                     f"（{len(daily) / span * 100:.0f}%）。")
        half = len(days_sorted) // 2
        first = sum(daily[d] for d in days_sorted[:half])
        last = sum(daily[d] for d in days_sorted[half:])
        if first:
            notes.append(f"前半 {first} 件 → 後半 {last} 件（{(last - first) / first * 100:+.0f}%）。"
                         + ("使われ方が伸びています。" if last > first else
                            "落ちています。失敗の内訳（errors）と合わせて見てください。"))
    if hour_rows:
        peak = max(hour_rows, key=lambda t: t[1])
        notes.append(f"最も使われる時間帯は {peak[0]}（{peak[1]} 件）。")

    return _usage_out("利用の推移",
                [_usage_table("日ごと", ["日付", "質問"], day_rows),
                 _usage_table("曜日", ["曜日", "質問"], dow_rows),
                 _usage_table("時間帯", ["時間", "質問"], hour_rows)],
                notes, {"active_days": len(daily), "questions": len(asked),
                        "median_wait_sec": round(mid) if waits else None})


def by_tool(records: list[dict]) -> dict:
    """呼ばれた機能。使われていない機能は、説明文が悪いか、要らないかのどちらか。"""
    if not records:
        return _empty(None)

    calls = Counter(t for r in records for t in r["tools"])
    total = sum(calls.values())
    rows = [(name, n, f"{n / total * 100:.1f}%") for name, n in calls.most_common()]

    notes = [_period_note(records)]
    if total:
        notes.append(f"ツール呼び出しは合計 {total:,} 回。"
                     f"種類は {len(calls)} 種です。")
        top = calls.most_common(3)
        notes.append("よく使われるのは " +
                     "、".join(f"{n}（{c}回）" for n, c in top) + "。")
    try:
        import tools as _tools
        unused = sorted(set(_tools._HANDLERS) - set(calls))
        if unused:
            notes.append(f"一度も呼ばれていない組み込みツールが {len(unused)} 種あります"
                         f"（{'、'.join(unused[:8])}"
                         f"{' ほか' if len(unused) > 8 else ''}）。"
                         "要らないなら「ツール」タブで無効にすると、AIの選択肢が減って"
                         "呼び分けが安定します。使ってほしいなら説明文を具体的に書き直します。")
    except Exception:
        pass
    return _usage_out("呼ばれた機能", [_usage_table("ツール別", ["ツール", "回数", "割合"], rows)],
                notes, {"total_calls": total, "kinds": len(calls)})


def by_database(records: list[dict]) -> dict:
    """DB別。会話の db_names は「実際にSQLが触ったDB」（選択UIは無い）。"""
    if not records:
        return _empty(None)

    per = Counter(name for r in records for name in set(r["dbs"]))
    rows = [(name, n, f"{n / len(records) * 100:.0f}%") for name, n in per.most_common()]
    widths = Counter(len(set(r["dbs"])) for r in records)
    width_rows = [(f"{k} DB", v) for k, v in sorted(widths.items())]

    notes = [_period_note(records),
             "実際にSQLが触ったDBで数えています（この仕組みを入れる前の会話は、"
             "当時選択されていたDBで数えます）。"]
    multi = sum(v for k, v in widths.items() if k >= 2)
    if records:
        notes.append(f"2つ以上のDBを使った会話は {multi}/{len(records)} 件"
                     f"（{multi / len(records) * 100:.0f}%）。"
                     + ("DBをまたぐ分析が実際に行われています。" if multi else
                        "横断分析がまだ使われていません。"
                        "またぎの結合定義と例文を足すと使われやすくなります。"))
    return _usage_out("対象データ別の利用",
                [_usage_table("DB別", ["DB", "会話数", "割合"], rows),
                 _usage_table("1会話で使ったDBの数", ["DB数", "会話"], width_rows)],
                notes, {"dbs": len(per)})


def errors(records: list[dict]) -> dict:
    """失敗の内訳。カタログを直して減るものと、そうでないものを分ける。"""
    if not records:
        return _empty(None)

    # 失敗は「どの質問で起きたか」まで対応づける。時刻と質問文が揃っていないと、
    # カタログの何を直せばよいかを後から辿れない。
    items = [(r["user"], a["at"], e, a["text"])
             for r in records for a in r["asked"] for e in a["errors"]]
    # 質問より前に出た失敗（会話を開いた直後など）。数を合わせるために拾っておく。
    items += [(r["user"], r["created"], e, "")
              for r in records for e in r["orphan_errors"]]
    items.sort(key=lambda t: t[1] or datetime.min)
    if not items:
        return _usage_out("失敗の内訳", [],
                    [_period_note(records),
                     "記録された失敗はありません。"], {"errors": 0})

    kinds: Counter = Counter()
    fixes: dict[str, str] = {}
    for _, _, msg, _ in items:
        kind, fix = classify_error(msg)
        kinds[kind] += 1
        fixes.setdefault(kind, fix)
    kind_rows = [(name, n, f"{n / len(items) * 100:.0f}%", fixes.get(name, ""))
                 for name, n in kinds.most_common()]

    detail_rows = [(f"{dt:%m-%d %H:%M}" if dt else "—", who, (q or "")[:40],
                    msg.splitlines()[0][:80])
                   for who, dt, msg, q in items[-MAX_LIST:]]

    notes = [_period_note(records),
             f"失敗 {len(items)} 件。分類は「誰が直せるか」で分けています。"]
    # 「カタログ画面で直せるもの」だけを足す。打ち手はいちばん多い分類のものを出す
    # （合計だけ言われても、何から手を付ければよいか分からないため）。
    fixable = {k: n for k, n in kinds.items()
               if k.startswith(("カタログ", "SQL", "ツール"))}
    if fixable:
        catalog_side = sum(fixable.values())
        top = max(fixable.items(), key=lambda kv: kv[1])
        notes.append(f"うち {catalog_side} 件（{catalog_side / len(items) * 100:.0f}%）は"
                     "カタログ画面での手当てで減らせます。"
                     f"いちばん多いのは「{top[0]}」{top[1]} 件で、{fixes.get(top[0], '')}")
    outside = kinds.get("LLM・API側の問題", 0)
    if outside:
        notes.append(f"{outside} 件はモデル・API側の問題で、カタログを直しても減りません。")
    return _usage_out("失敗の内訳",
                [_usage_table("分類", ["分類", "件数", "割合", "打ち手"], kind_rows),
                 _usage_table(f"直近の失敗（最大{MAX_LIST}件）",
                        ["日付", "利用者", "質問", "内容"], detail_rows)],
                notes, {"errors": len(items), "catalog_fixable": catalog_side})


def questions(records: list[dict]) -> dict:
    """実際に聞かれた質問。例文とカタログに反映するための材料。"""
    if not records:
        return _empty(None)

    asked = [(a["at"], r["user"], a["text"], a["failed"])
             for r in records for a in r["asked"] if a["text"]]
    if not asked:
        return _usage_out("聞かれた質問", [], [_period_note(records), "質問の記録がありません。"])
    asked.sort(key=lambda t: t[0] or datetime.min)

    rows = [(f"{dt:%m-%d %H:%M}" if dt else "—", who, q[:60], "×" if bad else "")
            for dt, who, q, bad in asked[-MAX_LIST:]]

    # 何を聞かれがちかを、語で大づかみに見る（形態素解析は入れない。傾向が分かれば足りる）
    words = Counter()
    for _, _, q, _ in asked:
        for w in re.findall(r"[一-龥ぁ-んァ-ヶa-zA-Z0-9]{2,}", q):
            if w not in ("教えて", "ください", "して", "この", "その", "どの", "です"):
                words[w] += 1
    word_rows = [(w, n) for w, n in words.most_common(20) if n > 1]

    notes = [_period_note(records),
             f"質問 {len(asked)} 件を記録しています。",
             "うまく答えられた質問は、チャットの⭐から例文としてカタログに登録できます。"
             "例文が増えるほど、同じ聞き方への精度が上がります。"]
    failed = sum(1 for *_, bad in asked if bad)
    if failed:
        notes.append(f"失敗を含む会話の質問が {failed} 件（× 印）。"
                     "この質問文がそのまま、カタログに足りない語彙の一覧になります。")
    return _usage_out("聞かれた質問",
                [_usage_table(f"直近の質問（最大{MAX_LIST}件）",
                        ["日付", "利用者", "質問", "失敗"], rows),
                 _usage_table("よく出る語", ["語", "回数"], word_rows)],
                notes, {"questions": len(asked)})


def imports(days: int | None = None) -> dict:
    """取り込みの実績。チャットとは別系統なので、まとめてここから見えるようにする。"""
    recs = history.recent_import_records(limit=2000)
    if not recs:
        return _usage_out("取り込みの実績", [], ["取り込みの記録がありません。"], {"runs": 0})

    limit = datetime.now() - timedelta(days=days) if days else None
    picked = []
    for r in recs:
        at = _usage_dt(r.get("at") or r.get("started"))
        if limit and at and at < limit:
            continue
        picked.append((at, r))
    if not picked:
        return _usage_out("取り込みの実績", [], [f"直近{days}日に取り込みの記録がありません。"],
                    {"runs": 0})

    ok = sum(1 for _, r in picked if r.get("ok"))
    kinds = Counter(history.IMPORT_RECORD_KINDS.get(str(r.get("kind")), str(r.get("kind")))
                    for _, r in picked)
    rows = [(k, n) for k, n in kinds.most_common()]
    tables = Counter(f"{r.get('db_file')} / {r.get('table')}" for _, r in picked)
    tbl_rows = [(name, n) for name, n in tables.most_common(20)]

    notes = [f"取り込み {len(picked)} 回。成功 {ok} 件 / 失敗 {len(picked) - ok} 件"
             f"（成功率 {ok / len(picked) * 100:.0f}%）。"]
    fails = [r for _, r in picked if not r.get("ok")]
    if fails:
        notes.append("直近の失敗: " + str(fails[-1].get("message", ""))[:100])
    return _usage_out("取り込みの実績",
                [_usage_table("実行のしかた別", ["種別", "回数"], rows),
                 _usage_table("テーブル別", ["DB / テーブル", "回数"], tbl_rows)],
                notes, {"runs": len(picked), "ok": ok})


#: メソッド名 -> 実処理。records を取らない imports だけ形が違う。
_METHOD_FUNCS = {
    "summary": summary,
    "users": by_user,
    "trend": trend,
    "tools": by_tool,
    "databases": by_database,
    "errors": errors,
    "questions": questions,
}


def analyze(method: str = "summary", days: int | None = None,
            user: str | None = None) -> dict:
    """入口。method はこのモジュールの METHODS のいずれか。"""
    method = (method or "summary").strip().lower()
    if method == "imports":
        return imports(days)
    fn = _METHOD_FUNCS.get(method)
    if fn is None:
        raise ValueError(f"method は {'、'.join([*METHODS, 'imports'])} "
                         f"のいずれかです（受け取った値: {method}）")
    records = collect(days=days, user=user)
    res = fn(records, days) if fn is summary else fn(records)
    if user:
        res["notes"] = [f"対象: {user} のみ", *res.get("notes", [])]
    return res


# ==========================================================================
# ===== 元 exports.py
# ダウンロードさせるファイル（CSV / テキスト / ZIP）の組み立てと、自動保存用HTML。
#
# xlsx の組み立ては excel.py。いずれもディスクには書かず、メモリ上のバイト列を返す。
# ==========================================================================
import csv
import datetime as _dt
import io
import re
import zipfile

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
CSV_MIME = "text/csv"
TEXT_MIME = "text/plain"
MD_MIME = "text/markdown"
ZIP_MIME = "application/zip"

# CSVの文字コード。Excelでそのまま開ける utf-8-sig を既定にする。
ENCODINGS = {
    "utf-8-sig": "UTF-8（BOM付き／Excelで文字化けしない・推奨）",
    "utf-8": "UTF-8（BOMなし）",
    "cp932": "Shift_JIS（cp932／古いWindows向け）",
}
DEFAULT_ENCODING = "utf-8-sig"
EXPORT_DELIMITERS = {"comma": ",", "tab": "\t", "semicolon": ";"}


def safe_filename(name: str | None, ext: str, default: str = "export") -> str:
    """ダウンロード用のファイル名を整える（末尾に日時、指定の拡張子）。"""
    s = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", str(name or "")).strip().strip(".")
    s = re.sub(r"\.(xlsx|csv|txt|md|zip)$", "", s, flags=re.IGNORECASE) or default
    stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M")
    return f"{s[:80]}_{stamp}.{ext.lstrip('.')}"


def _encode(text: str, encoding: str) -> bytes:
    enc = encoding if encoding in ENCODINGS else DEFAULT_ENCODING
    # Shift_JIS に無い文字（絵文字や一部の漢字）で落ちないよう置換する
    return text.encode(enc, errors="replace" if enc == "cp932" else "strict")


def build_csv(columns: list, rows: list, encoding: str = DEFAULT_ENCODING,
              delimiter: str = "comma") -> bytes:
    """1つの結果セットを CSV のバイト列にする。"""
    buf = io.StringIO()
    w = csv.writer(buf, delimiter=EXPORT_DELIMITERS.get(delimiter, ","),
                   lineterminator="\r\n", quoting=csv.QUOTE_MINIMAL)
    w.writerow([str(c) for c in columns])
    for r in rows:
        w.writerow(["" if v is None else v for v in r])
    return _encode(buf.getvalue(), encoding)


def build_zip(files: list[dict]) -> bytes:
    """[{"filename": str, "data": bytes}, ...] を1つのZIPにまとめる。"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        used = set()
        for f in files:
            name = str(f.get("filename") or "file")
            base, i = name, 2
            while name.lower() in used:
                stem, _, ext = base.rpartition(".")
                name = f"{stem}_{i}.{ext}" if stem else f"{base}_{i}"
                i += 1
            used.add(name.lower())
            z.writestr(name, f["data"])
    return buf.getvalue()


def table_to_text(columns: list, rows: list, style: str = "markdown") -> str:
    """結果セットを本文に埋め込めるテキスト表にする。"""
    cols = [str(c) for c in columns]
    body = [["" if v is None else str(v) for v in r] for r in rows]
    if style == "markdown":
        out = ["| " + " | ".join(cols) + " |",
               "| " + " | ".join("---" for _ in cols) + " |"]
        out += ["| " + " | ".join(r) + " |" for r in body]
        return "\n".join(out)
    if style == "tsv":
        return "\n".join(["\t".join(cols)] + ["\t".join(r) for r in body])
    # 等幅（プレーンテキスト用に桁を揃える）
    widths = [max(len(cols[i]), *(len(r[i]) for r in body)) if body else len(cols[i])
              for i in range(len(cols))]
    line = "-+-".join("-" * w for w in widths)
    out = [" | ".join(c.ljust(widths[i]) for i, c in enumerate(cols)), line]
    out += [" | ".join(v.ljust(widths[i]) for i, v in enumerate(r)) for r in body]
    return "\n".join(out)


def build_text(body: str, encoding: str = DEFAULT_ENCODING) -> bytes:
    return _encode(str(body or ""), encoding)




# ==========================================================================
# ===== 元 excel.py
# SELECT結果から Excel ブック(.xlsx)を組み立てる。
#
# ファイルはディスクに書かず、メモリ上のバイト列として返す。
#
# グラフはExcelネイティブのグラフとして入れる（画像ではない）。
# 受け取った側が範囲や種類を変えられるうえ、画像化ライブラリ（Chrome等）が
# 要らないので、サーバの環境に左右されない。
# ==========================================================================
import datetime as _dt  # noqa: F401  （シート値の型判定で使用）
import io
import re

from exports import XLSX_MIME  # noqa: F401  （既存の参照互換のため再公開）

from openpyxl import Workbook
from openpyxl.chart import (AreaChart, BarChart, LineChart, PieChart, Reference,
                            ScatterChart, Series)
from openpyxl.chart.label import DataLabelList
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

# Excelのシート名に使えない文字と長さ制限
_BAD_SHEET_CHARS = re.compile(r"[\[\]:*?/\\]")
_SHEET_NAME_MAX = 31
_MAX_WIDTH = 60          # 列幅の上限（文字数）
_HEADER_FILL = PatternFill("solid", fgColor="1F3B5C")
_HEADER_FONT = Font(bold=True, color="FFFFFF")
_BAND_FILL = PatternFill("solid", fgColor="F5F8FC")
_THIN = Side(style="thin", color="D5DBE2")
_BORDER = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)

# グラフの種類 -> (openpyxlのクラス, 積み上げ方)
EXCEL_CHART_TYPES = {
    "bar": (BarChart, "col", None),
    "bar_stacked": (BarChart, "col", "stacked"),
    "bar_percent": (BarChart, "col", "percentStacked"),
    "hbar": (BarChart, "bar", None),
    "hbar_stacked": (BarChart, "bar", "stacked"),
    "line": (LineChart, None, None),
    "line_stacked": (LineChart, None, "stacked"),
    "area": (AreaChart, None, None),
    "area_stacked": (AreaChart, None, "stacked"),
    "pie": (PieChart, None, None),
    "scatter": (ScatterChart, None, None),
}
_SERIES_COLORS = ["1F4E79", "F4B183", "70AD47", "C55A11", "7F7F7F",
                  "2E75B6", "A9D18E", "FFD966", "9DC3E6", "BFBFBF"]


def safe_sheet_name(name: str, used: set) -> str:
    """Excelの制約に合わせてシート名を整え、重複を避ける。"""
    s = _BAD_SHEET_CHARS.sub("_", str(name or "Sheet")).strip() or "Sheet"
    s = s[:_SHEET_NAME_MAX]
    base, i = s, 2
    while s.lower() in used:
        suffix = f"_{i}"
        s = base[: _SHEET_NAME_MAX - len(suffix)] + suffix
        i += 1
    used.add(s.lower())
    return s


def _cell_value(v):
    """openpyxl が扱えない型は文字列に落とす。"""
    if v is None or isinstance(v, (int, float, bool, str, _dt.datetime, _dt.date, _dt.time)):
        return v
    if isinstance(v, bytes):
        return f"<BLOB {len(v)} bytes>"
    return str(v)


def _autosize(ws, columns: list, rows: list):
    """見出しと先頭200行から列幅を決める。"""
    for ci, col in enumerate(columns, start=1):
        width = len(str(col))
        for r in rows[:200]:
            v = r[ci - 1] if ci - 1 < len(r) else None
            if v is not None:
                width = max(width, len(str(v)))
        ws.column_dimensions[get_column_letter(ci)].width = min(width + 2, _MAX_WIDTH)


def _add_chart(ws, spec: dict, columns: list, rows: list, header_row: int):
    """シートのデータ範囲からExcelネイティブのグラフを作って貼る。

    spec: {"type": 種類, "category_column": 横軸の列名, "value_columns": [系列の列名],
           "title": 見出し, "y_title": .., "x_title": .., "anchor": "H2",
           "data_labels": bool, "width": cm, "height": cm}
    """
    kind = str(spec.get("type") or "bar").lower()
    if kind not in EXCEL_CHART_TYPES:
        raise ValueError(f"未対応のグラフ種類です: {kind}。"
                         f"使えるのは {', '.join(EXCEL_CHART_TYPES)} です。")
    if not rows:
        raise ValueError("グラフにできる行がありません。")

    cat = spec.get("category_column") or (columns[0] if columns else None)
    if cat not in columns:
        raise ValueError(f"横軸の列 '{cat}' がありません。ある列: {', '.join(map(str, columns))}")
    vals = spec.get("value_columns") or [c for c in columns if c != cat]
    missing = [v for v in vals if v not in columns]
    if missing:
        raise ValueError(f"系列の列 {', '.join(map(str, missing))} がありません。"
                         f"ある列: {', '.join(map(str, columns))}")
    if not vals:
        raise ValueError("系列にする数値列がありません。")

    cls, direction, grouping = EXCEL_CHART_TYPES[kind]
    chart = cls()
    chart.title = spec.get("title") or None
    chart.style = 2
    if direction:
        chart.type = direction
    if grouping:
        chart.grouping = grouping
        chart.overlap = 100
    last = header_row + len(rows)
    cat_ref = Reference(ws, min_col=columns.index(cat) + 1, min_row=header_row + 1,
                        max_row=last)

    if kind == "scatter":
        # 散布図は x も数値列。1列目を x、残りを y にする。
        x_ref = Reference(ws, min_col=columns.index(vals[0]) + 1,
                          min_row=header_row + 1, max_row=last)
        for v in vals[1:] or vals[:1]:
            y_ref = Reference(ws, min_col=columns.index(v) + 1, min_row=header_row,
                              max_row=last)
            s = Series(y_ref, x_ref, title_from_data=True)
            s.marker.symbol = "circle"
            s.graphicalProperties.line.noFill = True
            chart.series.append(s)
    else:
        for i, v in enumerate(vals):
            ref = Reference(ws, min_col=columns.index(v) + 1, min_row=header_row,
                            max_row=last)
            chart.add_data(ref, titles_from_data=True)
        chart.set_categories(cat_ref)
        for i, s in enumerate(chart.series):
            color = _SERIES_COLORS[i % len(_SERIES_COLORS)]
            try:
                if kind.startswith("line"):
                    s.graphicalProperties.line.solidFill = color
                    s.smooth = False
                else:
                    s.graphicalProperties.solidFill = color
                    s.graphicalProperties.line.solidFill = color
            except AttributeError:
                pass

    if spec.get("data_labels") or (kind == "pie" and spec.get("data_labels") is not False):
        chart.dataLabels = DataLabelList()
        chart.dataLabels.showVal = kind != "pie"
        chart.dataLabels.showPercent = kind == "pie"
    if kind not in ("pie",):
        chart.y_axis.title = spec.get("y_title") or None
        chart.x_axis.title = spec.get("x_title") or None
        chart.y_axis.numFmt = spec.get("number_format") or "#,##0"
        chart.x_axis.delete = False      # これが無いとExcelで軸が消えることがある
        chart.y_axis.delete = False
    chart.width = float(spec.get("width") or 20)     # cm
    chart.height = float(spec.get("height") or 10)
    chart.legend.position = "b"
    if len(vals) <= 1 and kind != "pie":
        chart.legend = None

    anchor = spec.get("anchor") or f"{get_column_letter(len(columns) + 2)}{header_row}"
    ws.add_chart(chart, anchor)
    return chart


def build_excel(sheets: list[dict], title: str | None = None) -> bytes:
    """[{"name", "columns", "rows", "note"?, "charts"?}, ...] から xlsx を作る。

    charts は同じシートのデータから作るグラフの指定（複数可）。
    """
    if not sheets:
        raise ValueError("シートが1つもありません。")
    wb = Workbook()
    wb.remove(wb.active)
    used: set = set()

    for sh in sheets:
        columns = list(sh.get("columns") or [])
        rows = list(sh.get("rows") or [])
        ws = wb.create_sheet(safe_sheet_name(sh.get("name"), used))

        start = 1
        note = str(sh.get("note") or "").strip()
        if note:
            ws.cell(row=1, column=1, value=note).font = Font(italic=True, color="666666")
            start = 3

        for ci, col in enumerate(columns, start=1):
            c = ws.cell(row=start, column=ci, value=str(col))
            c.font = _HEADER_FONT
            c.fill = _HEADER_FILL
            c.alignment = Alignment(vertical="center", horizontal="center")
            c.border = _BORDER
        for ri, row in enumerate(rows, start=start + 1):
            banded = (ri - start) % 2 == 0
            for ci in range(1, len(columns) + 1):
                cell = ws.cell(row=ri, column=ci,
                               value=_cell_value(row[ci - 1] if ci - 1 < len(row) else None))
                cell.border = _BORDER
                if banded:
                    cell.fill = _BAND_FILL
                if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool):
                    cell.number_format = "#,##0.####"

        ws.freeze_panes = ws.cell(row=start + 1, column=1)
        if columns and rows:
            ws.auto_filter.ref = (f"A{start}:"
                                  f"{get_column_letter(len(columns))}{start + len(rows)}")
        _autosize(ws, columns, rows)

        charts = sh.get("charts")
        if isinstance(charts, dict):
            charts = [charts]
        for i, spec in enumerate(charts or []):
            spec = dict(spec or {})
            spec.setdefault("anchor",
                            f"{get_column_letter(len(columns) + 2)}"
                            f"{start + i * 21}")
            _add_chart(ws, spec, columns, rows, start)

    if title:
        wb.properties.title = str(title)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def excel_safe_filename(name: str | None, default: str = "export") -> str:
    """互換用。実体は exports.safe_filename（拡張子 .xlsx）。"""
    from exports import safe_filename as _sf
    return _sf(name, "xlsx", default)


# ==========================================================================
# ===== 元 charts.py
# SELECT結果からグラフを組み立てる。
#
# チャット画面とユーザー定義ツールの両方がここを通るので、
# 対応するグラフ種別を増やすときはこのファイルだけを直せばよい。
#
# 種別の追加手順:
#   1. CHART_SPECS に (説明, 必要な指定, 分類) を足す
#   2. _BUILDERS に組み立て関数を足す
# validate() と画面の説明文は CHART_SPECS から自動で作られる。
# ==========================================================================
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

# 列の指定をリストで受け取るもの（存在チェックの仕方が違う）
LIST_FIELDS = ("path", "dimensions")

# 種別 -> (日本語の説明, 必要な指定, 分類)
CHART_SPECS: dict[str, tuple[str, tuple, str]] = {
    # --- 比較 ---------------------------------------------------------------
    "bar":         ("棒。カテゴリ別の比較", ("x", "y"), "比較"),
    "hbar":        ("横棒。項目名が長いときや順位表に", ("x", "y"), "比較"),
    "stacked_bar": ("積み上げ棒。内訳つきの比較", ("x", "y"), "比較"),
    "percent_bar": ("100%積み上げ棒。構成比の比較", ("x", "y"), "比較"),
    "lollipop":    ("ロリポップ。棒より軽く順位を見せる", ("x", "y"), "比較"),
    "dumbbell":    ("ダンベル。2時点の差を1行で比べる", ("x", "y", "y2"), "比較"),
    "pareto":      ("パレート図。棒＋累積比率で重点を見つける", ("x", "y"), "比較"),
    "pyramid":     ("人口ピラミッド。左右に分けた横棒", ("x", "y", "color"), "比較"),
    "marimekko":   ("マリメッコ。幅も高さも意味を持つ積み上げ", ("x", "y", "size"), "比較"),
    "radar":       ("レーダー。複数指標のバランス", ("x", "y"), "比較"),
    "polar_bar":   ("極座標の棒。方位や時間帯の分布", ("x", "y"), "比較"),
    "bump":        ("バンプ。順位の入れ替わりを追う", ("x", "y", "color"), "比較"),
    # --- 推移 ---------------------------------------------------------------
    "line":        ("折れ線。時系列や推移", ("x", "y"), "推移"),
    "step":        ("階段。在庫や料金など段階的に変わる値", ("x", "y"), "推移"),
    "area":        ("面。積み上げの推移", ("x", "y"), "推移"),
    "area_percent": ("100%面。構成比の推移", ("x", "y"), "推移"),
    "range_area":  ("幅つき折れ線。予測の上下限や信頼区間", ("x", "y", "lower", "upper"), "推移"),
    "slope":       ("スロープ。2時点の順位・水準の変化", ("x", "y", "color"), "推移"),
    "candlestick": ("ローソク足。始値・高値・安値・終値", ("x", "open", "high", "low", "close"), "推移"),
    "ohlc":        ("OHLC。ローソク足の棒型", ("x", "open", "high", "low", "close"), "推移"),
    "gantt":       ("ガントチャート。作業や期間の並び", ("x", "start", "end"), "推移"),
    "calendar":    ("カレンダーヒートマップ。日ごとの多寡", ("x", "y"), "推移"),
    "control_chart": ("管理図。平均±3σを外れた点を見つける", ("x", "y"), "推移"),
    # --- 構成 ---------------------------------------------------------------
    "pie":         ("円。構成比", ("x", "y"), "構成"),
    "donut":       ("ドーナツ。構成比（中央に合計）", ("x", "y"), "構成"),
    "treemap":     ("ツリーマップ。階層つき構成比", ("path", "y"), "構成"),
    "sunburst":    ("サンバースト。階層つき構成比（円形）", ("path", "y"), "構成"),
    "icicle":      ("アイシクル。階層を短冊で並べる", ("path", "y"), "構成"),
    "funnel":      ("ファネル。段階ごとの減少", ("x", "y"), "構成"),
    "waterfall":   ("ウォーターフォール。増減の内訳", ("x", "y"), "構成"),
    "sankey":      ("サンキー。流れと量（どこからどこへ）", ("source", "target", "y"), "構成"),
    # --- 分布 ---------------------------------------------------------------
    "histogram":   ("ヒストグラム。1つの数値の分布", ("x",), "分布"),
    "density":     ("密度曲線。ヒストグラムをなめらかに", ("x",), "分布"),
    "ecdf":        ("累積分布。「〇〇以下が何%か」を読む", ("x",), "分布"),
    "box":         ("箱ひげ。カテゴリ別のばらつき", ("y",), "分布"),
    "violin":      ("バイオリン。分布の形まで見る", ("y",), "分布"),
    "strip":       ("ストリップ。個々の点を並べる", ("y",), "分布"),
    "ridgeline":   ("リッジライン。群ごとの分布を重ねる", ("x", "color"), "分布"),
    "qq":          ("Q-Qプロット。正規分布からのズレ", ("x",), "分布"),
    # --- 関係 ---------------------------------------------------------------
    "scatter":     ("散布。2つの数値の相関", ("x", "y"), "関係"),
    "bubble":      ("バブル。散布＋大きさで3指標", ("x", "y", "size"), "関係"),
    "histogram2d": ("2次元ヒストグラム。点が多すぎるときの散布図", ("x", "y"), "関係"),
    "contour":     ("等高線。2変数の密度", ("x", "y"), "関係"),
    "heatmap":     ("ヒートマップ。2軸の集計をマス目で", ("x", "y"), "関係"),
    "matrix":      ("行列ヒートマップ。集計済みのクロス表や相関行列をそのまま色で", (), "関係"),
    "scatter_matrix": ("散布図行列。数値列を総当たりで見る", ("dimensions",), "関係"),
    "parallel_coordinates": ("平行座標。多変量の傾向を線で追う", ("dimensions",), "関係"),
    "parallel_categories": ("平行カテゴリ。区分の組み合わせの多さ", ("dimensions",), "関係"),
    "scatter3d":   ("3D散布。3つの数値の関係", ("x", "y", "z"), "関係"),
    "surface":     ("3D曲面。集計済みのクロス表を立体で", (), "関係"),
    "network":     ("ネットワーク。つながりの図", ("source", "target"), "関係"),
    # --- 指標 ---------------------------------------------------------------
    "indicator":   ("数値の大写し。KPIを1つ見せる", ("value",), "指標"),
    "gauge":       ("ゲージ。目標に対する達成度", ("value",), "指標"),
    "bullet":      ("ブレット。実績と目標を並べる", ("value",), "指標"),
}

CHART_TYPES = tuple(CHART_SPECS)


def type_help(category: str | None = None) -> str:
    """LLMに見せる一覧。分類を指定するとその分だけ返す。"""
    items = [(k, v) for k, v in CHART_SPECS.items()
             if category is None or v[2] == category]
    return " / ".join(f"{k}={v[0]}" for k, v in items)


def types_in(category: str) -> list[str]:
    return [k for k, v in CHART_SPECS.items() if v[2] == category]


def required_fields(chart_type: str) -> tuple:
    spec = CHART_SPECS.get(chart_type)
    return spec[1] if spec else ("x", "y")


def validate(item: dict, columns: list) -> list[str]:
    """指定された列が結果に存在するか検証し、問題点を返す。"""
    ct = item.get("chart_type") or "bar"
    errs = []
    if ct not in CHART_SPECS:
        return [f"未対応のグラフ種別です: {ct} / 使えるのは {', '.join(CHART_TYPES)}"]
    for f in required_fields(ct):
        v = item.get(f)
        if f in LIST_FIELDS:
            cols = list(v or [])
            if not cols:
                errs.append(f"{ct} には {f}（列名のリスト）が必要です。")
            errs += [f"{f} の列 '{c}' が結果にありません。利用可能: {columns}"
                     for c in cols if c not in columns]
        elif not v:
            errs.append(f"{ct} には {f} の指定が必要です。")
        elif v not in columns:
            errs.append(f"指定列 '{v}' が結果にありません。利用可能: {columns}")
    # 任意指定も、指定されていれば存在チェック
    for f in ("color", "size", "text", "y2", "z", "lower", "upper", "target", "facet"):
        v = item.get(f)
        if not v:
            continue
        # target だけは列名ではなく目標値（数値）で来ることがある
        if f == "target" and isinstance(v, (int, float)) and not isinstance(v, bool):
            continue
        if v not in columns:
            errs.append(f"指定列 '{v}' が結果にありません。利用可能: {columns}")
    return errs


# =============================================================================
# 下ごしらえ
# =============================================================================

def _scale(name):
    """色スケール名を色のリストに直す。

    名前のまま渡すと、周辺分布つきのグラフで plotly が文字列を1文字ずつ
    色として読み、"Blues" が 'B' 扱いになって落ちる。
    """
    return getattr(px.colors.sequential, str(name or "Blues"), None) or "Blues"


def _numeric(df: pd.DataFrame, *cols):
    for c in cols:
        if c and c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(df[c])
    return df


def _num_series(df: pd.DataFrame, col) -> pd.Series:
    """必ず数値のSeriesにする（欠損は落とさず NaN のまま）。"""
    return pd.to_numeric(df[col], errors="coerce")


class _Ctx:
    """組み立て関数に渡す、よく使う値の詰め合わせ。"""

    def __init__(self, item: dict):
        self.item = item
        self.df = pd.DataFrame(item["rows"], columns=item["columns"])
        self.x, self.y = item.get("x"), item.get("y")
        self.title = item.get("title", "")
        for f in ("color", "size", "text", "y2", "z", "lower", "upper",
                  "target", "facet", "source", "start", "end",
                  "open", "high", "low", "close", "value"):
            setattr(self, f, item.get(f) if item.get(f) in self.df.columns else None)
        # target は列名でなく数値で来ることもある（目標値）
        self.target_value = item.get("target")
        self.path = [c for c in (item.get("path") or []) if c in self.df.columns]
        self.dimensions = [c for c in (item.get("dimensions") or []) if c in self.df.columns]
        _numeric(self.df, self.y, self.size, self.y2, self.z,
                 self.lower, self.upper, self.open, self.high, self.low, self.close)

    def get(self, key, default=None):
        return self.item.get(key, default)


# =============================================================================
# 比較
# =============================================================================

def _bar(c, orientation=None, barmode=None):
    return px.bar(c.df, x=c.x, y=c.y, color=c.color, text=c.text, title=c.title,
                  barmode=barmode or c.get("barmode") or "group",
                  orientation=orientation or c.get("orientation") or "v",
                  facet_col=c.facet)


def _hbar(c):
    # 横棒は「値が大きいものを上」に。並べ替えないと読みにくい
    d = c.df.sort_values(c.y) if c.y in c.df.columns else c.df
    return px.bar(d, x=c.y, y=c.x, color=c.color, text=c.text, title=c.title,
                  orientation="h", barmode=c.get("barmode") or "group")


def _stacked_bar(c):
    return _bar(c, barmode="stack")


def _percent_bar(c):
    d = c.df.copy()
    total = d.groupby(c.x)[c.y].transform("sum")
    d["_割合"] = _num_series(d, c.y) / total.replace(0, np.nan) * 100
    fig = px.bar(d, x=c.x, y="_割合", color=c.color, title=c.title, barmode="stack",
                 text=d["_割合"].round(1).astype(str) + "%")
    fig.update_yaxes(title_text="構成比(%)", range=[0, 100])
    return fig


def _lollipop(c):
    d = c.df.sort_values(c.y)
    fig = go.Figure()
    for _, r in d.iterrows():
        fig.add_shape(type="line", x0=0, x1=r[c.y], y0=r[c.x], y1=r[c.x],
                      line=dict(color="#9DC3E6", width=2))
    fig.add_trace(go.Scatter(x=d[c.y], y=d[c.x].astype(str), mode="markers",
                             marker=dict(size=12, color="#1F4E79"), name=c.y))
    fig.update_layout(title=c.title, xaxis_title=c.y, yaxis_title=c.x)
    return fig


def _dumbbell(c):
    d = c.df
    fig = go.Figure()
    for _, r in d.iterrows():
        fig.add_shape(type="line", x0=r[c.y], x1=r[c.y2], y0=r[c.x], y1=r[c.x],
                      line=dict(color="#BFBFBF", width=3))
    fig.add_trace(go.Scatter(x=d[c.y], y=d[c.x].astype(str), mode="markers",
                             name=str(c.y), marker=dict(size=12, color="#9DC3E6")))
    fig.add_trace(go.Scatter(x=d[c.y2], y=d[c.x].astype(str), mode="markers",
                             name=str(c.y2), marker=dict(size=12, color="#1F4E79")))
    fig.update_layout(title=c.title, xaxis_title="値", yaxis_title=c.x)
    return fig


def _pareto(c):
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    d = d.dropna(subset=[c.y]).sort_values(c.y, ascending=False)
    total = d[c.y].sum() or 1
    d["_累積"] = d[c.y].cumsum() / total * 100
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(go.Bar(x=d[c.x].astype(str), y=d[c.y], name=str(c.y),
                         marker_color="#2E75B6"), secondary_y=False)
    fig.add_trace(go.Scatter(x=d[c.x].astype(str), y=d["_累積"], name="累積構成比",
                             mode="lines+markers", line=dict(color="#C55A11")),
                  secondary_y=True)
    fig.add_hline(y=80, line_dash="dot", line_color="#C55A11", secondary_y=True,
                  annotation_text="80%")
    fig.update_yaxes(title_text=str(c.y), secondary_y=False)
    fig.update_yaxes(title_text="累積構成比(%)", range=[0, 105], secondary_y=True)
    fig.update_layout(title=c.title)
    return fig


def _pyramid(c):
    """人口ピラミッド。color の2種類を左右に振り分ける。"""
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    groups = list(pd.unique(d[c.color].dropna()))[:2]
    if len(groups) < 2:
        raise ValueError(f"人口ピラミッドには color 列に2種類の値が必要です"
                         f"（いま: {groups}）。")
    left, right = groups
    fig = go.Figure()
    dl, dr = d[d[c.color] == left], d[d[c.color] == right]
    fig.add_trace(go.Bar(y=dl[c.x].astype(str), x=-dl[c.y], name=str(left),
                         orientation="h", marker_color="#2E75B6"))
    fig.add_trace(go.Bar(y=dr[c.x].astype(str), x=dr[c.y], name=str(right),
                         orientation="h", marker_color="#F4B183"))
    fig.update_layout(title=c.title, barmode="overlay", bargap=0.1,
                      xaxis=dict(title=str(c.y),
                                 tickvals=None, ticktext=None))
    fig.update_xaxes(tickformat="~s")
    return fig


def _marimekko(c):
    """幅=size、高さ=y の積み上げ。x ごとの規模と内訳を同時に見せる。"""
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    d[c.size] = _num_series(d, c.size)
    widths = d.groupby(c.x, sort=False)[c.size].max()
    total_w = widths.sum() or 1
    fig = go.Figure()
    keys = list(widths.index)
    lefts, acc = {}, 0.0
    for k in keys:
        lefts[k] = acc
        acc += float(widths[k]) / total_w * 100
    groups = list(pd.unique(d[c.color].dropna())) if c.color else [None]
    for gi, g in enumerate(groups):
        sub = d if g is None else d[d[c.color] == g]
        xs, ys, ws = [], [], []
        for k in keys:
            row = sub[sub[c.x] == k]
            if row.empty:
                continue
            w = float(widths[k]) / total_w * 100
            xs.append(lefts[k] + w / 2)
            ws.append(w)
            ys.append(float(row[c.y].iloc[0]))
        fig.add_trace(go.Bar(x=xs, y=ys, width=ws, name=str(g) if g is not None else str(c.y),
                             marker_color=px.colors.qualitative.Set2[gi % 8]))
    fig.update_layout(title=c.title, barmode="stack", bargap=0,
                      xaxis_title=f"{c.x}（幅 = {c.size}）", yaxis_title=str(c.y))
    return fig


def _radar(c):
    fig = px.line_polar(c.df, r=c.y, theta=c.x, color=c.color, line_close=True,
                        title=c.title)
    fig.update_traces(fill="toself", opacity=0.5)
    return fig


def _polar_bar(c):
    return px.bar_polar(c.df, r=c.y, theta=c.x, color=c.color, title=c.title)


def _bump(c):
    """順位の推移。値が小さいほど上位なので、y軸を反転する。"""
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    fig = px.line(d, x=c.x, y=c.y, color=c.color, markers=True, title=c.title,
                  text=c.text)
    fig.update_traces(marker=dict(size=11))
    fig.update_yaxes(autorange="reversed", title_text=f"{c.y}（上が上位）",
                     dtick=1)
    return fig


# =============================================================================
# 推移
# =============================================================================

def _line(c):
    return px.line(c.df, x=c.x, y=c.y, color=c.color, text=c.text, title=c.title,
                   markers=True, facet_col=c.facet)


def _step(c):
    fig = px.line(c.df, x=c.x, y=c.y, color=c.color, title=c.title, markers=True)
    fig.update_traces(line_shape="hv")
    return fig


def _area(c):
    return px.area(c.df, x=c.x, y=c.y, color=c.color, title=c.title)


def _area_percent(c):
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    total = d.groupby(c.x)[c.y].transform("sum")
    d["_割合"] = d[c.y] / total.replace(0, np.nan) * 100
    fig = px.area(d, x=c.x, y="_割合", color=c.color, title=c.title)
    fig.update_yaxes(title_text="構成比(%)", range=[0, 100])
    return fig


def _range_area(c):
    d = c.df
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=d[c.x], y=d[c.upper], mode="lines", name="上限",
                             line=dict(width=0), showlegend=False))
    fig.add_trace(go.Scatter(x=d[c.x], y=d[c.lower], mode="lines", name="幅（95%）",
                             line=dict(width=0), fill="tonexty",
                             fillcolor="rgba(46,117,182,.18)"))
    fig.add_trace(go.Scatter(x=d[c.x], y=d[c.y], mode="lines+markers", name=str(c.y),
                             line=dict(color="#1F4E79", width=2)))
    fig.update_layout(title=c.title, xaxis_title=str(c.x), yaxis_title=str(c.y))
    return fig


def _slope(c):
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    fig = px.line(d, x=c.x, y=c.y, color=c.color, markers=True, title=c.title)
    fig.update_traces(line=dict(width=2))
    # 端に系列名を出す（凡例を目で追わなくて済む）
    first = str(d[c.x].iloc[0])
    for name, g in d.groupby(c.color):
        head = g[g[c.x].astype(str) == first]
        if len(head):
            fig.add_annotation(x=head[c.x].iloc[0], y=head[c.y].iloc[0], text=str(name),
                               xanchor="right", showarrow=False, xshift=-6, font_size=11)
    fig.update_layout(showlegend=False)
    return fig


def _candlestick(c):
    return go.Figure(go.Candlestick(
        x=c.df[c.x], open=c.df[c.open], high=c.df[c.high],
        low=c.df[c.low], close=c.df[c.close])).update_layout(
            title=c.title, xaxis_rangeslider_visible=False)


def _ohlc(c):
    return go.Figure(go.Ohlc(
        x=c.df[c.x], open=c.df[c.open], high=c.df[c.high],
        low=c.df[c.low], close=c.df[c.close])).update_layout(
            title=c.title, xaxis_rangeslider_visible=False)


def _gantt(c):
    fig = px.timeline(c.df, x_start=c.start, x_end=c.end, y=c.x, color=c.color,
                      text=c.text, title=c.title)
    fig.update_yaxes(autorange="reversed")     # 上から順に並べる
    return fig


def _calendar(c):
    """日付ごとの値を、週×曜日のマス目にする。"""
    d = c.df.copy()
    d[c.x] = pd.to_datetime(d[c.x], errors="coerce")
    d[c.y] = _num_series(d, c.y)
    d = d.dropna(subset=[c.x])
    if d.empty:
        raise ValueError(f"{c.x} を日付として読めませんでした。")
    d["_日付"] = d[c.x].dt.normalize()
    d = d.groupby("_日付", as_index=False)[c.y].sum().rename(columns={c.y: "_値"})
    d["_週"] = d["_日付"].dt.isocalendar().week.astype(int)
    d["_年"] = d["_日付"].dt.isocalendar().year.astype(int)
    d["_通週"] = (d["_年"] - d["_年"].min()) * 53 + d["_週"]
    names = ["月", "火", "水", "木", "金", "土", "日"]
    d["_曜日"] = d["_日付"].dt.weekday
    pivot = d.pivot_table(index="_曜日", columns="_通週", values="_値", aggfunc="sum")
    pivot = pivot.reindex(range(7))
    labels = (d.groupby("_通週")["_日付"].min().dt.strftime("%m/%d")
              .reindex(pivot.columns).tolist())
    fig = px.imshow(pivot.to_numpy(), x=labels, y=names, aspect="auto",
                    color_continuous_scale=_scale(c.get("colorscale")),
                    title=c.title, labels=dict(color=str(c.y)))
    fig.update_xaxes(title_text="週（週初の日付）", side="top")
    return fig


def _control_chart(c):
    """管理図。平均と±3σを引き、外れた点を赤くする。"""
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    d = d.dropna(subset=[c.y])
    m, sd = d[c.y].mean(), d[c.y].std(ddof=1)
    ucl, lcl = m + 3 * sd, m - 3 * sd
    out = (d[c.y] > ucl) | (d[c.y] < lcl)
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=d[c.x], y=d[c.y], mode="lines+markers", name=str(c.y),
                             line=dict(color="#2E75B6"),
                             marker=dict(size=8,
                                         color=np.where(out, "#B02A2A", "#2E75B6"))))
    for val, name, dash in ((m, "平均", "solid"), (ucl, "上方管理限界(+3σ)", "dash"),
                            (lcl, "下方管理限界(-3σ)", "dash")):
        fig.add_hline(y=val, line_dash=dash, line_color="#7F7F7F",
                      annotation_text=f"{name} {val:,.4g}", annotation_position="right")
    fig.update_layout(title=c.title or "管理図", xaxis_title=str(c.x),
                      yaxis_title=str(c.y))
    return fig


# =============================================================================
# 構成
# =============================================================================

def _pie(c):
    return px.pie(c.df, names=c.x, values=c.y, color=c.color, title=c.title,
                  hole=0.45 if c.get("chart_type") == "donut" else 0)


def _treemap(c):
    return px.treemap(c.df, path=c.path, values=c.y, color=c.color, title=c.title)


def _sunburst(c):
    return px.sunburst(c.df, path=c.path, values=c.y, color=c.color, title=c.title)


def _icicle(c):
    return px.icicle(c.df, path=c.path, values=c.y, color=c.color, title=c.title)


def _funnel(c):
    # plotly は x=値 / y=段階 なので入れ替える（x に段階、y に値を受け取る仕様）
    return px.funnel(c.df, x=c.y, y=c.x, color=c.color, title=c.title)


def _waterfall(c):
    d = c.df
    measure = ["relative"] * len(d)
    # 「合計」「total」で終わる行は合計として扱う
    for i, v in enumerate(d[c.x].astype(str)):
        if v.strip() in ("合計", "計", "total", "Total", "TOTAL"):
            measure[i] = "total"
    fig = go.Figure(go.Waterfall(
        x=d[c.x].astype(str), y=_num_series(d, c.y), measure=measure,
        text=d[c.y], textposition="outside"))
    fig.update_layout(title=c.title, waterfallgap=0.3)
    return fig


def _sankey(c):
    d = c.df.copy()
    d[c.y] = _num_series(d, c.y)
    labels = list(dict.fromkeys(d[c.source].astype(str).tolist()
                                + d[c.target].astype(str).tolist()))
    idx = {v: i for i, v in enumerate(labels)}
    fig = go.Figure(go.Sankey(
        node=dict(label=labels, pad=16, thickness=16,
                  line=dict(color="#BFBFBF", width=0.5)),
        link=dict(source=[idx[str(v)] for v in d[c.source]],
                  target=[idx[str(v)] for v in d[c.target]],
                  value=d[c.y].fillna(0).tolist())))
    fig.update_layout(title=c.title, font_size=12)
    return fig


# =============================================================================
# 分布
# =============================================================================

def _histogram(c):
    return px.histogram(c.df, x=c.x, color=c.color, title=c.title,
                        nbins=int(c.get("nbins")) if c.get("nbins") else None,
                        facet_col=c.facet, marginal=c.get("marginal"))


def _density(c):
    """ヒストグラム＋カーネル密度推定の曲線。"""
    from scipy import stats as sstats
    d = c.df.copy()
    d[c.x] = pd.to_numeric(d[c.x], errors="coerce")
    d = d.dropna(subset=[c.x])
    if len(d) < 3:
        raise ValueError("密度曲線には3行以上の数値が必要です。")
    fig = px.histogram(d, x=c.x, color=c.color, histnorm="probability density",
                       opacity=0.55, nbins=int(c.get("nbins") or 30), title=c.title)
    groups = [(None, d)] if not c.color else list(d.groupby(c.color))
    xs = np.linspace(d[c.x].min(), d[c.x].max(), 200)
    for name, g in groups:
        if g[c.x].nunique() < 2:
            continue
        try:
            kde = sstats.gaussian_kde(g[c.x].to_numpy())
        except np.linalg.LinAlgError:
            continue
        fig.add_trace(go.Scatter(x=xs, y=kde(xs), mode="lines",
                                 name=f"{name} 密度" if name is not None else "密度",
                                 line=dict(width=2)))
    fig.update_layout(bargap=0.02)
    return fig


def _ecdf(c):
    return px.ecdf(c.df, x=c.x, color=c.color, title=c.title, markers=False)


def _box(c):
    return px.box(c.df, x=c.x, y=c.y, color=c.color, title=c.title, points="outliers",
                  facet_col=c.facet)


def _violin(c):
    return px.violin(c.df, x=c.x, y=c.y, color=c.color, title=c.title, box=True,
                     points=False)


def _strip(c):
    return px.strip(c.df, x=c.x, y=c.y, color=c.color, title=c.title)


def _ridgeline(c):
    """群ごとの分布を少しずつずらして重ねる。"""
    d = c.df.copy()
    d[c.x] = pd.to_numeric(d[c.x], errors="coerce")
    d = d.dropna(subset=[c.x])
    fig = go.Figure()
    for name, g in d.groupby(c.color):
        fig.add_trace(go.Violin(x=g[c.x], name=str(name), side="positive",
                                width=2.2, points=False, meanline_visible=True,
                                orientation="h"))
    fig.update_layout(title=c.title, violingap=0, violinmode="overlay",
                      xaxis_title=str(c.x), showlegend=False)
    return fig


def _qq(c):
    """正規Q-Qプロット。点が直線に乗るほど正規分布に近い。"""
    from scipy import stats as sstats
    s = pd.to_numeric(c.df[c.x], errors="coerce").dropna().sort_values()
    if len(s) < 3:
        raise ValueError("Q-Qプロットには3行以上の数値が必要です。")
    theo = sstats.norm.ppf((np.arange(1, len(s) + 1) - 0.5) / len(s))
    theo = theo * s.std(ddof=1) + s.mean()
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=theo, y=s, mode="markers", name="実測",
                             marker=dict(color="#2E75B6", size=7)))
    lo, hi = float(min(theo.min(), s.min())), float(max(theo.max(), s.max()))
    fig.add_trace(go.Scatter(x=[lo, hi], y=[lo, hi], mode="lines", name="正規分布の直線",
                             line=dict(color="#C55A11", dash="dash")))
    fig.update_layout(title=c.title or f"{c.x} のQ-Qプロット",
                      xaxis_title="正規分布ならこうなる", yaxis_title="実測")
    return fig


# =============================================================================
# 関係
# =============================================================================

def _scatter(c):
    return px.scatter(c.df, x=c.x, y=c.y, color=c.color, text=c.text, title=c.title,
                      facet_col=c.facet,
                      trendline="ols" if c.get("trendline") else None)


def _bubble(c):
    return px.scatter(c.df, x=c.x, y=c.y, color=c.color, size=c.size, text=c.text,
                      title=c.title, size_max=50)


def _histogram2d(c):
    return px.density_heatmap(c.df, x=c.x, y=c.y, title=c.title,
                              nbinsx=int(c.get("nbins") or 30),
                              nbinsy=int(c.get("nbins") or 30),
                              color_continuous_scale=_scale(c.get("colorscale")),
                              marginal_x="histogram", marginal_y="histogram")


def _contour(c):
    fig = px.density_contour(c.df, x=c.x, y=c.y, color=c.color, title=c.title)
    fig.update_traces(contours_coloring="fill", contours_showlabels=True)
    return fig


def _heatmap(c):
    if c.color:
        return px.density_heatmap(c.df, x=c.x, y=c.y, z=c.color, histfunc="sum",
                                  title=c.title, text_auto=True,
                                  color_continuous_scale=_scale(c.get("colorscale")))
    return px.density_heatmap(c.df, x=c.x, y=c.y, histfunc="count", title=c.title,
                              text_auto=True,
                              color_continuous_scale=_scale(c.get("colorscale")))


def _matrix(c):
    """集計済みの表をそのまま行列として塗る。"""
    label = c.x if c.x in c.df.columns else c.df.columns[0]
    m = c.df.set_index(label)
    m = m.apply(lambda s: pd.to_numeric(s, errors="coerce")).dropna(axis=1, how="all")
    fig = px.imshow(m, text_auto=True, aspect="auto", title=c.title,
                    color_continuous_scale=_scale(c.get("colorscale")))
    fig.update_xaxes(side="top")
    return fig


def _scatter_matrix(c):
    d = c.df.copy()
    for col in c.dimensions:
        d[col] = pd.to_numeric(d[col], errors="coerce")
    fig = px.scatter_matrix(d, dimensions=c.dimensions, color=c.color, title=c.title)
    fig.update_traces(diagonal_visible=False, showupperhalf=False,
                      marker=dict(size=4, opacity=0.6))
    return fig


def _parallel_coordinates(c):
    d = c.df.copy()
    for col in c.dimensions:
        d[col] = pd.to_numeric(d[col], errors="coerce")
    d = d.dropna(subset=c.dimensions)
    color = c.color if (c.color and pd.api.types.is_numeric_dtype(
        pd.to_numeric(d[c.color], errors="coerce"))) else None
    if color:
        d[color] = pd.to_numeric(d[color], errors="coerce")
    return px.parallel_coordinates(d, dimensions=c.dimensions, color=color,
                                   title=c.title,
                                   color_continuous_scale=_scale(c.get("colorscale")))


def _parallel_categories(c):
    return px.parallel_categories(c.df, dimensions=c.dimensions, title=c.title,
                                  color=(pd.to_numeric(c.df[c.color], errors="coerce")
                                         if c.color else None))


def _scatter3d(c):
    return px.scatter_3d(c.df, x=c.x, y=c.y, z=c.z, color=c.color, size=c.size,
                         text=c.text, title=c.title)


def _surface(c):
    """集計済みのクロス表を立体にする（1列目が行ラベル）。"""
    label = c.x if c.x in c.df.columns else c.df.columns[0]
    m = c.df.set_index(label)
    m = m.apply(lambda s: pd.to_numeric(s, errors="coerce")).dropna(axis=1, how="all")
    fig = go.Figure(go.Surface(z=m.to_numpy(), x=list(m.columns),
                               y=[str(i) for i in m.index],
                               colorscale=_scale(c.get("colorscale"))))
    fig.update_layout(title=c.title, scene=dict(
        xaxis_title="列", yaxis_title=str(label), zaxis_title="値"))
    return fig


def _network(c):
    """つながりの図。円周上にノードを並べ、関係を線で結ぶ。"""
    d = c.df
    nodes = list(dict.fromkeys(d[c.source].astype(str).tolist()
                               + d[c.target].astype(str).tolist()))
    n = len(nodes)
    if not n:
        raise ValueError("つながりが1件もありません。")
    ang = {v: 2 * np.pi * i / n for i, v in enumerate(nodes)}
    pos = {v: (np.cos(a), np.sin(a)) for v, a in ang.items()}
    weights = _num_series(d, c.y) if c.y in d.columns else pd.Series([1] * len(d))
    wmax = float(weights.max() or 1)
    edge_x, edge_y = [], []
    for (_, r), w in zip(d.iterrows(), weights):
        x0, y0 = pos[str(r[c.source])]
        x1, y1 = pos[str(r[c.target])]
        edge_x += [x0, x1, None]
        edge_y += [y0, y1, None]
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=edge_x, y=edge_y, mode="lines", hoverinfo="skip",
                             line=dict(color="rgba(120,140,170,.45)", width=1.5),
                             showlegend=False))
    deg = pd.Series(d[c.source].astype(str).tolist()
                    + d[c.target].astype(str).tolist()).value_counts()
    fig.add_trace(go.Scatter(
        x=[pos[v][0] for v in nodes], y=[pos[v][1] for v in nodes],
        mode="markers+text", text=nodes, textposition="top center",
        marker=dict(size=[12 + 26 * deg.get(v, 1) / max(deg.max(), 1) for v in nodes],
                    color="#2E75B6"),
        hovertext=[f"{v}: {deg.get(v, 0)}件" for v in nodes], hoverinfo="text",
        showlegend=False))
    fig.update_layout(title=c.title, xaxis=dict(visible=False),
                      yaxis=dict(visible=False, scaleanchor="x"),
                      plot_bgcolor="rgba(0,0,0,0)")
    if wmax:
        fig.update_layout(margin=dict(l=20, r=20, t=50, b=20))
    return fig


# =============================================================================
# 指標
# =============================================================================

def _indicator_value(c) -> float:
    s = pd.to_numeric(c.df[c.value], errors="coerce").dropna()
    if s.empty:
        raise ValueError(f"{c.value} に数値がありません。")
    mode = (c.get("agg") or "sum").lower()
    return float({"sum": s.sum, "mean": s.mean, "max": s.max,
                  "min": s.min, "last": lambda: s.iloc[-1]}.get(mode, s.sum)())


def _target_of(c, value: float):
    if c.target and c.target in c.df.columns:
        t = pd.to_numeric(c.df[c.target], errors="coerce").dropna()
        return float(t.sum()) if len(t) else None
    try:
        return float(c.target_value) if c.target_value is not None else None
    except (TypeError, ValueError):
        return None


def _indicator(c):
    v = _indicator_value(c)
    t = _target_of(c, v)
    fig = go.Figure(go.Indicator(
        mode="number+delta" if t else "number", value=v,
        number=dict(valueformat=c.get("valueformat") or ",.4~f",
                    suffix=c.get("suffix") or ""),
        delta=dict(reference=t, relative=True, valueformat=".1%") if t else None,
        title=dict(text=c.title or str(c.value))))
    return fig


def _gauge(c):
    v = _indicator_value(c)
    t = _target_of(c, v)
    top = float(c.get("max") or (max(v, t or 0) * 1.25) or 1)
    fig = go.Figure(go.Indicator(
        mode="gauge+number+delta" if t else "gauge+number", value=v,
        number=dict(valueformat=c.get("valueformat") or ",.4~f",
                    suffix=c.get("suffix") or ""),
        delta=dict(reference=t, relative=True, valueformat=".1%") if t else None,
        title=dict(text=c.title or str(c.value)),
        gauge=dict(axis=dict(range=[0, top]), bar=dict(color="#2E75B6"),
                   steps=[dict(range=[0, top * 0.5], color="#F2F6FB"),
                          dict(range=[top * 0.5, top * 0.8], color="#D9E2F3")],
                   threshold=(dict(line=dict(color="#C55A11", width=3), value=t)
                              if t else None))))
    return fig


def _bullet(c):
    v = _indicator_value(c)
    t = _target_of(c, v)
    top = float(c.get("max") or (max(v, t or 0) * 1.25) or 1)
    fig = go.Figure(go.Indicator(
        mode="number+gauge+delta" if t else "number+gauge", value=v,
        delta=dict(reference=t) if t else None,
        number=dict(valueformat=c.get("valueformat") or ",.4~f"),
        title=dict(text=c.title or str(c.value)),
        gauge=dict(shape="bullet", axis=dict(range=[0, top]),
                   bar=dict(color="#1F4E79", thickness=0.6),
                   steps=[dict(range=[0, top * 0.6], color="#F2F6FB"),
                          dict(range=[top * 0.6, top * 0.85], color="#D9E2F3")],
                   threshold=(dict(line=dict(color="#C55A11", width=3), value=t)
                              if t else None))))
    fig.update_layout(height=190)
    return fig


# =============================================================================
# 組み立ての振り分け
# =============================================================================

_BUILDERS = {
    "bar": _bar, "hbar": _hbar, "stacked_bar": _stacked_bar,
    "percent_bar": _percent_bar, "lollipop": _lollipop, "dumbbell": _dumbbell,
    "pareto": _pareto, "pyramid": _pyramid, "marimekko": _marimekko,
    "radar": _radar, "polar_bar": _polar_bar, "bump": _bump,
    "line": _line, "step": _step, "area": _area, "area_percent": _area_percent,
    "range_area": _range_area, "slope": _slope, "candlestick": _candlestick,
    "ohlc": _ohlc, "gantt": _gantt, "calendar": _calendar,
    "control_chart": _control_chart,
    "pie": _pie, "donut": _pie, "treemap": _treemap, "sunburst": _sunburst,
    "icicle": _icicle, "funnel": _funnel, "waterfall": _waterfall, "sankey": _sankey,
    "histogram": _histogram, "density": _density, "ecdf": _ecdf, "box": _box,
    "violin": _violin, "strip": _strip, "ridgeline": _ridgeline, "qq": _qq,
    "scatter": _scatter, "bubble": _bubble, "histogram2d": _histogram2d,
    "contour": _contour, "heatmap": _heatmap, "matrix": _matrix,
    "scatter_matrix": _scatter_matrix,
    "parallel_coordinates": _parallel_coordinates,
    "parallel_categories": _parallel_categories,
    "scatter3d": _scatter3d, "surface": _surface, "network": _network,
    "indicator": _indicator, "gauge": _gauge, "bullet": _bullet,
}


def build_figure(item: dict):
    """render アイテム（kind="chart"）から plotly の figure を作る。"""
    ct = item.get("chart_type", "bar")
    builder = _BUILDERS.get(ct)
    if builder is None:
        raise ValueError(f"未対応のグラフ種別です: {ct} / "
                         f"使えるのは {', '.join(CHART_TYPES)}")
    fig = builder(_Ctx(item))
    fig.update_layout(margin=dict(l=55, r=20, t=50, b=50))
    return fig


def build_dual_figure(item: dict):
    """棒(左軸)+折れ線(右軸)の2軸グラフ。"""
    df = pd.DataFrame(item["rows"], columns=item["columns"])
    x = item["x"]
    bar_y = item.get("bar_y") or []
    line_y = item.get("line_y") or []
    _numeric(df, *bar_y, *line_y)

    fig = make_subplots(specs=[[{"secondary_y": True}]])
    for col in bar_y:
        fig.add_trace(go.Bar(x=df[x], y=df[col], name=col), secondary_y=False)
    for col in line_y:
        fig.add_trace(go.Scatter(x=df[x], y=df[col], name=col, mode="lines+markers"),
                      secondary_y=True)
    fig.update_layout(title=item.get("title", ""), barmode="group",
                      legend=dict(orientation="h", yanchor="bottom", y=1.02))
    fig.update_xaxes(title_text=x)
    fig.update_yaxes(title_text=item.get("left_title") or "（左軸）", secondary_y=False)
    fig.update_yaxes(title_text=item.get("right_title") or "（右軸）", secondary_y=True)
    return fig


# ==========================================================================
# ===== 元 figures.py
# グラフを画像（PNG）にする。Word や PowerPoint に貼るために使う。
#
# plotly の画像化は Chrome を裏で動かす（kaleido）。社内サーバに Chrome が
# 入っていないこともあるので、失敗したら None を返し、呼び出し側は
# 表や説明文だけで文書を作れるようにしてある。文書作成そのものは止めない。
#
# 画像は1文書内で何度も作るため、同じ図は使い回す（同じ処理を2回しない）。
# ==========================================================================
import hashlib
import threading

import config

_lock = threading.Lock()
_cache: dict[str, bytes] = {}
_MAX_CACHE = 40
# 一度失敗したら、その実行中は再挑戦しない（1枚あたり数秒待たされるため）
_broken: list[str] = []


def _figures_available() -> bool:
    """画像化できる環境かどうか（1回だけ実際に試して覚える）。"""
    if _broken:
        return False
    if _cache:
        return True
    import plotly.graph_objects as go
    return render(go.Figure(), width=80, height=60) is not None


def why_unavailable() -> str:
    return _broken[0] if _broken else ""


def render(fig, width: int | None = None, height: int | None = None,
           scale: float | None = None) -> bytes | None:
    """plotly の figure を PNG のバイト列にする。できなければ None。"""
    if _broken:
        return None
    w = int(width or config.REPORT_IMAGE_WIDTH)
    h = int(height or config.REPORT_IMAGE_HEIGHT)
    s = float(scale or config.REPORT_IMAGE_SCALE)
    try:
        key = hashlib.sha1(
            (fig.to_json() + f"|{w}x{h}@{s}").encode("utf-8")).hexdigest()
    except Exception:
        key = None
    if key:
        with _lock:
            hit = _cache.get(key)
        if hit is not None:
            return hit
    try:
        data = fig.to_image(format="png", width=w, height=h, scale=s)
    except Exception as e:
        msg = str(e).splitlines()[0][:200]
        _broken.append(f"グラフを画像にできませんでした（{type(e).__name__}: {msg}）。"
                       "文書には表と説明だけを入れます。"
                       "画像も入れたい場合は、サーバに Chrome/Chromium を用意して"
                       "kaleido が使える状態にしてください。")
        print(f"[figures] 画像化を無効にしました: {msg}")
        return None
    if key:
        with _lock:
            _cache[key] = data
            while len(_cache) > _MAX_CACHE:
                _cache.pop(next(iter(_cache)))
    return data


def for_print(fig, *, width=None, height=None):
    """紙・スライド向けに見た目を整えてから画像にする。

    画面はマウスで拡大できるが、紙とスライドはできない。
    文字を大きめに、余白を詰め、目盛りに桁区切りを入れる。
    """
    fig = _polish(fig)
    return render(fig, width=width, height=height)


def _polish(fig):
    """印刷向けの体裁に整える（元の figure は壊さない）。"""
    import copy
    fig = copy.deepcopy(fig)
    fig.update_layout(
        template="plotly_white",
        font=dict(family=config.REPORT_FONT_JA + ", sans-serif", size=15,
                  color="#1F1F1F"),
        title=dict(font=dict(size=17)),
        margin=dict(l=70, r=30, t=50, b=60),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0,
                    font=dict(size=13)),
        paper_bgcolor="white", plot_bgcolor="white",
        colorway=["#1F4E79", "#F4B183", "#70AD47", "#C55A11", "#7F7F7F",
                  "#2E75B6", "#A9D18E", "#FFD966", "#9DC3E6", "#BFBFBF"],
    )
    fig.update_xaxes(showgrid=False, linecolor="#BFBFBF", ticks="outside",
                     tickfont=dict(size=13))
    fig.update_yaxes(gridcolor="#E8E8E8", zerolinecolor="#BFBFBF",
                     tickfont=dict(size=13), tickformat=",")
    return fig


# ==========================================================================
# ===== 元 docx_report.py
# Wordレポートの生成。そのまま配布・回覧できる体裁で作る。
#
# 作りの方針:
#   - 表紙 → 目次 → 要約 → 本編 → 結論 → 付録 の順。報告書の型に合わせる。
#   - 図と表には通し番号とキャプションを付ける（「図3のとおり」と本文から呼べる）。
#   - 日本語フォントを明示的に当てる。指定しないと英字フォントが当たり、
#     開いた瞬間に体裁が崩れて見える。
#   - グラフは画像として貼る（Wordにネイティブのグラフが無いため）。
#     画像化できない環境では、同じ内容の表に自動で置き換える。
#
# 1セクション = 1つの dict:
#     {heading, body, bullets, table:{columns,rows}, image:bytes, caption,
#      note, callout, page_break}
# ==========================================================================
import io
import re
from datetime import datetime

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor

import config

DOCX_NAVY = RGBColor(0x1F, 0x3B, 0x5C)
DOCX_ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DOCX_HILITE = RGBColor(0xC5, 0x5A, 0x11)
DOCX_INK = RGBColor(0x22, 0x26, 0x2B)
DOCX_MUTED = RGBColor(0x6B, 0x72, 0x80)
DOCX_BAND = "F5F8FC"
HEADER_BG = "1F3B5C"
CALLOUT_BG = "FDF3E7"

DOCX_MAX_TABLE_ROWS = 40


class DocxReportError(Exception):
    """レポートを作れない理由（そのまま画面に出す）。"""


# =============================================================================
# 体裁の下ごしらえ
# =============================================================================

def _jp_font(run, size=None, bold=None, color=None, name=None):
    """日本語フォントを当てる（東アジア用は XML で直接指定する）。"""
    font = name or config.REPORT_FONT_JA
    run.font.name = font
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.insert(0, rFonts)
    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        rFonts.set(qn(attr), font)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _shade(cell_or_par, hex_color: str):
    el = cell_or_par._tc if hasattr(cell_or_par, "_tc") else cell_or_par._p
    pr = el.get_or_add_tcPr() if hasattr(el, "get_or_add_tcPr") else \
        el.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:fill"), hex_color)
    pr.append(shd)


def _border(par, *, size=6, color="C55A11", where="left"):
    pPr = par._p.get_or_add_pPr()
    borders = pPr.find(qn("w:pBdr"))
    if borders is None:
        borders = OxmlElement("w:pBdr")
        pPr.append(borders)
    b = OxmlElement(f"w:{where}")
    b.set(qn("w:val"), "single")
    b.set(qn("w:sz"), str(size * 4))
    b.set(qn("w:space"), "8")
    b.set(qn("w:color"), color)
    borders.append(b)


def _field(par, instr: str):
    """Wordのフィールド（ページ番号や目次）を入れる。開いたときに計算される。"""
    r1 = par.add_run()._element
    fld = OxmlElement("w:fldChar")
    fld.set(qn("w:fldCharType"), "begin")
    r1.append(fld)
    r2 = par.add_run()._element
    txt = OxmlElement("w:instrText")
    txt.set(qn("xml:space"), "preserve")
    txt.text = instr
    r2.append(txt)
    r3 = par.add_run()._element
    sep = OxmlElement("w:fldChar")
    sep.set(qn("w:fldCharType"), "separate")
    r3.append(sep)
    par.add_run("　")                        # 未計算のときに出る仮の文字
    r5 = par.add_run()._element
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    r5.append(end)


def _setup_styles(doc):
    """標準スタイルを日本語向けに整える。"""
    normal = doc.styles["Normal"]
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = DOCX_INK
    rPr = normal.element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.insert(0, rFonts)
    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        rFonts.set(qn(attr), config.REPORT_FONT_JA)
    normal.paragraph_format.line_spacing = 1.4
    normal.paragraph_format.space_after = Pt(6)

    for name, size, color, before in (("Heading 1", 16, DOCX_NAVY, 18),
                                      ("Heading 2", 13, DOCX_NAVY, 14),
                                      ("Heading 3", 11.5, DOCX_ACCENT, 10)):
        st = doc.styles[name]
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = color
        rPr = st.element.get_or_add_rPr()
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.insert(0, rFonts)
        for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
            rFonts.set(qn(attr), config.REPORT_FONT_JA)
        st.paragraph_format.space_before = Pt(before)
        st.paragraph_format.space_after = Pt(6)
        st.paragraph_format.keep_with_next = True


def _setup_page(doc, footer_text: str):
    sec = doc.sections[0]
    sec.top_margin = Cm(2.2)
    sec.bottom_margin = Cm(2.0)
    sec.left_margin = Cm(2.2)
    sec.right_margin = Cm(2.2)

    p = sec.footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if footer_text:
        r = p.add_run(footer_text + "　　")
        _jp_font(r, size=8.5, color=DOCX_MUTED)
    _field(p, "PAGE")
    r = p.add_run(" / ")
    _jp_font(r, size=8.5, color=DOCX_MUTED)
    _field(p, "NUMPAGES")
    for r in p.runs:
        _jp_font(r, size=8.5, color=DOCX_MUTED)


def _para(doc, text="", *, size=10.5, bold=False, color=DOCX_INK, align=None,
          space_after=6, style=None):
    p = doc.add_paragraph(style=style)
    if align is not None:
        p.alignment = align
    p.paragraph_format.space_after = Pt(space_after)
    for i, line in enumerate(str(text).split("\n")):
        if i:
            p.add_run().add_break()
        r = p.add_run(line)
        _jp_font(r, size=size, bold=bold, color=color)
    return p


def _docx_fmt(v) -> str:
    if v is None:
        return ""
    if isinstance(v, bool):
        return "はい" if v else "いいえ"
    if isinstance(v, float):
        if v == int(v) and abs(v) < 1e15:
            return f"{int(v):,}"
        return f"{v:,.2f}".rstrip("0").rstrip(".")
    if isinstance(v, int):
        return f"{v:,}"
    return str(v)


# =============================================================================
# 部品
# =============================================================================

def _cover(doc, args: dict):
    for _ in range(4):
        doc.add_paragraph()
    _para(doc, args.get("title", "レポート"), size=26, bold=True, color=DOCX_NAVY,
          align=WD_ALIGN_PARAGRAPH.CENTER, space_after=4)
    if args.get("subtitle"):
        _para(doc, args["subtitle"], size=13, color=DOCX_MUTED,
              align=WD_ALIGN_PARAGRAPH.CENTER, space_after=28)

    # 表紙の線
    line = doc.add_paragraph()
    line.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _border(line, size=3, color="1F3B5C", where="bottom")

    for _ in range(6):
        doc.add_paragraph()
    org = args.get("org") or config.REPORT_ORG
    for text in (args.get("date") or datetime.now().strftime("%Y年%m月%d日"),
                 org, args.get("author")):
        if text:
            _para(doc, text, size=11, color=DOCX_MUTED,
                  align=WD_ALIGN_PARAGRAPH.CENTER, space_after=2)
    doc.add_page_break()


def _toc(doc):
    _para(doc, "目次", size=15, bold=True, color=DOCX_NAVY, space_after=10)
    p = doc.add_paragraph()
    _field(p, r'TOC \o "1-2" \h \z \u')
    _para(doc, "※ 目次はWordで開いたあと、この部分を選んで F9 を押すと最新になります。",
          size=8.5, color=DOCX_MUTED, space_after=0)
    doc.add_page_break()


def _summary_box(doc, points: list):
    if not points:
        return
    _para(doc, "要点", size=12, bold=True, color=DOCX_HILITE, space_after=4)
    for s in points:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(str(s))
        _jp_font(r, size=11, bold=True, color=DOCX_INK)
    doc.add_paragraph()


def _docx_callout(doc, text: str, label="ポイント"):
    p = _para(doc, f"【{label}】{text}", size=10.5, color=DOCX_INK, space_after=10)
    _border(p, size=6, color="C55A11", where="left")
    _shade(p, CALLOUT_BG)
    p.paragraph_format.left_indent = Cm(0.4)
    p.paragraph_format.space_before = Pt(6)


def _table(doc, columns, rows, *, caption=None, number=None, note=None):
    limit = DOCX_MAX_TABLE_ROWS
    shown, cut = rows[:limit], max(0, len(rows) - limit)
    t = doc.add_table(rows=len(shown) + 1, cols=len(columns))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.autofit = True

    for j, c in enumerate(columns):
        cell = t.cell(0, j)
        cell.text = ""
        _shade(cell, HEADER_BG)
        p = cell.paragraphs[0]
        p.paragraph_format.space_after = Pt(2)
        r = p.add_run(str(c))
        _jp_font(r, size=9.5, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    for i, row in enumerate(shown, start=1):
        for j in range(len(columns)):
            v = row[j] if j < len(row) else ""
            cell = t.cell(i, j)
            cell.text = ""
            if i % 2 == 0:
                _shade(cell, DOCX_BAND)
            p = cell.paragraphs[0]
            p.paragraph_format.space_after = Pt(2)
            if isinstance(v, (int, float)) and not isinstance(v, bool):
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            r = p.add_run(_docx_fmt(v))
            _jp_font(r, size=9.5)

    tail = []
    if caption:
        cap = _para(doc, f"表{number} {caption}" if number else caption,
                    size=9, color=DOCX_MUTED, align=WD_ALIGN_PARAGRAPH.LEFT,
                    space_after=4)
        cap.paragraph_format.space_before = Pt(3)
    if cut:
        tail.append(f"全 {len(rows):,} 行のうち上位 {limit} 行を掲載")
    if note:
        tail.append(note)
    if tail:
        _para(doc, "　".join(tail), size=8.5, color=DOCX_MUTED, space_after=10)
    return t


def _image(doc, data: bytes, *, caption=None, number=None, width_cm=16.0):
    doc.add_picture(io.BytesIO(data), width=Cm(width_cm))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    if caption:
        _para(doc, f"図{number} {caption}" if number else caption, size=9,
              color=DOCX_MUTED, align=WD_ALIGN_PARAGRAPH.CENTER, space_after=10)


# =============================================================================
# 組み立て
# =============================================================================

def build_docx(sections: list[dict], *, title="レポート", subtitle="", summary=None,
          conclusion="", recommendations=None, caveats=None, footer="",
          org="", author="", toc=True, appendix=None) -> bytes:
    """セクションのリストから .docx のバイト列を作る。"""
    if not sections:
        raise DocxReportError("セクションが1つもありません。")

    doc = Document()
    _setup_styles(doc)
    _setup_page(doc, footer or config.REPORT_ORG or "")
    _cover(doc, {"title": title, "subtitle": subtitle, "org": org, "author": author})
    if toc:
        _toc(doc)

    if summary:
        doc.add_heading("要約", level=1)
        _summary_box(doc, summary)

    fig_no, tbl_no = 0, 0
    for i, s in enumerate(sections, 1):
        if not s.get("heading"):
            raise DocxReportError(f"{i}番目のセクションに heading がありません。")
        if s.get("page_break"):
            doc.add_page_break()
        doc.add_heading(s["heading"], level=int(s.get("level") or 1))
        if s.get("body"):
            _para(doc, s["body"])
        for b in (s.get("bullets") or []):
            text = b.get("text", "") if isinstance(b, dict) else str(b)
            level = int(b.get("level", 0)) if isinstance(b, dict) else 0
            p = doc.add_paragraph(style="List Bullet" if not level
                                  else "List Bullet 2")
            p.paragraph_format.space_after = Pt(3)
            r = p.add_run(text)
            _jp_font(r, size=10.5)
        if s.get("image"):
            fig_no += 1
            _image(doc, s["image"], caption=s.get("caption") or s["heading"],
                   number=fig_no)
        if s.get("table"):
            t = s["table"]
            if not t.get("columns"):
                raise DocxReportError(f"「{s['heading']}」の表に columns がありません。")
            tbl_no += 1
            _table(doc, t["columns"], t.get("rows") or [],
                   caption=s.get("table_caption") or s.get("caption") or s["heading"],
                   number=tbl_no, note=t.get("note"))
        if s.get("callout"):
            _docx_callout(doc, s["callout"])
        if s.get("note"):
            p = _para(doc, s["note"], size=10, color=DOCX_MUTED, space_after=10)
            _border(p, size=4, color="D5DBE2", where="left")
            p.paragraph_format.left_indent = Cm(0.4)

    if conclusion:
        doc.add_heading("結論", level=1)
        _para(doc, conclusion)
    if recommendations:
        doc.add_heading("推奨する打ち手", level=1)
        for i, r in enumerate(recommendations, 1):
            if isinstance(r, dict):
                text = r.get("text", "")
                extra = "　".join(x for x in
                                  (f"担当: {r['owner']}" if r.get("owner") else "",
                                   f"期限: {r['due']}" if r.get("due") else "") if x)
            else:
                text, extra = str(r), ""
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(4)
            run = p.add_run(f"{i}. ")
            _jp_font(run, size=10.5, bold=True, color=DOCX_HILITE)
            run = p.add_run(text)
            _jp_font(run, size=10.5, bold=True)
            if extra:
                run = p.add_run(f"（{extra}）")
                _jp_font(run, size=9.5, color=DOCX_MUTED)
    if caveats:
        doc.add_heading("前提・注意", level=1)
        for c in caveats:
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(3)
            r = p.add_run(str(c))
            _jp_font(r, size=10, color=DOCX_MUTED)

    for extra in (appendix or []):
        doc.add_page_break()
        doc.add_heading(extra.get("heading", "付録"), level=1)
        if extra.get("body"):
            _para(doc, extra["body"])
        if extra.get("table"):
            tbl_no += 1
            _table(doc, extra["table"]["columns"], extra["table"].get("rows") or [],
                   caption=extra.get("caption"), number=tbl_no)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def docx_safe_filename(name: str | None, default: str = "report") -> str:
    base = re.sub(r'[\\/:*?"<>|]', "_", str(name or default)).strip() or default
    return base if base.lower().endswith(".docx") else base + ".docx"


def outline_docx(sections: list[dict]) -> list[str]:
    out = []
    for i, s in enumerate(sections, 1):
        bits = []
        if s.get("image"):
            bits.append("図")
        if s.get("table"):
            bits.append(f"表{len(s['table'].get('rows') or [])}行")
        if s.get("bullets"):
            bits.append(f"箇条書き{len(s['bullets'])}件")
        out.append(f"{i}. {s.get('heading', '')}"
                   + (f"（{'・'.join(bits)}）" if bits else ""))
    return out


# ==========================================================================
# ===== 元 pptx_report.py
# PowerPointレポートの生成。会議でそのまま映せる体裁で作る。
#
# 作りの方針:
#   - 1スライド1メッセージ。上部の「キーメッセージ」に結論を1行で書き、
#     図表はその根拠として下に置く。読み手は上の1行だけで用が足りる。
#   - 日本語フォントを明示的に指定する。指定しないと英字フォントが当たり、
#     開いた瞬間に「ちゃんとしていない資料」に見える。
#   - グラフはPowerPointネイティブ（編集可）を既定にし、
#     ネイティブで表現できない種類だけ画像として貼る。
#
# 1スライド = 1つの dict。kind で中身が決まる:
#     title    表紙
#     agenda   目次
#     section  中扉
#     message  文字だけ（結論・考察）
#     table    表
#     chart    グラフ
#     kpi      数字を大きく並べる
#     compare  2つ並べて比較
#     closing  まとめ／次のアクション
# ==========================================================================
import io
import re
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor as _pptx_RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn as _pptx_qn
from pptx.util import Emu, Inches, Pt as _pptx_Pt

import config

SLIDE_W = Inches(13.333)          # 16:9
SLIDE_H = Inches(7.5)

PPTX_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_percent": XL_CHART_TYPE.COLUMN_STACKED_100,
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
    "hbar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "radar": XL_CHART_TYPE.RADAR_MARKERS,
}
SLIDE_KINDS = ("title", "agenda", "section", "message", "table", "chart",
               "kpi", "compare", "closing")

# 配色。1枚に何色も出さない。強調は1色だけ使う。
PPTX_NAVY = _pptx_RGBColor(0x1F, 0x3B, 0x5C)
PPTX_ACCENT = _pptx_RGBColor(0x2E, 0x75, 0xB6)
PPTX_HILITE = _pptx_RGBColor(0xC5, 0x5A, 0x11)
PPTX_INK = _pptx_RGBColor(0x22, 0x26, 0x2B)
PPTX_MUTED = _pptx_RGBColor(0x6B, 0x72, 0x80)
LINE = _pptx_RGBColor(0xD5, 0xDB, 0xE2)
PPTX_BAND = _pptx_RGBColor(0xF5, 0xF8, 0xFC)
WHITE = _pptx_RGBColor(0xFF, 0xFF, 0xFF)
GOOD = _pptx_RGBColor(0x1E, 0x7A, 0x3C)
BAD = _pptx_RGBColor(0xB0, 0x2A, 0x2A)
SERIES = ["1F4E79", "F4B183", "70AD47", "C55A11", "7F7F7F",
          "2E75B6", "A9D18E", "FFD966", "9DC3E6", "BFBFBF"]

PPTX_MAX_TABLE_ROWS = 12               # これを超えると字が小さくなって読めない
MAX_CATEGORIES = 24

# 余白（本文の左右端）
MARGIN = Inches(0.62)
BODY_W = SLIDE_W - MARGIN * 2


class PptxReportError(Exception):
    """レポートを作れない理由（そのまま画面に出す）。"""


# =============================================================================
# 文字まわり
# =============================================================================

def _jp(run):
    """日本語フォントを当てる。

    python-pptx は latin フォントしか設定しないので、
    日本語部分に別のフォントが当たってしまう。東アジア用を直接書く。
    """
    run.font.name = config.REPORT_FONT_JA
    rPr = run.font._element          # これ自体が rPr（文字の書式）
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(_pptx_qn(tag))
        if el is None:
            el = rPr.makeelement(_pptx_qn(tag), {})
            rPr.append(el)
        el.set("typeface", config.REPORT_FONT_JA)


def _text(frame, lines, *, size=18, bold=False, color=PPTX_INK, align=PP_ALIGN.LEFT,
          space_after=4, line_spacing=1.25):
    """text_frame に段落を流し込む。lines は文字列か (文字列, 上書き) の並び。"""
    frame.word_wrap = True
    items = lines if isinstance(lines, (list, tuple)) else [lines]
    first = True
    for item in items:
        opts = {}
        if isinstance(item, tuple):
            item, opts = item
        for line in str(item).split("\n"):
            p = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            p.text = line
            p.alignment = opts.get("align", align)
            p.space_after = _pptx_Pt(opts.get("space_after", space_after))
            p.line_spacing = opts.get("line_spacing", line_spacing)
            if opts.get("level"):
                p.level = opts["level"]
            for run in p.runs:
                run.font.size = _pptx_Pt(opts.get("size", size))
                run.font.bold = opts.get("bold", bold)
                run.font.color.rgb = opts.get("color", color)
                _jp(run)


def _pptx_box(slide, left, top, width, height, lines, **kw):
    shape = slide.shapes.add_textbox(left, top, width, height)
    _text(shape.text_frame, lines, **kw)
    return shape


def _rect(slide, left, top, width, height, fill=None, line=None,
          shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = _pptx_Pt(1)
    s.shadow.inherit = False
    return s


def _pptx_clean(v) -> str:
    return "" if v is None else str(v)


def _num(v):
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = re.sub(r"[,\s¥$%]", "", str(v))
    try:
        return float(s)
    except ValueError:
        return None


def _pptx_fmt(v) -> str:
    if isinstance(v, bool):
        return "はい" if v else "いいえ"
    if isinstance(v, float):
        if v == int(v) and abs(v) < 1e15:
            return f"{int(v):,}"
        # 12.30 ではなく 12.3 と出す（末尾の0は読み手には意味が無い）
        return f"{v:,.2f}".rstrip("0").rstrip(".")
    if isinstance(v, int):
        return f"{v:,}"
    return _pptx_clean(v)


# =============================================================================
# 共通の枠（ヘッダ・キーメッセージ・フッタ）
# =============================================================================

def _pptx_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, title: str, message: str = "") -> Emu:
    """見出しと、その下のキーメッセージ帯。戻り値は本文を始めてよい上端。"""
    _pptx_box(slide, MARGIN, Inches(0.30), BODY_W, Inches(0.55), title,
         size=25, bold=True, color=PPTX_NAVY)
    _rect(slide, MARGIN, Inches(0.92), BODY_W, Emu(12700), fill=PPTX_NAVY)

    if not message:
        return Inches(1.18)
    # 結論を1行で。ここだけ読めば分かるようにする。
    bar = _rect(slide, MARGIN, Inches(1.06), BODY_W, Inches(0.62), fill=PPTX_BAND)
    bar.line.color.rgb = LINE
    bar.line.width = _pptx_Pt(0.75)
    tf = bar.text_frame
    tf.margin_left, tf.margin_right = Inches(0.16), Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _text(tf, message, size=15, bold=True, color=PPTX_NAVY, space_after=0)
    return Inches(1.86)


def _footer(slide, text: str, page: int | None = None):
    _rect(slide, MARGIN, SLIDE_H - Inches(0.52), BODY_W, Emu(9525), fill=LINE)
    if text:
        _pptx_box(slide, MARGIN, SLIDE_H - Inches(0.46), Inches(9), Inches(0.32),
             text, size=9.5, color=PPTX_MUTED)
    if page:
        _pptx_box(slide, SLIDE_W - MARGIN - Inches(0.8), SLIDE_H - Inches(0.46),
             Inches(0.8), Inches(0.32), str(page), size=9.5, color=PPTX_MUTED,
             align=PP_ALIGN.RIGHT)


def _notes(slide, text: str):
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


def _source(slide, top, text: str):
    """出典・条件。数字の資料には必ず要る。"""
    if text:
        _pptx_box(slide, MARGIN, top, BODY_W, Inches(0.3), f"出所: {text}",
             size=9.5, color=PPTX_MUTED)


# =============================================================================
# スライドの種類ごと
# =============================================================================

def _slide_title(prs, spec):
    slide = _pptx_blank(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(3.05), fill=PPTX_NAVY)
    _rect(slide, 0, Inches(3.05), SLIDE_W, Inches(0.06), fill=PPTX_HILITE)
    _pptx_box(slide, Inches(0.9), Inches(1.05), SLIDE_W - Inches(1.8), Inches(1.2),
         spec.get("title", "レポート"), size=40, bold=True, color=WHITE)
    if spec.get("subtitle"):
        _pptx_box(slide, Inches(0.9), Inches(2.25), SLIDE_W - Inches(1.8), Inches(0.5),
             spec["subtitle"], size=18, color=_pptx_RGBColor(0xC5, 0xD5, 0xE8))

    y = Inches(3.55)
    for line in (spec.get("lines") or [])[:4]:
        _rect(slide, Inches(0.9), y + Inches(0.10), Inches(0.09), Inches(0.22),
              fill=PPTX_HILITE)
        _pptx_box(slide, Inches(1.15), y, SLIDE_W - Inches(2.2), Inches(0.42), line,
             size=16, color=PPTX_INK)
        y += Inches(0.52)

    org = spec.get("org") or config.REPORT_ORG
    foot = " ／ ".join(x for x in (org, spec.get("author")) if x)
    _pptx_box(slide, Inches(0.9), SLIDE_H - Inches(1.05), Inches(8), Inches(0.34),
         spec.get("date") or datetime.now().strftime("%Y年%m月%d日"),
         size=13, color=PPTX_MUTED)
    if foot:
        _pptx_box(slide, Inches(0.9), SLIDE_H - Inches(0.72), Inches(8), Inches(0.34),
             foot, size=13, color=PPTX_MUTED)
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_agenda(prs, spec):
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title") or "本日の内容", spec.get("message", ""))
    items = spec.get("items") or []
    y = top + Inches(0.18)
    step = min(Inches(0.72), (SLIDE_H - y - Inches(0.9)) / max(len(items), 1))
    for i, it in enumerate(items, 1):
        label = it.get("text") if isinstance(it, dict) else str(it)
        note = it.get("note", "") if isinstance(it, dict) else ""
        n = _rect(slide, MARGIN, y, Inches(0.44), Inches(0.44), fill=PPTX_NAVY,
                  shape=MSO_SHAPE.OVAL)
        tf = n.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(tf, str(i), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              space_after=0)
        _pptx_box(slide, MARGIN + Inches(0.62), y + Inches(0.02), BODY_W - Inches(0.8),
             Inches(0.4), label, size=17, bold=True, color=PPTX_INK)
        if note:
            _pptx_box(slide, MARGIN + Inches(0.62), y + Inches(0.34),
                 BODY_W - Inches(0.8), Inches(0.3), note, size=12, color=PPTX_MUTED)
        y += step
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_section(prs, spec):
    slide = _pptx_blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PPTX_NAVY)
    _rect(slide, MARGIN, Inches(3.18), Inches(0.9), Inches(0.07), fill=PPTX_HILITE)
    _pptx_box(slide, MARGIN, Inches(3.42), SLIDE_W - MARGIN * 2, Inches(0.9),
         spec.get("title", ""), size=32, bold=True, color=WHITE)
    if spec.get("subtitle"):
        _pptx_box(slide, MARGIN, Inches(4.35), SLIDE_W - MARGIN * 2, Inches(0.5),
             spec["subtitle"], size=15, color=_pptx_RGBColor(0xC5, 0xD5, 0xE8))
    _notes(slide, spec.get("notes", ""))
    return slide


def _bullets(slide, left, top, width, height, items, *, size=16):
    """箇条書き。dict なら {text, level, strong} を見る。"""
    y = top
    for item in items:
        if isinstance(item, dict):
            text, level = item.get("text", ""), int(item.get("level", 0))
            strong = bool(item.get("strong"))
        else:
            text, level, strong = str(item), 0, False
        if y > top + height - Inches(0.3):
            break
        mark_x = left + Inches(0.26) * level
        if level == 0:
            _rect(slide, mark_x, y + Inches(0.13), Inches(0.10), Inches(0.10),
                  fill=PPTX_HILITE if strong else PPTX_ACCENT)
        else:
            _pptx_box(slide, mark_x, y - Inches(0.02), Inches(0.2), Inches(0.3), "－",
                 size=13, color=PPTX_MUTED)
        _pptx_box(slide, mark_x + Inches(0.24), y - Inches(0.05),
             width - Inches(0.26) * level - Inches(0.24), Inches(0.42), text,
             size=size - 1.5 * level, bold=strong,
             color=PPTX_HILITE if strong else PPTX_INK)
        # 行数ぶん送る（おおよそでよい。重ならなければ十分）
        per = max(1, int(width / Inches(0.135) / max(size - 1.5 * level, 1) * 1.9))
        lines = max(1, (len(text) + per - 1) // per)
        y += Inches(0.34) * lines + Inches(0.12)
    return y


def _slide_message(prs, spec):
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    y = top + Inches(0.12)
    if spec.get("lead"):
        _pptx_box(slide, MARGIN, y, BODY_W, Inches(0.6), spec["lead"], size=17,
             color=PPTX_INK)
        y += Inches(0.72)
    items = spec.get("bullets") or []
    if items:
        y = _bullets(slide, MARGIN, y, BODY_W, SLIDE_H - y - Inches(1.2), items,
                     size=17)
    if spec.get("body"):
        _pptx_box(slide, MARGIN, y + Inches(0.1), BODY_W,
             SLIDE_H - y - Inches(1.1), spec["body"], size=14, color=PPTX_INK)
    if spec.get("callout"):
        _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.55), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _pptx_callout(slide, left, top, width, text, label="ポイント"):
    """強調枠。1枚に1つだけ置く。"""
    box = _rect(slide, left, top, width, Inches(0.92), fill=_pptx_RGBColor(0xFD, 0xF3, 0xE7))
    box.line.color.rgb = PPTX_HILITE
    box.line.width = _pptx_Pt(1.25)
    _rect(slide, left, top, Inches(0.07), Inches(0.92), fill=PPTX_HILITE)
    _pptx_box(slide, left + Inches(0.22), top + Inches(0.09), Inches(2), Inches(0.26),
         label, size=10.5, bold=True, color=PPTX_HILITE)
    _pptx_box(slide, left + Inches(0.22), top + Inches(0.34), width - Inches(0.44),
         Inches(0.52), text, size=14, color=PPTX_INK)


def _slide_kpi(prs, spec):
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    items = (spec.get("items") or [])[:4]
    if not items:
        raise PptxReportError("kpi スライドには items が必要です。")
    gap = Inches(0.32)
    width = (BODY_W - gap * (len(items) - 1)) / len(items)
    card_h = Inches(2.35)
    for i, it in enumerate(items):
        left = MARGIN + (width + gap) * i
        card = _rect(slide, left, top + Inches(0.25), width, card_h, fill=PPTX_BAND,
                     line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.04
        _rect(slide, left, top + Inches(0.25), width, Inches(0.07), fill=PPTX_ACCENT)
        _pptx_box(slide, left, top + Inches(0.48), width, Inches(0.34),
             _pptx_clean(it.get("label")), size=13, color=PPTX_MUTED, align=PP_ALIGN.CENTER)
        _pptx_box(slide, left, top + Inches(0.86), width, Inches(0.85),
             _pptx_fmt(it.get("value")) + _pptx_clean(it.get("unit")),
             size=34, bold=True, color=PPTX_NAVY, align=PP_ALIGN.CENTER)
        d = _num(it.get("delta"))
        if d is not None:
            mark = "▲" if d > 0 else ("▼" if d < 0 else "―")
            good = it.get("higher_is_better", True)
            col = PPTX_MUTED if d == 0 else (GOOD if (d > 0) == bool(good) else BAD)
            _pptx_box(slide, left, top + Inches(1.72), width, Inches(0.34),
                 f"{mark} {_pptx_fmt(abs(d))}{_pptx_clean(it.get('delta_unit'))}"
                 + (f"（{it['delta_label']}）" if it.get("delta_label") else ""),
                 size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        if it.get("note"):
            _pptx_box(slide, left, top + Inches(2.06), width, Inches(0.3),
                 it["note"], size=10.5, color=PPTX_MUTED, align=PP_ALIGN.CENTER)

    y = top + card_h + Inches(0.45)
    if spec.get("bullets"):
        y = _bullets(slide, MARGIN, y, BODY_W, SLIDE_H - y - Inches(1.0),
                     spec["bullets"], size=15)
    elif spec.get("comment"):
        _pptx_box(slide, MARGIN, y, BODY_W, Inches(1.0), spec["comment"], size=14)
    if spec.get("callout"):
        _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    _notes(slide, spec.get("notes", ""))
    return slide


def _add_table(slide, left, top, width, height, columns, rows, *,
               font=11.5, highlight_rows=()):
    shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top,
                                   width, height)
    table = shape.table
    for j, c in enumerate(columns):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTX_NAVY
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(cell.text_frame, _pptx_clean(c), size=font, bold=True, color=WHITE,
              space_after=0, line_spacing=1.0)
    for i, row in enumerate(rows, start=1):
        for j in range(len(columns)):
            v = row[j] if j < len(row) else ""
            cell = table.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                _pptx_RGBColor(0xFD, 0xF3, 0xE7) if (i - 1) in highlight_rows
                else (PPTX_BAND if i % 2 == 0 else WHITE))
            _text(cell.text_frame, _pptx_fmt(v), size=font, space_after=0,
                  line_spacing=1.0,
                  align=PP_ALIGN.RIGHT if isinstance(v, (int, float))
                  and not isinstance(v, bool) else PP_ALIGN.LEFT,
                  bold=(i - 1) in highlight_rows)
    return table


def _slide_table(prs, spec):
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    cols = [_pptx_clean(c) for c in (spec.get("columns") or [])]
    rows = spec.get("rows") or []
    if not cols:
        raise PptxReportError("table スライドには columns が必要です。")
    limit = int(spec.get("max_rows") or PPTX_MAX_TABLE_ROWS)
    shown, cut = rows[:limit], max(0, len(rows) - limit)

    comment = spec.get("comment")
    width = BODY_W if not comment else BODY_W - Inches(3.5)
    avail = SLIDE_H - top - Inches(1.05)
    height = min(Inches(0.36) * (len(shown) + 1), avail)
    _add_table(slide, MARGIN, top + Inches(0.1), width, height, cols, shown,
               font=11.5 if len(cols) <= 7 else 10,
               highlight_rows=set(spec.get("highlight_rows") or []))
    if comment:
        _pptx_box(slide, MARGIN + width + Inches(0.3), top + Inches(0.1),
             Inches(3.2), avail, comment, size=13.5)
    note = spec.get("source", "")
    if cut:
        note = (note + f"（全 {len(rows):,} 行のうち上位 {limit} 行）").strip()
    _source(slide, SLIDE_H - Inches(0.85), note)
    if spec.get("callout"):
        _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _style_chart(chart, spec, series_count):
    chart.has_title = False
    chart.font.size = _pptx_Pt(12)
    chart.font.name = config.REPORT_FONT_JA
    kind = str(spec.get("chart") or "bar").lower()

    if series_count > 1 or kind in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = _pptx_Pt(12)
    else:
        chart.has_legend = False

    try:
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = _pptx_RGBColor.from_string(
                SERIES[i % len(SERIES)])
            s.format.line.color.rgb = _pptx_RGBColor.from_string(SERIES[i % len(SERIES)])
    except (AttributeError, ValueError, NotImplementedError):
        pass

    # 数値軸は桁区切り。目盛り線は薄く。
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = _pptx_RGBColor(0xE8, 0xE8, 0xE8)
        va.format.line.color.rgb = LINE
        va.tick_labels.number_format = spec.get("number_format") or "#,##0"
        va.tick_labels.number_format_is_linked = False
        va.tick_labels.font.size = _pptx_Pt(11.5)
    except (AttributeError, ValueError, NotImplementedError):
        pass
    try:
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = LINE
        ca.tick_labels.font.size = _pptx_Pt(11.5)
    except (AttributeError, ValueError, NotImplementedError):
        pass

    # 値ラベル。多すぎると潰れるので、少ないときだけ。
    want = spec.get("data_labels")
    cats = len(spec.get("categories") or [])
    if want is None:
        want = kind in ("pie", "doughnut") or (cats and cats <= 8 and series_count <= 2)
    if want:
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = _pptx_Pt(11)
            dl.font.name = config.REPORT_FONT_JA
            if kind in ("pie", "doughnut"):
                dl.show_percentage = True
                dl.number_format = "0.0%"
                dl.number_format_is_linked = False
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            else:
                dl.number_format = spec.get("number_format") or "#,##0"
                dl.number_format_is_linked = False
                if kind in ("bar", "hbar"):
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except (AttributeError, ValueError, NotImplementedError):
            pass


def _slide_chart(prs, spec):
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    series = spec.get("series") or []
    cats = [_pptx_clean(c) for c in (spec.get("categories") or [])]

    # 画像として渡された図（PowerPointで描けない種類）はそのまま貼る
    if spec.get("image"):
        avail_h = SLIDE_H - top - Inches(1.05)
        has_side = bool(spec.get("comment"))
        w = BODY_W if not has_side else BODY_W - Inches(3.5)
        slide.shapes.add_picture(io.BytesIO(spec["image"]), MARGIN,
                                 top + Inches(0.08), width=w)
        if has_side:
            _pptx_box(slide, MARGIN + w + Inches(0.3), top + Inches(0.1), Inches(3.2),
                 avail_h, spec["comment"], size=13.5)
        _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
        if spec.get("callout"):
            _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
        _notes(slide, spec.get("notes", ""))
        return slide

    if not series:
        raise PptxReportError("chart スライドには series か image が必要です。")
    kind = str(spec.get("chart") or "bar").lower()
    if kind not in PPTX_CHART_TYPES:
        raise PptxReportError(f"未対応のグラフ種類です: {kind}。"
                          f"使えるのは {', '.join(PPTX_CHART_TYPES)} です。")
    if len(cats) > MAX_CATEGORIES:
        cats = cats[:MAX_CATEGORIES]
        series = [{**s, "values": (s.get("values") or [])[:MAX_CATEGORIES]}
                  for s in series]
        spec = {**spec, "source": (spec.get("source", "")
                                   + f"（上位{MAX_CATEGORIES}件）").strip()}

    if kind == "scatter":
        data = XyChartData()
        for s in series:
            sd = data.add_series(_pptx_clean(s.get("name") or "系列"))
            for x, y in zip(s.get("x") or [], s.get("values") or s.get("y") or []):
                if _num(x) is not None and _num(y) is not None:
                    sd.add_data_point(_num(x), _num(y))
    else:
        data = CategoryChartData()
        data.categories = cats or [str(i + 1) for i in
                                   range(len(series[0].get("values") or []))]
        for s in series:
            data.add_series(_pptx_clean(s.get("name") or "系列"),
                            [_num(v) for v in (s.get("values") or [])],
                            number_format=spec.get("number_format") or "#,##0")

    has_side = bool(spec.get("comment"))
    width = BODY_W if not has_side else BODY_W - Inches(3.5)
    height = SLIDE_H - top - Inches(1.05)
    frame = slide.shapes.add_chart(PPTX_CHART_TYPES[kind], MARGIN, top + Inches(0.08),
                                   width, height, data)
    _style_chart(frame.chart, {**spec, "categories": cats}, len(series))

    if has_side:
        _pptx_box(slide, MARGIN + width + Inches(0.3), top + Inches(0.15), Inches(3.2),
             height - Inches(0.2), spec["comment"], size=13.5)
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    if spec.get("callout"):
        _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_compare(prs, spec):
    """左右に並べて比べる（案A/案B、前年/今年 など）。"""
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    panes = (spec.get("panes") or [])[:2]
    if len(panes) != 2:
        raise PptxReportError("compare スライドには panes を2つ指定してください。")
    gap = Inches(0.4)
    w = (BODY_W - gap) / 2
    h = SLIDE_H - top - Inches(1.05)
    for i, pane in enumerate(panes):
        left = MARGIN + (w + gap) * i
        head = _rect(slide, left, top + Inches(0.05), w, Inches(0.46),
                     fill=PPTX_NAVY if i == 0 else PPTX_ACCENT)
        head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(head.text_frame, _pptx_clean(pane.get("title")), size=15, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, space_after=0)
        y = top + Inches(0.62)
        if pane.get("value") is not None:
            _pptx_box(slide, left, y, w, Inches(0.8),
                 _pptx_fmt(pane["value"]) + _pptx_clean(pane.get("unit")),
                 size=30, bold=True, color=PPTX_NAVY, align=PP_ALIGN.CENTER)
            y += Inches(0.9)
        if pane.get("image"):
            slide.shapes.add_picture(io.BytesIO(pane["image"]), left, y, width=w)
            y += Inches(2.6)
        if pane.get("bullets"):
            _bullets(slide, left + Inches(0.05), y, w - Inches(0.1),
                     top + h - y, pane["bullets"], size=14)
    if spec.get("callout"):
        _pptx_callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_closing(prs, spec):
    """まとめと次のアクション。担当と期限まで書けるようにする。"""
    slide = _pptx_blank(prs)
    top = _header(slide, spec.get("title") or "まとめと次のアクション",
                  spec.get("message", ""))
    y = top + Inches(0.1)
    if spec.get("summary"):
        _pptx_box(slide, MARGIN, y, BODY_W, Inches(0.3), "まとめ", size=12,
             bold=True, color=PPTX_MUTED)
        y = _bullets(slide, MARGIN, y + Inches(0.34), BODY_W, Inches(2.0),
                     spec["summary"], size=16) + Inches(0.15)

    actions = spec.get("actions") or []
    if actions:
        _pptx_box(slide, MARGIN, y, BODY_W, Inches(0.3), "次のアクション", size=12,
             bold=True, color=PPTX_MUTED)
        y += Inches(0.34)
        rows = []
        for a in actions:
            if isinstance(a, dict):
                rows.append([a.get("text", ""), a.get("owner", ""), a.get("due", "")])
            else:
                rows.append([str(a), "", ""])
        h = min(Inches(0.36) * (len(rows) + 1), SLIDE_H - y - Inches(0.8))
        _add_table(slide, MARGIN, y, BODY_W, h, ["やること", "担当", "期限"], rows,
                   font=12.5)
    _notes(slide, spec.get("notes", ""))
    return slide


_PPTX_BUILDERS = {"title": _slide_title, "agenda": _slide_agenda,
             "section": _slide_section, "message": _slide_message,
             "table": _slide_table, "chart": _slide_chart, "kpi": _slide_kpi,
             "compare": _slide_compare, "closing": _slide_closing,
             # 旧称
             "text": _slide_message}


# =============================================================================
# 組み立て
# =============================================================================

def build_pptx(slides: list[dict], title: str | None = None,
          subtitle: str | None = None, footer: str | None = None,
          agenda: bool = True) -> bytes:
    """スライド定義のリストから .pptx のバイト列を作る。"""
    if not slides:
        raise PptxReportError("スライドが1枚もありません。")
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    specs = list(slides)
    if title and (specs[0].get("kind") or "").lower() != "title":
        specs.insert(0, {"kind": "title", "title": title, "subtitle": subtitle or ""})

    # 中扉があれば、その並びから目次を自動で作る
    if agenda and not any((s.get("kind") or "") == "agenda" for s in specs):
        sections = [s.get("title", "") for s in specs
                    if (s.get("kind") or "") == "section"]
        if len(sections) >= 2:
            at = 1 if (specs[0].get("kind") or "") == "title" else 0
            specs.insert(at, {"kind": "agenda", "title": "本日の内容",
                              "items": sections})

    page = 0
    for i, spec in enumerate(specs):
        kind = str(spec.get("kind") or "message").lower()
        if kind not in _PPTX_BUILDERS:
            raise PptxReportError(f"{i + 1}枚目: 未対応の種類です: {kind}。"
                              f"使えるのは {', '.join(SLIDE_KINDS)} です。")
        try:
            slide = _PPTX_BUILDERS[kind](prs, spec)
        except PptxReportError:
            raise
        except Exception as e:
            raise PptxReportError(f"{i + 1}枚目（{kind}）の作成に失敗しました: {e}") from e
        if kind not in ("title", "section"):
            page += 1
            _footer(slide, footer or config.REPORT_ORG or "", page)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def pptx_safe_filename(name: str | None, default: str = "report") -> str:
    base = re.sub(r'[\\/:*?"<>|]', "_", str(name or default)).strip() or default
    return base if base.lower().endswith(".pptx") else base + ".pptx"


def outline_pptx(slides: list[dict]) -> list[str]:
    """何が入ったかの一覧（画面とLLMへの報告用）。"""
    labels = {"title": "表紙", "agenda": "目次", "section": "中扉",
              "message": "説明", "text": "説明", "table": "表", "chart": "グラフ",
              "kpi": "KPI", "compare": "比較", "closing": "まとめ"}
    out = []
    for i, s in enumerate(slides, 1):
        kind = str(s.get("kind") or "message").lower()
        extra = ""
        if kind == "chart":
            extra = (f"（{s.get('chart', 'bar')}・{len(s.get('series') or [])}系列）"
                     if not s.get("image") else "（画像）")
        elif kind == "table":
            extra = f"（{len(s.get('rows') or [])}行）"
        out.append(f"{i}. [{labels.get(kind, kind)}{extra}] {s.get('title', '')}")
    return out


# ==========================================================================
# ===== 元 business.py
# 業務でよく聞かれる分析。SQLでは書きにくく、pandas なら素直に書けるもの。
#
# advanced.py が統計の道具箱なのに対して、こちらは「現場の問い」に対応する。
#   期間比較    先月と比べてどうか。落ちた原因はどの区分か
#   ファネル    見積 → 受注 → 請求 → 入金 のどこで落ちているか
#   コホート    いつ始めた人が、どれだけ続いているか
#   併売        何と何が一緒に買われているか
#
# 戻り値の形は advanced.py と同じ {"title", "tables", "notes", "meta"}。
# 画面もLLMも同じ入れ物で受け取れるようにしてある。
# ==========================================================================
import numpy as np
import pandas as pd

# 表の作り方・数値の丸めは統計側と揃える（同じ見た目で出す）
from advanced import AnalysisError, _clean, _df, _out, _table as _advanced_table


def _business_numeric(df: pd.DataFrame, col: str) -> pd.Series:
    return pd.to_numeric(df[col], errors="coerce")


def _pct(a: float, b: float) -> float | None:
    """b に対する a の割合(%)。分母が0のときは出さない。"""
    return round(a / b * 100, 1) if b else None


def _delta_note(cur: float, prev: float, unit: str = "") -> str:
    diff = cur - prev
    rate = f"{diff / prev * 100:+.1f}%" if prev else "—"
    return f"{prev:,.4g}{unit} → {cur:,.4g}{unit}（{diff:+,.4g}{unit} / {rate}）"


# =============================================================================
# 期間比較（前月比・前年同月比）と寄与度分解
# =============================================================================

def compare_periods(columns: list, rows: list, period_col: str, value_col: str,
                    dimension_col: str | None = None,
                    current: str | None = None, previous: str | None = None,
                    qty_col: str | None = None, top: int = 15) -> dict:
    """2つの期間を比べ、差がどこから来たのかまで分解する。

    「先月と比べて売上が5%落ちた」で終わらせず、
    「どの区分が押し下げたのか」「数量が減ったのか単価が下がったのか」まで出す。
    period_col の値は文字列として比べるので、'2026-01' でも '2026年1月' でもよい。
    """
    df = _df(columns, rows)
    for c in (period_col, value_col):
        if c not in df.columns:
            raise AnalysisError(f"列が見つかりません: {c}"
                                f"（ある列: {', '.join(map(str, df.columns))}）")
    df[value_col] = _business_numeric(df, value_col)
    df = df.dropna(subset=[value_col])
    if df.empty:
        raise AnalysisError(f"{value_col} に数値がありません。")

    periods = [str(p) for p in sorted(df[period_col].astype(str).unique())]
    if len(periods) < 2 and not (current and previous):
        raise AnalysisError(
            f"比べるには期間が2つ以上必要です（いま {len(periods)} 個: {'、'.join(periods)}）。"
            "SQL側で2期間ぶんのデータを取ってください。")
    cur = str(current) if current else periods[-1]
    prev = str(previous) if previous else periods[-2]
    for p in (cur, prev):
        if p not in periods:
            raise AnalysisError(f"期間 '{p}' がデータにありません（ある期間: {'、'.join(periods)}）。")

    df["_p"] = df[period_col].astype(str)
    cur_df = df[df["_p"] == cur]
    prev_df = df[df["_p"] == prev]
    cur_total = float(cur_df[value_col].sum())
    prev_total = float(prev_df[value_col].sum())
    diff_total = cur_total - prev_total

    tables = [_advanced_table("全体", ["項目", prev, cur, "差分", "増減率(%)"],
                     [[value_col, round(prev_total, 4), round(cur_total, 4),
                       round(diff_total, 4), _pct(diff_total, abs(prev_total))]])]
    notes = [f"{value_col}: {_delta_note(cur_total, prev_total)}"]

    meta = {"current": cur, "previous": prev,
            "current_total": _clean(cur_total), "previous_total": _clean(prev_total)}

    if dimension_col and dimension_col in df.columns:
        a = prev_df.groupby(dimension_col)[value_col].sum()
        b = cur_df.groupby(dimension_col)[value_col].sum()
        seg = pd.DataFrame({prev: a, cur: b}).fillna(0.0)
        seg["差分"] = seg[cur] - seg[prev]
        seg["増減率(%)"] = np.where(seg[prev] != 0,
                                   (seg["差分"] / seg[prev].abs() * 100).round(1), np.nan)
        # 寄与度 = その区分の差分が、全体の差分のうち何割を占めるか
        seg["寄与度(%)"] = (seg["差分"] / abs(diff_total) * 100).round(1) if diff_total else np.nan
        seg = seg.sort_values("差分")
        show = pd.concat([seg.head(top), seg.tail(top)]).drop_duplicates()
        show = show.sort_values("差分", ascending=False).reset_index()
        cols, rws = _out(show.round(4))
        tables.append(_advanced_table(f"{dimension_col}別の内訳（増減の大きい順）", cols, rws))

        down = seg[seg["差分"] < 0].head(3)
        up = seg[seg["差分"] > 0].tail(3).iloc[::-1]
        if len(down):
            notes.append("押し下げた区分: " + "、".join(
                f"{i}（{r['差分']:+,.4g}"
                + (f" / 全体の変化の{abs(r['寄与度(%)']):.0f}%" if diff_total else "") + "）"
                for i, r in down.iterrows()))
        if len(up):
            notes.append("押し上げた区分: " + "、".join(
                f"{i}（{r['差分']:+,.4g}"
                + (f" / 全体の変化の{abs(r['寄与度(%)']):.0f}%" if diff_total else "") + "）"
                for i, r in up.iterrows()))
        # 新しく出てきた・消えた区分は、増減率だけ見ていると見落とす
        gone = [str(i) for i in seg.index[(seg[prev] > 0) & (seg[cur] == 0)]][:5]
        born = [str(i) for i in seg.index[(seg[prev] == 0) & (seg[cur] > 0)]][:5]
        if gone:
            notes.append(f"{cur} で無くなった{dimension_col}: {'、'.join(gone)}")
        if born:
            notes.append(f"{cur} で新たに出た{dimension_col}: {'、'.join(born)}")

    if qty_col and qty_col in df.columns:
        # 金額の変化を「数量が動いたぶん」と「単価が動いたぶん」に割る
        df[qty_col] = _business_numeric(df, qty_col)
        q0 = float(prev_df[qty_col].sum())
        q1 = float(cur_df[qty_col].sum())
        p0 = prev_total / q0 if q0 else 0.0
        p1 = cur_total / q1 if q1 else 0.0
        vol = (q1 - q0) * p0                      # 数量要因（単価は前期のまま）
        price = (p1 - p0) * q1                    # 単価要因（数量は当期）
        tables.append(_advanced_table("増減の要因分解", ["要因", "金額", "全体の変化に占める割合(%)"],
                             [["数量が変わったぶん", round(vol, 4), _pct(vol, abs(diff_total))],
                              ["単価が変わったぶん", round(price, 4), _pct(price, abs(diff_total))],
                              ["合計", round(vol + price, 4), None]]))
        notes.append(f"数量 {q0:,.4g} → {q1:,.4g}、平均単価 {p0:,.4g} → {p1:,.4g}。"
                     f"変化の内訳は数量 {vol:+,.4g}、単価 {price:+,.4g}。"
                     + ("数量の影響が大きいです。" if abs(vol) > abs(price) else
                        "単価の影響が大きいです。"))
        meta.update({"volume_effect": _clean(vol), "price_effect": _clean(price)})

    notes.append(f"比較した期間: {prev} と {cur}。"
                 "期間の長さや営業日数が違うと単純比較はできません。"
                 "日数が違う場合は1日あたりに直して比べてください。")
    return {"title": f"{value_col} の期間比較（{prev} → {cur}）",
            "tables": tables, "notes": notes, "meta": meta}


# =============================================================================
# ファネル（段階ごとの通過と滞留）
# =============================================================================

def _passed(s: pd.Series) -> pd.Series:
    """その段階を通過したか。日付なら「入っていれば通過」、数値なら0より大。"""
    num = pd.to_numeric(s, errors="coerce")
    if num.notna().mean() >= 0.8:
        return num.fillna(0) > 0
    return s.notna() & (s.astype(str).str.strip() != "")


def funnel_analysis(columns: list, rows: list, steps: list,
                    labels: list | None = None, group_col: str | None = None,
                    date_steps: bool = True) -> dict:
    """段階ごとの通過数・転換率・離脱と、段階間の滞留日数を出す。

    1行 = 1案件。steps には段階を表す列を順に並べる
    （例: 見積日, 受注日, 請求日, 入金日）。値が入っていればその段階を通過した扱い。
    """
    df = _df(columns, rows)
    steps = [s for s in (steps or []) if s]
    if len(steps) < 2:
        raise AnalysisError("steps に段階の列を2つ以上、順番に並べて指定してください。")
    missing = [s for s in steps if s not in df.columns]
    if missing:
        raise AnalysisError(f"列が見つかりません: {', '.join(missing)}"
                            f"（ある列: {', '.join(map(str, df.columns))}）")
    names = list(labels or []) + steps[len(labels or []):]

    flags = pd.DataFrame({s: _passed(df[s]) for s in steps})
    total = len(df)
    counts = [int(flags[s].sum()) for s in steps]

    frows = []
    for i, s in enumerate(steps):
        prev = counts[i - 1] if i else counts[0]
        frows.append([names[i], counts[i],
                      _pct(counts[i], counts[0]),
                      _pct(counts[i], prev) if i else None,
                      (prev - counts[i]) if i else 0])
    tables = [_advanced_table("段階ごとの通過",
                     ["段階", "件数", "最初からの通過率(%)", "直前からの転換率(%)", "離脱数"],
                     frows)]

    notes = [f"対象 {total:,} 件。{names[0]} {counts[0]:,} 件から "
             f"{names[-1]} {counts[-1]:,} 件まで、"
             f"通過率は {_pct(counts[-1], counts[0])}% です。"]
    # いちばん漏れている段階を名指しする。ここが改善の的になる
    drops = [(counts[i - 1] - counts[i], i) for i in range(1, len(steps))]
    if drops:
        worst = max(drops)
        if worst[0] > 0:
            i = worst[1]
            notes.append(f"最も落ちているのは {names[i - 1]} → {names[i]} で、"
                         f"{worst[0]:,} 件（{100 - (_pct(counts[i], counts[i - 1]) or 0):.1f}%）が"
                         "先へ進んでいません。")

    # 段階間の滞留日数。日付として読めるときだけ出す
    if date_steps:
        lag_rows = []
        for i in range(1, len(steps)):
            a = pd.to_datetime(df[steps[i - 1]], errors="coerce")
            b = pd.to_datetime(df[steps[i]], errors="coerce")
            days = (b - a).dt.total_seconds() / 86400
            days = days[days.notna() & (days >= 0)]
            if len(days) >= 3:
                lag_rows.append([f"{names[i - 1]} → {names[i]}", len(days),
                                 round(float(days.mean()), 1),
                                 round(float(days.median()), 1),
                                 round(float(days.quantile(0.9)), 1)])
        if lag_rows:
            tables.append(_advanced_table("段階間の日数", ["区間", "件数", "平均", "中央値", "90%点"],
                                 lag_rows))
            slow = max(lag_rows, key=lambda r: r[3])
            notes.append(f"最も時間がかかるのは {slow[0]} で、中央値 {slow[3]} 日"
                         f"（1割は {slow[4]} 日以上）。")

    if group_col and group_col in df.columns:
        grows = []
        for name, sub in df.groupby(group_col):
            f = pd.DataFrame({s: _passed(sub[s]) for s in steps})
            first, last = int(f[steps[0]].sum()), int(f[steps[-1]].sum())
            grows.append([str(name), first, last, _pct(last, first)])
        grows.sort(key=lambda r: (r[3] is None, r[3]))
        tables.append(_advanced_table(f"{group_col}別の通過率",
                             [group_col, names[0], names[-1], "通過率(%)"], grows))
        if len(grows) >= 2 and grows[0][3] is not None and grows[-1][3] is not None:
            notes.append(f"通過率が最も低いのは {grows[0][0]}（{grows[0][3]}%）、"
                         f"最も高いのは {grows[-1][0]}（{grows[-1][3]}%）。")

    return {"title": "ファネル分析", "tables": tables, "notes": notes,
            "meta": {"steps": names, "counts": counts, "total": total}}


# =============================================================================
# コホート（いつ始めた人が、どれだけ続いているか）
# =============================================================================

def cohort_analysis(columns: list, rows: list, id_col: str, period_col: str,
                    value_col: str | None = None, max_periods: int = 12) -> dict:
    """初回の期でグループ分けし、その後どれだけ残っているかを見る。

    期の並びはデータに出てくる値の昇順で決める。'2026-01' でも '第1四半期' でも、
    並べたときに正しい順になっていれば動く。
    """
    df = _df(columns, rows)
    for c in (id_col, period_col):
        if c not in df.columns:
            raise AnalysisError(f"列が見つかりません: {c}"
                                f"（ある列: {', '.join(map(str, df.columns))}）")
    df = df[df[id_col].notna() & df[period_col].notna()].copy()
    if df.empty:
        raise AnalysisError("対象データがありません。")
    df["_p"] = df[period_col].astype(str)

    order = {p: i for i, p in enumerate(sorted(df["_p"].unique()))}
    if len(order) < 2:
        raise AnalysisError(f"期が1つしかありません（{list(order)}）。"
                            "複数の期にまたがるデータを取ってください。")
    df["_i"] = df["_p"].map(order)
    first = df.groupby(id_col)["_i"].min().rename("_c")
    df = df.join(first, on=id_col)
    df["経過"] = df["_i"] - df["_c"]
    max_periods = max(1, min(int(max_periods or 12), len(order)))
    df = df[df["経過"] < max_periods]

    rev = {i: p for p, i in order.items()}
    people = df.pivot_table(index="_c", columns="経過", values=id_col,
                            aggfunc="nunique", fill_value=0)
    size = people[0] if 0 in people.columns else people.max(axis=1)

    keep = people.div(size, axis=0).mul(100).round(1)
    keep.index = [f"{rev[i]}（{int(size[i])}人)" for i in keep.index]
    keep.columns = [f"+{c}期" for c in keep.columns]
    k = keep.reset_index().rename(columns={"index": "コホート", "_c": "コホート"})
    cols, rws = _out(k)
    tables = [_advanced_table("継続率(%)", cols, rws)]

    p = people.copy()
    p.index = [f"{rev[i]}" for i in p.index]
    p.columns = [f"+{c}期" for c in p.columns]
    c2, r2 = _out(p.reset_index().rename(columns={"index": "コホート", "_c": "コホート"}))
    tables.append(_advanced_table("人数", c2, r2))

    notes = []
    if len(keep.columns) > 1:
        avg = keep.iloc[:, 1].mean()
        notes.append(f"初回の次の期に残っているのは平均 {avg:.1f}% です。")
    if len(keep.columns) > 3:
        avg3 = keep.iloc[:, 3].mean()
        notes.append(f"3期あとに残っているのは平均 {avg3:.1f}%。"
                     + ("落ち方が急なので、初期の定着に手を打つ余地があります。"
                        if avg3 < 30 else "比較的よく定着しています。"))
    # 新しいコホートほど良くなっているか（施策の効果が出ているか）
    if len(keep) >= 3 and len(keep.columns) > 1:
        early, late = keep.iloc[0, 1], keep.iloc[-1, 1]
        if abs(early - late) >= 5:
            notes.append(f"最初のコホート {early:.1f}% に対し、直近は {late:.1f}%。"
                         + ("改善しています。" if late > early else
                            "悪化しています。獲得の質か初期対応を確かめてください。"))

    if value_col and value_col in df.columns:
        df[value_col] = _business_numeric(df, value_col)
        amt = df.pivot_table(index="_c", columns="経過", values=value_col,
                             aggfunc="sum", fill_value=0).round(2)
        amt.index = [f"{rev[i]}" for i in amt.index]
        amt.columns = [f"+{c}期" for c in amt.columns]
        c3, r3 = _out(amt.reset_index().rename(columns={"index": "コホート", "_c": "コホート"}))
        tables.append(_advanced_table(f"{value_col}の合計", c3, r3))
        per = df.groupby("_c")[value_col].sum() / size
        notes.append("1人あたりの累計 " + value_col + ": " + "、".join(
            f"{rev[i]} {v:,.4g}" for i, v in per.items()))

    notes.append("直近のコホートは経過期間が短いぶん、右側のマスが空きます。"
                 "同じ経過期数どうし（縦ではなく列で）比べてください。")
    return {"title": "コホート分析（継続率）", "tables": tables, "notes": notes,
            "meta": {"cohorts": len(keep), "periods": len(keep.columns)}}


# =============================================================================
# 併売（何と何が一緒に買われているか）
# =============================================================================

def market_basket(columns: list, rows: list, transaction_col: str, item_col: str,
                  min_support: float = 1.0, top: int = 25,
                  max_items: int = 60) -> dict:
    """同じ伝票に一緒に入っている品目の組み合わせを見つける。

    リフト値は「たまたま一緒になる確率」に対して何倍かを表す。
    1.0 を大きく超える組み合わせが、置き場所や提案の手がかりになる。
    """
    df = _df(columns, rows)
    for c in (transaction_col, item_col):
        if c not in df.columns:
            raise AnalysisError(f"列が見つかりません: {c}"
                                f"（ある列: {', '.join(map(str, df.columns))}）")
    d = df[[transaction_col, item_col]].dropna().astype(str).drop_duplicates()
    n_tx = d[transaction_col].nunique()
    if n_tx < 10:
        raise AnalysisError(f"伝票が {n_tx} 件しかありません。10件以上必要です。")

    freq = d[item_col].value_counts()
    # 組み合わせの数は品目数の2乗で増える。よく出るものだけに絞って現実的な時間に収める
    keep = list(freq.head(max(2, int(max_items))).index)
    d = d[d[item_col].isin(keep)]
    cut = len(freq) - len(keep)

    baskets = d.groupby(transaction_col)[item_col].apply(set)
    baskets = baskets[baskets.map(len) >= 2]
    if baskets.empty:
        raise AnalysisError("2品目以上入っている伝票がありません。"
                            "1伝票1明細のデータになっていないか確認してください。")

    pair_count: dict = {}
    for items in baskets:
        picked = sorted(items)
        for i, a in enumerate(picked):
            for b in picked[i + 1:]:
                pair_count[(a, b)] = pair_count.get((a, b), 0) + 1

    out = []
    for (a, b), c in pair_count.items():
        support = c / n_tx * 100
        if support < float(min_support or 0):
            continue
        ca, cb = int(freq[a]), int(freq[b])
        conf_ab = c / ca * 100 if ca else 0.0
        conf_ba = c / cb * 100 if cb else 0.0
        lift = (c / n_tx) / ((ca / n_tx) * (cb / n_tx)) if ca and cb else 0.0
        out.append([a, b, c, round(support, 2), round(conf_ab, 1), round(conf_ba, 1),
                    round(lift, 2)])
    if not out:
        raise AnalysisError(f"支持度 {min_support}% 以上の組み合わせがありませんでした。"
                            "min_support を下げてください。")
    out.sort(key=lambda r: r[6], reverse=True)
    shown = out[: max(1, int(top))]

    tables = [_advanced_table("よく一緒に買われる組み合わせ",
                     ["品目A", "品目B", "同時件数", "支持度(%)",
                      "AならBも(%)", "BならAも(%)", "リフト"], shown)]
    tables.append(_advanced_table("よく出る品目", ["品目", "伝票数", "出現率(%)"],
                        [[i, int(c), round(c / n_tx * 100, 1)]
                         for i, c in freq.head(15).items()]))

    notes = [f"対象 {n_tx:,} 伝票、うち2品目以上入っているのは {len(baskets):,} 伝票です。"]
    if cut > 0:
        notes.append(f"品目が多いため、出現の多い上位 {len(keep)} 品目に絞って計算しました"
                     f"（{cut} 品目を除外）。")
    best = shown[0]
    notes.append(f"最も結びつきが強いのは「{best[0]}」と「{best[1]}」で、リフト {best[6]}倍。"
                 f"{best[0]}を買った人の {best[4]}% が{best[1]}も買っています。")
    notes.append("リフトは「たまたま一緒になる確率」に対する倍率です。1.0前後なら関係なし、"
                 "2.0を超えると強い結びつきと見ます。ただし件数が少ない組は偶然でも"
                 "大きな値になるので、同時件数も併せて見てください。")
    return {"title": "併売分析", "tables": tables, "notes": notes,
            "meta": {"transactions": int(n_tx), "pairs": len(out)}}


# ==========================================================================
# ===== 元 importer.py
# Excel / CSV を SQLite に取り込む層。
#
# このアプリで **唯一 DB に書き込む場所**。分析側（db.py）は読み取り専用のまま保つ。
# 書き込みをここに閉じ込めることで、「チャットからDBが書き換わることはない」という
# 保証を壊さずに、DB・テーブルの新規作成を提供する。
#
# 安全のための制約:
#   - 読むファイルは config.IMPORT_DIRS の中にあるものだけ（画面からパスは打たせない）
#   - シンボリックリンクや .. で許可フォルダの外に出ようとしたら拒否
#   - 拡張子・ファイルサイズ・行数の上限あり
#   - 作る .db は config.DATA_DIR の直下のみ。名前も英数字系に正規化する
# ==========================================================================
import re
import sqlite3
import unicodedata
from datetime import datetime
from pathlib import Path

import pandas as pd
import yaml

import config

# SQLiteの予約語のうち、テーブル名・列名に使われがちなもの
_RESERVED = {
    "abort", "action", "add", "all", "alter", "and", "as", "asc", "between", "by", "case",
    "check", "column", "commit", "create", "cross", "default", "delete", "desc", "distinct",
    "drop", "else", "end", "escape", "except", "exists", "for", "from", "full", "group",
    "having", "in", "index", "inner", "insert", "into", "is", "join", "key", "left", "like",
    "limit", "not", "null", "offset", "on", "or", "order", "outer", "primary", "references",
    "right", "select", "set", "table", "then", "to", "transaction", "union", "unique",
    "update", "using", "values", "view", "when", "where", "with",
}


class ImportError_(Exception):
    """取り込みに失敗したときに投げる（画面にそのまま出せる日本語メッセージ）。"""


# =============================================================================
# 取り込み元ファイルの列挙（許可フォルダの中だけ）
# =============================================================================

def _read_extra() -> list[str]:
    p = config.IMPORT_DIRS_FILE
    if not p.exists():
        return []
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}
    except Exception as e:
        print(f"[importer] 追加フォルダの設定を読めませんでした: {p} ({e})")
        return []
    items = data.get("dirs") if isinstance(data, dict) else data
    return [str(x) for x in (items or []) if str(x).strip()]


def extra_dirs() -> list[Path]:
    """画面から追加されたフォルダ。"""
    return [Path(s).expanduser() for s in _read_extra()]


def configured_dirs() -> list[dict]:
    """許可フォルダの一覧（どこで設定されたかつき）。"""
    out = [{"path": d, "source": "env", "removable": False} for d in config.IMPORT_DIRS]
    for d in extra_dirs():
        out.append({"path": d, "source": "ui", "removable": True})
    return out


def add_dir(raw: str) -> Path:
    """画面からフォルダを追加する。存在と読み取り可否をその場で確かめる。"""
    if not config.IMPORT_DIRS_EDITABLE:
        raise ImportError_("画面からのフォルダ追加は無効化されています（IMPORT_DIRS_EDITABLE）。")
    text = (raw or "").strip().strip('"')
    if not text:
        raise ImportError_("フォルダのパスを入力してください。")
    p = Path(text).expanduser()
    try:
        real = p.resolve(strict=True)
    except OSError:
        raise ImportError_(f"見つかりません: {text}（マウントされているか確認してください）") from None
    if not real.is_dir():
        raise ImportError_(f"フォルダではありません: {real}")
    # ルート直下を丸ごと許可すると、走査が終わらないうえ事故のもとになる
    if real.parent == real:
        raise ImportError_("ドライブ/ファイルシステムのルートは指定できません。"
                           "取り込み用のフォルダを切って指定してください。")
    try:
        next(real.iterdir(), None)
    except PermissionError:
        raise ImportError_(f"読み取り権限がありません: {real}") from None
    except OSError as e:
        raise ImportError_(f"アクセスできません: {real}（{e.strerror or e}）") from None

    current = _read_extra()
    if any(Path(s).expanduser().resolve() == real for s in current
           if Path(s).expanduser().exists()):
        raise ImportError_("そのフォルダは既に登録されています。")
    if any(d.resolve() == real for d in config.IMPORT_DIRS if d.exists()):
        raise ImportError_("env の IMPORT_DIRS に既に入っています。")

    current.append(str(real))
    _write_extra(current)
    return real


def remove_dir(raw: str) -> bool:
    if not config.IMPORT_DIRS_EDITABLE:
        raise ImportError_("画面からのフォルダ変更は無効化されています。")
    current = _read_extra()
    left = [s for s in current if s != raw]
    if len(left) == len(current):
        return False
    _write_extra(left)
    return True


def _write_extra(items: list[str]) -> None:
    p = config.IMPORT_DIRS_FILE
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(yaml.safe_dump({"dirs": items}, allow_unicode=True, sort_keys=False),
                 encoding="utf-8")


def allowed_dirs() -> list[Path]:
    """実際に読み込みを許可するフォルダ（env + 画面から追加した分）。"""
    out, seen = [], set()
    for d in list(config.IMPORT_DIRS) + extra_dirs():
        try:
            real = d.resolve()
        except OSError:
            continue
        if real not in seen:
            seen.add(real)
            out.append(real)
    return out


def dir_status() -> list[dict]:
    """許可フォルダごとの状態。

    本番ではネットワークマウント（/mnt/... など）を指すため、
    「ファイルが無い」のか「マウントされていない・権限がない」のかを
    画面で切り分けられるようにする。
    """
    rows = []
    for entry in configured_dirs():
        d = entry["path"]
        info = {"設定値": str(d), "実際のパス": "", "状態": "", "ok": False,
                "source": entry["source"], "removable": entry["removable"]}
        try:
            real = d.resolve()
            info["実際のパス"] = str(real)
            if not real.exists():
                info["状態"] = "見つかりません（マウントされていない可能性があります）"
            elif not real.is_dir():
                info["状態"] = "フォルダではありません"
            else:
                next(real.iterdir(), None)      # 読めるかどうかを実際に試す
                info["状態"] = "利用できます"
                info["ok"] = True
        except PermissionError:
            info["状態"] = "読み取り権限がありません"
        except OSError as e:
            info["状態"] = f"アクセスできません（{e.strerror or e}）"
        rows.append(info)
    return rows


def _importer_walk(root: Path, depth: int = 0, only_supported: bool = True):
    """root 以下を走査する。depth=0 なら階層の制限なし。

    権限の無いフォルダは黙って飛ばす（共有フォルダには必ずあるため、
    そこで止まると他のファイルまで見えなくなる）。

    only_supported=False にすると拡張子で絞らない。「何が置いてあるか」を
    調べる用で、読み込みの可否は別に判断する。
    """
    stack = [(root, 0)]
    seen_dirs: set = set()
    while stack:
        cur, level = stack.pop()
        try:
            real = cur.resolve()
            if real in seen_dirs:        # リンクの輪でぐるぐる回らないように
                continue
            seen_dirs.add(real)
            entries = list(cur.iterdir())
        except OSError:
            continue
        for p in entries:
            try:
                if p.is_dir():
                    if not depth or level < depth:
                        stack.append((p, level + 1))
                elif is_noise(p.name):
                    continue
                elif not only_supported or p.suffix.lower() in config.IMPORT_EXTENSIONS:
                    yield p
            except OSError:
                continue


def is_noise(name: str) -> bool:
    """Windows共有フォルダに必ず混ざる、開いても意味の無いファイル。

    ~$売上.xlsx … Excelで開いている間だけできるロックファイル（開くと壊れて見える）
    .DS_Store / desktop.ini / Thumbs.db … OSが勝手に作るもの
    """
    low = name.lower()
    return (name.startswith("~$") or name.startswith(".")
            or low in ("desktop.ini", "thumbs.db"))


def is_allowed(path: Path) -> bool:
    """許可フォルダの中にある実ファイルかどうか（.. やリンク経由の脱出を防ぐ）。"""
    try:
        real = path.resolve(strict=True)
    except (OSError, RuntimeError):
        return False
    if not real.is_file() or is_noise(real.name):
        return False
    # 拡張子は大小を無視する（Windows側では .CSV と .csv が混在しうる）
    if real.suffix.lower() not in config.IMPORT_EXTENSIONS:
        return False
    return any(real == d or d in real.parents for d in allowed_dirs())


def list_source_files() -> list[Path]:
    """取り込める候補ファイル（許可フォルダ配下を IMPORT_SCAN_DEPTH 階層まで探す）。

    件数は IMPORT_MAX_FILES で打ち切る（そこに達したかは呼び出し側で
    len() を見て判断する）。
    """
    seen: set = set()
    found: list[Path] = []
    for d in allowed_dirs():
        try:
            if not d.is_dir():
                continue
        except OSError:
            continue
        for p in _importer_walk(d, config.IMPORT_SCAN_DEPTH):
            try:
                real = p.resolve()
            except OSError:
                continue
            # 同じファイルが複数の許可フォルダから見えることがあるので重複を除く
            if real in seen or not is_allowed(real):
                continue
            seen.add(real)
            found.append(real)
            if len(found) >= config.IMPORT_MAX_FILES:
                return sorted(found)
    return sorted(found)


def is_within_allowed(path: Path) -> bool:
    """許可フォルダの中にある実ファイルか（拡張子は問わない）。

    is_allowed() は「取り込めるファイルか」まで見るので拡張子で弾く。
    こちらは置き場所だけを見る。何が置いてあるかを一覧するためのもので、
    「読んでよいか」は呼び出し側が別に判断すること。
    """
    try:
        real = path.resolve(strict=True)
    except (OSError, RuntimeError):
        return False
    if not real.is_file() or is_noise(real.name):
        return False
    return any(real == d or d in real.parents for d in allowed_dirs())


def list_all_files(depth: int = 0, limit: int | None = None) -> list[Path]:
    """許可フォルダ配下のファイル（拡張子を問わない）。調査用。"""
    cap = limit or config.IMPORT_MAX_FILES
    seen: set = set()
    found: list[Path] = []
    for d in allowed_dirs():
        try:
            if not d.is_dir():
                continue
        except OSError:
            continue
        for p in _importer_walk(d, depth, only_supported=False):
            try:
                real = p.resolve()
            except OSError:
                continue
            if real in seen:
                continue
            seen.add(real)
            found.append(real)
            if len(found) >= cap:
                return sorted(found)
    return sorted(found)


def display_name(path: Path) -> str:
    """画面に出す相対パス（許可フォルダからの位置）。"""
    for d in allowed_dirs():
        try:
            return str(path.relative_to(d))
        except ValueError:
            continue
    return path.name


def is_allowed_dir(path: Path) -> bool:
    """許可フォルダ自身か、その配下のフォルダか。"""
    try:
        real = path.resolve(strict=True)
    except (OSError, RuntimeError):
        return False
    if not real.is_dir():
        return False
    return any(real == d or d in real.parents for d in allowed_dirs())


def browse(path: str | None = None) -> dict:
    """フォルダの中身を1階層ぶん返す（エクスプローラ風の画面用）。

    path を省略すると許可フォルダの一覧を返す。
    許可フォルダの外は、パスを直接渡されても開かない。
    """
    roots = allowed_dirs()
    if not path:
        return {
            "path": "", "label": "取り込み元フォルダ", "parent": None,
            "dirs": [{"path": str(d), "name": str(d)} for d in roots if d.is_dir()],
            "files": [], "crumbs": [],
        }

    here = Path(path)
    if not is_allowed_dir(here):
        raise ImportError_("そのフォルダは開けません（許可されたフォルダの外です）。")
    here = here.resolve()

    dirs, files = [], []
    try:
        for p in sorted(here.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower())):
            try:
                if p.is_dir():
                    if not is_noise(p.name):
                        dirs.append({"path": str(p), "name": p.name})
                elif p.suffix.lower() in config.IMPORT_EXTENSIONS and not is_noise(p.name):
                    files.append({"path": str(p), "name": p.name,
                                  "size": p.stat().st_size,
                                  "mtime": datetime.fromtimestamp(
                                      p.stat().st_mtime).strftime("%Y-%m-%d %H:%M")})
            except OSError:
                continue
    except PermissionError:
        raise ImportError_(f"読み取り権限がありません: {here}") from None
    except OSError as e:
        raise ImportError_(f"開けませんでした: {here}（{e.strerror or e}）") from None

    # パンくず。許可フォルダより上には遡らせない
    root = next((d for d in roots if here == d or d in here.parents), None)
    crumbs, cur = [], here
    while root is not None and cur != root:
        crumbs.append({"path": str(cur), "name": cur.name})
        cur = cur.parent
    crumbs.append({"path": str(root), "name": str(root)})
    crumbs.reverse()
    parent = str(here.parent) if (root is not None and here != root) else ""
    return {"path": str(here), "label": here.name or str(here), "parent": parent,
            "dirs": dirs, "files": files, "crumbs": crumbs}


def check_readable(path: Path) -> None:
    if not is_allowed(path):
        raise ImportError_("許可されたフォルダの中のファイルではありません。")
    mb = path.stat().st_size / (1024 * 1024)
    if mb > config.IMPORT_MAX_FILE_MB:
        raise ImportError_(
            f"ファイルが大きすぎます（{mb:.1f}MB / 上限 {config.IMPORT_MAX_FILE_MB}MB）。")


# =============================================================================
# 読み込み
# =============================================================================

def sheet_names(path: Path) -> list[str]:
    """Excelのシート名。CSVなら空リスト。"""
    if path.suffix.lower() not in (".xlsx", ".xlsm"):
        return []
    check_readable(path)
    return _sheet_names_of(path)


def _sheet_names_of(src) -> list[str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
        try:
            return list(wb.sheetnames)
        finally:
            wb.close()
    except Exception as e:
        raise ImportError_(f"Excelを開けませんでした: {e}") from e


# CSVの文字コード。上から順に試す。
CSV_ENCODINGS = ["utf-8-sig", "cp932", "utf-8", "shift_jis", "euc_jp"]


# --- アップロードされたファイル（サーバのフォルダには置かない） --------------------

def check_upload(data: bytes, filename: str) -> str:
    """アップロードの受け入れ判定。戻り値は正規化した拡張子。"""
    if not config.IMPORT_ALLOW_UPLOAD:
        raise ImportError_("アップロードからの取り込みは無効化されています（IMPORT_ALLOW_UPLOAD）。")
    ext = Path(filename or "").suffix.lower()
    if ext not in config.IMPORT_EXTENSIONS:
        raise ImportError_(
            f"扱えない形式です（{ext or '拡張子なし'}）。{'、'.join(config.IMPORT_EXTENSIONS)} のみ対応です。")
    mb = len(data) / (1024 * 1024)
    if mb > config.IMPORT_MAX_FILE_MB:
        raise ImportError_(
            f"ファイルが大きすぎます（{mb:.1f}MB / 上限 {config.IMPORT_MAX_FILE_MB}MB）。")
    return ext


def upload_sheet_names(data: bytes, filename: str) -> list[str]:
    if check_upload(data, filename) not in (".xlsx", ".xlsm"):
        return []
    import io
    return _sheet_names_of(io.BytesIO(data))


def read_upload(data: bytes, filename: str, sheet: str | None = None, header_row: int = 0,
                delimiter: str | None = None, nrows: int | None = None) -> pd.DataFrame:
    """アップロードされたバイト列を DataFrame として読む。ディスクには書かない。"""
    import io
    ext = check_upload(data, filename)
    try:
        if ext in (".xlsx", ".xlsm"):
            return pd.read_excel(io.BytesIO(data), sheet_name=sheet or 0,
                                 header=header_row, nrows=nrows, dtype=object)
        sep = delimiter if delimiter else ("\t" if ext == ".tsv" else None)
        last = None
        for enc in CSV_ENCODINGS:
            try:
                return pd.read_csv(io.BytesIO(data), header=header_row, nrows=nrows,
                                   dtype=object, sep=sep, engine="python", encoding=enc)
            except UnicodeDecodeError as e:
                last = e
        raise ImportError_(
            "文字コードを判定できませんでした。UTF-8 か Shift_JIS で保存し直してください。"
            f"（{last}）")
    except ImportError_:
        raise
    except Exception as e:
        raise ImportError_(f"ファイルを読めませんでした: {e}") from e


# 画面に出す区切り文字の選択肢（.txt は区切りがまちまちなので選べるようにする）
DELIMITERS = {
    "自動判定": None,
    "カンマ ( , )": ",",
    "タブ": "\t",
    "パイプ ( | )": "|",
    "セミコロン ( ; )": ";",
    "空白（連続もまとめる）": r"\s+",
}


def _explain_read_error(e: Exception, path: Path, sheet: str | None) -> str:
    """pandas / OS の例外を、管理者がそのまま対処できる日本語にする。

    定期取り込みの失敗はメール・⚠マーク・AIの注記にこの文がそのまま載るので、
    'No columns to parse from file' のような英語のままでは何をすればよいか分からない。
    """
    name = path.name
    if isinstance(e, PermissionError):
        return (f"{name} を開けません。他のプログラム（Excel など）で開かれているか、"
                "読み取り権限がありません。閉じてから、次回の実行を待つか「今すぐ更新」してください。")
    if isinstance(e, FileNotFoundError):
        return f"{name} が見つかりません（移動・削除された可能性）。"
    msg = str(e)
    if "No columns to parse" in msg or isinstance(e, pd.errors.EmptyDataError):
        return f"{name} の中身が空です（0バイト、または見出し行がありません）。"
    if "Worksheet named" in msg or "Worksheet index" in msg:
        try:
            names = _sheet_names_of(path)
            have = "、".join(names) if names else "（なし）"
        except Exception:
            have = "（不明）"
        return (f"シート「{sheet}」が {name} にありません（シート名が変わった可能性）。"
                f"いまあるシート: {have}。設定のシートを直してください。")
    if "BadZipFile" in type(e).__name__ or "not a zip file" in msg.lower() or "File is not a zip file" in msg:
        return (f"{name} を Excel ファイルとして開けません（壊れているか、拡張子だけ .xlsx の別形式）。"
                "Excel で開いて保存し直してください。")
    if isinstance(e, pd.errors.ParserError):
        return f"{name} を表として読めませんでした（行ごとの列数が揃っていない等）: {msg}"
    return f"{name} を読めませんでした: {msg}"


def read_table(path: Path, sheet: str | None = None, header_row: int = 0,
               encoding: str | None = None, nrows: int | None = None,
               delimiter: str | None = None) -> pd.DataFrame:
    """ファイルを DataFrame として読む。

    header_row は0始まり。見出しが2行目にあるなら 1 を渡す。
    delimiter は CSV/TSV/TXT 用。None なら .tsv はタブ、それ以外は自動判定。
    """
    check_readable(path)
    ext = path.suffix.lower()
    try:
        if ext in (".xlsx", ".xlsm"):
            df = pd.read_excel(path, sheet_name=sheet or 0, header=header_row,
                               nrows=nrows, dtype=object)
        else:
            if path.stat().st_size == 0:
                raise ImportError_(f"{path.name} の中身が空です（0バイト）。")
            sep = delimiter if delimiter else ("\t" if ext == ".tsv" else None)
            last = None
            for enc in ([encoding] if encoding else CSV_ENCODINGS):
                try:
                    df = pd.read_csv(path, header=header_row, nrows=nrows, dtype=object,
                                     sep=sep, engine="python", encoding=enc)
                    break
                except UnicodeDecodeError as e:
                    last = e
            else:
                raise ImportError_(
                    f"{path.name} の文字コードを判定できませんでした。テキスト（CSV）ではないか、"
                    "壊れている可能性があります。UTF-8 か Shift_JIS で保存し直してください。"
                    f"（{last}）")
    except ImportError_:
        raise
    except Exception as e:
        raise ImportError_(_explain_read_error(e, path, sheet)) from e

    if df.empty and not len(df.columns):
        raise ImportError_(f"{path.name} の中身が空のようです（見出し行の指定を確認してください）。")
    return df


# =============================================================================
# 名前と型の正規化
# =============================================================================

def safe_name(name: str, fallback: str = "col") -> str:
    """SQLiteで扱いやすい識別子にする（日本語はそのまま残す）。

    記号と空白を _ にし、数字始まり・予約語・空文字を避ける。
    """
    s = unicodedata.normalize("NFKC", str(name)).strip()
    s = re.sub(r"[^\w]", "_", s, flags=re.UNICODE)   # \w は日本語も含む
    s = re.sub(r"_+", "_", s).strip("_")
    if not s:
        return fallback
    if s[0].isdigit():
        s = "_" + s
    if s.lower() in _RESERVED:
        s = s + "_"
    return s[:64]


def unique_names(names: list[str]) -> list[str]:
    """列名の重複を _2, _3 … で解消する。"""
    out, used = [], {}
    for i, n in enumerate(names):
        base = safe_name(n, fallback=f"col{i + 1}")
        if base in used:
            used[base] += 1
            base = f"{base}_{used[base]}"
        else:
            used[base] = 1
        out.append(base)
    return out


def infer_type(series: pd.Series) -> str:
    """列の中身から SQLite の型を決める（判断できなければ TEXT）。"""
    s = series.dropna()
    s = s[s.astype(str).str.strip() != ""]
    if s.empty:
        return "TEXT"
    num = pd.to_numeric(s, errors="coerce")
    if num.notna().all():
        # 小数点を含まず整数で表せるなら INTEGER
        if (num == num.round()).all() and num.abs().max() < 2 ** 63:
            return "INTEGER"
        return "REAL"
    return "TEXT"


def plan_columns(df: pd.DataFrame) -> list[dict]:
    """列ごとの「元の名前 / 使う名前 / 型」の一覧を作る。"""
    names = unique_names([str(c) for c in df.columns])
    return [{"元の列名": str(orig), "列名": name, "型": infer_type(df[orig])}
            for orig, name in zip(df.columns, names)]


def _cast(series: pd.Series, sqlite_type: str):
    """SQLiteに渡せる素のPython値（int / float / str / None）に揃える。

    numpy の int64 などをそのまま渡すと、sqlite3 がバッファとみなして
    BLOB で保存してしまう（数値として比較も集計もできなくなる）。
    """
    if sqlite_type == "INTEGER":
        num = pd.to_numeric(series, errors="coerce")
        return num.map(lambda v: None if pd.isna(v) else int(v))
    if sqlite_type == "REAL":
        num = pd.to_numeric(series, errors="coerce")
        return num.map(lambda v: None if pd.isna(v) else float(v))
    return series.map(lambda v: None if v is None or pd.isna(v) else str(v))


def prepare_frame(df: pd.DataFrame, columns: list[dict]):
    """型を当てはめた DataFrame と、TEXTに落とした列名の一覧を返す。

    型は先頭数千行から推定するので、後ろの行に数値でない値が混ざることがある。
    そのまま数値に変換すると黙ってNULLになって値が消えるため、
    1件でも変換できない値があればその列は TEXT に落とす。
    """
    out, degraded = {}, []
    for c in columns:
        src = df[c["元の列名"]]
        if c["型"] in ("INTEGER", "REAL"):
            num = pd.to_numeric(src, errors="coerce")
            filled = src.notna() & (src.astype(str).str.strip() != "")
            if bool((filled & num.isna()).any()):
                c["型"] = "TEXT"
                degraded.append(c["列名"])
        out[c["列名"]] = _cast(src, c["型"])
    return pd.DataFrame(out), degraded


# =============================================================================
# 書き込み（このアプリで唯一DBに書く場所）
# =============================================================================

def db_path_for(name: str) -> Path:
    """新しいDBファイルのパス。data/ の直下に限定する。"""
    stem = safe_name(name, fallback="")
    if not stem:
        raise ImportError_("DB名を入力してください（英数字・かな・漢字が使えます）。")
    return (config.DATA_DIR / f"{stem}.db")


def _importer_qi(name: str) -> str:
    return '"' + str(name).replace('"', '""') + '"'


def existing_tables(db_path: Path) -> list[str]:
    if not db_path.exists():
        return []
    conn = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    try:
        rows = conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table' "
            "AND name NOT LIKE 'sqlite_%' ORDER BY name").fetchall()
        return [r[0] for r in rows]
    finally:
        conn.close()


def table_columns(db_path: Path, table: str) -> list[str]:
    if not db_path.exists():
        return []
    conn = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    try:
        return [r[1] for r in conn.execute(f"PRAGMA table_info({_importer_qi(table)})")]
    except sqlite3.Error:
        return []
    finally:
        conn.close()


def import_dataframe(db_path: Path, table: str, df: pd.DataFrame, columns: list[dict],
                     mode: str = "create", timestamp_col: str | None = None,
                     timestamp_value: str | None = None):
    """DataFrame を1テーブルとして書き込む。

    mode: create=新規作成 / replace=作り直す / append=既存に追記
    timestamp_col: 指定すると、その名前の列に取り込み日時を入れて一緒に書く。
                   追記を重ねたとき「いつ取り込んだ分か」を後から絞れるようにするため。
    戻り値: (書き込んだ行数, 型をTEXTに落とした列名の一覧)
    """
    if db_path.parent.resolve() != config.DATA_DIR.resolve():
        raise ImportError_("DBファイルは data/ の直下にしか作れません。")
    table = safe_name(table, fallback="")
    if not table:
        raise ImportError_("テーブル名を入力してください。")
    if len(df) > config.IMPORT_MAX_ROWS:
        raise ImportError_(
            f"行数が多すぎます（{len(df):,}行 / 上限 {config.IMPORT_MAX_ROWS:,}行）。")
    if not columns:
        raise ImportError_("取り込む列がありません。")

    data, degraded = prepare_frame(df, columns)

    write_cols = [{"列名": c["列名"], "型": c["型"]} for c in columns]
    if timestamp_col:
        ts_name = safe_name(timestamp_col, fallback="取得日時")
        if ts_name in {c["列名"] for c in write_cols}:
            raise ImportError_(
                f"取得日時の列名 '{ts_name}' が元データの列名とぶつかっています。別の名前にしてください。")
        stamp = timestamp_value or datetime.now().isoformat(timespec="seconds")
        data[ts_name] = stamp
        write_cols.append({"列名": ts_name, "型": "TEXT"})

    config.DATA_DIR.mkdir(parents=True, exist_ok=True)
    # timeout: 裏のスケジューラと画面からの手動更新が重なっても、即エラーにせず順番待ちする
    conn = sqlite3.connect(db_path, timeout=30)
    try:
        have = table in existing_tables(db_path) if db_path.exists() else False
        if mode == "create" and have:
            raise ImportError_(f"テーブル '{table}' は既にあります。"
                               "「作り直す」か「追記する」を選ぶか、別の名前にしてください。")
        if mode == "replace":
            conn.execute(f"DROP TABLE IF EXISTS {_importer_qi(table)}")
            have = False
        if not have:
            cols_sql = ", ".join(f"{_importer_qi(c['列名'])} {c['型']}" for c in write_cols)
            conn.execute(f"CREATE TABLE {_importer_qi(table)} ({cols_sql})")
        else:
            # 追記先に無い列があると INSERT が落ちるので、先に照合して分かる形で止める。
            # 取得日時だけは後から足せるので ALTER で追加する。
            have_cols = set(table_columns(db_path, table))
            missing = [c["列名"] for c in write_cols if c["列名"] not in have_cols]
            if timestamp_col and safe_name(timestamp_col, "取得日時") in missing:
                ts_name = safe_name(timestamp_col, "取得日時")
                conn.execute(f"ALTER TABLE {_importer_qi(table)} ADD COLUMN {_importer_qi(ts_name)} TEXT")
                missing.remove(ts_name)
            if missing:
                raise ImportError_(
                    f"追記先の '{table}' に無い列があります: {', '.join(missing)}。"
                    "列名を合わせるか、「作り直す」を選んでください。")

        placeholders = ", ".join("?" for _ in write_cols)
        cols_list = ", ".join(_importer_qi(c["列名"]) for c in write_cols)
        rows = list(data[[c["列名"] for c in write_cols]]
                    .itertuples(index=False, name=None))   # _cast で素の値に揃え済み
        conn.executemany(
            f"INSERT INTO {_importer_qi(table)} ({cols_list}) VALUES ({placeholders})", rows)
        conn.commit()
        return len(rows), degraded
    except sqlite3.Error as e:
        conn.rollback()
        raise ImportError_(f"書き込みに失敗しました: {e}") from e
    finally:
        conn.close()


def prune_runs(db_path: Path, table: str, timestamp_col: str, keep: int) -> int:
    """取得日時の新しい keep 回分だけ残し、それより古い回を削除する。

    「回」は取得日時の値の種類で数える（1回の取り込みで入った行は同じ値を持つ）。
    取得日時が NULL の行 ―― この仕組みを入れる前から入っていた行 ―― は消さない。
    戻り値は削除した行数。
    """
    keep = int(keep)
    if keep < 1 or not timestamp_col:
        return 0
    table, ts = safe_name(table), safe_name(timestamp_col, "取得日時")
    if ts not in table_columns(db_path, table):
        return 0
    conn = sqlite3.connect(db_path, timeout=30)
    try:
        cur = conn.execute(
            f"DELETE FROM {_importer_qi(table)} "
            f"WHERE {_importer_qi(ts)} IS NOT NULL AND {_importer_qi(ts)} NOT IN "
            f"(SELECT {_importer_qi(ts)} FROM {_importer_qi(table)} WHERE {_importer_qi(ts)} IS NOT NULL "
            f" GROUP BY {_importer_qi(ts)} ORDER BY {_importer_qi(ts)} DESC LIMIT ?)", (keep,))
        removed = cur.rowcount or 0
        conn.commit()
        return removed
    except sqlite3.Error as e:
        conn.rollback()
        raise ImportError_(f"古い取り込み分の削除に失敗しました: {e}") from e
    finally:
        conn.close()


def table_info(db_path: Path, table: str, timestamp_col: str | None = None) -> dict:
    """1テーブルの中身の要約。「DBの管理」画面で状態を確かめるために使う。

    取得日時の列は、ジョブに設定があればそれを、無ければ既定の列名を探す
    （画面から手で取り込んだテーブルにも付いているため）。
    """
    cols = table_columns(db_path, table)
    ts = None
    for cand in (timestamp_col, config.IMPORT_TIMESTAMP_COLUMN):
        if cand and safe_name(cand, "") in cols:
            ts = safe_name(cand, "")
            break

    info = {"name": table, "columns": cols, "column_count": len(cols),
            "rows": 0, "timestamp_column": ts, "runs": None,
            "latest": None, "oldest": None}
    if not db_path.exists():
        return info
    conn = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    try:
        info["rows"] = conn.execute(f"SELECT COUNT(*) FROM {_importer_qi(table)}").fetchone()[0]
        if ts:
            row = conn.execute(
                f"SELECT COUNT(DISTINCT {_importer_qi(ts)}), MIN({_importer_qi(ts)}), MAX({_importer_qi(ts)}) "
                f"FROM {_importer_qi(table)} WHERE {_importer_qi(ts)} IS NOT NULL").fetchone()
            info["runs"], info["oldest"], info["latest"] = row[0], row[1], row[2]
    except sqlite3.Error as e:
        info["error"] = str(e)
    finally:
        conn.close()
    return info


def sample_rows(db_path: Path, table: str, limit: int | None = None,
                timestamp_col: str | None = None) -> dict:
    """テーブルの中身を数行だけ覗く（読み取り専用）。

    取得日時の列があれば新しい順に取る。「さっきの取り込みがちゃんと入ったか」を
    確かめるのが主な用途なので、先頭から取ると古い行しか見えず役に立たない。
    """
    table = safe_name(table)
    limit = int(limit or config.IMPORT_SAMPLE_ROWS)
    cols = table_columns(db_path, table)
    out = {"table": table, "columns": cols, "rows": [], "order_by": None, "limit": limit}
    if not db_path.exists() or not cols:
        out["error"] = "テーブルが見つかりません。"
        return out

    ts = None
    for cand in (timestamp_col, config.IMPORT_TIMESTAMP_COLUMN):
        if cand and safe_name(cand, "") in cols:
            ts = safe_name(cand, "")
            break
    order = f" ORDER BY {_importer_qi(ts)} DESC" if ts else ""
    out["order_by"] = ts

    conn = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    try:
        cur = conn.execute(f"SELECT * FROM {_importer_qi(table)}{order} LIMIT ?", (limit,))
        out["columns"] = [d[0] for d in cur.description]
        # BLOB はそのままだとJSONに載らないので、見える形に潰しておく
        out["rows"] = [[v.hex()[:32] if isinstance(v, (bytes, bytearray)) else v
                        for v in row] for row in cur.fetchall()]
    except sqlite3.Error as e:
        out["error"] = str(e)
    finally:
        conn.close()
    return out


def run_count(db_path: Path, table: str, timestamp_col: str) -> int:
    """いま何回分の取り込みが入っているか。"""
    ts = safe_name(timestamp_col or "", "")
    if not ts or ts not in table_columns(db_path, safe_name(table)):
        return 0
    conn = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    try:
        row = conn.execute(
            f"SELECT COUNT(DISTINCT {_importer_qi(ts)}) FROM {_importer_qi(safe_name(table))} "
            f"WHERE {_importer_qi(ts)} IS NOT NULL").fetchone()
        return int(row[0]) if row else 0
    except sqlite3.Error:
        return 0
    finally:
        conn.close()


def drop_table(db_path: Path, table: str) -> None:
    if db_path.parent.resolve() != config.DATA_DIR.resolve():
        raise ImportError_("data/ の外は操作できません。")
    conn = sqlite3.connect(db_path)
    try:
        conn.execute(f"DROP TABLE IF EXISTS {_importer_qi(table)}")
        conn.commit()
    finally:
        conn.close()


# ==========================================================================
# ===== 元 jobs.py
# 定期取り込み（ジョブ）の定義と実行。
#
# 1ジョブ = 「どのファイルを / どう読んで / どのテーブルへ / どの方式で / どの間隔で」。
# 定義は data/import_jobs.yaml に置く（DBファイル自体が全ユーザー共通なのでジョブも共通）。
#
# 実行の入口は3つ。中身はすべて run_job() に集約してある。
#   - 画面の「▶ 今すぐ更新」
#   - 画面を開いたときの自動実行（config.IMPORT_AUTO_REFRESH が true のとき）
#   - cron から refresh.py    ← 本番はこれが確実（誰も画面を開かなくても動く）
# ==========================================================================
import threading
import os
import uuid
from datetime import datetime, timedelta
from pathlib import Path

import yaml

import config
import history
import importer

# 定義ファイルの書き換えとジョブ実行を直列化する。
# 裏で回るスケジューラと、画面からの「▶ 今すぐ更新」が同時に走りうるため。
_jobs_lock = threading.RLock()

# 画面に出す更新間隔。値は分。0 は「手動のみ」。
INTERVALS = {
    "手動のみ": 0,
    "15分ごと": 15,
    "1時間ごと": 60,
    "3時間ごと": 180,
    "6時間ごと": 360,
    "1日ごと": 1440,
    "1週間ごと": 10080,
}
MODES = {
    "replace": "全件入れ替え（毎回すべて削除して入れ直す）",
    "append": "追記（前回までのデータを残して足す）",
}

# 追記のとき「何回分の取り込みを残すか」。これを超えた古い回は消す。
# 上限を決めておかないと、日次で回すだけでもテーブルが際限なく膨らむ。
MAX_KEEP_RUNS = 800
# 既定値は置かない。何回分残すかは業務ごとに違うので、必ず自分で決めてもらう。
DEFAULT_KEEP_RUNS = None
# 開始日時の判定に使う許容。送信のタイムラグで「今」が過去扱いになるのを防ぐ。
START_GRACE_MINUTES = 2


def parse_dt(value) -> datetime | None:
    """画面から来る日時文字列（'2026-08-10T09:00' など）を datetime に。"""
    text = str(value or "").strip()
    if not text:
        return None
    try:
        return datetime.fromisoformat(text)
    except ValueError:
        return None


def validate_job(job: dict, check_start: bool = True) -> list[str]:
    """保存前の点検。画面にそのまま出せる日本語で返す。

    check_start=False にすると開始日時が過去でも通す（登録済みジョブの
    間隔変更や停止/再開など、開始日時を触らない更新のため）。
    """
    errors = []
    if not (job.get("db_file") and job.get("table")):
        errors.append("取り込み先のDBとテーブルを指定してください。")
    if not job.get("source"):
        errors.append("取り込み元のファイルを選んでください。")

    raw_start = str(job.get("start_at") or "").strip()
    if raw_start:
        start = parse_dt(raw_start)
        if start is None:
            errors.append("開始日時の形式が正しくありません。")
        elif check_start and start < datetime.now() - timedelta(minutes=START_GRACE_MINUTES):
            errors.append(f"開始日時に過去の時刻は指定できません"
                          f"（指定: {start:%Y-%m-%d %H:%M}）。今より後の日時にしてください。")

    # 取得日時列は更新の仕方によらず必須。全件入れ替えでも「いつ時点のデータか」が
    # 分からないと、取り込み後の分析で断面を説明できない。
    if not str(job.get("timestamp_column") or "").strip():
        errors.append("取得日時の列名が必須です。")

    if job.get("mode") == "append":
        keep = job.get("keep_runs")
        if keep in (None, ""):
            errors.append("追記のときは保存回数が必須です。")
        else:
            try:
                keep = int(keep)
            except (TypeError, ValueError):
                errors.append("保存回数は数値で指定してください。")
            else:
                if not (1 <= keep <= MAX_KEEP_RUNS):
                    errors.append(f"保存回数は 1〜{MAX_KEEP_RUNS} の範囲で指定してください。")
    return errors


def manual_run_blocked(job: dict) -> str | None:
    """手で走らせてはいけない設定なら、その理由を返す（問題なければ None）。

    定期実行 × 追記 の組み合わせだけは止める。この2つが重なると、
      ・次回予定が「前回実行＋間隔」で決まるので、手で走らせた分だけ後ろにずれる
      ・保存回数を1回ぶん余計に使い、その回だけ間隔の違うデータが混ざる
    となって、せっかく決めた更新頻度が崩れる。
    全件入れ替えや「手動のみ」の設定は、何度走らせても頻度の意味が変わらないので通す。
    """
    if int(job.get("interval_minutes") or 0) <= 0:
        return None
    if (job.get("mode") or "replace") != "append":
        return None
    label = interval_label(job.get("interval_minutes", 0))
    return (f"「{job.get('name') or 'この設定'}」は定期実行（{label}）＋追記です。"
            "手動で動かすと次回の実行時刻がずれ、保存回数も1回ぶん余計に使うため、"
            "手動実行はできません。どうしても今すぐ入れたいときは、"
            "更新の頻度を「手動のみ」に変えてから実行してください。")


def interval_label(minutes: int) -> str:
    for k, v in INTERVALS.items():
        if v == int(minutes or 0):
            return k
    return f"{minutes}分ごと"


# =============================================================================
# 保存と読み出し
# =============================================================================

def _read() -> list[dict]:
    p = config.IMPORT_JOBS_FILE
    if not p.exists():
        return []
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"[jobs] 読めませんでした: {p} ({e})")
        return []
    items = data.get("jobs") if isinstance(data, dict) else data
    return [j for j in (items or []) if isinstance(j, dict) and j.get("id")]


def _write(items: list[dict]) -> None:
    p = config.IMPORT_JOBS_FILE
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(yaml.safe_dump({"jobs": items}, allow_unicode=True, sort_keys=False),
                 encoding="utf-8")


def list_jobs() -> list[dict]:
    return sorted(_read(), key=lambda j: (j.get("name") or ""))


def get_job(job_id: str) -> dict | None:
    return next((j for j in _read() if j.get("id") == job_id), None)


def _same_target(a: dict, b: dict) -> bool:
    """同じ取り込み元（ファイル＋シート）を同じDBの同じテーブルへ入れる設定か。"""
    def norm_path(p):
        return os.path.normcase(os.path.normpath(str(p or "").strip()))
    return (norm_path(a.get("source")) == norm_path(b.get("source"))
            and (a.get("sheet") or None) == (b.get("sheet") or None)
            and (a.get("db_file") or "") == (b.get("db_file") or "")
            and (a.get("table") or "") == (b.get("table") or ""))


def find_duplicate(job: dict) -> dict | None:
    """同じ取り込み元→同じテーブルの設定がすでにあれば、それを返す（自分自身は除く）。

    同じ設定が2つあると、同じ時刻に2回追記されて全行が二重になる
    （「保持N回」は取得日時で数えるので、同時刻の2バッチを1回分とみなして両方残す）。
    登録時に止めるためのもの。
    """
    for j in _read():
        if j.get("id") != job.get("id") and _same_target(j, job):
            return j
    return None


def save_job(job: dict) -> dict:
    with _jobs_lock:
        job = dict(job)
        # setdefault では駄目。呼び出し側が id=None を明示的に入れてくることがあり、
        # そのまま保存すると読み出し時に落とされて「保存したのに消える」ことになる。
        if not job.get("id"):
            job["id"] = uuid.uuid4().hex[:12]
        if not job.get("created_at"):
            job["created_at"] = datetime.now().isoformat(timespec="seconds")
        items = [j for j in _read() if j.get("id") != job["id"]]
        items.append(job)
        _write(items)
        return job


def delete_job(job_id: str) -> bool:
    with _jobs_lock:
        items = _read()
        left = [j for j in items if j.get("id") != job_id]
        if len(left) == len(items):
            return False
        _write(left)
        return True


# =============================================================================
# 実行タイミング
# =============================================================================

def next_run_at(job: dict) -> datetime | None:
    """次に動く予定の時刻。手動のみなら None。

    開始日時が設定されていれば、それより前には動かさない。
    """
    minutes = int(job.get("interval_minutes") or 0)
    if minutes <= 0:
        return None
    start = parse_dt(job.get("start_at"))
    last = parse_dt(job.get("last_run"))
    if last is None:
        # 一度も動いていない。開始日時があればその時刻、無ければすぐ対象。
        return start or datetime.now()
    nxt = last + timedelta(minutes=minutes)
    return max(nxt, start) if start else nxt


def is_due(job: dict, now: datetime | None = None) -> bool:
    if not job.get("enabled", True):
        return False
    nxt = next_run_at(job)
    return nxt is not None and nxt <= (now or datetime.now())


def due_jobs(now: datetime | None = None) -> list[dict]:
    return [j for j in list_jobs() if is_due(j, now)]


# =============================================================================
# 「設定どおりに更新できていない」ジョブ
#
# 失敗は履歴を見に行かないと分からず、日次の取り込みが月曜から失敗して金曜まで
# 誰も気づかない、が起こり得る。そこで「いま健全でないジョブ」を1か所で判定し、
#   ・チャットのサイドバー（DB名・テーブル名に警告マーク）
#   ・AIの回答（そのテーブルを使う質問に、データが古い可能性を添える）
#   ・管理者へのメール通知
# の3つが同じ判断を使う。
# =============================================================================

def problems() -> list[dict]:
    """設定どおりに更新できていない定期取り込み。

    3種類ある:
      failed   … 前回の実行が失敗した（ファイルが無い・シート名や列が変わった等）
      degraded … 取り込めたが、数値列に文字が混ざって文字として保存した
                 （合計・平均がずれる。元ファイルの値を直すべき）
      overdue  … 有効な自動実行なのに、予定の2周期ぶん以上動いていない
                 （スケジューラが止まっている・アプリが落ちていた等）
    戻り値: [{id, name, db_file, table, kind, since, message}, ...]
    """
    now = datetime.now()
    out = []
    for j in list_jobs():
        if not j.get("enabled", True):
            continue
        if j.get("last_status") == "error":
            out.append({"id": j.get("id"), "name": j.get("name"),
                        "db_file": j.get("db_file"), "table": j.get("table"),
                        "kind": "failed", "since": j.get("last_run") or "",
                        "message": j.get("last_message") or "前回の実行が失敗しました。"})
            continue
        if j.get("last_degraded"):
            cols = "、".join(str(c) for c in j["last_degraded"])
            out.append({"id": j.get("id"), "name": j.get("name"),
                        "db_file": j.get("db_file"), "table": j.get("table"),
                        "kind": "degraded", "since": j.get("last_run") or "",
                        "message": (f"前回の取り込みで、数値の列（{cols}）に数値でない値が混ざり、"
                                    "文字として保存しました。合計や平均がずれる可能性があります。"
                                    "元ファイルの値を確認してください。")})
            continue
        minutes = int(j.get("interval_minutes") or 0)
        last = parse_dt(j.get("last_run"))
        if minutes > 0 and last and (now - last) > timedelta(minutes=minutes * 2):
            out.append({"id": j.get("id"), "name": j.get("name"),
                        "db_file": j.get("db_file"), "table": j.get("table"),
                        "kind": "overdue", "since": j.get("last_run") or "",
                        "message": (f"{interval_label(minutes)}の予定ですが、"
                                    f"{last:%m/%d %H:%M} から更新されていません。"
                                    "自動実行が止まっている可能性があります。")})
    return out


def problems_by_table() -> dict:
    """{(db_file, table): [problem, ...]}。画面やAIの注記で引きやすい形。"""
    out: dict = {}
    for p in problems():
        out.setdefault((p["db_file"], p["table"]), []).append(p)
    return out


# =============================================================================
# 実行
# =============================================================================

def source_path(job: dict) -> Path:
    return Path(job.get("source", ""))


def run_job(job: dict, kind: str = "auto", user: str | None = None) -> dict:
    """1ジョブを実行して、結果を定義ファイルに書き戻す。

    kind は履歴に残す実行のきっかけ。"auto"=スケジューラ、"job"=画面の「▶ 今すぐ更新」。

    例外は投げず、結果を dict で返す（1本こけても他を止めないため）。
      {"ok": bool, "rows": int, "message": str, "degraded": [...]}
    """
    with _jobs_lock:                       # 同じテーブルへ同時に書かないように直列化する
        return _run_job_locked(job, kind, user)


def _run_job_locked(job: dict, kind: str = "auto", user: str | None = None) -> dict:
    started = datetime.now()
    result = {"ok": False, "rows": 0, "message": "", "degraded": []}
    removed = 0
    kept = None
    try:
        path = source_path(job)
        if not importer.is_allowed(path):
            raise importer.ImportError_(
                "取り込み元のファイルが見つかりません（移動・削除、または許可フォルダの設定変更）。")

        df = importer.read_table(
            path,
            sheet=job.get("sheet") or None,
            header_row=int(job.get("header_row") or 0),
            delimiter=job.get("delimiter") or None,
        )
        cols = [dict(c) for c in (job.get("columns") or [])]
        if not cols:
            cols = importer.plan_columns(df)
        missing = [c["元の列名"] for c in cols if c["元の列名"] not in df.columns]
        if missing:
            if len(missing) == len(cols):
                # 1列も合わない＝列名の変更ではなく、区切り文字か見出し行の位置が変わった
                found = "、".join(str(c) for c in list(df.columns)[:5])
                raise importer.ImportError_(
                    f"設定した列が1つも見つかりません（ファイル側の見出し: {found}"
                    f"{' …' if len(df.columns) > 5 else ''}）。"
                    "区切り文字・見出し行の位置・シートが変わった可能性があります。"
                    "取り込み画面で開き直して設定を作り直してください。")
            raise importer.ImportError_(
                f"ファイル側に無い列があります: {', '.join(missing)}。"
                "列構成が変わった可能性があります。設定を作り直してください。")

        db_path = config.DATA_DIR / job["db_file"]
        mode = job.get("mode") or "replace"
        ts_col = job.get("timestamp_column") or config.IMPORT_TIMESTAMP_COLUMN
        if len(df) == 0:
            # 見出しだけのファイル（上流の出力が失敗した等）で全件入れ替えすると、
            # テーブルが空になり「成功」で終わる。前回の内容を残して止める。
            # 本当に0件にしたいときは、取り込み画面から手で入れ替える。
            raise importer.ImportError_(
                f"{path.name} にデータ行がありません（見出しだけ）。"
                f"{'テーブルを空にしないため、前回の内容を残しました。' if mode != 'append' else '追記する行が無いため何もしていません。'}"
                "本当に0件なら、取り込み画面から手で入れ替えてください。")
        n, degraded = importer.import_dataframe(
            db_path, job["table"], df, cols, mode=mode,
            timestamp_col=ts_col,
            timestamp_value=started.isoformat(timespec="seconds"),
        )
        message = f"{n:,}行を{'追記' if mode == 'append' else '全件入れ替え'}しました。"
        if degraded:
            # 数値列に文字が混ざった。取り込み自体は通るが集計がずれるので、黙って通さない
            message += (f" ⚠ 数値にできない値があったため文字として保存した列: {', '.join(degraded)}"
                        "（元ファイルの値を確認してください）")

        # 追記のときは、保存回数を超えた古い取り込み分を落とす
        if mode == "append" and ts_col and job.get("keep_runs"):
            removed = importer.prune_runs(db_path, job["table"], ts_col, job["keep_runs"])
            kept = importer.run_count(db_path, job["table"], ts_col)
            result["removed"], result["kept"] = removed, kept
            message += f" 保持 {kept}/{job['keep_runs']}回"
            if removed:
                message += f"（古い {removed:,}行を削除）"
        result.update(ok=True, rows=n, degraded=degraded, message=message)
    except importer.ImportError_ as e:
        result["message"] = str(e)
    except Exception as e:                        # 想定外でもジョブ一覧は壊さない
        result["message"] = f"想定外のエラー: {e}"

    saved = get_job(job.get("id", "")) or dict(job)
    saved.update({
        "last_run": started.isoformat(timespec="seconds"),
        "last_status": "ok" if result["ok"] else "error",
        "last_message": result["message"],
        "last_rows": result["rows"],
        "last_degraded": list(result.get("degraded") or []),
    })
    save_job(saved)
    history.add_import_record(job.get("db_file", ""), job.get("table", ""),
                result["ok"], result["message"], kind=kind,
                mode=job.get("mode") or "replace", rows=result["rows"],
                removed=removed, kept=kept, keep=job.get("keep_runs"),
                source=job.get("source", ""), sheet=job.get("sheet"),
                job_id=job.get("id"), job_name=job.get("name"),
                user=user, started=started)
    return result


def run_due(now: datetime | None = None, kind: str = "auto",
            user: str | None = None) -> list[tuple[dict, dict]]:
    """期限が来たジョブをまとめて実行する。戻り値は (ジョブ, 結果) の一覧。"""
    return [(j, run_job(j, kind, user)) for j in due_jobs(now)]


# ==========================================================================
# ===== 元 scheduler.py
# アプリ内スケジューラ。cron や常駐サービスを別に用意せず、Pythonだけで定期実行する。
#
# アプリ起動時にデーモンスレッドを1本立て、一定間隔で「期限が来たジョブ」を実行する。
# Streamlit のスクリプトは操作のたびに再実行されるが、スレッドはプロセスに1本だけ。
# 画面を誰も開いていなくても、アプリのプロセスが生きていれば動く。
#
#   app.py ──start()──▶ [aiagent-import-scheduler スレッド]
#                           └─ 60秒ごと: jobs.due_jobs() → jobs.run_job()
#
# このスレッドから Streamlit の API（st.*）は呼ばない。
# 画面の描画コンテキストが無いので、状態は _state に置いて画面側から読む。
# ==========================================================================
import atexit
import threading
import traceback
from datetime import datetime

import config
import jobs

# 名前でスレッドの生存を確認する。モジュールが再読込されてもフラグに頼らず
# 二重起動を防げる（開発中にファイルを保存するとStreamlitが読み直すため）。
_THREAD_NAME = "aiagent-import-scheduler"

# 終了時に眠っているスレッドを起こして片付けるための合図。
# sleep() で寝かせたままプロセスを終わらせると、後始末中に標準出力を掴んだままになり
# 「_enter_buffered_busy ... daemon threads」で異常終了することがある。
_stop = threading.Event()

_state: dict = {
    "started_at": None,
    "last_tick": None,
    "tick_count": 0,
    "last_ran": [],        # 直近に実行したジョブ [{name, ok, message, at}]
    "last_error": None,
}


def is_running() -> bool:
    return any(t.name == _THREAD_NAME and t.is_alive() for t in threading.enumerate())


def scheduler_status() -> dict:
    return {**_state, "running": is_running(),
            "tick_sec": config.IMPORT_SCHEDULER_TICK_SEC,
            "enabled": config.IMPORT_SCHEDULER}


def _scheduler_log(msg: str) -> None:
    if _stop.is_set():                 # 終了処理中は標準出力に触らない
        return
    try:
        print(f"[scheduler {datetime.now():%Y-%m-%d %H:%M:%S}] {msg}", flush=True)
    except (ValueError, OSError):      # 出力先が既に閉じられている
        pass


#: 前回の周回で「設定どおりに更新できていなかった」ジョブ。変化を見るために持つ。
_prev_problems: list = []


def tick() -> list:
    """期限が来たジョブを実行する（スレッドの1周分。テストからも呼べる）。"""
    global _prev_problems
    ran = []
    for job in jobs.due_jobs():
        res = jobs.run_job(job)
        ran.append({"name": job.get("name") or job.get("id"), "ok": res["ok"],
                    "message": res["message"],
                    "at": datetime.now().isoformat(timespec="seconds")})
        _scheduler_log(("OK  " if res["ok"] else "NG  ") + f"{job.get('name')}: {res['message']}")
    _state["last_tick"] = datetime.now().isoformat(timespec="seconds")
    _state["tick_count"] += 1
    if ran:
        _state["last_ran"] = ran[-10:]
    # 「健全→失敗」「失敗→復旧」の変わり目だけ管理者に知らせる
    try:
        import mailer
        cur = jobs.problems()
        if _state["tick_count"] > 1:          # 起動直後の1周目は「変化」ではないので送らない
            r = mailer.alert_import_problems(cur, _prev_problems)
            if r:
                _scheduler_log(f"管理者に通知: {r.get('message', '')}")
        _prev_problems = cur
    except Exception as e:
        _scheduler_log(f"通知の判定に失敗（続行）: {e}")
    return ran


def _loop() -> None:
    _scheduler_log(f"開始（{config.IMPORT_SCHEDULER_TICK_SEC}秒ごとに確認）")
    while not _stop.is_set():
        try:
            tick()
            _state["last_error"] = None
        except Exception as e:
            # 1周こけても止めない。止まると以後ずっと更新されなくなるため。
            _state["last_error"] = f"{e}"
            _scheduler_log("巡回でエラー: " + traceback.format_exc(limit=3).replace("\n", " "))
        # sleep ではなく wait。停止の合図が来たら即座に抜ける。
        _stop.wait(max(5, config.IMPORT_SCHEDULER_TICK_SEC))


def stop(timeout: float = 2.0) -> None:
    """スレッドを止める（プロセス終了時に自動で呼ばれる）。"""
    _stop.set()
    for t in threading.enumerate():
        if t.name == _THREAD_NAME and t.is_alive() and t is not threading.current_thread():
            t.join(timeout=timeout)


def start() -> bool:
    """スケジューラを起動する。何度呼んでも1本しか立たない。"""
    if not config.IMPORT_SCHEDULER:
        return False
    if is_running():
        return False
    _stop.clear()
    t = threading.Thread(target=_loop, name=_THREAD_NAME, daemon=True)
    t.start()
    atexit.register(stop)
    _state["started_at"] = datetime.now().isoformat(timespec="seconds")
    return True


# ==========================================================================
# ===== 元 cleanup.py
# テーブル・DBを消したときの後片付け。
#
# 消すこと自体は DROP TABLE とファイル移動で済む。面倒なのはその後で、
# 参照は方々に散らばっている。
#
#   ・そのDBの .meta.yaml  … 説明・関連・用語・例文・検算ルール・ER図の配置
#   ・他のDBの .meta.yaml  … DBをまたぐ関連、ER図に借りたテーブル
#   ・定期取り込みの設定    … 残すと、消したテーブルが次の実行で復活する
#   ・利用者ごとの選択      … 「対象データ」に消えたDBが残り続ける
#
# 放っておくと、AIには存在しないテーブルの説明が渡り続け、例文の検証は
# 「no such table」で落ちる。掃除をここに集めて、消し忘れが出ないようにする。
#
# 消す前の「何が巻き添えになるか」も同じ規則で数える（table_impact / db_impact）。
# 数えるだけの関数は何も書き換えない。
# ==========================================================================
import re
from pathlib import Path

import catalog
import config
import db
import jobs
import prefs
import verify


# =============================================================================
# SQLがそのテーブルを触っているか
# =============================================================================

def uses_table(sql: str, table: str, alias: str | None = None) -> bool:
    """SQL文字列がそのテーブル名を参照しているか。

    "orders" と "demo_sales.orders" の両方を見る。素の名前だけを探すと
    DB名で修飾された書き方（例文はたいていこちら）を取りこぼし、
    修飾ありだけを探すと単一DBの例文を取りこぼす。
    構文解析まではしない。掃除の判定なので、取りこぼすより拾いすぎるほうがまし。
    ただし何を消したかは呼び出し側で必ず報告すること。
    """
    if not table:
        return False
    text = str(sql or "")
    t = r'"?' + re.escape(table) + r'"?(?![\w])'
    if re.search(r'(?<![\w."])' + t, text, re.IGNORECASE):
        return True
    if alias and re.search(r'(?<![\w."])' + re.escape(alias) + r'\s*\.\s*' + t,
                           text, re.IGNORECASE):
        return True
    return False


def _rel_text(rel: dict) -> str:
    return f"{rel.get('from', '')} → {rel.get('to', '')}"


def _ep_hits(rel: dict, own_alias: str, alias: str, table: str | None) -> bool:
    """関連の端点が (alias, table) を指しているか。table=None ならDB丸ごと。"""
    for key in ("from", "to"):
        ep = catalog.parse_endpoint(rel.get(key, ""), own_alias)
        if ep and ep[0] == alias and (table is None or ep[1] == table):
            return True
    return False


# =============================================================================
# メタの掃除（1ファイルぶん）
# =============================================================================

def _scrub_meta(meta: dict, own_alias: str, alias: str, table: str | None) -> dict:
    """meta から (alias, table) への参照を落とす。落としたものを返す。

    meta はその場で書き換える。table=None なら、そのDBへの参照すべて。
    自分自身のDB（own_alias == alias）かどうかで消す範囲が変わる:
      自分   … テーブルの説明・用語・例文・検算も消す
      他所   … 関連とER図の置き場所だけ（例文は他DBのテーブルを引くこともある）
    """
    hit: dict = {"relationships": [], "glossary": [], "examples": [],
                 "checks": [], "tables": [], "er_external": [], "er_layout": []}
    mine = own_alias == alias

    rels = meta.get("relationships") or []
    keep = [r for r in rels if not _ep_hits(r, own_alias, alias, table)]
    if len(keep) != len(rels):
        hit["relationships"] = [_rel_text(r) for r in rels
                                if _ep_hits(r, own_alias, alias, table)]
        meta["relationships"] = keep
    if not meta.get("relationships"):
        meta.pop("relationships", None)

    # ER図に借りているテーブル（"alias.table" の並び）
    ext = [str(x) for x in (meta.get("er_external") or [])]
    gone = [x for x in ext
            if x.split(".")[0] == alias and (table is None or x.split(".")[-1] == table)]
    if gone:
        hit["er_external"] = gone
        left = [x for x in ext if x not in gone]
        if left:
            meta["er_external"] = left
        else:
            meta.pop("er_external", None)

    # ER図の配置（"alias.table": [x, y]）
    layout = meta.get("er_layout") or {}
    lgone = [k for k in layout
             if str(k).split(".")[0] == alias
             and (table is None or str(k).split(".")[-1] == table)]
    if lgone:
        hit["er_layout"] = lgone
        for k in lgone:
            layout.pop(k, None)
        if not layout:
            meta.pop("er_layout", None)

    if not mine:
        return hit

    # --- ここから先は自分のDBのときだけ -------------------------------------
    if table is None:
        return hit                      # DBごと消えるのでファイルごと処分される

    tables = meta.get("tables") or {}
    if table in tables:
        hit["tables"] = [table]
        tables.pop(table, None)
        if not tables:
            meta.pop("tables", None)

    # DB全体の用語。SQL式がそのテーブルを引いているものだけ消す。
    # 説明文だけの用語は、文中にテーブル名が出てきても残す（文章なので）
    gl = meta.get("glossary") or {}
    gone_terms = [t for t, v in gl.items()
                  if isinstance(v, dict) and uses_table(v.get("sql"), table, alias)]
    if gone_terms:
        hit["glossary"] = gone_terms
        for t in gone_terms:
            gl.pop(t, None)
        if not gl:
            meta.pop("glossary", None)

    exs = meta.get("examples") or []
    left = [e for e in exs if not uses_table(e.get("sql"), table, alias)]
    if len(left) != len(exs):
        hit["examples"] = [str(e.get("q") or e.get("sql") or "")[:60]
                           for e in exs if uses_table(e.get("sql"), table, alias)]
        meta["examples"] = left
        if not left:
            meta.pop("examples", None)

    cks = verify.normalize(meta.get("checks"))
    def ck_hits(c):
        return any(uses_table(s, table, alias) for s in
                   (c["left"]["sql"], c["right"]["sql"], c.get("drilldown") or ""))
    left = [c for c in cks if not ck_hits(c)]
    if len(left) != len(cks):
        hit["checks"] = [c["name"] for c in cks if ck_hits(c)]
        meta["checks"] = left
        if not left:
            meta.pop("checks", None)

    return hit


def _merge(into: dict, where: str, hit: dict) -> None:
    """掃除の結果を「どのDBで何を消したか」の形で積む。"""
    for key, items in hit.items():
        for it in items:
            into.setdefault(key, []).append({"db": where, "text": it})


# =============================================================================
# 下見（消す前に見せる。何も書き換えない）
# =============================================================================

def _cleanup_walk(alias: str, table: str | None, apply: bool, skip: Path | None = None) -> dict:
    """全DBのメタを見て、(alias, table) への参照を数える／消す。

    apply=False なら保存しない（load_meta は毎回ファイルから読み直すので、
    その場の dict を書き換えても他に影響しない）。
    skip はそのDB自身のファイル（消すので触らない）。
    """
    found: dict = {}
    for f in db.list_db_files():
        if skip is not None and f == skip:
            continue
        own = db.alias_for(f)
        meta = catalog.load_meta(f)
        hit = _scrub_meta(meta, own, alias, table)
        if any(hit.values()):
            _merge(found, own, hit)
            if apply:
                catalog.save_meta(f, meta)
    return found


def _jobs_for(db_name: str, table: str | None) -> list[dict]:
    return [j for j in jobs.list_jobs()
            if j.get("db_file") == db_name and (table is None or j.get("table") == table)]


def _job_text(j: dict) -> str:
    """定期取り込みの1行分。何がどの間隔で入ってくる設定かが分かればよい。"""
    label = jobs.interval_label(j.get("interval_minutes"))
    name = j.get("name") or j.get("table") or "（無題）"
    stopped = "・停止中" if j.get("enabled") is False else ""
    return f"{name}（{j.get('table')} / {label}{stopped}）"


def table_impact(path: Path, table: str) -> dict:
    """テーブルを消したときに巻き添えになるもの。数えるだけ。"""
    alias = db.alias_for(path)
    out = _cleanup_walk(alias, table, apply=False)
    out["jobs"] = [{"id": j.get("id"), "name": j.get("name") or j.get("table"),
                    "text": _job_text(j)}
                   for j in _jobs_for(path.name, table)]
    return out


def db_impact(path: Path) -> dict:
    """DBを消したときに巻き添えになるもの。数えるだけ。"""
    alias = db.alias_for(path)
    out = _cleanup_walk(alias, None, apply=False, skip=path)
    out["jobs"] = [{"id": j.get("id"), "name": j.get("name") or j.get("table"),
                    "text": _job_text(j)}
                   for j in _jobs_for(path.name, None)]
    out["users"] = [{"db": path.name, "text": u} for u in _users_selecting(path.name)]
    try:
        profile = catalog.profile_db(path)
        out["own_tables"] = [{"db": alias,
                              "text": f"{t}（{(info.get('row_count') or 0):,}行）"}
                             for t, info in (profile.get("tables") or {}).items()]
    except Exception as e:
        out["own_tables"] = [{"db": alias, "text": f"（読めませんでした: {e}）"}]
    return out


def _users_selecting(db_name: str) -> list[str]:
    """「対象データ」にそのDBを入れている利用者。"""
    root = config.USER_META_DIR
    if not root.exists():
        return []
    out = []
    for d in sorted(root.iterdir()):
        if not (d / "prefs.yaml").is_file():
            continue
        if db_name in prefs.get_selection(d.name):
            out.append(d.name)
    return out


# =============================================================================
# 実際に消す
# =============================================================================

def clean_table(path: Path, table: str, drop_jobs: bool = True) -> dict:
    """テーブルを消したあとの掃除。DROP TABLE 自体は importer 側で済ませておく。"""
    alias = db.alias_for(path)
    done = _cleanup_walk(alias, table, apply=True)
    done["jobs"] = []
    if drop_jobs:
        for j in _jobs_for(path.name, table):
            if jobs.delete_job(j.get("id")):
                done["jobs"].append({"db": path.name,
                                     "text": j.get("name") or j.get("table")})
    catalog.forget(path)
    return done


def delete_db(path: Path, drop_jobs: bool = True) -> dict:
    """DBを消す。ファイルごと削除するので元には戻せない。

    押す前に、巻き添えになるものの一覧とファイル名の入力で二重に確認している。
    """
    if path.parent.resolve() != config.DATA_DIR.resolve():
        raise ValueError("data/ の外は操作できません。")
    if not path.is_file():
        raise FileNotFoundError(f"DBが見つかりません: {path.name}")

    alias = db.alias_for(path)

    # 先にファイルを消す。順番が逆だと、削除に失敗したときに
    # 他のDBの関連だけが消えて、DBは残るという半端な状態になる。
    removed = []
    try:
        for p in (path, Path(f"{path}.meta.yaml")):
            if p.is_file():
                p.unlink()
                removed.append(p.name)
    except OSError as e:
        raise ValueError(f"{path.name} を削除できませんでした（{e}）。"
                         "このDBを使っている処理が終わってから、もう一度試してください。") from e

    done = _cleanup_walk(alias, None, apply=True, skip=path)

    done["jobs"] = []
    if drop_jobs:
        for j in _jobs_for(path.name, None):
            if jobs.delete_job(j.get("id")):
                done["jobs"].append({"db": path.name,
                                     "text": j.get("name") or j.get("table")})

    # 利用者の「対象データ」から外す（残すと、次に開いたとき選択が壊れて見える）
    done["users"] = []
    for u in _users_selecting(path.name):
        sel = prefs.get_selection(u)
        sel.pop(path.name, None)
        prefs.set_selection(u, sel)
        done["users"].append({"db": path.name, "text": u})

    catalog.forget(path)
    done["removed"] = [{"db": path.name, "text": n} for n in removed]
    print(f"[cleanup] {path.name} を削除しました（{'、'.join(removed)}）")
    return done


#: 画面に出すときの見出し。キーの順にそのまま並べる。
LABELS = {
    "own_tables": "テーブルと中のデータ",
    "tables": "テーブルの説明",
    "relationships": "関連（ER図の線）",
    "glossary": "業務用語",
    "examples": "質問とSQLの例文",
    "checks": "検算ルール",
    "er_external": "ER図に借りているテーブル",
    "er_layout": "ER図の配置",
    "jobs": "定期取り込みの設定",
    "users": "利用者の「対象データ」の選択",
    "removed": "削除したファイル",
}


def summarize(impact: dict) -> list[dict]:
    """{キー: [{db, text}]} を画面用の並びにする。空の項目は落とす。"""
    return [{"key": k, "label": LABELS.get(k, k), "items": impact[k]}
            for k in LABELS if impact.get(k)]


# ==========================================================================
# ===== 元 mailer.py
# メールの下書きと送信（SMTP）。
#
# 方針: 「作る」と「送る」を必ず分ける。
# LLMは compose（下書き作成）までしかできず、実際の送信は
# 画面でユーザーが本文と宛先を見て承認したときだけ実行される。
# 宛先を間違えた1通は取り消せないので、AIの判断だけでは外に出さない。
#
# 宛先はDBのテーブルから探す。人の情報がどのテーブルにあるかは
# DBごとに違うので、列名と実際の値（@を含むか等）から推測する。
# ==========================================================================
import mimetypes
import re
import smtplib
import ssl
import threading
from dataclasses import dataclass, field
from datetime import datetime
from email.header import Header
from email.message import EmailMessage
from email.utils import formataddr, formatdate, make_msgid, parseaddr

import yaml

import config
import db

# ざっくりだが実用上これで十分。厳密なRFC準拠より、明らかな入力ミスを弾く方が大事。
EMAIL_RE = re.compile(r"^[^@\s,;:<>\"]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}$")

# 宛先探しで手がかりにする列名（部分一致・大文字小文字は無視）
MAIL_HINTS = ("mail", "メール", "eメール", "address", "アドレス", "宛先", "email")
NAME_HINTS = ("name", "氏名", "名前", "担当", "社員名", "person", "user", "顧客名", "得意先")
DEPT_HINTS = ("部署", "部門", "所属", "課", "dept", "department", "division", "組織", "拠点")

_mailer_lock = threading.Lock()
_sent_log: list[dict] = []          # 直近の送信記録（画面表示用）
_MAX_LOG = 200


class MailError(Exception):
    """送信できない理由（そのまま画面に出す）。"""


# =============================================================================
# 設定
# =============================================================================

@dataclass
class SmtpSettings:
    host: str = ""
    port: int = 25
    security: str = "none"          # none / starttls / ssl
    user: str = ""
    password: str = ""
    sender: str = ""
    sender_name: str = ""
    timeout: int = 20
    allow_addresses: list = field(default_factory=list)
    senders: list = field(default_factory=list)      # 画面で選べる差出人の候補
    max_recipients: int = 20
    dry_run: bool = True
    alert_to: list = field(default_factory=list)     # 定期取り込みの失敗を知らせる管理者

    @property
    def configured(self) -> bool:
        return bool(self.host and self.sender)

    def problems(self) -> list[str]:
        out = []
        if not self.host:
            out.append("送信サーバのホスト名が未設定です。「メール設定」画面で登録してください。")
        if not self.sender:
            out.append("差出人アドレスが未設定です。「メール設定」画面で登録してください。")
        if not self.allow_addresses:
            out.append("送信できる宛先が1件も登録されていません。"
                       "登録するまでメールは送れません（「メール設定」画面で追加してください）。")
        elif not EMAIL_RE.match(self.sender):
            out.append(f"差出人アドレスの形式が正しくありません: {self.sender}")
        if self.security not in ("none", "starttls", "ssl"):
            out.append(f"暗号化は none / starttls / ssl のいずれかです: {self.security}")
        if self.user and not self.password:
            out.append("認証ユーザー名を設定したのに、パスワードが空です。")
        return out

    def allows(self, address: str) -> bool:
        """この宛先に送ってよいか。

        許可リストに載っているアドレスだけに送れる。**空のときは誰にも送れない。**
        「未設定なら全員に送れる」だと、設定を忘れたまま社外へ出てしまうため。
        """
        if not self.allow_addresses:
            return False
        return str(address).strip().lower() in [a.lower() for a in self.allow_addresses]

    @property
    def restricted(self) -> bool:
        """常に True。許可リスト方式なので、制限が外れることはない。"""
        return True


# =============================================================================
# 画面から変えられる設定（data/mail_settings.yaml）
#
# env の値を初期値として、このファイルの内容を上書きで重ねる。
# この設定ファイルは中身をそのまま画面に出すので、秘密は入れないこと。
# =============================================================================

# 送信サーバ（接続先）。暗号化と認証は社内リレー前提で画面に出さないので、
# 変えたい環境では env の SMTP_SECURITY / SMTP_USER / SMTP_PASSWORD を使う。
SERVER_KEYS = ("host", "port", "timeout")
# 差出人と宛先まわり
EDITABLE_KEYS = SERVER_KEYS + ("sender", "sender_name", "senders",
                               "allow_addresses", "max_recipients", "dry_run",
                               "alert_to")


def _read_overrides() -> dict:
    p = config.SMTP_SETTINGS_FILE
    if not p.exists():
        return {}
    try:
        data = yaml.safe_load(p.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"[mailer] 設定を読めませんでした: {p} ({e})")
        return {}
    return {k: v for k, v in (data or {}).items() if k in EDITABLE_KEYS} \
        if isinstance(data, dict) else {}


def _write_overrides(data: dict) -> None:
    p = config.SMTP_SETTINGS_FILE
    p.parent.mkdir(parents=True, exist_ok=True)
    with _mailer_lock:
        p.write_text(yaml.safe_dump({k: data[k] for k in EDITABLE_KEYS if k in data},
                                    allow_unicode=True, sort_keys=False),
                     encoding="utf-8")


def settings() -> SmtpSettings:
    """env の値に、画面から保存した設定を重ねて返す。"""
    ov = _read_overrides()
    senders = [str(s).strip() for s in (ov.get("senders") or []) if str(s).strip()]
    sender = str(ov.get("sender") or config.SMTP_SENDER or "").strip()
    if not sender and senders:
        sender = senders[0]
    return SmtpSettings(
        host=str(ov.get("host", config.SMTP_HOST) or "").strip(),
        port=int(ov.get("port", config.SMTP_PORT) or 25),
        security=str(config.SMTP_SECURITY or "none").lower(),
        user=str(config.SMTP_USER or "").strip(),
        password=config.SMTP_PASSWORD,
        sender=sender,
        sender_name=str(ov.get("sender_name", config.SMTP_SENDER_NAME) or "").strip(),
        timeout=int(ov.get("timeout", config.SMTP_TIMEOUT) or 20),
        allow_addresses=[str(a).strip() for a in (ov.get("allow_addresses") or [])
                         if str(a).strip()],
        senders=senders,
        max_recipients=int(ov.get("max_recipients", config.SMTP_MAX_RECIPIENTS) or 20),
        dry_run=bool(ov["dry_run"]) if "dry_run" in ov else config.SMTP_DRY_RUN,
        alert_to=[str(a).strip() for a in (ov.get("alert_to") or []) if str(a).strip()],
    )


# --- 宛先に登録してよいドメイン（env の SEND_OK_MAIL_DOMAIN）------------------------

def domain_ok(address: str) -> bool:
    """このアドレスを許可リストに登録してよいか。

    env が空なら制限なし（それでも許可リストへの登録自体は必要）。
    サブドメイン（sales.example.co.jp）も対象に含める。
    """
    allowed = config.SEND_OK_MAIL_DOMAIN
    if not allowed:
        return True
    dom = str(address).strip().lower().rsplit("@", 1)[-1]
    return any(dom == d or dom.endswith("." + d) for d in allowed)


def allowed_domains_label() -> str:
    """画面に出す「登録できるドメイン」の表記。"""
    return "、".join("@" + d for d in config.SEND_OK_MAIL_DOMAIN) or "すべてのドメイン"


_HOSTNAME_RE = re.compile(r"^[A-Za-z0-9]([A-Za-z0-9.-]*[A-Za-z0-9])?$")


def _validate_server(data: dict) -> list[str]:
    """送信サーバ（接続情報）の点検。"""
    errors = []
    host = str(data.get("host") or "").strip()
    if not host:
        errors.append("SMTPサーバのホスト名を入力してください。")
    elif not _HOSTNAME_RE.match(host):
        errors.append(f"ホスト名の形式が正しくありません: {host}")

    try:
        port = int(data.get("port") or 0)
    except (TypeError, ValueError):
        errors.append("ポート番号は数値で指定してください。")
    else:
        if not (1 <= port <= 65535):
            errors.append("ポート番号は 1〜65535 で指定してください。")

    try:
        t = int(data.get("timeout") or 0)
    except (TypeError, ValueError):
        errors.append("タイムアウトは数値で指定してください。")
    else:
        if not (1 <= t <= 300):
            errors.append("タイムアウトは 1〜300 秒で指定してください。")
    return errors


def validate_settings(data: dict) -> list[str]:
    """画面から来た設定の点検。1つでも返ったら保存しない。"""
    errors = _validate_server(data)
    senders = [str(s).strip() for s in (data.get("senders") or []) if str(s).strip()]
    for s in senders:
        if not EMAIL_RE.match(s):
            errors.append(f"差出人アドレスの形式が正しくありません: {s}")
    sender = str(data.get("sender") or "").strip()
    if not sender:
        errors.append("使用する差出人アドレスを選んでください。")
    elif not EMAIL_RE.match(sender):
        errors.append(f"差出人アドレスの形式が正しくありません: {sender}")
    elif senders and sender not in senders:
        errors.append(f"{sender} は差出人の候補に入っていません。先に候補へ追加してください。")

    for a in (data.get("allow_addresses") or []):
        addr = str(a).strip()
        if not EMAIL_RE.match(addr):
            errors.append(f"宛先として登録できない形式です: {a}")
        elif not domain_ok(addr):
            errors.append(f"{addr} は登録できません。"
                          f"登録できるのは {allowed_domains_label()} のアドレスだけです"
                          "（env の SEND_OK_MAIL_DOMAIN）。")
    # 通知先の管理者も同じ縛り（許可ドメイン）。社外へは飛ばさない
    for a in (data.get("alert_to") or []):
        addr = str(a).strip()
        if not EMAIL_RE.match(addr):
            errors.append(f"通知先として登録できない形式です: {a}")
        elif not domain_ok(addr):
            errors.append(f"{addr} は通知先に登録できません。"
                          f"登録できるのは {allowed_domains_label()} のアドレスだけです。")

    try:
        n = int(data.get("max_recipients") or 0)
    except (TypeError, ValueError):
        errors.append("一度に送れる宛先数は数値で指定してください。")
    else:
        if not (1 <= n <= 500):
            errors.append("一度に送れる宛先数は 1〜500 で指定してください。")
    return errors


def _with_current(data: dict) -> dict:
    """省略されたキーを現在の値で埋める。

    画面は毎回すべて送ってくるが、一部だけ変えたい呼び出し方もできるように
    しておく。埋めずに検証すると「送っていない項目」で弾かれてしまう。
    """
    s = settings()
    merged = {"host": s.host, "port": s.port, "timeout": s.timeout,
              "sender": s.sender, "sender_name": s.sender_name,
              "senders": s.senders, "allow_addresses": s.allow_addresses,
              "max_recipients": s.max_recipients, "dry_run": s.dry_run,
              "alert_to": s.alert_to}
    # None は「指定なし」。空文字や空リストは「消したい」なので通す。
    merged.update({k: v for k, v in (data or {}).items()
                   if k in merged and v is not None})
    return merged


def save_settings(data: dict, user: str | None = None) -> SmtpSettings:
    """画面からの保存。検証してから書く。"""
    merged = _with_current(data)
    errors = validate_settings(merged)
    if errors:
        raise MailError(" / ".join(errors))
    keep = _read_overrides()
    keep.update({
        "host": str(merged["host"] or "").strip(),
        "port": int(merged["port"] or 25),
        "timeout": int(merged["timeout"] or 20),
        "sender": str(merged["sender"] or "").strip(),
        "sender_name": str(merged["sender_name"] or "").strip(),
        "senders": [str(s).strip() for s in (merged["senders"] or []) if str(s).strip()],
        "allow_addresses": [str(a).strip() for a in (merged["allow_addresses"] or [])
                            if str(a).strip()],
        "max_recipients": int(merged["max_recipients"] or 20),
        "dry_run": bool(merged["dry_run"]),
        "alert_to": [str(a).strip() for a in (merged.get("alert_to") or [])
                     if str(a).strip()],
    })
    _write_overrides(keep)
    print(f"[mailer] 設定を更新しました（{user or '不明'}）: "
          f"サーバ={keep['host']}:{keep['port']} / "
          f"差出人={keep['sender']} / 宛先制限="
          f"{len(keep['allow_addresses'])}アドレス / "
          f"テスト送信={keep['dry_run']}")
    return settings()


def mail_status() -> dict:
    """画面に渡す現在の設定。認証情報は含めない。"""
    s = settings()
    return {"configured": s.configured, "host": s.host, "port": s.port,
            "sender": s.sender,
            "sender_name": s.sender_name, "dry_run": s.dry_run,
            "senders": s.senders,
            "allow_addresses": s.allow_addresses,
            "alert_to": s.alert_to,
            "restricted": s.restricted, "timeout": s.timeout,
            "allowed_domains": list(config.SEND_OK_MAIL_DOMAIN),
            "allowed_domains_label": allowed_domains_label(),
            "max_recipients": s.max_recipients, "problems": s.problems(),
            "settings_file": str(config.SMTP_SETTINGS_FILE)}


# =============================================================================
# 宛先を探す
# =============================================================================

def _hit(name: str, hints) -> bool:
    low = str(name).lower()
    return any(h.lower() in low for h in hints)


def _looks_like_email_column(conn, table: str, column: str) -> bool:
    """列名で分からないときは、実際の値に @ が入っているかで判断する。"""
    try:
        cur = conn.execute(f'SELECT "{column}" FROM "{table}" '
                           f'WHERE "{column}" IS NOT NULL LIMIT 20')
        vals = [str(r[0]) for r in cur.fetchall()]
    except Exception:
        return False
    if not vals:
        return False
    return sum(1 for v in vals if "@" in v and "." in v.split("@")[-1]) >= max(1, len(vals) // 2)


def address_tables(scope: list[dict]) -> list[dict]:
    """選択中のDBから「人とメールアドレスが載っていそうな表」を探す。"""
    found = []
    for s in scope:
        try:
            conn = db.connect_ro(s["path"])
        except Exception:
            continue
        try:
            tables = [r[0] for r in conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' "
                "AND name NOT LIKE 'sqlite_%' ORDER BY name")]
            for t in tables:
                if s.get("tables") and t not in s["tables"]:
                    continue
                try:
                    cols = [r[1] for r in conn.execute(f'PRAGMA table_info("{t}")')]
                except Exception:
                    continue
                mail_cols = [c for c in cols if _hit(c, MAIL_HINTS)]
                if not mail_cols:
                    mail_cols = [c for c in cols if _looks_like_email_column(conn, t, c)]
                if not mail_cols:
                    continue
                found.append({
                    "alias": s["alias"], "table": t, "columns": cols,
                    "mail_columns": mail_cols,
                    "name_columns": [c for c in cols if _hit(c, NAME_HINTS)],
                    "dept_columns": [c for c in cols if _hit(c, DEPT_HINTS)],
                })
        finally:
            conn.close()
    return found


def find_recipients(scope: list[dict], query: str = "", *, limit: int = 50,
                    table: str | None = None) -> dict:
    """名前・部署・アドレスの断片から宛先候補を探す。

    どの表を見ればよいかは address_tables() が推測する。
    query が空なら、その表の先頭から候補を出す（一覧確認のため）。
    """
    sources = address_tables(scope)
    if table:
        sources = [s for s in sources
                   if s["table"] == table or f"{s['alias']}.{s['table']}" == table]
    if not sources:
        return {"ok": False, "candidates": [], "sources": [],
                "message": "メールアドレスが入っていそうな表が、選択中のDBに見つかりません。"
                           "サイドバーで名簿のあるDBを選ぶか、"
                           "アドレスを直接指定してください。"}

    conf = settings()
    q = str(query or "").strip()
    out, seen = [], set()
    for src in sources:
        cols = src["columns"]
        search_cols = list(dict.fromkeys(
            src["mail_columns"] + src["name_columns"] + src["dept_columns"]
            + [c for c in cols if c not in src["mail_columns"]]))[:12]
        where, params = "", {}
        if q:
            # 値はプレースホルダで渡す。列名は実在するものだけを使うので識別子として安全。
            where = " WHERE " + " OR ".join(f'CAST("{c}" AS TEXT) LIKE :q'
                                            for c in search_cols)
            params["q"] = f"%{q}%"
        sql = (f'SELECT {", ".join(chr(34) + c + chr(34) for c in cols)} '
               f'FROM "{src["table"]}"{where} LIMIT {int(limit)}')
        try:
            conn = db.connect_ro(next(s["path"] for s in scope
                                      if s["alias"] == src["alias"]))
        except Exception:
            continue
        try:
            rows = conn.execute(sql, params).fetchall()
        except Exception:
            rows = []
        finally:
            conn.close()

        for r in rows:
            rec = dict(zip(cols, r))
            mail = next((str(rec[c]).strip() for c in src["mail_columns"]
                         if rec.get(c) and "@" in str(rec[c])), "")
            if not mail or mail.lower() in seen:
                continue
            seen.add(mail.lower())
            out.append({
                "email": mail,
                "name": next((str(rec[c]) for c in src["name_columns"] if rec.get(c)), ""),
                "dept": next((str(rec[c]) for c in src["dept_columns"] if rec.get(c)), ""),
                "source": f"{src['alias']}.{src['table']}",
                "valid": bool(EMAIL_RE.match(mail)),
                # 許可リストの外なら、下書きを作る前に分かるようにしておく
                "allowed": conf.allows(mail),
                "row": {k: (str(v) if v is not None else "") for k, v in list(rec.items())[:8]},
            })
            if len(out) >= limit:
                break
        if len(out) >= limit:
            break

    msg = (f"{len(out)}件の宛先候補が見つかりました。" if out
           else f"「{q}」に一致する宛先が見つかりませんでした。"
                f"探した表: {', '.join(s['alias'] + '.' + s['table'] for s in sources)}")
    return {"ok": bool(out), "candidates": out, "message": msg,
            "sources": [{"table": f"{s['alias']}.{s['table']}",
                         "mail_columns": s["mail_columns"],
                         "name_columns": s["name_columns"],
                         "dept_columns": s["dept_columns"]} for s in sources]}


# =============================================================================
# 下書きの検証
# =============================================================================

def _norm_addresses(value) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        parts = re.split(r"[,;\n]", value)
    else:
        parts = list(value)
    out = []
    for p in parts:
        p = str(p).strip()
        if not p:
            continue
        name, addr = parseaddr(p)
        out.append(addr or p)
    return out


def validate_draft(draft: dict, *, system: bool = False) -> list[str]:
    """送る前の点検。1つでも返ったら送信できない。

    system=True は、アプリ自身が管理者へ送る通知（定期取り込みの失敗など）。
    宛先は「メール設定」の通知先（alert_to）そのものなので、利用者向けの
    許可リスト（allow_addresses）とは独立に通す。サーバ・差出人の設定は同じく必要。
    """
    s = settings()
    errors = [e for e in s.problems()
              if not (system and "送信できる宛先" in e)]
    to = _norm_addresses(draft.get("to"))
    cc = _norm_addresses(draft.get("cc"))
    bcc = _norm_addresses(draft.get("bcc"))
    if not to:
        errors.append("宛先(To)が空です。")
    for addr in to + cc + bcc:
        if not EMAIL_RE.match(addr):
            errors.append(f"アドレスの形式が正しくありません: {addr}")
    total = len(to) + len(cc) + len(bcc)
    if total > s.max_recipients:
        errors.append(f"宛先が {total} 件あります。一度に送れるのは "
                      f"{s.max_recipients} 件までです（SMTP_MAX_RECIPIENTS）。")
    if system:
        # 通知先として登録したアドレスにだけ送る
        ok = {a.lower() for a in s.alert_to}
        for addr in to + cc + bcc:
            if addr.lower() not in ok:
                errors.append(f"{addr} は通知先に登録されていません。")
    elif not s.allow_addresses:
        errors.append("送信できる宛先が1件も登録されていません。"
                      "「メール設定」画面で登録するまで、どこにも送信できません。")
    else:
        for addr in to + cc + bcc:
            if not s.allows(addr):
                errors.append(f"{addr} は送信が許可されていません"
                              f"（許可済みは {len(s.allow_addresses)}件）。"
                              "メール設定で追加してください。")
    if not str(draft.get("subject") or "").strip():
        errors.append("件名が空です。")
    if not str(draft.get("body") or "").strip():
        errors.append("本文が空です。")
    return errors


def build_message(draft: dict, attachments: list[dict] | None = None) -> EmailMessage:
    """EmailMessage を組み立てる（送信せずに中身を確認するのにも使う）。"""
    s = settings()
    msg = EmailMessage()
    msg["From"] = formataddr((str(Header(s.sender_name, "utf-8")), s.sender)) \
        if s.sender_name else s.sender
    msg["To"] = ", ".join(_norm_addresses(draft.get("to")))
    if _norm_addresses(draft.get("cc")):
        msg["Cc"] = ", ".join(_norm_addresses(draft.get("cc")))
    if draft.get("reply_to"):
        msg["Reply-To"] = str(draft["reply_to"])
    msg["Subject"] = str(draft.get("subject") or "")
    msg["Date"] = formatdate(localtime=True)
    msg["Message-ID"] = make_msgid()
    msg.set_content(str(draft.get("body") or ""))

    for a in (attachments or []):
        data, name = a.get("data"), a.get("filename") or "attachment"
        if not data:
            continue
        mime = a.get("mime") or mimetypes.guess_type(name)[0] or "application/octet-stream"
        maintype, _, subtype = mime.partition("/")
        msg.add_attachment(data, maintype=maintype, subtype=subtype or "octet-stream",
                           filename=name)
    return msg


def preview(draft: dict, attachments: list[dict] | None = None) -> dict:
    """送信せずに、送られる内容をそのまま見せる。"""
    s = settings()
    to = _norm_addresses(draft.get("to"))
    cc = _norm_addresses(draft.get("cc"))
    bcc = _norm_addresses(draft.get("bcc"))
    body = str(draft.get("body") or "")
    return {
        "from": (f"{s.sender_name} <{s.sender}>" if s.sender_name else s.sender),
        "to": to, "cc": cc, "bcc": bcc,
        "subject": str(draft.get("subject") or ""),
        "body": body,
        "body_lines": len(body.splitlines()),
        "attachments": [{"filename": a.get("filename"),
                         "size": len(a.get("data") or b"")}
                        for a in (attachments or [])],
        "errors": validate_draft(draft),
        "dry_run": s.dry_run,
        "smtp": f"{s.host}:{s.port} ({s.security})",
    }


# =============================================================================
# 送信
# =============================================================================

def _connect(s: SmtpSettings):
    if s.security == "ssl":
        server = smtplib.SMTP_SSL(s.host, s.port, timeout=s.timeout,
                                  context=ssl.create_default_context())
    else:
        server = smtplib.SMTP(s.host, s.port, timeout=s.timeout)
        if s.security == "starttls":
            server.starttls(context=ssl.create_default_context())
    if s.user:
        server.login(s.user, s.password)
    return server


def send(draft: dict, attachments: list[dict] | None = None,
         user: str | None = None, *, system: bool = False) -> dict:
    """実際に送る。呼ぶ前に必ずユーザーの承認を取ること。

    SMTP_DRY_RUN=true のあいだは接続せず、組み立てた内容だけ返す
    （本番のSMTPを教えてもらう前に画面を試せるようにするため）。
    system=True はアプリ自身からの管理者通知（validate_draft 参照）。
    """
    errors = validate_draft(draft, system=system)
    if errors:
        raise MailError(" / ".join(errors))
    s = settings()
    msg = build_message(draft, attachments)
    recipients = (_norm_addresses(draft.get("to")) + _norm_addresses(draft.get("cc"))
                  + _norm_addresses(draft.get("bcc")))

    record = {"at": datetime.now().isoformat(timespec="seconds"),
              "to": _norm_addresses(draft.get("to")),
              "cc": _norm_addresses(draft.get("cc")),
              "bcc_count": len(_norm_addresses(draft.get("bcc"))),
              "subject": msg["Subject"], "user": user,
              "attachments": [a.get("filename") for a in (attachments or [])],
              "dry_run": s.dry_run, "ok": False, "message": ""}

    if s.dry_run:
        record.update(ok=True, message="下書きの確認のみ（SMTP_DRY_RUN=true のため送信していません）")
        # 送らない代わりに、文面をサーバのログに残す（通知の中身を SMTP 無しで確かめるため）
        body_lines = "\n".join("  | " + l for l in str(draft.get("body") or "").splitlines())
        print(f"[mailer] DRY RUN → {', '.join(record['to'])}\n件名: {msg['Subject']}\n{body_lines}",
              flush=True)
    else:
        try:
            server = _connect(s)
        except Exception as e:
            record["message"] = f"SMTPサーバに接続できません（{s.host}:{s.port}）: {e}"
            _mailer_log(record)
            raise MailError(record["message"]) from e
        try:
            server.send_message(msg, from_addr=s.sender, to_addrs=recipients)
            record.update(ok=True, message=f"{len(recipients)}件の宛先に送信しました。")
        except Exception as e:
            record["message"] = f"送信に失敗しました: {e}"
            _mailer_log(record)
            raise MailError(record["message"]) from e
        finally:
            try:
                server.quit()
            except Exception:
                pass
    _mailer_log(record)
    return record


def test_connection() -> dict:
    """設定の疎通確認だけ行う（メールは送らない）。"""
    s = settings()
    problems = s.problems()
    if problems:
        return {"ok": False, "message": " / ".join(problems)}
    try:
        server = _connect(s)
    except Exception as e:
        return {"ok": False, "message": f"接続できませんでした（{s.host}:{s.port}）: {e}"}
    try:
        server.noop()
        return {"ok": True, "message": f"{s.host}:{s.port} に接続できました"
                                       f"（{s.security}{'・認証あり' if s.user else ''}）。"}
    finally:
        try:
            server.quit()
        except Exception:
            pass


def _mailer_log(record: dict) -> None:
    with _mailer_lock:
        _sent_log.append(record)
        while len(_sent_log) > _MAX_LOG:
            _sent_log.pop(0)


def sent_log(limit: int = 50) -> list[dict]:
    with _mailer_lock:
        return list(reversed(_sent_log[-limit:]))


# =============================================================================
# 定期取り込みの失敗を管理者に知らせる
#
# 状態が「健全 → 失敗」に変わった瞬間に1回だけ送る。失敗が続くあいだ毎周期
# 送ると（15分ごとの設定なら1日96通）読まれなくなるので、直るまで黙る。
# 直ったら「復旧しました」を1回送る。宛先は「メール設定」の通知先（管理者）。
# =============================================================================

def alert_import_problems(current: list[dict], previous: list[dict]) -> dict | None:
    """定期取り込みの状態変化を管理者に送る。送らなかったときは None。

    current / previous は jobs.problems() の結果（今回と前回）。
    """
    s = settings()
    if not s.alert_to:
        return None
    now_ids = {p["id"]: p for p in current}
    prev_ids = {p["id"]: p for p in previous}
    newly = [now_ids[i] for i in now_ids if i not in prev_ids]
    fixed = [prev_ids[i] for i in prev_ids if i not in now_ids]
    if not newly and not fixed:
        return None

    lines = []
    if newly:
        lines.append("■ 設定どおりに更新できなくなった定期取り込み")
        for p in newly:
            tag = {"failed": "失敗", "degraded": "値の型がずれた", "overdue": "動いていない"}.get(p.get("kind"), "")
            lines.append(f"  ・{p['name']}（{p['db_file']} / {p['table']}）{'［' + tag + '］' if tag else ''}")
            lines.append(f"     {p['message']}")
        lines.append("")
        kinds = {p.get("kind") for p in newly}
        if "failed" in kinds:
            lines.append("  失敗: 取り込み元のファイル・シート名・列構成を確認してください。"
                         "失敗している間、そのテーブルは前回の内容のまま変わりません。")
        if "degraded" in kinds:
            lines.append("  値の型がずれた: 取り込みはできていますが、数値の列に文字が混ざったため"
                         "文字として保存しました。元ファイルの値を直して次回の実行を待ってください。")
        if "overdue" in kinds:
            lines.append("  動いていない: アプリが起動しているか、自動実行が止まっていないか確認してください。")
        lines.append("")
    if fixed:
        lines.append("■ 復旧した定期取り込み")
        for p in fixed:
            lines.append(f"  ・{p['name']}（{p['db_file']} / {p['table']}）")
        lines.append("")
    lines.append("確認: データカタログ > DB・テーブル > 各テーブルの「管理」")
    subject = ("[DB分析アシスタント] 定期取り込みが失敗しています"
               if newly else "[DB分析アシスタント] 定期取り込みが復旧しました")
    draft = {"to": list(s.alert_to), "subject": subject, "body": "\n".join(lines)}
    try:
        return send(draft, [], user="scheduler", system=True)
    except Exception as e:
        print(f"[mailer] 通知を送れませんでした: {e}")
        return None


# ==========================================================================
# ===== 元 custom_tools.py
# ユーザーがUIから定義するツール（SQLテンプレート型）。
#
# ツール1つは「名前 + 説明 + パラメータ定義 + SQLテンプレート + 出力形式」で表す。
# Pythonコードは書かせない。SQLは既存の SELECT専用ガード（db.run_select）を通し、
# パラメータは SQLite のバインド変数として渡すのでSQLインジェクションは起こらない。
#
# 保存先は各DBの .meta.yaml の `tools:`。どのDBに置かれていても全DB共通で使える
# （collect_everywhere が全DBから集める。置き場はSQLが主に見ているDBに自動で決まる）。
#
#   tools:
#     - name: monthly_sales
#       description: 指定年の月別売上を返す。「今年の売上推移」などで使う。
#       parameters:
#         - name: year
#           type: string
#           description: "対象年 'YYYY'"
#           required: true
#       sql: |
#         SELECT strftime('%Y-%m', o.order_date) AS 月, ... WHERE ... = :year ...
#       render: chart          # table | chart | chart_dual | none
#       chart: {chart_type: line, x: 月, y: 売上, title: 月別売上}
#       enabled: true
#
# 組み込みツールの有効/無効と説明文の上書きは .meta.yaml の `builtin_tools:` に持つ。
#
#   builtin_tools:
#     plot_dual_axis: {enabled: false}
#     run_sql_query: {description: "…独自の言い回しに差し替え…"}
# ==========================================================================
import re

RENDER_KINDS = ("table", "chart", "chart_dual", "excel", "csv", "none")
PARAM_TYPES = ("string", "integer", "number", "boolean")

# ツール名はOpenAIのfunction名の制約に合わせる（英数字とアンダースコア）
_NAME_RE = re.compile(r"^[A-Za-z][A-Za-z0-9_]{0,47}$")
# SQL中のバインド変数 :name を拾う（:: は型キャストなので除外）
_BIND_RE = re.compile(r"(?<!:):([A-Za-z_][A-Za-z0-9_]*)")

# 組み込みツール名（ユーザー定義ツールと衝突させない）。
# 以前はここに4つだけ列挙していたが、実際の組み込みは38個ある。
# 漏れた名前（forecast など）でツールを作れてしまい、AIに同じ名前の関数が
# 2つ渡って、しかも実行されるのは組み込み側だけ、という不整合が起きていた。
# 一覧は tools.schemas が持っているので、そちらから引く（循環importを避けて遅延）。
def builtin_names() -> set:
    from tools import BUILTIN_TOOLS
    return {t["function"]["name"] for t in BUILTIN_TOOLS}


def custom_tool_safe_name(candidate: str, taken=()) -> str:
    """どんな文字列からでも、使えるツール名を作る。

    AIが起こした名前や日本語の説明が元でも、function名の制約
    （英字始まり・英数字と_・48文字以内）に収め、既存とも組み込みとも
    衝突しない名前にする。ユーザーに名前で悩ませないための道具。
    """
    ascii_ = re.sub(r"[^a-z0-9_]+", "_", str(candidate or "").lower()).strip("_")
    base = ascii_[:40] if re.match(r"^[a-z]", ascii_) else ""
    if not base:
        base = "tool"
    used = {str(t).lower() for t in taken} | {n.lower() for n in builtin_names()}
    name, n = base, 2
    while name.lower() in used:
        name = f"{base}_{n}"
        n += 1
    return name


def bind_names(sql: str) -> list[str]:
    """SQLテンプレートに現れるバインド変数名（重複なし・出現順）。"""
    seen, out = set(), []
    for m in _BIND_RE.finditer(sql or ""):
        if m.group(1) not in seen:
            seen.add(m.group(1))
            out.append(m.group(1))
    return out


def validate_custom_tool(tool: dict, existing_names: set = frozenset()) -> list[str]:
    """ツール定義を検証し、問題点のリストを返す（空なら妥当）。"""
    errs: list[str] = []
    name = str(tool.get("name") or "").strip()
    if not name:
        errs.append("ツール名は必須です。")
    elif not _NAME_RE.match(name):
        errs.append("ツール名は英字で始まる英数字とアンダースコアのみ（48文字以内）にしてください。")
    elif name in builtin_names():
        errs.append(f"'{name}' は組み込みツールと同じ名前です。別の名前にしてください。")
    elif name in existing_names:
        errs.append(f"'{name}' は既に存在します。")

    if not str(tool.get("description") or "").strip():
        errs.append("説明は必須です（AIがこのツールを使うかどうかの判断材料になります）。")

    sql = str(tool.get("sql") or "").strip()
    if not sql:
        errs.append("SQLは必須です。")

    params = tool.get("parameters") or []
    pnames = []
    for i, p in enumerate(params, start=1):
        pn = str((p or {}).get("name") or "").strip()
        if not pn:
            errs.append(f"パラメータ{i}: 名前が空です。")
            continue
        if not re.match(r"^[A-Za-z_][A-Za-z0-9_]*$", pn):
            errs.append(f"パラメータ '{pn}': 英数字とアンダースコアのみ使えます。")
        if pn in pnames:
            errs.append(f"パラメータ '{pn}' が重複しています。")
        pnames.append(pn)
        if (p or {}).get("type") not in PARAM_TYPES:
            errs.append(f"パラメータ '{pn}': 型は {', '.join(PARAM_TYPES)} のいずれかにしてください。")

    # SQL中の :name と パラメータ定義の対応
    if sql:
        binds = set(bind_names(sql))
        for miss in sorted(binds - set(pnames)):
            errs.append(f"SQLに :{miss} がありますが、パラメータが定義されていません。")
        for unused in sorted(set(pnames) - binds):
            errs.append(f"パラメータ '{unused}' がSQL中で使われていません（:{unused} と書きます）。")

    render = tool.get("render") or "table"
    if render not in RENDER_KINDS:
        errs.append(f"出力形式は {', '.join(RENDER_KINDS)} のいずれかにしてください。")
    if render == "chart":
        import charts
        c = tool.get("chart") or {}
        ct = str(c.get("chart_type") or "").strip()
        if not ct:
            errs.append("グラフ出力には chart.chart_type が必要です。")
        elif ct not in charts.CHART_TYPES:
            errs.append(f"未対応のグラフ種別です: {ct}")
        else:
            for k in charts.required_fields(ct):
                v = c.get(k)
                if (not v) if k != "path" else (not list(v or [])):
                    errs.append(f"{ct} には chart.{k} が必要です。")
    if render == "chart_dual":
        c = tool.get("chart") or {}
        if not str(c.get("x") or "").strip():
            errs.append("2軸グラフには chart.x が必要です。")
        if not (c.get("bar_y") or []):
            errs.append("2軸グラフには chart.bar_y（棒にする列）が1つ以上必要です。")
        if not (c.get("line_y") or []):
            errs.append("2軸グラフには chart.line_y（折れ線にする列）が1つ以上必要です。")
    return errs


def to_schema(tool: dict) -> dict:
    """ユーザー定義ツール → OpenAI function calling のJSON Schema。"""
    props, required = {}, []
    for p in (tool.get("parameters") or []):
        pn = str(p.get("name") or "").strip()
        if not pn:
            continue
        props[pn] = {"type": p.get("type") or "string",
                     "description": str(p.get("description") or "")}
        if p.get("required", True):
            required.append(pn)
    return {
        "type": "function",
        "function": {
            "name": tool["name"],
            "description": str(tool.get("description") or ""),
            "parameters": {"type": "object", "properties": props, "required": required},
        },
    }


def coerce_params(tool: dict, args: dict) -> dict:
    """LLMが渡してきた引数を、定義された型に寄せてバインド用の辞書にする。"""
    out = {}
    for p in (tool.get("parameters") or []):
        pn = str(p.get("name") or "").strip()
        if not pn:
            continue
        v = args.get(pn)
        if v is None:
            out[pn] = None
            continue
        t = p.get("type") or "string"
        try:
            if t == "integer":
                out[pn] = int(v)
            elif t == "number":
                out[pn] = float(v)
            elif t == "boolean":
                out[pn] = 1 if (v is True or str(v).lower() in ("true", "1", "yes")) else 0
            else:
                out[pn] = str(v)
        except (TypeError, ValueError):
            raise ValueError(f"パラメータ '{pn}' を {t} として解釈できません: {v!r}")
    return out


def collect_everywhere(selected: list[dict] | None = None) -> list[dict]:
    """全DBのユーザー定義ツールを集める。置き場のDBを選んでいなくても拾う。

    ツールは作るときにDBを意識させない（SQLがどのDBに入るかはAIが決める）ので、
    「置き場のDBを選んでいないと存在しないことになる」のは作った人の意図と食い違う。
    組み込みツールと同じで、DBの選択に関係なく在ることにする。

    selected を渡すと、そのSQLが名指ししているDBが1つも選ばれていないツールは外す。
    いま見ている範囲と関係のないツールまで並べると、AIの選び分けが鈍るため。
    """
    import catalog                       # 循環importを避けるため、使うときに読む
    import db as dbmod

    picked = {str(s.get("name") or "") for s in (selected or [])}
    out, seen = [], set()
    for p in dbmod.list_db_files():
        alias = dbmod.alias_for(p)
        for t in (catalog.load_meta(p).get("tools") or []):
            if not isinstance(t, dict):
                continue
            name = str(t.get("name") or "").strip()
            if not name or name in seen or t.get("enabled") is False:
                continue
            if picked:
                needs = set(dbmod.dbs_named_in(str(t.get("sql") or "")))
                # どのDBも名指ししていないSQLは、置き場のDBのものとして扱う
                if not (needs or {p.name}) & picked:
                    continue
            seen.add(name)
            # owner_file は編集画面が保存先を知るためのもの（aliasはファイル名と別物）
            out.append({**t, "owner": alias, "owner_file": p.name})
    return out


def builtin_overrides(entries: list[dict]) -> dict:
    """組み込みツールの有効/無効・説明上書きを合成する。

    無効化はどれか1つのDBで無効なら無効（安全側）。説明は最初に見つかったものを採用。
    """
    merged: dict[str, dict] = {}
    for e in entries:
        for name, ov in (e.get("meta", {}).get("builtin_tools") or {}).items():
            if not isinstance(ov, dict):
                continue
            cur = merged.setdefault(name, {})
            if ov.get("enabled") is False:
                cur["enabled"] = False
            desc = str(ov.get("description") or "").strip()
            if desc and not cur.get("description"):
                cur["description"] = desc
    return merged


# ==========================================================================
# ===== 元 tools/results.py
# ツールが取ったデータを短いあいだ覚えておく置き場。
#
# 同じSQLを何度も流し直さないための仕組み。「集計 → グラフ → レポート」と
# 進むとき、以前は各ツールが自分でSQLを実行していたので、1つの問いに対して
# 同じSQLが3回走っていた。往復の上限（config.MAX_AGENT_STEPS）も、そのぶん
# 無駄に消える。結果に名前（result_id）を付けて返し、後続のツールはSQLの
# 代わりにその名前を指せるようにする。
#
# 副産物として、表とグラフが必ず同じデータを見ることになる。
# 実行し直す方式では、その間にデータが入れ替わると数字がずれ得た。
#
# 置き方の約束:
#   * プロセス内の辞書に持つ。ワーカーは1つで運用する前提（run.py参照）。
#   * 古いものから捨てる。件数と総セル数の両方に上限を設ける。
#   * 取り出すときは、預けたときと同じDBの組み合わせかを確かめる。
#     IDを当てずっぽうで指されても、選んでいないDBの中身は出さない。
# ==========================================================================
import uuid
from collections import OrderedDict

#: 覚えておく結果の数。会話1本で使う量に対して十分な余裕を見た数。
MAX_ENTRIES = 40
#: 総セル数の上限（行×列の合計）。これを超えたら古いものから捨てる。
MAX_CELLS = 400_000

_store: "OrderedDict[str, dict]" = OrderedDict()


def scope_key(scope: list[dict]) -> str:
    """どのDBの組み合わせで取ったデータかを表す文字列。"""
    return "|".join(sorted(str((s or {}).get("path") or "") for s in (scope or [])))


def _cells(entry: dict) -> int:
    return len(entry["rows"]) * max(1, len(entry["columns"]))


def _evict() -> None:
    """上限を超えたぶんを、古い順に捨てる。"""
    while len(_store) > MAX_ENTRIES:
        _store.popitem(last=False)
    total = sum(_cells(e) for e in _store.values())
    while total > MAX_CELLS and len(_store) > 1:
        _, old = _store.popitem(last=False)
        total -= _cells(old)


def put(scope: list[dict], columns: list, rows: list, truncated: bool = False,
        sql: str | None = None, label: str | None = None) -> str:
    """結果を預けて result_id を返す。"""
    rid = "r_" + uuid.uuid4().hex[:8]
    _store[rid] = {
        "scope": scope_key(scope),
        "columns": list(columns),
        "rows": [tuple(r) for r in rows],
        "truncated": bool(truncated),
        "sql": sql,
        "label": label,
    }
    _evict()
    return rid


def get(scope: list[dict], rid: str) -> dict | None:
    """預けた結果を取り出す。無い・別のDBの組み合わせ、のときは None。"""
    entry = _store.get(str(rid or ""))
    if entry is None or entry["scope"] != scope_key(scope):
        return None
    _store.move_to_end(rid)          # 使ったものは新しい扱いにして残す
    return entry


def describe(rid: str) -> str:
    """LLMに返す一言。何のデータなのかを思い出せるようにする。"""
    entry = _store.get(str(rid or ""))
    if entry is None:
        return ""
    return entry.get("label") or (entry.get("sql") or "")[:80]


def clear() -> None:
    """テスト用。"""
    _store.clear()


# ==========================================================================
# ===== 元 tools/common.py
# ツールの実処理が共通で使う小道具。
#
# LLMへ返すJSONの組み立てと、データを用意して advanced.py に渡す定型。
# データは sql から取ることも、前のツールが返した result_id を指すこともできる
# （results.py 参照）。
# ==========================================================================
import json

import advanced
import config
import db


def _json(obj) -> str:
    return json.dumps(obj, ensure_ascii=False, default=str)


def _err(message: str) -> dict:
    return {
        "ok": False,
        "llm_content": _json({"error": message}),
        "render": {"role": "assistant", "kind": "error", "message": message},
    }


def _total_rows(sql: str, scope: list[dict]) -> int | None:
    """上限で切り詰められたとき、本当は何行あるのかを数える。

    「全部見た」と誤解したまま結論を書かせないための添え物。
    重いSQLだと数え直しも失敗し得るので、そのときは黙って諦める。
    """
    try:
        _, rows, _ = db.run_select(f"SELECT COUNT(*) FROM ({sql.strip().rstrip(';')})",
                                   scope, max_rows=1)
        return int(rows[0][0])
    except Exception:
        return None


def fetch(spec: dict, scope: list[dict], *, label: str | None = None,
          max_rows: int | None = None):
    """ツールが使うデータを用意する。

    spec に result_id があれば前の結果を使い、無ければ sql を実行する。
    どちらの道でも result_id を返すので、呼び出し側はそれをLLMに伝えて
    次のツールで使い回せるようにする。

    max_rows は取得の上限。省略時は画面向けの MAX_RESULT_ROWS(2,000)。
    ファイル出力は EXPORT_MAX_ROWS を渡してくる。「表とグラフは2,000行で足りるが、
    CSVは全行欲しい」ので、道具ごとに上限が違う。

    戻り値: (columns, rows, truncated, result_id, total_rows)
      total_rows は切り詰めが起きたときだけ入る（本当の総件数）。
    """
    rid = str((spec or {}).get("result_id") or "").strip()
    if rid:
        entry = results.get(scope, rid)
        if entry is None:
            raise advanced.AnalysisError(
                f"result_id '{rid}' のデータが見つかりません。"
                "古くなって捨てられたか、別の会話の結果です。"
                "sql を指定して取り直してください。")
        # 預かっているのは2,000行に切り詰めた結果のことがある。ファイル出力のように
        # もっと大きな上限で全行欲しい呼び出しなら、預けたときのSQLで取り直す。
        # （「集計→CSVに」の流れで、CSVだけ2,000行で欠ける事故を防ぐ）
        if (max_rows and entry["truncated"] and entry.get("sql")
                and max_rows > len(entry["rows"])):
            sql = entry["sql"]
            wide = db.widen_scope(sql, scope)
            columns, rows, truncated = db.run_select(sql, wide, max_rows=max_rows)
            return columns, rows, truncated, rid, (
                _total_rows(sql, wide) if truncated else None)
        return entry["columns"], entry["rows"], entry["truncated"], rid, None

    sql = str((spec or {}).get("sql") or "").strip()
    if not sql:
        raise advanced.AnalysisError(
            "sql と result_id のどちらも指定されていません。"
            "新しくデータを取るなら sql を、前のツールの結果を使うなら result_id を指定してください。")
    # スコープは質問ごとの自動判定なので、例文由来のSQLなどが範囲外のDBを
    # 名指しすることがある。必要なぶんは繋いで実行する（読み取り専用のまま）。
    scope = db.widen_scope(sql, scope)
    columns, rows, truncated = db.run_select(sql, scope, max_rows=max_rows)
    # 使い回し用の預かりは、後続の表・グラフに足りる行数まで。
    # 100万行をそのまま預けると、これ1つで置き場(MAX_CELLS)を食い潰すため。
    keep = rows[: config.MAX_RESULT_ROWS]
    rid = results.put(scope, columns, keep, truncated or len(rows) > len(keep),
                      sql=sql, label=label)
    return columns, rows, truncated, rid, (_total_rows(sql, scope) if truncated else None)


def source_note(row_count: int, truncated: bool, total: int | None,
                cap: int | None = None) -> dict:
    """LLMに渡す「元データの規模」。切り詰めのときは実際の件数も添える。"""
    out = {"source_row_count": row_count, "source_truncated": bool(truncated)}
    if truncated:
        out["source_total_rows"] = total
        out["warning"] = (
            f"上限 {cap or config.MAX_RESULT_ROWS:,} 行で切り詰めました。"
            + (f"実際は {total:,} 行あります。" if total else "")
            + "この結果は全体の一部です。全体を語るなら、SQL側で"
              "GROUP BY で集計するか、条件を絞って取り直してください。")
    return out


def _select_for(args: dict, scope: list[dict]):
    """分析ツール共通の入口。データを用意して (columns, rows) を返す。"""
    columns, rows, truncated, rid, total = fetch(args, scope)
    if not rows:
        raise advanced.AnalysisError("データが0行でした。抽出条件を見直してください。")
    return columns, rows, truncated, rid, total


def _report_result(res: dict, *, source_rows: int | None = None,
                   truncated: bool = False, total: int | None = None,
                   result_id: str | None = None,
                   scope: list[dict] | None = None,
                   extra: dict | None = None) -> dict:
    """advanced.py の戻り値を、画面用アイテムとLLM用の要約に変換する。

    LLMには表を丸ごと渡さない。所見(notes)と各表の先頭数行があれば
    十分に説明でき、トークンも節約できる。
    表そのものを次のツールへ渡せるよう、1つ目の表は result_id を付けて預ける。
    """
    tables = res.get("tables") or []
    llm_tables = []
    for i, t in enumerate(tables):
        rows = t.get("rows") or []
        item = {
            "name": t.get("name"), "columns": t.get("columns"),
            "row_count": len(rows),
            "rows": [list(r) for r in rows[: config.SAMPLE_ROWS_FOR_LLM]],
        }
        # 分析結果の表もグラフやレポートの材料になる。指せるようにしておく。
        if scope is not None and rows:
            item["result_id"] = results.put(scope, t.get("columns") or [], rows,
                                            label=f"{res.get('title')} / {t.get('name')}")
        llm_tables.append(item)
    payload = {"status": "analysis_ready", "title": res.get("title"),
               "notes": res.get("notes") or [], "tables": llm_tables,
               "meta": res.get("meta") or {}}
    if source_rows is not None:
        payload.update(source_note(source_rows, truncated, total))
    if result_id:
        payload["source_result_id"] = result_id
    payload.update(extra or {})
    return {
        "ok": True,
        "llm_content": _json(payload),
        "render": {"role": "assistant", "kind": "report", "title": res.get("title"),
                   "tables": tables, "notes": res.get("notes") or []},
    }


def _analysis_tool(fn):
    """データを用意して advanced.py の関数に渡す、共通のかたち。"""
    def run(args: dict, scope: list[dict]) -> dict:
        try:
            columns, rows, truncated, rid, total = _select_for(args, scope)
            res = fn(args, columns, rows)
        except advanced.AnalysisError as e:
            return _err(str(e))
        except Exception as e:
            return _err(f"分析に失敗しました: {e}")
        if args.get("title"):
            res["title"] = args["title"]
        return _report_result(res, source_rows=len(rows), truncated=truncated,
                              total=total, result_id=rid, scope=scope)
    return run


# ==========================================================================
# ===== 元 tools/schemas.py
# LLMに渡すツールの定義（JSON Schema）。
#
# ここは「何ができるか」の宣言だけを書く場所で、処理は置かない。
# 実処理は query / stats / reports / mail にある。
# ==========================================================================
import advanced
import analysis
import charts
import excel
import exports
import pptx_report
import usage

#: 前のツールが返したデータを、SQLを書き直さずに使い回すための指定。
#: 「集計 → グラフ → レポート」で同じSQLが何度も走るのを避ける。
_RESULT_ID = {
    "type": "string",
    "description": "前のツールが返した result_id。これを指定すると sql は不要で、"
                   "同じデータをそのまま使う（同じSQLを書き直さないこと）。",
}


# 指定の名前 -> スキーマ（同じ説明を何度も書かないためのまとめ）
_CHART_ARGS = {
    "x": {"type": "string", "description": "横軸／カテゴリにする列名。"},
    "y": {"type": "string", "description": "値にする列名（数値）。"},
    "y2": {"type": "string", "description": "もう一方の値の列（dumbbell の比較先）。"},
    "z": {"type": "string", "description": "3つ目の数値の列（scatter3d の高さ）。"},
    "color": {"type": "string",
              "description": "色分けに使う列名。積み上げや群分けにも使う。"},
    "size": {"type": "string", "description": "大きさ／幅に使う数値列。"},
    "text": {"type": "string", "description": "点や棒に添えるラベルの列名。"},
    "facet": {"type": "string", "description": "この列の値ごとに小さく分割して並べる。"},
    "path": {"type": "array", "items": {"type": "string"},
             "description": "階層。大きい分類から順に列名を並べる。"},
    "dimensions": {"type": "array", "items": {"type": "string"},
                   "description": "対象にする列名のリスト（3〜6列が読みやすい）。"},
    "lower": {"type": "string", "description": "下限の列（信頼区間や予測の幅）。"},
    "upper": {"type": "string", "description": "上限の列。"},
    "source": {"type": "string", "description": "流れの起点になる列。"},
    "target": {"type": ["string", "number"],
               "description": "流れの終点の列。指標では目標値（数値そのもの、"
                              "または目標が入っている列名）。"},
    "start": {"type": "string", "description": "開始日時の列。"},
    "end": {"type": "string", "description": "終了日時の列。"},
    "open": {"type": "string", "description": "始値の列。"},
    "high": {"type": "string", "description": "高値の列。"},
    "low": {"type": "string", "description": "安値の列。"},
    "close": {"type": "string", "description": "終値の列。"},
    "value": {"type": "string", "description": "指標にする数値の列。"},
    "agg": {"type": "string", "enum": ["sum", "mean", "max", "min", "last"],
            "description": "value をどうまとめるか。既定は sum。"},
    "max": {"type": "number", "description": "ゲージの上限値。省略すると自動。"},
    "suffix": {"type": "string", "description": "数値の後ろに付ける単位（円・%など）。"},
    "nbins": {"type": "integer", "description": "階級の数。既定は自動。"},
    "orientation": {"type": "string", "enum": ["v", "h"],
                    "description": "棒の向き。横棒は h。"},
    "barmode": {"type": "string", "enum": ["group", "stack", "relative"],
                "description": "棒の積み方。既定は group。"},
    "marginal": {"type": "string", "enum": ["box", "violin", "rug"],
                 "description": "ヒストグラムの上に添える分布。任意。"},
    "trendline": {"type": "boolean", "description": "散布図に回帰直線を重ねる。"},
    "colorscale": {"type": "string",
                   "description": "色の濃淡（Blues / Reds / Greens など）。"},
}


# ツール名 -> (分類, 説明, 使う指定, 必須)
_CHART_TOOLS = {
    "plot_comparison": (
        "比較",
        "項目どうしを比べるグラフ。「部署別」「商品別」「順位」「ランキング」"
        "「前年と比べて」「重点管理」を見せたいときに使う。",
        ("x", "y", "y2", "color", "size", "text", "facet", "orientation", "barmode"),
        ("sql", "chart_type", "x", "y", "title")),
    "plot_trend": (
        "推移",
        "時間とともにどう変わったかを見せるグラフ。「推移」「時系列」「予測の幅」"
        "「工程の期間」「日ごとの多寡」「異常な回」を扱うときに使う。",
        ("x", "y", "color", "text", "lower", "upper", "start", "end",
         "open", "high", "low", "close", "facet"),
        ("sql", "chart_type", "title")),
    "plot_composition": (
        "構成",
        "全体が何でできているかを見せるグラフ。「内訳」「構成比」「シェア」"
        "「階層」「増減の要因」「どこからどこへ流れたか」を扱うときに使う。",
        ("x", "y", "color", "path", "source", "target", "text"),
        ("sql", "chart_type", "title")),
    "plot_distribution": (
        "分布",
        "ばらつきの形を見せるグラフ。「分布」「ヒストグラム」「箱ひげ」"
        "「偏り」「正規分布か」「群ごとの散らばり」を扱うときに使う。",
        ("x", "y", "color", "nbins", "facet", "marginal"),
        ("sql", "chart_type", "title")),
    "plot_relationship": (
        "関係",
        "2つ以上の項目の関係を見せるグラフ。「相関」「散布図」「密度」"
        "「多変量」「総当たり」「つながり」を扱うときに使う。",
        ("x", "y", "z", "color", "size", "text", "dimensions", "nbins",
         "source", "target", "colorscale", "trendline", "facet"),
        ("sql", "chart_type", "title")),
    "plot_kpi": (
        "指標",
        "数字を1つ大きく見せるグラフ。「KPI」「達成率」「目標に対して」"
        "「今いくら」を見せたいときに使う。",
        ("value", "target", "agg", "max", "suffix", "colorscale"),
        ("sql", "chart_type", "value", "title")),
}


def _chart_tools() -> list[dict]:
    """用途ごとのグラフツール定義を作る。"""
    out = []
    for name, (cat, desc, fields, required) in _CHART_TOOLS.items():
        props = {
            "sql": {"type": "string",
                    "description": "グラフに使うデータを取る SELECT 文。"
                                   "集計が要るものは GROUP BY 済みにすること。"},
            "chart_type": {"type": "string", "enum": charts.types_in(cat),
                           "description": f"グラフ種別。{charts.type_help(cat)}"},
            "title": {"type": "string", "description": "グラフのタイトル。"},
            "purpose": {"type": "string",
                        "description": "このグラフで示したいことの短い説明。"},
        }
        props.update({f: _CHART_ARGS[f] for f in fields})
        out.append({"type": "function", "function": {
            "name": name,
            "description": (f"{desc}使える種別: {charts.type_help(cat)}"),
            "parameters": {"type": "object", "properties": props,
                           "required": list(required)},
        }})
    return out


BUILTIN_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "run_sql_query",
            "description": (
                "選択中の SQLite DB群 に対して読み取り専用の SELECT 文を実行し、結果テーブルを取得する。"
                "SELECT(または WITH ... SELECT)以外は実行不可。"
                "複数DBが選択されている場合、テーブル名は必ず『エイリアス.テーブル名』で修飾すること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {
                        "type": "string",
                        "description": "実行する SQLite 用 SELECT 文。SELECT または WITH で始めること。",
                    },
                    "purpose": {
                        "type": "string",
                        "description": "このクエリで何を確認したいかの短い説明(日本語)。",
                    },
                },
                "required": ["sql", "purpose"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "plot_chart",
            "description": (
                "SELECT 文の結果をグラフ化する。時系列の推移や分布を可視化したいときに使用。"
                "内部で SELECT を実行し、指定の x / y / color 列でグラフを描画する。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "グラフ用データを取得する SELECT 文(GROUP BYで集計済みにする)。"},
                    "chart_type": {
                        "type": "string",
                        "enum": list(charts.CHART_TYPES),
                        "description": "グラフ種別。" + charts.type_help(),
                    },
                    "x": {"type": "string",
                          "description": "x軸に使う列名。pie/donut はカテゴリ、funnel は段階、"
                                         "radar は指標名、histogram は対象の数値列。"},
                    "y": {"type": "string",
                          "description": "y軸に使う列名。pie/donut/treemap/sunburst/funnel は値、"
                                         "radar は値。histogram では不要。"},
                    "color": {"type": "string",
                              "description": "系列(色分け)に使う列名。任意。heatmap ではマス目の値になる。"},
                    "size": {"type": "string", "description": "bubble の大きさに使う数値列名。"},
                    "text": {"type": "string", "description": "点や棒に表示するラベルの列名。任意。"},
                    "path": {
                        "type": "array", "items": {"type": "string"},
                        "description": "treemap/sunburst の階層。大きい分類から順に列名を並べる。",
                    },
                    "nbins": {"type": "integer", "description": "histogram の階級数。任意。"},
                    "orientation": {
                        "type": "string", "enum": ["v", "h"],
                        "description": "棒の向き。横棒にしたいときは h。bar以外では無視。",
                    },
                    "barmode": {
                        "type": "string",
                        "enum": ["group", "stack", "relative"],
                        "description": "棒グラフの積み方。積み上げ=stack、横並び比較=group(既定)。bar以外では無視。",
                    },
                    "title": {"type": "string", "description": "グラフのタイトル。"},
                    "purpose": {"type": "string", "description": "このグラフで示したいことの短い説明。"},
                },
                "required": ["sql", "chart_type", "title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "plot_dual_axis",
            "description": (
                "棒グラフ(左軸)と折れ線グラフ(右軸)を組み合わせた2軸グラフを描く。"
                "件数(棒)と比率など単位の異なる指標(折れ線)を同時に見せたいときに使用。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "x列と数値の指標列を返す SELECT 文(GROUP BYで集計)。"},
                    "x": {"type": "string", "description": "x軸に使う列名。"},
                    "bar_y": {
                        "type": "array", "items": {"type": "string"},
                        "description": "左軸に棒で表示する数値列名のリスト(1つ以上)。例: 件数。",
                    },
                    "line_y": {
                        "type": "array", "items": {"type": "string"},
                        "description": "右軸に折れ線で表示する数値列名のリスト(1つ以上)。例: 比率(%)。",
                    },
                    "left_title": {"type": "string", "description": "左軸のラベル(任意)。"},
                    "right_title": {"type": "string", "description": "右軸のラベル(任意)。"},
                    "title": {"type": "string", "description": "グラフのタイトル。"},
                    "purpose": {"type": "string", "description": "このグラフで示したいことの短い説明。"},
                },
                "required": ["sql", "x", "bar_y", "line_y", "title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pivot_table",
            "description": (
                "クロス集計表（ピボットテーブル）を作る。"
                "「AとBのマトリクスで」「行に○○、列に△△」「クロス集計」「表形式で比較」"
                "などと言われたら使う。"
                "SQLite には PIVOT 構文が無く CASE WHEN を列の数だけ手書きする必要があるため、"
                "列に展開したい集計はSQLで書かずにこのツールを使うこと。"
                "sql では集計せず、明細または index/columns/values の3列を返すだけでよい。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "元データを取得する SELECT 文。集計はこのツールが行う。"},
                    "index": {"type": "array", "items": {"type": "string"},
                              "description": "行にする列名（複数可）。"},
                    "columns": {"type": "string",
                                "description": "列に展開する列名。省略すると行ごとの集計表になる。"},
                    "values": {"type": "string", "description": "集計する値の列名。"},
                    "aggfunc": {"type": "string", "enum": list(analysis.AGG_FUNCS),
                                "description": "集計方法。既定は sum。"},
                    "margins": {"type": "boolean", "description": "総計の行と列を付けるか。既定は false。"},
                    "percent": {"type": "string", "enum": list(analysis.PERCENT_MODES),
                                "description": "実数の代わりに構成比(%)で出す。"
                                               + " / ".join(f"{k}={v}" for k, v
                                                            in analysis.PERCENT_MODES.items())},
                    "rank": {"type": "string",
                             "description": "大きい順に並べて順位を付ける。列名を書くとその列で、"
                                            "'total' と書くと行の合計で並べる。"},
                    "render": {"type": "string", "enum": ["table", "heatmap"],
                               "description": "表示方法。heatmap にすると色付きの行列で見せる。既定は table。"},
                    "title": {"type": "string", "description": "見出し。"},
                    "purpose": {"type": "string", "description": "この集計で確認したいことの短い説明。"},
                },
                "required": ["sql", "index", "values"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "analyze_stats",
            "description": (
                "統計的な分析を行う。"
                "SQLite には STDDEV / MEDIAN / CORR / PERCENTILE が無いため、"
                "「相関」「中央値」「ばらつき」「四分位」「外れ値」「異常値」を聞かれたら"
                "SQLで計算しようとせず必ずこのツールを使うこと。"
                "sql は集計せずに明細を返す（1行1件）ようにする。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "分析対象の明細を取得する SELECT 文。集計はしない。"},
                    "method": {
                        "type": "string", "enum": ["describe", "correlation", "outliers"],
                        "description": "describe=基本統計量(件数/平均/標準偏差/最小/四分位/中央値/最大) / "
                                       "correlation=相関行列 / outliers=外れ値の抽出",
                    },
                    "columns": {"type": "array", "items": {"type": "string"},
                                "description": "対象の数値列。省略すると数値列を自動判定する。"},
                    "group_by": {"type": "string",
                                 "description": "describe のとき、この列ごとに分けて統計を出す。任意。"},
                    "target": {"type": "string",
                               "description": "outliers のとき、外れ値を調べる数値列。必須。"
                                              "mahalanobis のときは列名をカンマ区切りで複数。"},
                    "outlier_method": {"type": "string",
                                       "enum": list(advanced.OUTLIER_METHODS_EXT),
                                       "description": " / ".join(
                                           f"{k}={v}" for k, v
                                           in advanced.OUTLIER_METHODS_EXT.items())},
                    "threshold": {"type": "number",
                                  "description": "外れ値の閾値。iqr は既定1.5、zscore は既定3。"},
                    "corr_method": {"type": "string", "enum": list(analysis.CORR_METHODS),
                                    "description": "pearson=直線的な関係(既定) / spearman=順位の関係"},
                    "lag": {"type": "integer",
                            "description": "correlation のとき、何期先までずらして相関を見るか。"
                                           "「広告費は翌月の売上に効くか」のような遅れて出る効果を"
                                           "調べたいときに指定する。sql は時点の昇順で1行1期にすること。"},
                    "partial": {"type": "boolean",
                                "description": "correlation のとき true にすると偏相関にする。"
                                               "control で指定した列の影響を取り除いてから相関を見るので、"
                                               "「第3の変数のせいで関係して見えるだけ」を切り分けられる。"},
                    "control": {"type": "array", "items": {"type": "string"},
                                "description": "partial=true のとき、影響を取り除きたい列。"},
                    "title": {"type": "string", "description": "見出し。"},
                    "purpose": {"type": "string", "description": "この分析で確認したいことの短い説明。"},
                },
                "required": ["sql", "method"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "export_excel",
            "description": (
                "SELECT の結果を Excel ファイル(.xlsx)にまとめ、ユーザーがダウンロードできる状態にする。"
                "「エクセルで」「xlsxで」「ファイルにして」「ダウンロードしたい」"
                "などと言われたら使う。sheets に複数の SELECT を渡すと複数シートのブックになる。"
                "chart を書くと、そのシートのデータからExcelのグラフを作って貼る"
                "（画像ではないので、受け取った側が範囲や種類を変えられる）。"
                "数字を並べるだけのシートより、グラフを1つ付けた方が伝わる。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sheets": {
                        "type": "array",
                        "description": "ブックに入れるシート。1要素につき1シート。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "シート名(31文字以内)。"},
                                "sql": {"type": "string", "description": "このシートに書き出す SELECT 文。"},
                                "note": {"type": "string", "description": "シート先頭に入れる補足(任意)。"},
                                "chart": {
                                    "type": "object",
                                    "description": "このシートのデータから作るグラフ（任意）。",
                                    "properties": {
                                        "type": {"type": "string",
                                                 "enum": list(excel.EXCEL_CHART_TYPES),
                                                 "description": "グラフの種類。"},
                                        "category_column": {"type": "string",
                                                            "description": "横軸にする列名。"},
                                        "value_columns": {"type": "array",
                                                          "items": {"type": "string"},
                                                          "description": "系列にする数値列。"},
                                        "title": {"type": "string", "description": "グラフの題名。"},
                                        "y_title": {"type": "string", "description": "縦軸の名前。"},
                                        "x_title": {"type": "string", "description": "横軸の名前。"},
                                        "data_labels": {"type": "boolean",
                                                        "description": "値ラベルを出すか。"},
                                    },
                                    "required": ["type"],
                                },
                                "charts": {
                                    "type": "array",
                                    "items": {"type": "object", "additionalProperties": True},
                                    "description": "グラフを複数貼るときはこちらに並べる。",
                                },
                            },
                            "required": ["name", "sql"],
                        },
                    },
                    "filename": {"type": "string", "description": "ファイル名(拡張子不要)。例: 月別売上"},
                    "purpose": {"type": "string", "description": "何のためのファイルかの短い説明。"},
                },
                "required": ["sheets", "filename"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "export_csv",
            "description": (
                "SELECT の結果を CSV ファイルとして書き出し、ユーザーがダウンロードできる状態にする。"
                "「CSVで」「csvにして」「取り込み用のファイル」などと言われたら使う。"
                "files に複数指定すると、まとめてZIPで渡す。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "files": {
                        "type": "array",
                        "description": "書き出すファイル。1要素につき1CSV。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "ファイル名(拡張子不要)。"},
                                "sql": {"type": "string", "description": "書き出す SELECT 文。"},
                            },
                            "required": ["name", "sql"],
                        },
                    },
                    "encoding": {
                        "type": "string", "enum": list(exports.ENCODINGS),
                        "description": "文字コード。既定は utf-8-sig（Excelで開いても文字化けしない）。"
                                       "Shift_JIS が要るときだけ cp932。",
                    },
                    "delimiter": {
                        "type": "string", "enum": list(exports.EXPORT_DELIMITERS),
                        "description": "区切り文字。既定は comma。TSVにしたいときは tab。",
                    },
                    "purpose": {"type": "string", "description": "何のためのファイルかの短い説明。"},
                },
                "required": ["files"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "export_text",
            "description": (
                "文章（レポート・要約・メモ）をテキストファイルとして書き出し、"
                "ユーザーがダウンロードできる状態にする。"
                "「テキストで」「レポートにして」「議事録に」「まとめを文書で」などと言われたら使う。"
                "body に本文を自分で書き、必要なら sections に SELECT を指定して集計表を差し込む。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "ファイル名(拡張子不要)。"},
                    "body": {
                        "type": "string",
                        "description": "本文。あなたが書いた文章をそのまま入れる。"
                                       "sections を使う場合、差し込みたい位置に {{見出し}} と書く。",
                    },
                    "sections": {
                        "type": "array",
                        "description": "本文に差し込む集計表。省略可。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string",
                                            "description": "見出し。本文の {{この文字列}} が表に置き換わる。"
                                                           "本文に無ければ末尾に追記される。"},
                                "sql": {"type": "string", "description": "表にする SELECT 文。"},
                            },
                            "required": ["heading", "sql"],
                        },
                    },
                    "format": {
                        "type": "string", "enum": ["md", "txt"],
                        "description": "md=Markdown（表が罫線付き）/ txt=プレーンテキスト。既定は md。",
                    },
                    "encoding": {
                        "type": "string", "enum": list(exports.ENCODINGS),
                        "description": "文字コード。既定は utf-8-sig。",
                    },
                },
                "required": ["filename", "body"],
            },
        },
    },
    # ---- 統計 ---------------------------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "hypothesis_test",
            "description": (
                "統計的仮説検定。「差があると言えるか」「偶然ではないか」「有意か」"
                "「A/Bどちらが良いか」「効果があったか」を判断したいときに使う。"
                "平均の差・比率の差・分布の偏り・相関の有無を、p値と効果量つきで判定する。"
                "sql は集計せず明細（1行1件）を返すこと。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "検定対象の明細を取る SELECT 文。"},
                    "method": {"type": "string", "enum": list(advanced.TEST_METHODS),
                               "description": " / ".join(f"{k}={v}" for k, v
                                                         in advanced.TEST_METHODS.items())},
                    "value_col": {"type": "string", "description": "測定値の列（数値）。"},
                    "group_col": {"type": "string",
                                  "description": "群を表す列。2群比較・分散分析・カイ二乗で使う。"},
                    "value_col2": {"type": "string",
                                   "description": "対応のある検定や相関で、もう一方の列。"},
                    "popmean": {"type": "number", "description": "1標本t検定で比較する基準値。"},
                    "expected": {"type": "array", "items": {"type": "number"},
                                 "description": "適合度検定で期待する比率や度数。"},
                    "alternative": {"type": "string",
                                    "enum": ["two-sided", "less", "greater"],
                                    "description": "対立仮説。既定は両側(two-sided)。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "method"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "regression",
            "description": (
                "回帰分析。「何が効いているか」「要因分析」「どの変数が影響するか」"
                "「予測式を作りたい」ときに使う。係数・p値・寄与の大きさ・"
                "あてはまり(R²)・多重共線性まで返す。"
                "目的変数が0/1なら logistic、件数なら poisson を選ぶ。"
                "文字列の説明変数（部署など）は自動でダミー変数にする。"
                "sql は集計せず明細（1行1件）を返すこと。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "分析対象の明細を取る SELECT 文。"},
                    "target": {"type": "string", "description": "目的変数（説明したい列）。"},
                    "features": {"type": "array", "items": {"type": "string"},
                                 "description": "説明変数の列名。"},
                    "method": {"type": "string", "enum": list(advanced.REGRESSION_METHODS),
                               "description": " / ".join(f"{k}={v}" for k, v
                                                         in advanced.REGRESSION_METHODS.items())},
                    "predict": {"type": "array", "items": {"type": "object"},
                                "description": "予測したい入力の一覧。例 [{\"広告費\":100}]",
                                "additionalProperties": True},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "target", "features"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "distribution_analysis",
            "description": (
                "分布の形を調べる。「ばらつき」「ヒストグラム」「偏り」「どんな分布か」"
                "「上位下位の広がり」を見たいときに使う。度数分布・要約統計・"
                "正規分布などへの当てはめ判定を返す。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "明細を取る SELECT 文。"},
                    "target": {"type": "string", "description": "調べる数値列。"},
                    "bins": {"type": "integer", "description": "階級の数。既定20。"},
                    "group_col": {"type": "string", "description": "群ごとに比べるときの列。"},
                    "fit": {"type": "array", "items": {"type": "string",
                                                       "enum": list(advanced.DISTRIBUTIONS)},
                            "description": "当てはめを試す分布。既定は norm と lognorm。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "target"],
            },
        },
    },
    # ---- 時系列 -------------------------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "forecast",
            "description": (
                "将来の値を予測する。「来月はいくら」「このままいくと」「着地見込み」"
                "「予測」を聞かれたら使う。予測値と95%の幅、"
                "過去データで試した誤差率(MAPE)を返す。"
                "sql は時点ごとに1行（例: 月ごとの売上合計）にすること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "時点と値を返す SELECT 文（時点の昇順で1行1期）。"},
                    "time_col": {"type": "string", "description": "時点の列（年月など）。"},
                    "value_col": {"type": "string", "description": "予測する数値の列。"},
                    "periods": {"type": "integer", "description": "何期先まで予測するか。既定6。"},
                    "method": {"type": "string",
                               "enum": ["auto"] + list(advanced.FORECAST_METHODS),
                               "description": "auto=データ量から自動選択。" + " / ".join(
                                   f"{k}={v}" for k, v in advanced.FORECAST_METHODS.items())},
                    "season_length": {"type": "integer",
                                      "description": "季節の周期。月次で1年なら12、曜日なら7。"},
                    "exog": {
                        "type": "object", "additionalProperties": True,
                        "description": "説明変数つきで予測する場合の指定。"
                                       "{\"columns\": [\"広告費\"], \"future\": [[120],[130]]} の形で、"
                                       "future には予測する期数と同じ数だけ将来の値を並べる。"
                                       "「広告費をこう置いたら売上はどうなるか」に答えられる。",
                    },
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "time_col", "value_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "timeseries_analysis",
            "description": (
                "時系列の見方をまとめる。「推移」「トレンド」「季節性」「前年同月比」"
                "「移動平均」「周期」を聞かれたときに使う。"
                "sql は時点ごとに1行にすること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "時点と値を返す SELECT 文。"},
                    "time_col": {"type": "string", "description": "時点の列。"},
                    "value_col": {"type": "string", "description": "値の列。"},
                    "window": {"type": "integer", "description": "移動平均の期間。既定3。"},
                    "season_length": {"type": "integer",
                                      "description": "季節の周期（月次なら12）。指定すると季節分解する。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "time_col", "value_col"],
            },
        },
    },
    # ---- 試算・シミュレーション ---------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "monte_carlo_simulation",
            "description": (
                "モンテカルロ・シミュレーション。「もし〜だったら」「リスク」"
                "「確率」「見込みの幅」「何%の確率で」を扱うときに使う。"
                "不確実な入力を分布で与え、式を何万回も試して結果の分布を出す。"
                "実データのばらつきをそのまま使いたい変数は dist=empirical と column を指定し、"
                "その列を返す sql も併せて渡すこと。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "formula": {"type": "string",
                                "description": "変数名を使った計算式。例 (単価 - 原価) * 数量 - 固定費"},
                    "variables": {
                        "type": "object", "additionalProperties": True,
                        "description": "{変数名: {dist, ...}}。dist は normal(mean,std) / "
                                       "uniform(min,max) / triangular(min,mode,max) / "
                                       "lognormal(mean,std) / poisson(lam) / binomial(n,p) / "
                                       "empirical(column) / fixed(value)。",
                    },
                    "trials": {"type": "integer", "description": "試行回数。既定10000。"},
                    "sql": {"type": "string",
                            "description": "empirical を使うときに、その列を含む明細を取る SELECT 文。"},
                    "targets": {"type": "array", "items": {"type": "number"},
                                "description": "「この値を超える確率」を知りたいしきい値。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["formula", "variables"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "scenario_analysis",
            "description": (
                "シナリオ比較（楽観・標準・悲観など）。前提を数パターン置いて"
                "結果を並べ、どの変数の影響が大きいかも出す。"
                "確率分布まで置く必要がないときは、こちらの方が説明しやすい。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "formula": {"type": "string", "description": "変数名を使った計算式。"},
                    "scenarios": {"type": "object", "additionalProperties": True,
                                  "description": "{シナリオ名: {変数名: 値}}"},
                    "base": {"type": "object", "additionalProperties": True,
                             "description": "共通の前提値。感度分析の基準にもなる。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["formula", "scenarios"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bootstrap_estimate",
            "description": (
                "ブートストラップ法で平均などの信頼区間を出す。"
                "「この差は誤差の範囲か」「どのくらい確からしいか」を、"
                "分布の形を仮定せずに示せる。件数が少ないときにも使える。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "明細を取る SELECT 文。"},
                    "target": {"type": "string", "description": "対象の数値列。"},
                    "statistic": {"type": "string",
                                  "enum": ["mean", "median", "std", "sum", "p90"],
                                  "description": "推定する統計量。既定 mean。"},
                    "group_col": {"type": "string", "description": "群ごとに出すときの列。"},
                    "trials": {"type": "integer", "description": "再抽出の回数。既定5000。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "target"],
            },
        },
    },
    # ---- 分ける -------------------------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "clustering",
            "description": (
                "k-meansでグループ分けする。「セグメント」「タイプ分け」「似ている順に分類」"
                "「顧客を分けたい」ときに使う。列ごとの尺度差は標準化して吸収する。"
                "分け方が分からないときは k に \"auto\" を指定すると、"
                "最も素直に分かれる数を自動で選び、各グループの特徴も言葉で返す。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "明細を取る SELECT 文。"},
                    "features": {"type": "array", "items": {"type": "string"},
                                 "description": "分類に使う数値列。"},
                    "k": {"type": ["integer", "string"],
                          "description": "グループ数。既定3。\"auto\" にすると"
                                         "シルエット係数が最も高い数（2〜8）を自動で選ぶ。"},
                    "categorical": {"type": "array", "items": {"type": "string"},
                                    "description": "分類に使いたい区分の列（地域・会員区分など）。"
                                                   "0/1に開いてから一緒に分ける。"},
                    "label_col": {"type": "string", "description": "行の名前になる列（顧客名など）。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "features"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "abc_analysis",
            "description": (
                "ABC分析（パレート分析）。「売上の8割を占める商品」「重点管理」"
                "「上位集中度」を見るときに使う。累計構成比でA/B/Cに区分する。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "品目と値を返す SELECT 文。"},
                    "label_col": {"type": "string", "description": "品目の列（商品名など）。"},
                    "value_col": {"type": "string", "description": "金額や数量の列。"},
                    "thresholds": {"type": "array", "items": {"type": "number"},
                                   "description": "A/Bの境目。既定 [70, 90]（累計%）。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "label_col", "value_col"],
            },
        },
    },
    # ---- レポート -----------------------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "export_pptx",
            "description": (
                "PowerPointのレポートを作る。「パワポで」「スライドにして」「報告資料」"
                "「プレゼン」と言われたら使う。会議でそのまま映せる体裁で出力する。"
                "グラフは編集できるPowerPointのグラフとして入る。"
                "各スライドに sql を書くと、その場でSQLを実行して中身を埋める。"
                "\n"
                "重要: 各スライドに message（そのページで言いたいこと1行）を必ず書く。"
                "見出しの下に帯で表示され、聞き手はここだけ読めば分かる。"
                "「売上推移」ではなく「3月の落ち込みは期ずれで、実勢は右肩上がり」と書く。"
                "\n"
                "構成の目安: title（表紙）→ kpi（数字の要約）→ "
                "section（章の区切り）→ chart / table（根拠）→ compare（案の比較）→ "
                "closing（まとめと次のアクション）。中扉が2つ以上あれば目次は自動で入る。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "ファイル名（拡張子不要）。"},
                    "title": {"type": "string", "description": "レポート全体の題名。"},
                    "subtitle": {"type": "string", "description": "副題（対象期間など）。"},
                    "footer": {"type": "string", "description": "各ページ下部に入れる文字。"},
                    "slides": {
                        "type": "array",
                        "description": "スライドの並び。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "kind": {"type": "string", "enum": list(pptx_report.SLIDE_KINDS),
                                         "description": "title=表紙 / agenda=目次 / "
                                                        "section=中扉 / message=文字だけ / "
                                                        "table=表 / chart=グラフ / "
                                                        "kpi=数字を大きく / compare=2案の比較 / "
                                                        "closing=まとめと次のアクション"},
                                "title": {"type": "string"},
                                "subtitle": {"type": "string"},
                                "message": {"type": "string",
                                            "description": "そのページで言いたいこと1行。"
                                                           "見出しの下に帯で出る。必ず書く。"},
                                "sql": {"type": "string",
                                        "description": "table/chart のとき、中身を取る SELECT 文。"},
                                "chart": {"type": "string", "enum": list(pptx_report.PPTX_CHART_TYPES),
                                          "description": "グラフの種類。"},
                                "category_column": {"type": "string",
                                                    "description": "chart のとき横軸にする列。"},
                                "value_columns": {"type": "array", "items": {"type": "string"},
                                                  "description": "chart のとき系列にする数値列。"},
                                "bullets": {"type": "array",
                                            "items": {"type": "object",
                                                      "additionalProperties": True},
                                            "description": "箇条書き。文字列でも "
                                                           "{text, level, strong} でもよい。"},
                                "lead": {"type": "string",
                                         "description": "message のとき、箇条書きの前に置く導入文。"},
                                "body": {"type": "string", "description": "本文。"},
                                "items": {"type": "array", "items": {"type": "object",
                                                                     "additionalProperties": True},
                                          "description": "kpi のとき "
                                                         "[{label, value, unit, delta, "
                                                         "delta_unit, delta_label, "
                                                         "higher_is_better, note}]。"},
                                "panes": {"type": "array", "items": {"type": "object",
                                                                     "additionalProperties": True},
                                          "description": "compare のとき、左右2つ "
                                                         "[{title, value, unit, bullets}]。"},
                                "summary": {"type": "array", "items": {"type": "string"},
                                            "description": "closing のときのまとめ。"},
                                "actions": {"type": "array",
                                            "items": {"type": "object",
                                                      "additionalProperties": True},
                                            "description": "closing のとき "
                                                           "[{text, owner, due}]。"},
                                "callout": {"type": "string",
                                            "description": "下部の囲みで強調する一文。1枚に1つまで。"},
                                "comment": {"type": "string",
                                            "description": "図表の右に添える所見。"},
                                "source": {"type": "string",
                                           "description": "出所・集計条件。数字の資料には入れる。"},
                                "notes": {"type": "string", "description": "発表者ノート。"},
                                "data_labels": {"type": "boolean",
                                                "description": "グラフに数値ラベルを出す。"
                                                               "既定は自動判断。"},
                                "highlight_rows": {"type": "array",
                                                   "items": {"type": "integer"},
                                                   "description": "table で強調する行（0始まり）。"},
                                "max_rows": {"type": "integer",
                                             "description": "table で載せる最大行数。既定12。"},
                            },
                            "required": ["kind"],
                        },
                    },
                },
                "required": ["slides"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "build_report",
            "description": (
                "分析の結果を1つのレポートにまとめる。"
                "「レポートにして」「まとめて」「報告書」「分析結果を整理して」"
                "「結論と根拠を示して」と言われたら使う。"
                "画面に読みやすい形で出しつつ、ダウンロードできるファイルも作る。"
                "各セクションに sql を書けば表が入り、chart を書けばグラフも入る。"
                "要点(summary)と結論(conclusion)は必ず自分の言葉で書くこと。"
                "数字を並べるだけでなく『だから何か』を書く。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "レポートの題名。"},
                    "subtitle": {"type": "string", "description": "対象期間や条件。"},
                    "summary": {
                        "type": "array", "items": {"type": "string"},
                        "description": "要点。最初に読む人が3行で分かるように書く。",
                    },
                    "sections": {
                        "type": "array",
                        "description": "本編。1セクション＝1つの論点。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string", "description": "見出し。"},
                                "body": {"type": "string",
                                         "description": "説明の文章。何が言えるかを書く。"},
                                "sql": {"type": "string",
                                        "description": "根拠として載せる表の SELECT 文。任意。"},
                                "chart": {
                                    "type": "object", "additionalProperties": True,
                                    "description": "グラフの指定。chart_type と x / y などを"
                                                   "入れる。sql の結果を使う。任意。",
                                },
                                "note": {"type": "string",
                                         "description": "この節の所見・注意点。任意。"},
                                "max_rows": {"type": "integer",
                                             "description": "表に載せる最大行数。既定20。"},
                            },
                            "required": ["heading"],
                        },
                    },
                    "conclusion": {"type": "string", "description": "結論。"},
                    "recommendations": {
                        "type": "array", "items": {"type": "string"},
                        "description": "推奨する打ち手。実行できる粒度で書く。",
                    },
                    "caveats": {
                        "type": "array", "items": {"type": "string"},
                        "description": "前提・制約・数字の読み方の注意。",
                    },
                    "format": {
                        "type": "string",
                        "enum": ["md", "docx", "pptx", "xlsx", "none"],
                        "description": "ダウンロード用ファイルの形式。"
                                       "md=軽い文書(既定) / docx=Word報告書（図表つき）/ "
                                       "pptx=スライド / xlsx=表ごとにシート / "
                                       "none=画面表示だけ",
                    },
                    "filename": {"type": "string", "description": "ファイル名（拡張子不要）。"},
                    "org": {"type": "string", "description": "表紙に入れる部署名など。"},
                    "footer": {"type": "string",
                               "description": "各ページ下部の文字（「社外秘」など）。"},
                },
                "required": ["title", "sections"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "export_docx",
            "description": (
                "Word文書（.docx）を作る。「Wordで」「docxで」「報告書にして」"
                "「配布資料」「回覧」と言われたら使う。"
                "表紙・目次・図表番号つきのキャプション・ページ番号が入り、"
                "そのまま配布できる体裁になる。"
                "各セクションに sql を書くと表が入り、chart も書くとグラフが図として入る。"
                "本文(body)は必ず自分の言葉で書くこと。表を貼っただけの文書は読まれない。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "文書の題名。"},
                    "subtitle": {"type": "string", "description": "副題（対象期間など）。"},
                    "org": {"type": "string", "description": "表紙に入れる部署名。"},
                    "author": {"type": "string", "description": "作成者名。"},
                    "footer": {"type": "string",
                               "description": "各ページ下部の文字（「社外秘」など）。"},
                    "toc": {"type": "boolean", "description": "目次を入れるか。既定は入れる。"},
                    "summary": {"type": "array", "items": {"type": "string"},
                                "description": "冒頭の要約。3点程度。"},
                    "sections": {
                        "type": "array",
                        "description": "本編。1セクション＝1つの見出し。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string", "description": "見出し。"},
                                "level": {"type": "integer",
                                          "description": "見出しの階層（1が大見出し）。"},
                                "body": {"type": "string", "description": "本文。"},
                                "bullets": {"type": "array", "items": {"type": "string"},
                                            "description": "箇条書き。"},
                                "sql": {"type": "string",
                                        "description": "表・グラフの元になる SELECT 文。"},
                                "chart": {"type": "object", "additionalProperties": True,
                                          "description": "グラフの指定（chart_type と x / y など）。"
                                                         "図として貼られる。"},
                                "table": {"type": "boolean",
                                          "description": "表も載せるか。既定は載せる。"},
                                "caption": {"type": "string", "description": "図のキャプション。"},
                                "table_caption": {"type": "string",
                                                  "description": "表のキャプション。"},
                                "note": {"type": "string", "description": "補足・注記。"},
                                "callout": {"type": "string",
                                            "description": "囲みで強調したい一文。"},
                                "max_rows": {"type": "integer",
                                             "description": "表に載せる最大行数。既定40。"},
                                "page_break": {"type": "boolean",
                                               "description": "このセクションの前で改ページする。"},
                            },
                            "required": ["heading"],
                        },
                    },
                    "conclusion": {"type": "string", "description": "結論。"},
                    "recommendations": {
                        "type": "array",
                        "items": {"type": "object", "additionalProperties": True},
                        "description": "推奨する打ち手。[{text, owner, due}] または文字列の並び。",
                    },
                    "caveats": {"type": "array", "items": {"type": "string"},
                                "description": "前提・注意。"},
                    "filename": {"type": "string", "description": "ファイル名（拡張子不要）。"},
                },
                "required": ["title", "sections"],
            },
        },
    },
    # ---- メール -------------------------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "find_mail_recipients",
            "description": (
                "宛先をDBから探す。名前・部署・アドレスの一部で検索できる。"
                "「〇〇部に送って」「田中さんに送って」と言われたら、"
                "まずこれで実在するアドレスを確認してから compose_email を呼ぶこと。"
                "アドレスを推測で作ってはいけない。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string",
                              "description": "検索語（氏名・部署名・アドレスの一部）。空なら一覧。"},
                    "table": {"type": "string", "description": "探す表を絞るとき。省略可。"},
                    "limit": {"type": "integer", "description": "最大件数。既定50。"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "compose_email",
            "description": (
                "メールの下書きを作る。画面に確認カードが出て、"
                "ユーザーが「送信」を押したときだけ実際に送られる（自動送信はしない）。"
                "宛先は find_mail_recipients で確認した実在のアドレスを使うこと。"
                "直前に作った Excel / CSV / PowerPoint を添付できる。"
                "本文は挨拶・要点・詳細・結びの順で、日本語のビジネスメールとして書く。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "to": {"type": "array", "items": {"type": "string"},
                           "description": "宛先アドレス。"},
                    "to_query": {"type": "string",
                                 "description": "アドレスの代わりに検索語で指定する場合"
                                                "（例: 営業部）。DBから引いて宛先にする。"},
                    "cc": {"type": "array", "items": {"type": "string"}},
                    "bcc": {"type": "array", "items": {"type": "string"}},
                    "subject": {"type": "string", "description": "件名。"},
                    "body": {"type": "string", "description": "本文（プレーンテキスト）。"},
                    "attach_filenames": {
                        "type": "array", "items": {"type": "string"},
                        "description": "この会話で作ったファイル名。省略時は添付なし。"
                                       "'all' を入れると直近に作ったファイルを全部添付する。"},
                    "reply_to": {"type": "string", "description": "返信先アドレス。"},
                },
                "required": ["subject", "body"],
            },
        },
    },
    # ---- 業務でよく聞かれる分析 ---------------------------------------------
    {
        "type": "function",
        "function": {
            "name": "compare_periods",
            "description": (
                "2つの期間を比べ、差がどこから来たのかまで分解する。"
                "「先月と比べて」「前年同月比」「前期からどう変わったか」"
                "「なぜ落ちたのか」を聞かれたら使う。"
                "全体の増減だけでなく、どの区分が押し下げ／押し上げたか（寄与度）を出す。"
                "qty_col を渡すと、金額の変化を「数量が動いたぶん」と「単価が動いたぶん」に分ける。"
                "sql は期間の列・値の列（あれば区分の列）を含む形で、2期間ぶんまとめて取ること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "2期間ぶんのデータを取る SELECT 文。"
                                           "期間の列・値の列・（任意で）区分の列を返す。"},
                    "period_col": {"type": "string", "description": "期間を表す列（'2026-01' など）。"},
                    "value_col": {"type": "string", "description": "比べる数値の列（売上など）。"},
                    "dimension_col": {"type": "string",
                                      "description": "増減の内訳を見る区分の列（部門・商品など）。任意。"},
                    "qty_col": {"type": "string",
                                "description": "数量の列。指定すると数量要因と単価要因に分解する。"},
                    "current": {"type": "string", "description": "当期。省略すると最後の期。"},
                    "previous": {"type": "string", "description": "前期。省略すると最後から2番目の期。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "period_col", "value_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "data_quality",
            "description": (
                "分析の前に、データそのものの異常を洗い出す。"
                "「数字が合わない」「件数がおかしい」「このデータは信用できるか」"
                "と言われたとき、また重要な集計を出す前の確認に使う。"
                "行数・主キーの重複・空の列・親に存在しない外部キー・日付の範囲を調べ、"
                "深刻な順に並べて返す。SQLは要らない（対象のDBを直接見る）。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "tables": {"type": "array", "items": {"type": "string"},
                               "description": "調べるテーブル名（'stocks' でも 'demo_inventory.stocks' でもよい）。省略すると対象DBのテーブルを順に見る。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "detect_anomalies",
            "description": (
                "時系列から「いつもと違う時点」と「いつから変わったか」を見つける。"
                "「異常」「急に増えた」「おかしい日」「いつから悪化したか」を聞かれたら使う。"
                "前後の期間と比べるので、右肩上がりのデータでも直近を全部異常とは言わない。"
                "静的な外れ値（analyze_stats の outliers）とは用途が違う。"
                "sql は時点ごとに1行にすること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "時点と値を返す SELECT 文。"},
                    "time_col": {"type": "string", "description": "時点の列。"},
                    "value_col": {"type": "string", "description": "監視する数値の列。"},
                    "window": {"type": "integer", "description": "比べる前後の期間。既定7。"},
                    "threshold": {"type": "number",
                                  "description": "何倍離れたら異常とするか。既定3。小さくすると多く拾う。"},
                    "season_length": {"type": "integer",
                                      "description": "曜日や月の周期。指定すると季節変動を除いてから判定する。"},
                    "changepoints": {"type": "boolean",
                                     "description": "水準が変わった時点も探すか。既定true。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "time_col", "value_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "funnel_analysis",
            "description": (
                "段階ごとの通過・離脱・滞留を出す。"
                "「見積から受注までの転換率」「どこで落ちているか」「滞留」"
                "「リードタイム」を聞かれたら使う。"
                "sql は 1行=1案件 にして、各段階の日付（または通過フラグ）の列を並べること。"
                "例: 見積日・受注日・請求日・入金日を1行に持つSELECT。"
                "値が入っていればその段階を通過した扱いになる。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "1行=1案件で、各段階の日付列を並べた SELECT 文。"},
                    "steps": {"type": "array", "items": {"type": "string"},
                              "description": "段階を表す列名を、順番に並べる。"},
                    "labels": {"type": "array", "items": {"type": "string"},
                               "description": "画面に出す段階の名前。省略すると列名を使う。"},
                    "group_col": {"type": "string",
                                  "description": "区分ごとに通過率を比べるときの列（担当・地域など）。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "steps"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "cohort_analysis",
            "description": (
                "いつ始めた人がどれだけ続いているかを見る。"
                "「継続率」「定着」「リピート」「離脱」「初回からの推移」を聞かれたら使う。"
                "初回の期でグループ分けし、経過期ごとの残存率をマトリクスで返す。"
                "sql は 1行=(対象, 期) の明細にすること（同じ人が複数期に出てよい）。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "対象と期の明細を返す SELECT 文。"},
                    "id_col": {"type": "string", "description": "対象の列（顧客ID・社員IDなど）。"},
                    "period_col": {"type": "string",
                                   "description": "期の列（'2026-01' など、並べて正しい順になる形）。"},
                    "value_col": {"type": "string",
                                  "description": "金額なども見るときの数値列。任意。"},
                    "max_periods": {"type": "integer", "description": "何期先まで見るか。既定12。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "id_col", "period_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "market_basket",
            "description": (
                "一緒に買われている（使われている）品目の組み合わせを見つける。"
                "「併売」「セット販売」「一緒に買われる」「関連商品」を聞かれたら使う。"
                "支持度・確信度・リフトを返す。SQLでは実質書けない分析。"
                "sql は 1行=(伝票, 品目) の明細にすること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string",
                            "description": "伝票と品目の明細を返す SELECT 文。"},
                    "transaction_col": {"type": "string",
                                        "description": "伝票を表す列（受注ID・レシートIDなど）。"},
                    "item_col": {"type": "string", "description": "品目の列（商品名など）。"},
                    "min_support": {"type": "number",
                                    "description": "全伝票に占める最低の出現率(%)。既定1.0。"},
                    "top": {"type": "integer", "description": "返す組み合わせの数。既定25。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "transaction_col", "item_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "survival_analysis",
            "description": (
                "「どれだけ持つか」「いつ辞めるか」を扱う。"
                "設備の故障間隔(MTBF)・部品の寿命・社員の在籍期間・顧客の継続期間に使う。"
                "まだ起きていない分（在籍中・稼働中）を捨てずに計算するので、"
                "単純平均のように短く見積もることがない。"
                "Weibull分布の形から、劣化型か初期不良型か偶発型かも判定する。"
                "sql は 1行=1対象 にして、期間の列と（あれば）発生フラグの列を返すこと。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "sql": {"type": "string", "description": "1行=1対象の明細を返す SELECT 文。"},
                    "duration_col": {"type": "string",
                                     "description": "期間の列（稼働時間・在籍日数など）。"},
                    "event_col": {"type": "string",
                                  "description": "起きたか(1)まだか(0)の列。省略すると全件で起きた扱い。"},
                    "group_col": {"type": "string",
                                  "description": "群ごとに比べるときの列（機種・部署など）。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": ["sql", "duration_col"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "explore_import_files",
            "description": (
                "「データ取り込み」の取り込み元フォルダにあるファイルを調べる。読むだけで、"
                "取り込み・変更・削除はしない（取り込みは画面の操作でしか行われない）。"
                "「取り込み元に何がある？」「新しく届いたファイルは？」「このExcelは取り込める？」"
                "「まだ取り込んでいないファイルは？」と聞かれたら使う。"
                "まだDBに入っていないファイルの話なので、SQLでは答えられない。"
                "\n"
                "file を指定しなければ一覧（拡張子を問わず、場所・サイズ・更新日時・"
                "取り込み済みか）、指定すればそのファイルの下見になる。"
                "一覧で得たパスをそのまま file に渡すこと（パスを推測して組み立てないこと）。"
                "\n"
                "下見では「そのまま取り込める / 手直しが要る / 取り込みに向かない」を判定し、"
                "理由と直し方を返す。セル結合・多段見出し・月が横に並んだクロス表・"
                "合計行の混入・見出しが1行目にない、といった"
                "「取り込めるが正しく使えない」形を見つけられる。"
                "取り込みを勧める前に、この判定を確認すること。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string",
                             "description": "見たいフォルダ。省略すると許可フォルダの直下を見る。"},
                    "recursive": {"type": "boolean",
                                  "description": "下の階層もまとめて見るか。既定 false。"},
                    "pattern": {"type": "string",
                                "description": "名前に含まれる文字での絞り込み（例: 売上）。"},
                    "only_not_imported": {"type": "boolean",
                                          "description": "まだ取り込んでいないファイルだけに絞る。"},
                    "check": {"type": "boolean",
                              "description": "一覧の各ファイルについて、表として使える形かどうかも"
                                             "判定する。1件ずつ開くので20件までにしている。"},
                    "file": {"type": "string",
                             "description": "中身を下見するファイルのパス。一覧で得たものを使う。"},
                    "sheet": {"type": "string",
                              "description": "Excelのとき、見たいシート名。省略すると先頭のシート。"},
                    "header_row": {"type": "integer",
                                   "description": "見出しの行（0始まり）。2行目が見出しなら1。"},
                    "rows": {"type": "integer",
                             "description": "下見で読む行数。既定5、最大20。"},
                    "title": {"type": "string", "description": "見出し。"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "propose_glossary_term",
            "description": (
                "業務用語をデータカタログの用語集に登録する「登録カード」をチャットに出す。"
                "実際に登録するかはユーザーがカードのボタンで決める（勝手には登録されない）。"
                "\n使いどころ:"
                "\n- ユーザーが言葉の定義を教えてくれたとき"
                "（「有効な受注とはキャンセル以外のこと」など）"
                "\n- あいまいな用語をあなたが解釈し、その解釈をユーザーが認めたとき"
                "\n- 同じ言葉の意味を何度も聞き直していると気づいたとき"
                "\nまず「この定義で用語集に登録しますか？」と一言確認し、"
                "前向きな返事があったらこのツールでカードを出す。"
                "会話のたびに毎回は出さない（うるさくなる）。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "term": {"type": "string", "description": "用語（例: 有効な受注）。"},
                    "description": {"type": "string",
                                    "description": "日本語の定義。ユーザーの言い回しを活かす。"},
                    "sql": {"type": "string",
                            "description": "SQLの条件式・計算式（任意）。"
                                           "例: orders.status != '9' AND orders.kbn = '1'。"
                                           "会話で実際に使って正しかった式を入れる。"},
                    "how": {"type": "string",
                            "description": "どのデータをどこから取り、どう絞る/計算するのかを、"
                                           "SQLを知らない人に伝わる日本語で書く。テーブルや列は"
                                           "業務の言葉で呼ぶ。例:「受注データから、キャンセル(9)"
                                           "以外で取引区分が通常(1)の行を数える」。必須。"},
                    "table": {"type": "string",
                              "description": "用語を置くテーブル。その用語が主に関わるテーブル名。"
                                             "複数テーブルにまたがる用語のときは省略し db を指定。"},
                    "db": {"type": "string",
                           "description": "DB全体の用語にするときのDBファイル名（例: demo_sales.db）。"},
                },
                "required": ["term", "description", "how"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "propose_example",
            "description": (
                "いまの質問と実行済みSQLを「例文」としてカタログに登録する登録カードを出す。"
                "例文はAIのお手本になり、似た質問への精度が上がる。"
                "実際に登録するかはユーザーがカードのボタンで決める。"
                "\n使いどころ:"
                "\n- ユーザーが「これを例文にして」「この答えを覚えて」と言ったとき"
                "\n- ユーザーが回答を「合っている」と認め、その質問が今後もよく出そうなとき"
                "\nsql には、この会話で実際に実行して正しかったSQLをそのまま入れる（書き直さない）。"
                "毎回は提案しない。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "question": {"type": "string",
                                 "description": "ユーザーの質問文（言い回しを変えない）。"},
                    "sql": {"type": "string",
                            "description": "実行して正しかったSELECT文そのまま。"},
                    "summary": {"type": "string",
                                "description": "どのデータをどこから取り、どう集計したかを、"
                                               "SQLを知らない人に伝わる日本語で。"
                                               "例:「受注データと顧客マスタをつなぎ、"
                                               "ランクごとに売上金額を合計した」。"},
                },
                "required": ["question", "sql", "summary"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "show_er_diagram",
            "description": (
                "DBのER図（テーブル同士の関係図）をチャット画面に表示する。"
                "「ER図を見せて」「テーブルの関係を図で」「データ構造を見たい」"
                "と言われたら使う。図は読み取り専用で、利用者が拡大縮小・全画面表示できる。"
                "表示と同時に結合の一覧も返るので、それを踏まえて補足してよい。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "db": {"type": "string",
                           "description": "対象のDBファイル名（例: demo_sales.db）。"
                                          "質問がどのDBの話かはカタログの説明から判断する。"},
                    "purpose": {"type": "string",
                                "description": "何を確かめたくて表示するかの短い説明。"},
                },
                "required": ["db"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "open_table",
            "description": (
                "テーブルの中身（全行）を見る画面へのリンクをチャットに出す"
                "（利用者が「テーブル全体を開く」を押すと別タブで開く。勝手には開かない）。"
                "「〇〇テーブルを見せて」「中身を全部見たい」「データそのものを確認したい」"
                "と言われたら、SELECT * を書くのではなくこれを使う。"
                "画面では列ごとのフィルター・並べ替え・ページ送りができ、全行を辿れる。"
                "出したあとは、そのテーブルが何かを1〜2文で補足するだけでよい"
                "（中身を表で貼り直さない）。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "table": {"type": "string", "description": "開くテーブル名。"},
                    "db": {"type": "string",
                           "description": "そのテーブルがあるDBファイル名（例: demo_sales.db）。"
                                          "同じ表名が複数のDBにあるときは必須。"},
                    "purpose": {"type": "string",
                                "description": "何を確かめたくて開くかの短い説明。"},
                },
                "required": ["table"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "analyze_usage",
            "description": (
                "このアプリ自身の使われ方（利用状況）を調べる。分析対象のDBの中身ではなく、"
                "チャット履歴と取り込みの記録が材料なので、SQLでは答えられない。"
                "「このアプリはどれくらい使われている？」「誰が使っている？」"
                "「よく使われる機能は？」「どんな質問が多い？」「どこで失敗している？」"
                "「カタログのどこを直せばいい？」と聞かれたら使う。"
                "\n"
                "method で見る角度を選ぶ。まず summary で全体像を出し、"
                "気になった点を errors や users で掘るとよい。"
                "特に errors は失敗を「カタログを直せば減るもの」と"
                "「モデル・API側の問題」に分けて返すので、改善の打ち手を答えるときに使う。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "method": {
                        "type": "string",
                        "enum": [*usage.METHODS, "imports"],
                        "description": " / ".join(
                            f"{k}={v}" for k, v in usage.METHODS.items())
                        + " / imports=取り込みの実行実績",
                    },
                    "days": {"type": "integer",
                             "description": "直近何日ぶんを見るか。省略すると全期間。"},
                    "user": {"type": "string",
                             "description": "特定の利用者だけに絞る（ユーザー名）。任意。"},
                    "title": {"type": "string", "description": "見出し。"},
                    "purpose": {"type": "string",
                                "description": "この集計で確認したいことの短い説明。"},
                },
                "required": ["method"],
            },
        },
    },
    # ---- グラフ（用途別） ---------------------------------------------------
    *_chart_tools(),
    {
        "type": "function",
        "function": {
            "name": "describe_table",
            "description": (
                "選択中DBのテーブル詳細（列・型・説明・コード値の意味・実値の分布・サンプル行）を取得する。"
                "初めて使うテーブルでSQLを書く前に呼んで、列名や値の実体を確認する。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "db": {"type": "string", "description": "DBのエイリアス名。"},
                    "table": {"type": "string", "description": "テーブル名。"},
                },
                "required": ["db", "table"],
            },
        },
    },
]


def _allow_result_id(node) -> None:
    """sql を受け取る所すべてに result_id を足し、sql を必須から外す。

    レポートの節やExcelのシートのように、入れ子の中にも sql がある。
    1つずつ手で書き足すと必ず抜けるので、木をたどって機械的に付ける。
    """
    if isinstance(node, list):
        for v in node:
            _allow_result_id(v)
        return
    if not isinstance(node, dict):
        return
    props = node.get("properties")
    if isinstance(props, dict) and "sql" in props and "result_id" not in props:
        props["result_id"] = dict(_RESULT_ID)
        req = node.get("required")
        if isinstance(req, list) and "sql" in req:
            node["required"] = [r for r in req if r != "sql"]
    for v in node.values():
        _allow_result_id(v)


_allow_result_id(BUILTIN_TOOLS)

# plot_chart は用途別のグラフツール（plot_comparison など）で完全に置き換えられる。
# 同じことが2通りでできると、AIはどちらを使うか毎回迷い、定義の文字数も倍かかる。
# 実処理は残してあるので、過去の会話やユーザー定義の上書きが壊れることはない。
_RETIRED = {"plot_chart"}
BUILTIN_TOOLS = [t for t in BUILTIN_TOOLS if t["function"]["name"] not in _RETIRED]


# ==========================================================================
# ===== 元 tools/business.py
# 業務でよく聞かれる分析のツール。
#
# business.py（期間比較・ファネル・コホート・併売）と
# advanced.py（異常検知・生存時間）をツールとして公開する。
# データ品質チェックだけは、SQLを受け取るのではなくDBそのものを見に行くので
# ここに実処理を置く。
# ==========================================================================
import advanced
import business
import catalog
import db

_compare_periods = _analysis_tool(lambda a, c, r: business.compare_periods(
    c, r, a.get("period_col"), a.get("value_col"),
    dimension_col=a.get("dimension_col"), current=a.get("current"),
    previous=a.get("previous"), qty_col=a.get("qty_col")))


_funnel_analysis = _analysis_tool(lambda a, c, r: business.funnel_analysis(
    c, r, a.get("steps") or [], labels=a.get("labels"),
    group_col=a.get("group_col")))


_cohort_analysis = _analysis_tool(lambda a, c, r: business.cohort_analysis(
    c, r, a.get("id_col"), a.get("period_col"), value_col=a.get("value_col"),
    max_periods=int(a.get("max_periods") or 12)))


_market_basket = _analysis_tool(lambda a, c, r: business.market_basket(
    c, r, a.get("transaction_col"), a.get("item_col"),
    min_support=float(a.get("min_support") or 1.0), top=int(a.get("top") or 25)))


_detect_anomalies = _analysis_tool(lambda a, c, r: advanced.detect_anomalies(
    c, r, a.get("time_col"), a.get("value_col"),
    window=int(a.get("window") or 7), threshold=float(a.get("threshold") or 3.0),
    season_length=int(a["season_length"]) if a.get("season_length") else None,
    changepoints=a.get("changepoints", True)))


_survival_analysis = _analysis_tool(lambda a, c, r: advanced.survival_analysis(
    c, r, a.get("duration_col"), event_col=a.get("event_col"),
    group_col=a.get("group_col")))


# =============================================================================
# データ品質チェック（DBを直接見る）
# =============================================================================

#: 1回のチェックで見るテーブルの上限。全部見ると時間がかかりすぎる。
_MAX_TABLES = 12
#: これより大きいテーブルでは、重い集計（種類数）を省く。
_HEAVY_ROWS = 200_000


def _q(alias: str, table: str) -> str:
    return f'{alias}."{table}"'


def _column_stats(scope: list[dict], alias: str, table: str, cols: list[dict],
                  rowcount: int | None) -> tuple:
    """1テーブルぶんの欠損・種類数を1本のSQLで数える。"""
    heavy = (rowcount or 0) <= _HEAVY_ROWS
    parts = ["COUNT(*) AS n"]
    for i, c in enumerate(cols):
        name = c["name"].replace('"', '""')
        parts.append(f'SUM(CASE WHEN "{name}" IS NULL THEN 1 ELSE 0 END) AS nul{i}')
        if heavy:
            parts.append(f'COUNT(DISTINCT "{name}") AS uni{i}')
        if str(c.get("type") or "").upper().startswith(("TEXT", "VARCHAR", "CHAR")):
            parts.append(f"SUM(CASE WHEN TRIM(\"{name}\") = '' THEN 1 ELSE 0 END) AS emp{i}")
        else:
            parts.append(f"0 AS emp{i}")
    sql = f"SELECT {', '.join(parts)} FROM {_q(alias, table)}"
    _, rows, _ = db.run_select(sql, scope, max_rows=1)
    got = rows[0]
    n = int(got[0] or 0)
    out, pos = [], 1
    for c in cols:
        nul = int(got[pos] or 0)
        pos += 1
        uni = int(got[pos] or 0) if heavy else None
        pos += 1 if heavy else 0
        emp = int(got[pos] or 0)
        pos += 1
        out.append({"column": c["name"], "type": c.get("type") or "",
                    "nulls": nul, "empty": emp, "unique": uni})
    return n, out


def _pk_duplicates(scope: list[dict], alias: str, table: str, pk: list) -> int | None:
    """主キーが重複している組み合わせの数。"""
    if not pk:
        return None
    keys = ", ".join(f'"{c}"' for c in pk)
    sql = (f"SELECT COUNT(*) FROM (SELECT {keys} FROM {_q(alias, table)} "
           f"GROUP BY {keys} HAVING COUNT(*) > 1)")
    try:
        _, rows, _ = db.run_select(sql, scope, max_rows=1)
        return int(rows[0][0] or 0)
    except Exception:
        return None


def _orphans(scope: list[dict], child: tuple, parent: tuple) -> int | None:
    """親に居ない子（孤立した外部キー）の件数。"""
    ca, ct, cc = child
    pa, pt, pc = parent
    sql = (f'SELECT COUNT(*) FROM {_q(ca, ct)} c WHERE c."{cc}" IS NOT NULL '
           f'AND NOT EXISTS (SELECT 1 FROM {_q(pa, pt)} p WHERE p."{pc}" = c."{cc}")')
    try:
        _, rows, _ = db.run_select(sql, scope, max_rows=1)
        return int(rows[0][0] or 0)
    except Exception:
        return None


def _date_range(scope: list[dict], alias: str, table: str, col: str) -> tuple | None:
    """日付列の最小・最大。データがいつまで入っているかを見る。"""
    try:
        _, rows, _ = db.run_select(
            f'SELECT MIN("{col}"), MAX("{col}") FROM {_q(alias, table)}',
            scope, max_rows=1)
        return rows[0][0], rows[0][1]
    except Exception:
        return None


def _looks_like_date(col: dict, sample: str = "") -> bool:
    name = str(col.get("name") or "").lower()
    return (str(col.get("type") or "").upper().startswith("DATE")
            or any(k in name for k in ("date", "日", "_at", "time", "月")))


def _data_quality(args: dict, scope: list[dict]) -> dict:
    """選択中のDBを見て、分析の前に気づいておくべき異常を洗い出す。"""
    if not scope:
        return _err("対象のDBがありません。")

    # テーブル名は 'stocks' でも 'demo_inventory.stocks' でもよい
    # （プロンプトが『DB名.テーブル名』で書くよう求めているので、後者で来ることが多い）
    want = [str(t) for t in (args.get("tables") or [])]
    def _wanted(alias, tname):
        return (not want) or tname in want or f"{alias}.{tname}" in want
    issues, tbl_rows, col_rows, ref_rows = [], [], [], []
    checked = 0

    # 結合定義はDBをまたぐので、スコープ全体で一度だけ組み立てる
    try:
        entries = [{"alias": e["alias"], "profile": catalog.profile_db(e["path"]),
                    "meta": e.get("meta") or catalog.load_meta(e["path"])} for e in scope]
        edges = catalog.collect_edges(entries)
    except Exception:
        edges = []

    for s in scope:
        alias = s["alias"]
        try:
            profile = catalog.profile_db(s["path"])
        except Exception as e:
            issues.append(("高", f"{alias}: プロファイルを読めませんでした（{e}）"))
            continue
        meta = s.get("meta") or catalog.load_meta(s["path"])
        allowed = set(s.get("tables") or profile.get("tables") or {})

        for tname, t in (profile.get("tables") or {}).items():
            if tname not in allowed:
                continue
            if not _wanted(alias, tname):
                continue
            if checked >= _MAX_TABLES:
                break
            checked += 1
            cols = t.get("columns") or []
            try:
                n, stats = _column_stats(scope, alias, tname, cols, t.get("row_count"))
            except Exception as e:
                issues.append(("中", f"{alias}.{tname}: 集計できませんでした（{e}）"))
                continue

            pk, pk_src = catalog.effective_pk(profile, meta, tname)
            dup = _pk_duplicates(scope, alias, tname, pk)
            tbl_rows.append([f"{alias}.{tname}", n, len(cols),
                             "、".join(pk) if pk else "（無し）",
                             dup if dup is not None else "—"])
            if n == 0:
                issues.append(("高", f"{alias}.{tname} は0行です。取り込みが済んでいない可能性があります。"))
                continue
            if dup:
                issues.append(("高", f"{alias}.{tname} は主キー（{'、'.join(pk)}）が "
                                     f"{dup} 組で重複しています。件数や金額が二重に数えられます。"))
            if not pk:
                issues.append(("中", f"{alias}.{tname} に主キーがありません。"
                                     "重複を検出できないので、カタログで指定してください。"))

            for st in stats:
                nul_pct = round(st["nulls"] / n * 100, 1) if n else 0.0
                emp_pct = round(st["empty"] / n * 100, 1) if n else 0.0
                col_rows.append([f"{alias}.{tname}", st["column"], st["type"],
                                 nul_pct, emp_pct,
                                 st["unique"] if st["unique"] is not None else "—"])
                if st["nulls"] == n:
                    issues.append(("中", f"{alias}.{tname}.{st['column']} は全て空です。"
                                         "この列は集計に使えません。"))
                elif nul_pct >= 30:
                    issues.append(("中", f"{alias}.{tname}.{st['column']} は "
                                         f"{nul_pct}% が空です。平均を取ると母数がずれます。"))
                elif emp_pct >= 10:
                    issues.append(("低", f"{alias}.{tname}.{st['column']} は "
                                         f"{emp_pct}% が空文字です。NULLと混在しています。"))
                if st["unique"] == 1 and n > 1:
                    issues.append(("低", f"{alias}.{tname}.{st['column']} は1種類の値しかありません。"))

            # データがいつまで入っているか（古いまま気づかないのを防ぐ）
            for c in cols:
                if not _looks_like_date(c):
                    continue
                rng = _date_range(scope, alias, tname, c["name"])
                if rng and rng[1]:
                    tbl_rows[-1].append(f"{c['name']}: {rng[0]} 〜 {rng[1]}")
                    break

        # 参照整合性。カタログの結合定義とDBのFK宣言の両方を見る。
        # 子と親は保存順ではなく主キーの位置から決める（手書きのYAMLが
        # 逆向きでも、「親に居ない子」を正しい向きで数えるため）
        for edge in edges:
            (ca, ct, cc), (pa, pt, pc) = catalog.child_parent(entries, edge)
            if ca != alias or ct not in allowed:
                continue
            miss = _orphans(scope, (ca, ct, cc), (pa, pt, pc))
            if miss is None:
                continue
            kind = "FK宣言" if edge.get("kind") == "fk" else "カタログの結合定義"
            ref_rows.append([f"{ca}.{ct}.{cc}", f"{pa}.{pt}.{pc}", miss, kind])
            if miss:
                issues.append(("高", f"{ca}.{ct}.{cc} の {miss} 件が "
                                     f"{pa}.{pt} に存在しません。"
                                     "内部結合すると、この件数ぶん落ちます。"))

    if not tbl_rows:
        return _err("調べられるテーブルがありませんでした。"
                    "tables に指定した名前が合っているか確認してください"
                    "（例: 'stocks' または 'demo_inventory.stocks'）。")

    head = ["テーブル", "行数", "列数", "主キー", "主キー重複"]
    if any(len(r) > 5 for r in tbl_rows):
        head.append("日付の範囲")
    tbl_rows = [r + [""] * (len(head) - len(r)) for r in tbl_rows]

    rank = {"高": 0, "中": 1, "低": 2}
    issues.sort(key=lambda x: rank.get(x[0], 3))
    tables = [_table_of("テーブル", head, tbl_rows),
              _table_of("見つかった問題", ["深刻度", "内容"],
                        [[lv, msg] for lv, msg in issues[:80]] or [["—", "問題は見つかりませんでした。"]]),
              _table_of("列ごとの状態", ["テーブル", "列", "型", "空の割合(%)",
                                        "空文字の割合(%)", "値の種類数"], col_rows)]
    if ref_rows:
        tables.append(_table_of("参照整合性", ["子", "親", "親に無い件数", "定義元"], ref_rows))

    high = [m for lv, m in issues if lv == "高"]
    notes = [f"{checked} テーブルを調べました。"
             + (f"深刻な問題が {len(high)} 件あります。" if high
                else "分析を止めるような問題は見つかりませんでした。")]
    notes += high[:5]
    if checked >= _MAX_TABLES:
        notes.append(f"テーブルが多いため {_MAX_TABLES} 件までにしています。"
                     "続きは tables で対象を指定してください。")
    notes.append("行数が0・主キーの重複・親に無い外部キーは、集計結果を直接ゆがめます。"
                 "先にここを直してから数字を読んでください。")

    return _report_result({"title": "データ品質チェック", "tables": tables, "notes": notes,
                           "meta": {"tables_checked": checked,
                                    "issues": len(issues), "critical": len(high)}},
                          scope=scope)


def _table_of(name: str, columns: list, rows: list) -> dict:
    return {"name": name, "columns": columns, "rows": [tuple(r) for r in rows]}


HANDLERS_business = {
    "compare_periods": _compare_periods,
    "funnel_analysis": _funnel_analysis,
    "cohort_analysis": _cohort_analysis,
    "market_basket": _market_basket,
    "detect_anomalies": _detect_anomalies,
    "survival_analysis": _survival_analysis,
    "data_quality": _data_quality,
}

# data_quality は自分でDBを見に行くので、SQLプレビューの対象外
SQL_TOOLS_business = {"compare_periods", "funnel_analysis", "cohort_analysis",
             "market_basket", "detect_anomalies", "survival_analysis"}


# ==========================================================================
# ===== 元 tools/files.py
# 取り込み元フォルダを調べるツール（読むだけ）。
#
# 「取り込み元に何が来ているか」「このCSVはどんな列か」「まだ取り込んでいない
# ファイルはあるか」に答えるためのもの。DBに入る前のファイルの話なので、
# SQLでは答えられない。
#
# 安全のうえで大事なところは、すべて importer.py の既存の仕組みに任せる。
#   allowed_dirs()   … 読んでよいフォルダ（env の IMPORT_DIRS ＋画面で追加した分）
#   is_allowed()     … .. やリンクで許可フォルダの外へ出ようとしても弾く
#   check_readable() … 読む直前にもう一度確かめる
# ここで新しくパスの判定を書かない（守りの仕組みを二重に持つと必ずズレる）。
#
# できるのは一覧と下見だけ。取り込み・作成・変更・削除は一切しない。
# 取り込みの実行は今までどおり「データ取り込み」画面の操作に限る。
# ==========================================================================
from datetime import datetime
from pathlib import Path

import config
import filecheck
import history
import importer

#: 一覧で返す最大件数。多すぎるとLLMに渡すだけで無駄になる。
_MAX_ROWS = 300
#: 下見で見せる行数の上限。
_MAX_PREVIEW_ROWS = 20
#: 一覧でまとめて形を判定するときの上限。1件ずつ開くので数を抑える。
_MAX_CHECK = 20


def _files_table(name: str, columns: list, rows: list) -> dict:
    return {"name": name, "columns": columns, "rows": [tuple(r) for r in rows]}


def _size(n: int | None) -> str:
    if n is None:
        return ""
    if n < 1024:
        return f"{n} B"
    if n < 1024 * 1024:
        return f"{n / 1024:.1f} KB"
    return f"{n / 1024 / 1024:.1f} MB"


def _mtime(p: Path) -> str:
    try:
        return datetime.fromtimestamp(p.stat().st_mtime).strftime("%Y-%m-%d %H:%M")
    except OSError:
        return ""


def _imported_index() -> dict:
    """取り込み済みファイルの索引。パスの表記ゆれに備えて小文字で引く。"""
    out = {}
    for src, rec in history.latest_by_source().items():
        out[src.strip().lower()] = rec
    return out


def _import_state(p: Path, index: dict) -> str:
    """そのファイルが取り込み済みかどうかの一言。"""
    rec = index.get(str(p).lower())
    if rec is None:
        # パスが変わっていても、ファイル名が一致すれば手がかりにはなる
        rec = next((r for key, r in index.items()
                    if Path(key).name == p.name.lower()), None)
        if rec is None:
            return "未取り込み"
        return (f"同名を取り込み済み（{rec.get('db_file')} / {rec.get('table')}"
                f"・{rec.get('at', '')[:16]}）")
    mark = "" if rec.get("ok") else "／前回は失敗"
    return f"{rec.get('db_file')} / {rec.get('table')}（{rec.get('at', '')[:16]}{mark}）"


def _roots_table() -> dict:
    """許可フォルダの状態。マウント切れや権限なしをここで切り分ける。"""
    rows = [[r["設定値"], r["状態"], r["source"]] for r in importer.dir_status()]
    return _files_table("取り込み元フォルダ", ["フォルダ", "状態", "設定元"], rows)


def _supported(p: Path) -> bool:
    return p.suffix.lower() in config.IMPORT_EXTENSIONS


def _listing(args: dict) -> dict:
    path = str(args.get("path") or "").strip()
    recursive = bool(args.get("recursive"))
    pattern = str(args.get("pattern") or "").strip().lower()
    only_new = bool(args.get("only_not_imported"))
    check = bool(args.get("check"))
    index = _imported_index()

    roots = importer.allowed_dirs()
    if not roots:
        return _err("取り込み元フォルダが設定されていません。"
                    "「データ取り込み」画面で追加するか、env の IMPORT_DIRS を設定してください。")

    if path and not importer.is_allowed_dir(Path(path)):
        return _err(f"そのフォルダは見られません（許可フォルダの外です）: {path}。"
                    f"見られるのは {'、'.join(str(d) for d in roots)} の中だけです。")

    here = Path(path).resolve() if path else None
    dirs: list = []
    files: list[Path] = []

    # 一覧は拡張子で絞らない。「何が置いてあるか」を知るのが目的なので、
    # 取り込めない形式（PDFなど）も見せて、可否は列で示す。
    if recursive:
        for p in importer.list_all_files():
            if here is None or here == p.parent or here in p.parents:
                files.append(p)
    else:
        for d in ([here] if here is not None else roots):
            try:
                entries = sorted(d.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower()))
            except OSError:
                continue
            for p in entries:
                try:
                    if p.is_dir():
                        if not importer.is_noise(p.name):
                            dirs.append(str(p.relative_to(d)) if here is not None
                                        else f"{d.name}/{p.name}")
                    elif importer.is_within_allowed(p):
                        files.append(p.resolve())
                except OSError:
                    continue

    if pattern:
        files = [f for f in files if pattern in f.name.lower()]
        dirs = [d for d in dirs if pattern in str(d).lower()]

    rows, checked = [], 0
    for p in files:
        state = _import_state(p, index)
        if only_new and state != "未取り込み":
            continue
        if not _supported(p):
            shape = "取り込み対象外の形式"
        elif check and checked < _MAX_CHECK:
            checked += 1
            try:
                shape = filecheck.summary_line(filecheck.inspect(p))
            except Exception as e:
                shape = f"判定できず（{type(e).__name__}）"
        else:
            shape = "未判定"
        rows.append([importer.display_name(p), p.suffix.lower().lstrip(".") or "（なし）",
                     _size(p.stat().st_size if p.exists() else None), _mtime(p),
                     state, shape])
    rows.sort(key=lambda r: r[3], reverse=True)      # 新しいものから
    shown = rows[:_MAX_ROWS]

    tables = [_roots_table()]
    if dirs:
        tables.append(_files_table("フォルダ", ["名前"], [[d] for d in sorted(set(dirs))]))
    tables.append(_files_table("ファイル",
                         ["場所", "種類", "サイズ", "更新日時", "取り込み状況", "表として使えるか"],
                         shown))

    supported = [r for r in rows if r[5] != "取り込み対象外の形式"]
    notes = [f"{len(rows)} 件のファイルが見つかりました"
             + (f"（多いので新しい順に {len(shown)} 件だけ載せています）"
                if len(rows) > len(shown) else "")
             + f"。うち取り込みに対応した形式は {len(supported)} 件です"
             + f"（対応: {'、'.join(config.IMPORT_EXTENSIONS)}）。"]
    if not recursive and not path:
        notes.append("いまは許可フォルダの直下だけを見ています。"
                     "下の階層も見るなら recursive=true、"
                     "特定のフォルダを見るなら path を指定してください。")
    if not check:
        notes.append("「表として使えるか」は check=true を付けると調べます"
                     "（1件ずつ開くので、多いときは絞ってから）。")
    elif checked >= _MAX_CHECK:
        notes.append(f"判定は {_MAX_CHECK} 件までにしています。"
                     "pattern や path で絞ると続きを見られます。")
    bad = [r for r in rows if r[5].startswith(("取り込みに向かない", "手直しが要る"))]
    if bad:
        notes.append("そのままでは取り込めないものがあります: "
                     + "、".join(f"{r[0]}（{r[5]}）" for r in bad[:3]))
    fresh = [r for r in rows if r[4] == "未取り込み"]
    if fresh:
        notes.append(f"まだ取り込んでいないファイルが {len(fresh)} 件あります: "
                     + "、".join(r[0] for r in fresh[:5])
                     + ("（ほか）" if len(fresh) > 5 else ""))
    if not rows and not dirs:
        notes.append("ファイルはありませんでした。フォルダが空か、権限が無い可能性があります。")
    notes.append("中身と形を詳しく見るには file にパスを指定してください。"
                 "実際に取り込むのは「データ取り込み」画面の操作です。ここでは読むだけです。")

    return _report_result({"title": "取り込み元フォルダの中身",
                           "tables": tables, "notes": notes,
                           "meta": {"files": len(rows), "supported": len(supported),
                                    "not_imported": len(fresh)}})


def _issues_table(res: dict) -> dict:
    return _files_table("見つかった問題", ["深刻度", "内容", "直し方"],
                  [[i["level"], i["text"], i["fix"]] for i in res["issues"]]
                  or [["—", "気になる点はありませんでした。", ""]])


def _preview(args: dict) -> dict:
    raw = Path(str(args.get("file") or "").strip())
    # 場所の判定と「読める形式か」の判定を分ける。対象外の形式でも
    # 「なぜ取り込めないか」は答えられるようにする。
    if not importer.is_within_allowed(raw):
        return _err(f"そのファイルは見られません: {raw}。"
                    "許可フォルダの中のファイルだけを指定してください。"
                    "パスは一覧（file を指定しない呼び方）で得たものを使ってください。")
    target = raw.resolve()

    if not _supported(target):
        res = filecheck.inspect(target)
        return _report_result({
            "title": f"{target.name} は取り込みに対応していない形式です",
            "tables": [_issues_table(res)],
            "notes": [f"{importer.display_name(target)}"
                      f"（{_size(target.stat().st_size)}・更新 {_mtime(target)}）",
                      f"扱えるのは {'、'.join(config.IMPORT_EXTENSIONS)} です。"
                      "中身は読んでいません。"],
            "meta": {"verdict": res["verdict"]}})

    rows_want = max(1, min(int(args.get("rows") or 5), _MAX_PREVIEW_ROWS))

    try:
        sheets = importer.sheet_names(target)
    except importer.ImportError_ as e:
        return _err(str(e))
    sheet = args.get("sheet") or (sheets[0] if sheets else None)

    # まず形を見る。見出しが1行目に無ければ、その行で読み直す
    try:
        res = filecheck.inspect(target, sheet=sheet)
    except Exception as e:
        res = {"verdict": "判定できず", "issues": [
            {"level": "低", "text": f"形を調べられませんでした: {e}", "fix": ""}],
            "header_row": 0, "shape": {}}
    header_row = (int(args["header_row"]) if args.get("header_row") is not None
                  else int(res.get("header_row") or 0))

    try:
        df = importer.read_table(target, sheet=sheet, header_row=header_row,
                                 nrows=rows_want)
    except importer.ImportError_ as e:
        return _err(f"{target.name}: {e}")
    except Exception as e:
        return _err(f"{target.name} を読めませんでした: {e}")

    plan = importer.plan_columns(df)
    tables = [
        _files_table("判定", ["項目", "内容"],
               [["そのまま取り込めるか", res["verdict"]],
                ["見出しの行", f"{header_row + 1} 行目"],
                *[[k, v] for k, v in (res.get("shape") or {}).items() if k != "見出し行"]]),
        _issues_table(res),
        _files_table("列", ["元の列名", "取り込み後の列名", "推定される型"],
               [[c["元の列名"], c["列名"], c["型"]] for c in plan]),
        _files_table(f"先頭 {len(df)} 行", [str(c) for c in df.columns],
               [[("" if v is None else v) for v in r] for r in df.values.tolist()]),
    ]
    if sheets:
        tables.insert(0, _files_table("シート", ["シート名"], [[s] for s in sheets]))

    index = _imported_index()
    notes = [f"{importer.display_name(target)}（{_size(target.stat().st_size)}"
             f"・更新 {_mtime(target)}）",
             f"判定: {res['verdict']}",
             f"取り込み状況: {_import_state(target, index)}"]
    for i in res.get("issues", []):
        if i["level"] == "高":
            notes.append(f"{i['text']} → {i['fix']}")
    if sheets:
        notes.append(f"シートは {len(sheets)} 枚あります（いま見ているのは「{sheet}」）。"
                     "シートごとに形が違うので、使いたいシートを sheet で指定して確かめてください。")
    if res.get("encoding"):
        notes.append(f"文字コード: {res['encoding']} / 区切り: {res.get('delimiter')}")
    if header_row:
        notes.append(f"見出しが1行目ではないので、{header_row + 1} 行目を見出しとして読みました。"
                     f"取り込み画面でも「見出しの行」に {header_row + 1} を指定してください。")
    if df.empty:
        notes.append("中身の行が読めませんでした。header_row を変えて試してください。")
    notes.append("型は中身からの推定で、取り込み画面で直せます。"
                 "ここでは読むだけで、取り込みはしていません。")

    return _report_result({"title": f"{target.name} の下見（{res['verdict']}）",
                           "tables": tables, "notes": notes,
                           "meta": {"verdict": res["verdict"], "columns": len(plan),
                                    "header_row": header_row, "sheets": sheets}})


def _explore_import_files(args: dict, scope: list[dict]) -> dict:
    try:
        return _preview(args) if str(args.get("file") or "").strip() else _listing(args)
    except importer.ImportError_ as e:
        return _err(str(e))
    except PermissionError:
        return _err("読み取り権限がありません。共有フォルダの権限を確認してください。")
    except OSError as e:
        return _err(f"フォルダにアクセスできませんでした: {e.strerror or e}")


HANDLERS_files = {"explore_import_files": _explore_import_files}

# SQLは受け取らない
SQL_TOOLS_files: set = set()

# 管理者だけに渡すツール。「データ取り込み」画面が管理者専用なので、
# AI経由なら誰でも中身が見られる、という抜け道を作らない。
ADMIN_TOOLS_files = {"explore_import_files"}


# ==========================================================================
# ===== 元 tools/mail.py
# メールの宛先探しと下書き。送信そのものは画面のボタンからだけ。
# ==========================================================================
import mailer


def _find_mail_recipients(args: dict, scope: list[dict]) -> dict:
    res = mailer.find_recipients(scope, args.get("query") or "",
                                 limit=int(args.get("limit") or 50),
                                 table=args.get("table"))
    cands = res["candidates"]
    return {
        "ok": res["ok"],
        "llm_content": _json({
            "status": "recipients", "message": res["message"],
            "sources": res["sources"],
            "count": len(cands),
            "candidates": [{k: c[k] for k in ("email", "name", "dept", "source")}
                           for c in cands[:50]],
            "note": "ここに出たアドレスだけを宛先に使うこと。推測で作らない。",
        }),
        "render": {"role": "assistant", "kind": "table",
                   "columns": ["メールアドレス", "氏名", "部署", "出所"],
                   "rows": [[c["email"], c["name"], c["dept"], c["source"]]
                            for c in cands]} if cands else
                  {"role": "assistant", "kind": "text", "content": res["message"]},
    }


def _compose_email(args: dict, scope: list[dict]) -> dict:
    to = list(args.get("to") or [])
    matched = []
    if args.get("to_query"):
        res = mailer.find_recipients(scope, args["to_query"], limit=50)
        matched = res["candidates"]
        to += [c["email"] for c in matched if c["valid"]]
        if not matched:
            return _err(f"「{args['to_query']}」に一致する宛先が見つかりませんでした。"
                        "find_mail_recipients で候補を確認してください。")
    to = list(dict.fromkeys(a for a in to if a))
    draft = {"to": to, "cc": args.get("cc") or [], "bcc": args.get("bcc") or [],
             "subject": args.get("subject") or "", "body": args.get("body") or "",
             "reply_to": args.get("reply_to") or "",
             "attach_filenames": list(args.get("attach_filenames") or [])}
    view = mailer.preview(draft)
    # 添付は会話ログから web 側が解決する。ここでは名前だけ持たせる。
    return {
        "ok": not view["errors"],
        "llm_content": _json({
            "status": "mail_draft" if not view["errors"] else "mail_draft_invalid",
            "to": view["to"], "cc": view["cc"], "bcc_count": len(view["bcc"]),
            "subject": view["subject"], "body_lines": view["body_lines"],
            "attach_filenames": draft["attach_filenames"],
            "matched_from_db": [{"email": c["email"], "name": c["name"],
                                 "dept": c["dept"]} for c in matched[:20]],
            "problems": view["errors"],
            "note": ("下書きを画面に出した。送信するかはユーザーが画面のボタンで決める。"
                     "こちらから送信することはできないので、"
                     "『内容を確認して送信ボタンを押してください』と伝えること。"),
        }),
        "render": {"role": "assistant", "kind": "mail_draft", "draft": draft,
                   "preview": view},
    }

HANDLERS_mail = {
    "find_mail_recipients": _find_mail_recipients,
    "compose_email": _compose_email,
}

SQL_TOOLS_mail: set[str] = set()


# ==========================================================================
# ===== 元 tools/query.py
# 調べる・集計する・描く・出す。SQLの結果をそのまま扱うツール。
# ==========================================================================
from pathlib import Path

import advanced
import analysis
import catalog
import charts
import config
import excel
import exports


def _example_exists(sql: str) -> bool:
    """このSQLと同じ例文がどこかのDBに登録済みか。"""
    if not sql.strip():
        return False
    import db as dbmod
    for f in dbmod.list_db_files():
        if catalog.find_example(catalog.load_meta(f).get("examples") or [], sql):
            return True
    return False


def _run_sql_query(args: dict, scope: list[dict]) -> dict:
    try:
        columns, rows, truncated, rid, total = fetch(args, scope,
                                                     label=args.get("purpose"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"SQL実行エラー: {e}")

    sample = rows[: config.SAMPLE_ROWS_FOR_LLM]
    llm_content = _json({
        "columns": columns,
        "row_count": len(rows),
        "rows": [list(r) for r in sample],
        "result_id": rid,
        "note": (f"全{len(rows)}行中 先頭{len(sample)}行を表示。"
                 f"この結果は result_id '{rid}' で他のツールから使い回せます"
                 "（同じSQLを書き直さなくてよい）。"
                 if len(rows) > len(sample) else
                 f"この結果は result_id '{rid}' で他のツールから使い回せます。"),
        **source_note(len(rows), truncated, total),
        # AIが「例文に登録しますか？」と聞くべきかの手がかり。
        # 既に同じSQLの例文があれば聞かない（同じ提案の繰り返しは邪魔になる）。
        "example_registered": _example_exists(str(args.get("sql") or "")),
    })
    return {
        "ok": True,
        "llm_content": llm_content,
        "render": {
            "role": "assistant", "kind": "table",
            "columns": columns, "rows": rows, "truncated": truncated,
        },
    }


_CHART_FIELDS = ("chart_type", "x", "y", "color", "size", "text", "path",
                 "nbins", "orientation", "barmode", "title",
                 # 種別ごとに使う指定
                 "y2", "z", "lower", "upper", "facet", "dimensions",
                 "source", "target", "start", "end",
                 "open", "high", "low", "close",
                 "value", "agg", "max", "suffix", "valueformat",
                 "colorscale", "marginal", "trendline")


def _plot_chart(args: dict, scope: list[dict]) -> dict:
    try:
        columns, rows, truncated, rid, total = fetch(args, scope,
                                                     label=args.get("title"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"グラフ用SQLの実行エラー: {e}")

    item = {k: args.get(k) for k in _CHART_FIELDS}
    item["chart_type"] = item.get("chart_type") or "bar"
    errs = charts.validate(item, columns)
    if errs:
        return _err(" / ".join(errs))

    return {
        "ok": True,
        "llm_content": _json({
            "status": "chart_rendered",
            "chart_type": item["chart_type"],
            "columns": columns,
            "row_count": len(rows),
            "result_id": rid,
            **source_note(len(rows), truncated, total),
        }),
        "render": {
            "role": "assistant", "kind": "chart",
            "columns": columns, "rows": rows, **item,
            "title": args.get("title", ""),
        },
    }


def _plot_dual_axis(args: dict, scope: list[dict]) -> dict:
    try:
        columns, rows, truncated, rid, total = fetch(args, scope,
                                                     label=args.get("title"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"2軸グラフ用SQLの実行エラー: {e}")

    x = args.get("x")
    bar_y = args.get("bar_y") or []
    line_y = args.get("line_y") or []
    needed = [x] + list(bar_y) + list(line_y)
    missing = [c for c in needed if c and c not in columns]
    if missing or not bar_y or not line_y:
        msg = (f"指定列が結果に存在しません: {missing} / 利用可能な列: {columns}"
               if missing else "bar_y と line_y にはそれぞれ1つ以上の数値列を指定してください。")
        return _err(msg)

    return {
        "ok": True,
        "llm_content": _json({
            "status": "dual_axis_chart_rendered",
            "columns": columns, "row_count": len(rows),
            "bar_y": bar_y, "line_y": line_y, "result_id": rid,
            **source_note(len(rows), truncated, total),
        }),
        "render": {
            "role": "assistant", "kind": "chart_dual",
            "columns": columns, "rows": rows,
            "x": x, "bar_y": bar_y, "line_y": line_y,
            "left_title": args.get("left_title"), "right_title": args.get("right_title"),
            "title": args.get("title", ""),
        },
    }


def _describe_table(args: dict, scope: list[dict]) -> dict:
    text = catalog.describe_table_text(scope, args.get("db", ""), args.get("table", ""))
    return {"ok": not text.startswith("エラー"), "llm_content": text, "render": None}


def _pivot_table(args: dict, scope: list[dict]) -> dict:
    try:
        columns, rows, truncated, rid, total = fetch(args, scope,
                                                     label=args.get("title"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"クロス集計用SQLの実行エラー: {e}")
    if not rows:
        return _err("データが0行でした。抽出条件を見直してください。")

    try:
        cols, prows = analysis.pivot(
            columns, rows,
            index=args.get("index") or [], cols=args.get("columns") or None,
            values=args.get("values"), aggfunc=args.get("aggfunc") or "sum",
            margins=bool(args.get("margins")),
            percent=args.get("percent"), rank_by=args.get("rank"),
        )
    except Exception as e:
        return _err(f"クロス集計に失敗しました: {e}")

    title = args.get("title") or "クロス集計"
    # 集計後の表もグラフやレポートの材料になるので、指せるようにして返す
    out_rid = results.put(scope, cols, prows, label=f"{title}（クロス集計の結果）")
    llm_content = _json({
        "status": "pivot_ready", "columns": cols, "row_count": len(prows),
        "rows": [list(r) for r in prows[: config.SAMPLE_ROWS_FOR_LLM]],
        "result_id": out_rid, "source_result_id": rid,
        "note": f"集計後の表は result_id '{out_rid}' でグラフやレポートに渡せます。",
        **source_note(len(rows), truncated, total),
    })
    if (args.get("render") or "table") == "heatmap":
        render = {"role": "assistant", "kind": "chart", "columns": cols, "rows": prows,
                  "chart_type": "matrix", "x": cols[0], "title": title}
    else:
        render = {"role": "assistant", "kind": "table", "columns": cols, "rows": prows,
                  "truncated": False}
    return {"ok": True, "llm_content": llm_content, "render": render}


def _analyze_stats(args: dict, scope: list[dict]) -> dict:
    method = args.get("method") or "describe"
    try:
        columns, rows, truncated, rid, total = fetch(args, scope,
                                                     label=args.get("title"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"分析用SQLの実行エラー: {e}")
    if not rows:
        return _err("データが0行でした。抽出条件を見直してください。")

    title = args.get("title") or {"describe": "基本統計量", "correlation": "相関",
                                  "outliers": "外れ値"}.get(method, "分析")
    note = f"（元データ {len(rows):,} 行"
    note += "／上限で切り詰め済み）" if truncated else "）"

    try:
        if method == "describe":
            cols, srows = analysis.describe(columns, rows, args.get("columns"),
                                            args.get("group_by"))
            extra = {}
            render = {"role": "assistant", "kind": "table", "columns": cols, "rows": srows}
        elif method == "correlation":
            cm = args.get("corr_method") or "pearson"
            lag = int(args.get("lag") or 0)
            if args.get("partial"):
                # 交絡を取り除いた相関。「効いて見えるのは第3の変数のせい」を切り分ける
                res = advanced.partial_correlation(
                    columns, rows, args.get("columns"), args.get("control") or [],
                    method=cm)
                return _report_result(res, source_rows=len(rows), truncated=truncated,
                                      total=total, result_id=rid, scope=scope,
                                      extra={"method": "partial_correlation"})
            if lag:
                # 時差相関。「広告費は翌月の売上に効く」を見るための道具
                res = advanced.lag_correlation(
                    columns, rows, args.get("target"), args.get("columns"),
                    max_lag=lag, method=cm)
                return _report_result(res, source_rows=len(rows), truncated=truncated,
                                      total=total, result_id=rid, scope=scope,
                                      extra={"method": "lag_correlation"})
            cols, srows = analysis.correlation(columns, rows, args.get("columns"), cm)
            extra = {"method": cm, "strong_pairs": analysis.correlation_pairs(columns, rows, cm)[:8],
                     "caution": "相関は因果ではありません。効いている理由を確かめるには "
                                "partial=true で交絡を除くか、regression を使ってください。"}
            render = {"role": "assistant", "kind": "chart", "columns": cols, "rows": srows,
                      "chart_type": "matrix", "x": cols[0], "title": f"{title}{note}",
                      "colorscale": "RdBu"}
        else:
            target = args.get("target")
            if not target:
                return _err("outliers には target（外れ値を調べる数値列）が必要です。")
            om = args.get("outlier_method") or "iqr"
            # mahalanobis は複数列をまとめて見る。"売上, 客数" のような指定も許す。
            cols_in = [t.strip() for t in str(target).split(",")] if isinstance(target, str) \
                else list(target)
            res = advanced.outliers_ext(columns, rows,
                                        cols_in if len(cols_in) > 1 else cols_in[0],
                                        method=om, threshold=args.get("threshold"))
            return _report_result(res, source_rows=len(rows), truncated=truncated,
                                  total=total, result_id=rid, scope=scope,
                                  extra={"method": "outliers", "outlier_method": om})
    except Exception as e:
        return _err(f"{title}の計算に失敗しました: {e}")

    out_rid = results.put(scope, cols, srows, label=title)
    return {
        "ok": True,
        "llm_content": _json({
            "status": "stats_ready", "method": method, "columns": cols,
            "row_count": len(srows),
            "rows": [list(r) for r in srows[: config.SAMPLE_ROWS_FOR_LLM]],
            "result_id": out_rid, "source_result_id": rid,
            **source_note(len(rows), truncated, total), **extra,
        }),
        "render": render,
    }


def _propose_glossary_term(args: dict, scope: list[dict]) -> dict:
    """業務用語の登録カードをチャットに出す（提案まで。保存は人がボタンで確定）。

    メールと同じ型: 作るのはAI、確定するのは人。カタログは全員共通の土台なので、
    会話の「はい」だけで書き換えず、カードのボタンという明示の操作を挟む。
    SQL式があれば実データで確かめてから見せる（通らない式を提案しない）。
    """
    import db as dbmod

    term = str(args.get("term") or "").strip()
    desc = str(args.get("description") or "").strip()
    sql = str(args.get("sql") or "").strip()
    table = str(args.get("table") or "").strip()
    if not term or not desc:
        return _err("term（用語）と description（説明）は必須です。")

    # 置き場のDBを決める。table を持つDBを探す（複数あれば db で指定させる）
    name = str(args.get("db") or "").strip()
    files = dbmod.list_db_files()
    target = None
    if name:
        low = name.lower()
        target = next((f for f in files
                       if low in (f.name.lower(), f.stem.lower(),
                                  dbmod.alias_for(f).lower())), None)
        if target is None:
            return _err(f"DB '{name}' が見つかりません。")
    elif table:
        owners = [f for f in files
                  if table in (catalog.profile_db(f).get("tables") or {})]
        if len(owners) == 1:
            target = owners[0]
        elif not owners:
            return _err(f"テーブル '{table}' を持つDBが見つかりません。")
        else:
            return _err(f"テーブル '{table}' は複数のDBにあります。"
                        "db でどのDBか指定してください: "
                        + "、".join(f.name for f in owners))
    else:
        return _err("table（用語を置くテーブル）か db を指定してください。")

    alias = dbmod.alias_for(target)
    if table and table not in (catalog.profile_db(target).get("tables") or {}):
        return _err(f"{target.name} にテーブル '{table}' がありません。")

    # SQL式の検証。条件式→該当件数 / 計算式→計算例 / 通らない→エラーで差し戻し
    verdict, detail = "", ""
    if sql:
        ref = f"{alias}.{table}" if table else None
        wide = dbmod.widen_scope(f"{alias}. {sql}", [
            {"path": str(target), "alias": alias, "name": target.name, "tables": None}])
        wide = dbmod.widen_scope(sql, wide)
        try:
            if ref:
                _, rows, _ = dbmod.run_select(
                    f"SELECT COUNT(*) AS n, (SELECT COUNT(*) FROM {ref}) AS t "
                    f"FROM {ref} WHERE {sql}", wide, max_rows=1)
                n, t = rows[0]
                verdict = "条件式"
                detail = f"該当 {n:,} 行 / 全 {t:,} 行"
            else:
                _, rows, _ = dbmod.run_select(f"SELECT {sql} AS v", wide, max_rows=1)
                verdict, detail = "計算式", f"計算結果: {rows[0][0]}"
        except Exception:
            try:
                src = ref or f"{alias}.{(list(catalog.profile_db(target)['tables']) or [''])[0]}"
                _, rows, _ = dbmod.run_select(f"SELECT {sql} AS v FROM {src}",
                                              wide, max_rows=1)
                verdict, detail = "計算式", f"計算結果の例: {rows[0][0]}"
            except Exception as e:
                return _err(f"SQL式が実データで通りませんでした: {str(e).splitlines()[0][:120]} "
                            "式を直して提案し直すか、SQL式なし（説明だけ）で提案してください。")

    meta_now = catalog.load_meta(target)
    current = (catalog.table_glossary(meta_now, table) if table
               else catalog.db_glossary(meta_now)).get(term)
    return {
        "ok": True,
        "llm_content": _json({
            "status": "proposed", "db": target.name, "table": table or "(DB全体)",
            "term": term, "verdict": verdict or "説明のみ", "detail": detail,
            "already_exists": current is not None,
            "note": "登録カードをユーザーの画面に出しました。登録するかはユーザーが"
                    "カードのボタンで決めます。あなたはこれ以上の操作をしなくてよい。",
        }),
        "render": {"role": "assistant", "kind": "glossary_term",
                   "db": target.name, "table": table, "term": term,
                   "description": desc, "sql": sql,
                   "how": str(args.get("how") or "").strip(),
                   "verdict": verdict, "detail": detail,
                   "exists": current is not None,
                   "old": current or None},
    }


def _propose_example(args: dict, scope: list[dict]) -> dict:
    """例文（質問とSQLのペア）の登録カードをチャットに出す。保存は人がボタンで確定。

    カードにはSQLは出さない。代わりに「どのデータをどこから取って、どう集計したか」の
    日本語（summary。AIが書く）と、実際に動かした結果の先頭数行を載せる。
    SQLを読めない人でも、中身を見て正しさを判断できるようにするため。
    """
    import db as dbmod

    q = str(args.get("question") or "").strip()
    sql = str(args.get("sql") or "").strip()
    summary = str(args.get("summary") or "").strip()
    if not q or not sql or not summary:
        return _err("question（質問文）・sql・summary（何をどう集計したかの日本語）は必須です。")

    wide = dbmod.widen_scope(sql, scope)
    try:
        columns, rows, truncated = dbmod.run_select(sql, wide, max_rows=5)
    except Exception as e:
        return _err(f"SQLが実データで通りませんでした: {str(e).splitlines()[0][:120]} "
                    "実際に成功したSQLをそのまま渡してください。")
    total = _total_rows(sql, wide) if truncated else len(rows)

    # 置き場のDB = SQLが最初に名指ししているDB（例文はDBごとのファイルに残る）
    hits = []
    import re as _re
    for f in dbmod.list_db_files():
        m = _re.search(r'(?<![\w."])' + _re.escape(dbmod.alias_for(f)) + r"\s*\.",
                       sql, _re.IGNORECASE)
        if m:
            hits.append((m.start(), f))
    hits.sort(key=lambda t: t[0])
    home = hits[0][1] if hits else (Path(scope[0]["path"]) if len(scope) == 1 else None)
    if home is None:
        return _err("このSQLがどのDBのものか判断できませんでした。"
                    "テーブル名を『DB名.テーブル名』の形で書いたSQLを渡してください。")

    same = catalog.find_example(catalog.load_meta(home).get("examples") or [], sql)
    return {
        "ok": True,
        "llm_content": _json({
            "status": "proposed", "db": home.name, "question": q,
            "already_exists": same is not None,
            "note": "登録カードをユーザーの画面に出しました。登録するかはユーザーが"
                    "カードのボタンで決めます。あなたはこれ以上の操作をしなくてよい。",
        }),
        "render": {"role": "assistant", "kind": "example_proposal",
                   "db": home.name, "question": q, "sql": sql, "summary": summary,
                   "columns": columns, "rows": [list(r) for r in rows],
                   "total": total, "exists": same is not None,
                   "old_q": (same or {}).get("q", "")},
    }


def _show_er_diagram(args: dict, scope: list[dict]) -> dict:
    """ER図をチャットに出す（読み取り専用）。

    描くのはカタログ画面と同じキャンバス（er.js）。データも同じ er_payload なので、
    画面で見える図とAIの理解（結合定義）は必ず一致する。
    """
    import db as dbmod

    files = dbmod.list_db_files()
    if not files:
        return _err("data/ にDBがありません。")
    name = str(args.get("db") or "").strip()
    target = None
    if name:
        low = name.lower()
        target = next((f for f in files
                       if low in (f.name.lower(), f.stem.lower(),
                                  dbmod.alias_for(f).lower())), None)
        if target is None:
            return _err(f"DB '{name}' が見つかりません。指定できるのは: "
                        + "、".join(f.name for f in files))
    elif len(scope) == 1:
        target = Path(scope[0]["path"])
    else:
        return _err("どのDBのER図かを db で指定してください。候補: "
                    + "、".join(s_["name"] for s_ in scope))

    try:
        payload = catalog.er_payload(target)
    except Exception as e:
        return _err(f"ER図データの組み立てに失敗しました: {e}")

    own = [n for n in payload["nodes"] if not n.get("external")]
    rels = [{"from": ".".join(str(x) for x in e["from"]),
             "to": ".".join(str(x) for x in e["to"]),
             "cardinality": e.get("cardinality") or "",
             "kind": "FOREIGN KEY宣言" if e.get("kind") == "fk" else "カタログ登録"}
            for e in payload["edges"]]
    return {
        "ok": True,
        "llm_content": _json({
            "status": "er_ready", "db": target.name,
            "tables": [n["table"] for n in own],
            "borrowed_tables": payload.get("extra") or [],
            "relationships": rels[:60],
            "note": "ER図はユーザーの画面に表示済み。関係を文章で説明し直す必要はない。"
                    "結合の一覧は上の relationships のとおり。",
        }),
        "render": {"role": "assistant", "kind": "er", "db": target.name,
                   "title": f"{target.name} のER図", "er": payload},
    }


def _export_excel(args: dict, scope: list[dict]) -> dict:
    sheets_in = args.get("sheets") or []
    if not sheets_in:
        return _err("sheets が空です。少なくとも1つ SELECT を指定してください。")

    built, summary = [], []
    for i, sh in enumerate(sheets_in, start=1):
        name = (sh or {}).get("name") or f"Sheet{i}"
        # ファイルには全行入れる（画面向けの2,000行とは別枠）。
        # Excelのシートは仕様上 1,048,576 行までなので、見出しぶんを引いて丸める。
        cap = min(config.EXPORT_MAX_ROWS, 1_048_575)
        try:
            columns, rows, truncated, _, _ = fetch(sh or {}, scope, label=name,
                                                   max_rows=cap)
        except advanced.AnalysisError as e:
            return _err(f"シート '{name}': {e}")
        except Exception as e:
            return _err(f"シート '{name}' のSQL実行エラー: {e}")
        note = (sh or {}).get("note") or ""
        if truncated:
            note = (note + f"（{cap:,}行で切り詰め）").strip()
        charts = (sh or {}).get("charts") or (sh or {}).get("chart")
        if isinstance(charts, dict):
            charts = [charts]
        built.append({"name": name, "columns": columns, "rows": rows, "note": note,
                      "charts": charts or []})
        summary.append({"sheet": name, "columns": columns, "row_count": len(rows),
                        "truncated": truncated,
                        "charts": [c.get("type") for c in (charts or [])]})

    try:
        data = excel.build_excel(built, title=args.get("purpose") or args.get("filename"))
    except ValueError as e:
        return _err(f"Excelのグラフを作れませんでした: {e}")
    except Exception as e:
        return _err(f"Excelの作成に失敗しました: {e}")

    filename = exports.safe_filename(args.get("filename"), "xlsx")
    return {
        "ok": True,
        "llm_content": _json({
            "status": "file_ready",
            "filename": filename,
            "sheets": summary,
            "note": "ユーザーの画面に保存済み。ファイルの中身を再度説明する必要はない。",
        }),
        "render": {
            "role": "assistant", "kind": "file", "filename": filename,
            "mime": exports.XLSX_MIME, "data": data, "sheets": built,
            "note": f"{len(built)}シート",
        },
    }


def _export_csv(args: dict, scope: list[dict]) -> dict:
    files_in = args.get("files") or []
    if not files_in:
        return _err("files が空です。少なくとも1つ SELECT を指定してください。")
    enc = args.get("encoding") or exports.DEFAULT_ENCODING
    delim = args.get("delimiter") or "comma"

    made, summary, preview = [], [], []
    for i, f in enumerate(files_in, start=1):
        name = (f or {}).get("name") or f"data{i}"
        try:
            columns, rows, truncated, _, _ = fetch(f or {}, scope, label=name,
                                                   max_rows=config.EXPORT_MAX_ROWS)
        except advanced.AnalysisError as e:
            return _err(f"'{name}': {e}")
        except Exception as e:
            return _err(f"'{name}' のSQL実行エラー: {e}")
        try:
            data = exports.build_csv(columns, rows, enc, delim)
        except Exception as e:
            return _err(f"'{name}' のCSV作成に失敗しました（文字コード {enc}）: {e}")
        made.append({"filename": exports.safe_filename(name, "csv"), "data": data})
        summary.append({"file": name, "columns": columns, "row_count": len(rows),
                        "truncated": truncated})
        preview.append({"name": name, "columns": columns, "rows": rows})

    if len(made) == 1:
        filename, data, mime = made[0]["filename"], made[0]["data"], exports.CSV_MIME
    else:
        filename = exports.safe_filename(args.get("purpose") or "csv_files", "zip")
        data, mime = exports.build_zip(made), exports.ZIP_MIME

    return {
        "ok": True,
        "llm_content": _json({
            "status": "file_ready", "filename": filename,
            "encoding": enc, "delimiter": delim, "files": summary,
            "note": "ユーザーの画面に保存済み。中身を再度全部説明する必要はない。",
        }),
        "render": {
            "role": "assistant", "kind": "file", "filename": filename,
            "mime": mime, "data": data, "sheets": preview,
            "note": f"文字コード {enc} / 区切り {delim}",
        },
    }


def _export_text(args: dict, scope: list[dict]) -> dict:
    body = str(args.get("body") or "")
    fmt = args.get("format") or "md"
    enc = args.get("encoding") or exports.DEFAULT_ENCODING
    style = "markdown" if fmt == "md" else "plain"

    summary, preview = [], []
    for sec in (args.get("sections") or []):
        heading = str((sec or {}).get("heading") or "")
        try:
            columns, rows, truncated, _, _ = fetch(sec or {}, scope, label=heading,
                                                   max_rows=config.EXPORT_MAX_ROWS)
        except advanced.AnalysisError as e:
            return _err(f"セクション '{heading}': {e}")
        except Exception as e:
            return _err(f"セクション '{heading}' のSQL実行エラー: {e}")
        table = exports.table_to_text(columns, rows, style)
        block = (f"## {heading}\n\n{table}\n" if fmt == "md"
                 else f"■ {heading}\n\n{table}\n")
        placeholder = "{{" + heading + "}}"
        if placeholder in body:
            body = body.replace(placeholder, block)
        else:
            body = body.rstrip() + "\n\n" + block
        summary.append({"heading": heading, "columns": columns, "row_count": len(rows),
                        "truncated": truncated})
        preview.append({"name": heading, "columns": columns, "rows": rows})

    try:
        data = exports.build_text(body, enc)
    except Exception as e:
        return _err(f"テキストの書き出しに失敗しました（文字コード {enc}）: {e}")

    filename = exports.safe_filename(args.get("filename"), fmt)
    return {
        "ok": True,
        "llm_content": _json({
            "status": "file_ready", "filename": filename,
            "format": fmt, "encoding": enc, "sections": summary,
            "chars": len(body),
            "note": "ユーザーの画面に保存済み。本文を再度全部繰り返す必要はない。",
        }),
        "render": {
            "role": "assistant", "kind": "file", "filename": filename,
            "mime": exports.MD_MIME if fmt == "md" else exports.TEXT_MIME,
            "data": data, "text": body, "sheets": preview,
            "note": f"{fmt} / 文字コード {enc} / {len(body):,} 文字",
        },
    }

def _open_table(args: dict, scope: list[dict]) -> dict:
    """テーブルの中身（全行）を別タブのビューアで開く。

    「テーブルを見せて」に対して、AIが SELECT * を打って先頭数行を貼るのは
    (1) 行数の上限で切れる (2) 会話が表で埋まる、の2つで具合が悪い。
    全行を辿れる画面（/table）を開いて、AIは要約に専念する。
    """
    import db as dbmod

    files = dbmod.list_db_files()
    if not files:
        return _err("data/ にDBがありません。")
    name = str(args.get("db") or "").strip()
    table = str(args.get("table") or "").strip()
    if not table:
        return _err("どのテーブルを開くかを table で指定してください。")

    # DB名は 'sales.db' でも 'sales' でもよい。省略時は表名から探す
    target = None
    if name:
        low = name.lower()
        target = next((f for f in files
                       if low in (f.name.lower(), f.stem.lower(), dbmod.alias_for(f).lower())), None)
        if target is None:
            return _err(f"DB '{name}' が見つかりません。指定できるのは: "
                        + "、".join(f.name for f in files))
    else:
        hits = [f for f in files if table in (catalog.profile_db(f).get("tables") or {})]
        if not hits:
            return _err(f"テーブル '{table}' が見つかりません。db も指定してください。")
        if len(hits) > 1:
            return _err(f"テーブル '{table}' が複数のDBにあります（"
                        + "、".join(f.name for f in hits) + "）。db で指定してください。")
        target = hits[0]

    profile = catalog.profile_db(target)
    info = (profile.get("tables") or {}).get(table)
    if info is None:
        return _err(f"テーブル '{table}' が {target.name} にありません。"
                    "このDBのテーブル: " + "、".join(list((profile.get("tables") or {}).keys())[:20]))

    meta = catalog.load_meta(target)
    tmeta = (meta.get("tables") or {}).get(table) or {}
    cols = [c["name"] for c in (info.get("columns") or [])]
    rows = info.get("row_count")
    return {
        "ok": True,
        "llm_content": _json({
            "status": "table_view_opened", "db": target.name, "table": table,
            "rows": rows, "columns": cols,
            "note": "テーブルの全行を見る画面へのリンク（「テーブル全体を開く」ボタン）を"
                    "利用者のチャットに出した。利用者がそれを押すと別タブで開く。"
                    "中身を SELECT * で貼り直す必要はない。"
                    "何のテーブルか・何に使えるかを1〜2文で補足し、"
                    "「テーブル全体を開く」から見られることを伝えるだけでよい。",
        }),
        "render": {
            "role": "assistant", "kind": "table_link",
            "db": target.name, "table": table,
            "title": f"{table}（{target.name}）",
            "rows": rows, "columns": cols,
            "description": tmeta.get("description") or "",
        },
    }


# このモジュールが受け持つツール
HANDLERS_query = {
    "run_sql_query": _run_sql_query,
    "describe_table": _describe_table,
    "show_er_diagram": _show_er_diagram,
    "open_table": _open_table,
    "propose_glossary_term": _propose_glossary_term,
    "propose_example": _propose_example,
    "pivot_table": _pivot_table,
    "analyze_stats": _analyze_stats,
    "plot_chart": _plot_chart,
    "plot_dual_axis": _plot_dual_axis,
    # 用途別のグラフツールは、中身はどれも同じ組み立てを通る
    **{name: _plot_chart for name in _CHART_TOOLS},
    "export_excel": _export_excel,
    "export_csv": _export_csv,
    "export_text": _export_text,
}

# SQLを受け取るツール（実行前プレビュー表示の対象）
SQL_TOOLS_query = {"run_sql_query", "plot_chart", "plot_dual_axis", "pivot_table",
             "analyze_stats", *_CHART_TOOLS}


# ==========================================================================
# ===== 元 tools/reports.py
# レポート出力（PowerPoint / Word / 画面用のレポート）。
# ==========================================================================
import re
from datetime import datetime

import advanced
import charts
import docx_report
import excel
import exports
import figures
import pptx_report


def _slide_from_sql(spec: dict, scope: list[dict], index: int) -> dict:
    """slides の1枚ぶん。sql か result_id があればここで中身に変える。"""
    out = dict(spec)
    kind = str(spec.get("kind") or "message").lower()
    if not (spec.get("sql") or spec.get("result_id")) or kind not in ("table", "chart"):
        return out

    columns, rows, truncated, _, _ = fetch(spec, scope, label=spec.get("title"))
    if not rows:
        raise pptx_report.PptxReportError(f"{index}枚目「{spec.get('title', '')}」の"
                                      "SQLが0行でした。抽出条件を見直してください。")
    if kind == "table":
        out["columns"], out["rows"] = columns, [list(r) for r in rows]
        if truncated:
            out["comment"] = (out.get("comment", "") + "　※ 上限で切り詰め済み").strip()
        return out

    cat = spec.get("category_column") or columns[0]
    if cat not in columns:
        raise pptx_report.PptxReportError(
            f"{index}枚目: 横軸の列 '{cat}' がSQLの結果にありません"
            f"（ある列: {', '.join(columns)}）。")
    vals = spec.get("value_columns") or [c for c in columns if c != cat]
    missing = [v for v in vals if v not in columns]
    if missing:
        raise pptx_report.PptxReportError(
            f"{index}枚目: 系列の列 {', '.join(missing)} がSQLの結果にありません"
            f"（ある列: {', '.join(columns)}）。")
    ci = columns.index(cat)
    out["categories"] = [r[ci] for r in rows]
    if str(spec.get("chart") or "bar").lower() == "scatter":
        xi = columns.index(vals[0])
        out["series"] = [{"name": vals[1] if len(vals) > 1 else vals[0],
                          "x": [r[xi] for r in rows],
                          "values": [r[columns.index(vals[1] if len(vals) > 1 else vals[0])]
                                     for r in rows]}]
    else:
        out["series"] = [{"name": v, "values": [r[columns.index(v)] for r in rows]}
                         for v in vals]
    return out


def _export_pptx(args: dict, scope: list[dict]) -> dict:
    slides_in = args.get("slides") or []
    if not slides_in:
        return _err("slides が空です。少なくとも1枚は指定してください。")
    built = []
    for i, spec in enumerate(slides_in, start=1):
        try:
            built.append(_slide_from_sql(spec or {}, scope, i))
        except (pptx_report.PptxReportError, advanced.AnalysisError) as e:
            return _err(f"{i}枚目: {e}" if isinstance(e, advanced.AnalysisError) else str(e))
        except Exception as e:
            return _err(f"{i}枚目のSQL実行エラー: {e}")
    try:
        data = pptx_report.build_pptx(built, title=args.get("title"),
                                 subtitle=args.get("subtitle"),
                                 footer=args.get("footer"))
    except pptx_report.PptxReportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"PowerPointの作成に失敗しました: {e}")

    filename = pptx_report.pptx_safe_filename(args.get("filename") or args.get("title"))
    outline = pptx_report.outline_pptx(built)
    return {
        "ok": True,
        "llm_content": _json({
            "status": "file_ready", "filename": filename,
            "slides": outline,
            "note": "ユーザーの画面に保存済み。中身を再度説明する必要はない。",
        }),
        "render": {"role": "assistant", "kind": "file", "filename": filename,
                   "mime": PPTX_MIME, "data": data,
                   "note": f"{len(built)}スライド", "outline": outline},
    }


PPTX_MIME = ("application/vnd.openxmlformats-officedocument."
             "presentationml.presentation")


DOCX_MIME = ("application/vnd.openxmlformats-officedocument."
             "wordprocessingml.document")


def _export_docx(args: dict, scope: list[dict]) -> dict:
    """Word文書。図表つきで、そのまま配布できる体裁にする。"""
    secs_in = args.get("sections") or []
    if not secs_in:
        return _err("sections が空です。少なくとも1つの見出しを入れてください。")

    sections, figs, tbls = [], 0, 0
    for i, s in enumerate(secs_in, start=1):
        s = s or {}
        if not s.get("heading"):
            return _err(f"{i}番目のセクションに heading がありません。")
        sec = {k: s.get(k) for k in ("heading", "body", "bullets", "note",
                                     "callout", "level", "page_break")}
        # 図表のキャプションは、見出し頭の「1. 」を落として重複を避ける
        label = re.sub(r"^\s*\d+[.．)、]\s*", "", s["heading"])
        if s.get("sql") or s.get("result_id"):
            try:
                columns, rows, truncated, _, _ = fetch(s, scope, label=s["heading"])
            except advanced.AnalysisError as e:
                return _err(f"「{s['heading']}」: {e}")
            except Exception as e:
                return _err(f"「{s['heading']}」のSQL実行エラー: {e}")
            if not rows:
                return _err(f"「{s['heading']}」のデータが0行でした。")
            limit = int(s.get("max_rows") or 40)
            if s.get("chart"):
                chart = {**(s["chart"] or {}), "columns": columns,
                         "rows": [list(r) for r in rows],
                         "title": (s["chart"] or {}).get("title") or s["heading"]}
                chart.setdefault("chart_type", "bar")
                errs = charts.validate(chart, columns)
                if errs:
                    return _err(f"「{s['heading']}」のグラフ指定: {' / '.join(errs)}")
                img = _chart_image(chart)
                if img:
                    sec["image"] = img
                    sec["caption"] = s.get("caption") or label
                    figs += 1
            if s.get("table", True):        # 既定で表も載せる（根拠として残す）
                sec["table"] = {"columns": columns,
                                "rows": [list(r) for r in rows[:limit]]}
                if len(rows) > limit or truncated:
                    sec["table"]["note"] = f"全 {len(rows):,} 行から抜粋"
                sec["table_caption"] = s.get("table_caption") or label
                tbls += 1
        sections.append(sec)

    try:
        data = docx_report.build_docx(
            sections, title=args.get("title", "レポート"),
            subtitle=args.get("subtitle", ""),
            summary=args.get("summary") or [],
            conclusion=args.get("conclusion", ""),
            recommendations=args.get("recommendations") or [],
            caveats=args.get("caveats") or [],
            footer=args.get("footer", ""), org=args.get("org", ""),
            author=args.get("author", ""), toc=bool(args.get("toc", True)))
    except docx_report.DocxReportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"Word文書の作成に失敗しました: {e}")

    filename = docx_report.docx_safe_filename(args.get("filename") or args.get("title"))
    note = ""
    if figs == 0 and any(s.get("chart") for s in secs_in):
        note = figures.why_unavailable()
    return {
        "ok": True,
        "llm_content": _json({
            "status": "file_ready", "filename": filename,
            "sections": docx_report.outline_docx(sections),
            "figures": figs, "tables": tbls,
            "warning": note or None,
            "note": "ユーザーの画面に保存済み。中身を再度説明する必要はない。",
        }),
        "render": {"role": "assistant", "kind": "file", "filename": filename,
                   "mime": DOCX_MIME, "data": data,
                   "note": f"{len(sections)}セクション / 図{figs} 表{tbls}"
                           + (f" ／ {note}" if note else ""),
                   "outline": docx_report.outline_docx(sections)},
    }


def _report_markdown(args: dict, sections: list[dict]) -> str:
    """レポートを1本の Markdown にする（ダウンロード用と、表示の下敷き）。"""
    out = [f"# {args.get('title', 'レポート')}"]
    if args.get("subtitle"):
        out.append(f"*{args['subtitle']}*")
    out.append(f"*作成: {datetime.now():%Y-%m-%d %H:%M}*")
    if args.get("summary"):
        out += ["", "## 要点"] + [f"- {s}" for s in args["summary"]]
    for i, s in enumerate(sections, 1):
        out += ["", f"## {i}. {s['heading']}"]
        if s.get("body"):
            out += ["", s["body"]]
        if s.get("table"):
            t = s["table"]
            out += ["", exports.table_to_text(t["columns"], t["rows"], "markdown")]
            if t.get("truncated"):
                out.append(f"（全 {t['total']:,} 行のうち上位 {len(t['rows'])} 行）")
        if s.get("chart"):
            out.append(f"（グラフ: {s['chart'].get('chart_type')} — 画面で確認できます）")
        if s.get("note"):
            out += ["", f"> {s['note']}"]
    if args.get("conclusion"):
        out += ["", "## 結論", "", args["conclusion"]]
    if args.get("recommendations"):
        out += ["", "## 推奨する打ち手"] + [f"{i}. {r}" for i, r
                                            in enumerate(args["recommendations"], 1)]
    if args.get("caveats"):
        out += ["", "## 前提・注意"] + [f"- {c}" for c in args["caveats"]]
    return "\n".join(out) + "\n"


# PowerPointのネイティブグラフは種類が限られる。近いもので描けるならそれを使い、
# 描けない種類（サンキー・箱ひげ等）は画像として貼る。
_PPTX_CHART_MAP = {
    "bar": "bar", "hbar": "hbar", "stacked_bar": "bar_stacked",
    "percent_bar": "bar_percent", "lollipop": "bar", "pareto": "bar",
    "line": "line", "step": "line", "bump": "line", "slope": "line",
    "area": "area", "area_percent": "area_stacked",
    "pie": "pie", "donut": "doughnut", "funnel": "hbar", "radar": "radar",
    "scatter": "scatter", "bubble": "scatter", "polar_bar": "radar",
}


# Excelのグラフも種類が限られる。近いものに寄せ、無理なものは表だけにする。
_XLSX_CHART_MAP = {
    "bar": "bar", "hbar": "hbar", "stacked_bar": "bar_stacked",
    "percent_bar": "bar_percent", "lollipop": "bar", "pareto": "bar",
    "line": "line", "step": "line", "bump": "line", "slope": "line",
    "area": "area", "area_percent": "area_stacked",
    "pie": "pie", "donut": "pie", "funnel": "hbar",
    "scatter": "scatter", "bubble": "scatter",
}


def _chart_image(chart: dict):
    """グラフを印刷向けの画像にする。できなければ None。"""
    try:
        return figures.for_print(charts.build_figure(chart))
    except Exception as e:
        print(f"[report] 画像化に失敗: {e}")
        return None


def _series_from(chart: dict, table: dict):
    """表からPowerPointのネイティブグラフ用の系列を組み立てる。"""
    cols = table["columns"]
    cat = chart.get("x") if chart.get("x") in cols else cols[0]
    ci = cols.index(cat)
    if chart.get("chart_type") in ("scatter", "bubble"):
        xs = [r[cols.index(chart["x"])] for r in table["rows"]]
        ys = [r[cols.index(chart["y"])] for r in table["rows"]]
        return [], [{"name": chart.get("y", "値"), "x": xs, "values": ys}]
    wanted = [chart["y"]] if chart.get("y") in cols else \
        [c for c in cols if c != cat][:3]
    return ([r[ci] for r in table["rows"]],
            [{"name": v, "values": [r[cols.index(v)] for r in table["rows"]]}
             for v in wanted])


def _split_message(s: dict) -> tuple[str, str]:
    """1ページの「言いたいこと1行」と、横に添える残りの文章に分ける。

    所見(note)があればそれを1行目にする。無ければ本文の最初の一文を使い、
    その場合は本文の残りだけを横に置く（同じ文を2回出さない）。
    """
    body = (s.get("body") or "").strip()
    if s.get("note"):
        return s["note"], body
    if not body:
        return "", ""
    head, sep, rest = body.partition("。")
    if not sep:
        return body[:90], ""
    return head + "。", rest.strip()


def _report_slides(args: dict, sections: list[dict]) -> list[dict]:
    """レポートの内容を、会議で映せるスライドの並びに翻訳する。"""
    summary = args.get("summary") or []
    slides = [{"kind": "title", "title": args.get("title", "レポート"),
               "subtitle": args.get("subtitle", ""),
               "lines": summary[:4], "org": args.get("org", "")}]
    if len(sections) >= 2:
        slides.append({"kind": "agenda", "title": "本日の内容",
                       "items": [s["heading"] for s in sections]})
    if summary:
        # 帯に出した1点目は繰り返さない（同じ文が2回出ると雑に見える）
        slides.append({"kind": "message", "title": "要約",
                       "message": summary[0],
                       "bullets": summary[1:] or summary,
                       "callout": args.get("conclusion", "")})

    for s in sections:
        message, comment = _split_message(s)
        base = {"title": s["heading"],
                # 見出しの下に置く1行。結論を先に言う。
                "message": message, "comment": comment,
                "notes": s.get("body") or "",
                "source": args.get("source", "")}
        chart, table = s.get("chart"), s.get("table")
        if chart and table:
            kind = chart.get("chart_type")
            native = _PPTX_CHART_MAP.get(kind)
            if native:
                cats, series = _series_from(chart, table)
                slides.append({**base, "kind": "chart", "chart": native,
                               "categories": cats, "series": series})
            else:
                img = _chart_image(chart)
                if img:
                    slides.append({**base, "kind": "chart", "image": img})
                else:
                    slides.append({**base, "kind": "table", **table})
        elif table:
            slides.append({**base, "kind": "table", **table})
        else:
            slides.append({**base, "kind": "message", "comment": None,
                           "lead": base.pop("comment", "") or "",
                           "bullets": s.get("bullets") or []})

    if args.get("conclusion") or args.get("recommendations"):
        slides.append({"kind": "closing", "title": "まとめと次のアクション",
                       "message": args.get("conclusion", ""),
                       "summary": (args.get("summary") or [])[:3],
                       "actions": args.get("recommendations") or []})
    return slides


def _report_docx_sections(args: dict, sections: list[dict]) -> list[dict]:
    """レポートの内容を、Wordのセクションに翻訳する（図は画像で貼る）。"""
    out = []
    for s in sections:
        label = re.sub(r"^\s*\d+[.．)、]\s*", "", s["heading"])
        sec = {"heading": s["heading"], "body": s.get("body", ""),
               "note": s.get("note", ""), "caption": label}
        if s.get("chart"):
            img = _chart_image(s["chart"])
            if img:
                sec["image"] = img
        if s.get("table"):
            sec["table"] = {"columns": s["table"]["columns"],
                            "rows": s["table"]["rows"]}
            sec["table_caption"] = label
            if s["table"].get("truncated"):
                sec["table"]["note"] = f"全 {s['table']['total']:,} 行から抜粋"
        out.append(sec)
    return out


def _build_report(args: dict, scope: list[dict]) -> dict:
    sections_in = args.get("sections") or []
    if not sections_in:
        return _err("sections が空です。少なくとも1つの論点を入れてください。")

    sections, dropped = [], []
    for i, s in enumerate(sections_in, 1):
        s = s or {}
        if not s.get("heading"):
            return _err(f"{i}番目のセクションに heading がありません。")
        out = {"heading": s["heading"], "body": s.get("body", ""),
               "note": s.get("note", "")}
        if s.get("sql") or s.get("result_id"):
            try:
                columns, rows, truncated, _, _ = fetch(s, scope, label=s["heading"])
            except advanced.AnalysisError as e:
                return _err(f"「{s['heading']}」: {e}")
            except Exception as e:
                return _err(f"「{s['heading']}」のSQL実行エラー: {e}")
            limit = int(s.get("max_rows") or 20)
            out["table"] = {"columns": columns, "rows": [list(r) for r in rows[:limit]],
                            "total": len(rows), "truncated": len(rows) > limit or truncated}
            out["sql"] = s.get("sql")
            if s.get("chart"):
                chart = {k: v for k, v in (s["chart"] or {}).items()}
                chart.setdefault("chart_type", "bar")
                errs = charts.validate(chart, columns)
                if errs:
                    # グラフ1つのためにレポート全体を捨てない。
                    # 表は根拠として残るので、図を落として作り切る方が役に立つ。
                    dropped.append(f"「{s['heading']}」のグラフ: {' / '.join(errs)}")
                else:
                    # グラフは全行を使う（表は読みやすさのために切っている）
                    out["chart"] = {**chart, "columns": columns,
                                    "rows": [list(r) for r in rows],
                                    "title": chart.get("title") or s["heading"]}
        elif s.get("chart"):
            dropped.append(f"「{s['heading']}」のグラフ: 元になる sql / result_id がありません。")
        sections.append(out)

    md = _report_markdown(args, sections)
    fmt = (args.get("format") or "md").lower()
    name = args.get("filename") or args.get("title") or "report"
    data = filename = mime = None
    try:
        if fmt == "md":
            data = exports.build_text(md)
            filename = exports.safe_filename(name, "md")
            mime = exports.TEXT_MIME
        elif fmt == "xlsx":
            sheets = []
            for s in sections:
                if not s.get("table"):
                    continue
                sheet = {"name": s["heading"], "columns": s["table"]["columns"],
                         "rows": s["table"]["rows"], "note": s.get("note", "")}
                # 画面のグラフ指定を、そのままExcelのグラフに読み替える
                if s.get("chart"):
                    ch = s["chart"]
                    kind = _XLSX_CHART_MAP.get(ch.get("chart_type"))
                    if kind:
                        cat = ch.get("x") if ch.get("x") in sheet["columns"] \
                            else sheet["columns"][0]
                        vals = ([ch["y"]] if ch.get("y") in sheet["columns"]
                                else [c for c in sheet["columns"] if c != cat][:3])
                        sheet["charts"] = [{"type": kind, "category_column": cat,
                                            "value_columns": vals,
                                            "title": s["heading"]}]
                sheets.append(sheet)
            if not sheets:
                return _err("xlsx にするには、表（sql）のあるセクションが1つ以上必要です。")
            data = excel.build_excel(sheets, title=args.get("title"))
            filename = exports.safe_filename(name, "xlsx")
            mime = exports.XLSX_MIME
        elif fmt == "pptx":
            data = pptx_report.build_pptx(_report_slides(args, sections),
                                     title=args.get("title"),
                                     subtitle=args.get("subtitle"),
                                     footer=args.get("footer", ""))
            filename = pptx_report.pptx_safe_filename(name)
            mime = PPTX_MIME
        elif fmt == "docx":
            data = docx_report.build_docx(
                _report_docx_sections(args, sections),
                title=args.get("title", "レポート"),
                subtitle=args.get("subtitle", ""),
                summary=args.get("summary") or [],
                conclusion=args.get("conclusion", ""),
                recommendations=args.get("recommendations") or [],
                caveats=args.get("caveats") or [],
                footer=args.get("footer", ""), org=args.get("org", ""))
            filename = docx_report.docx_safe_filename(name)
            mime = DOCX_MIME
    except Exception as e:
        return _err(f"ファイルの作成に失敗しました: {e}")

    render = {"role": "assistant", "kind": "report_doc",
              "title": args.get("title", "レポート"),
              "subtitle": args.get("subtitle", ""),
              "summary": args.get("summary") or [],
              "sections": sections,
              "conclusion": args.get("conclusion", ""),
              "recommendations": args.get("recommendations") or [],
              "caveats": args.get("caveats") or [],
              "markdown": md}
    if data:
        render.update(data=data, filename=filename, mime=mime)

    return {
        "ok": True,
        "llm_content": _json({
            "status": "report_ready", "title": args.get("title"),
            "sections": [{"heading": s["heading"],
                          "rows": (s.get("table") or {}).get("total"),
                          "chart": (s.get("chart") or {}).get("chart_type")}
                         for s in sections],
            "filename": filename,
            "dropped_charts": dropped or None,
            "note": "レポートは画面に表示済み。内容をもう一度書き出す必要はない。"
                    "次に何をするか（送付・追加分析など）だけ短く伝えること。"
                    + ("　※ 一部のグラフは指定が合わず省いた。作り直すなら、"
                       "dropped_charts の指摘どおりに列名を直して呼ぶこと"
                       "（同じ引数で呼び直さない）。" if dropped else ""),
        }),
        "render": render,
    }

HANDLERS_reports = {
    "export_pptx": _export_pptx,
    "export_docx": _export_docx,
    "build_report": _build_report,
}

SQL_TOOLS_reports: set[str] = set()


# ==========================================================================
# ===== 元 tools/stats.py
# 統計と試算。advanced.py の分析関数をツールとして公開する。
# ==========================================================================
import advanced


_hypothesis_test = _analysis_tool(lambda a, c, r: advanced.hypothesis_test(
    c, r, a.get("method"), value_col=a.get("value_col"), group_col=a.get("group_col"),
    value_col2=a.get("value_col2"), popmean=float(a.get("popmean") or 0),
    expected=a.get("expected"), alternative=a.get("alternative") or "two-sided"))


_regression = _analysis_tool(lambda a, c, r: advanced.regression(
    c, r, a.get("target"), a.get("features") or [], method=a.get("method") or "ols",
    predict=a.get("predict")))


_distribution_analysis = _analysis_tool(lambda a, c, r: advanced.distribution(
    c, r, a.get("target"), bins=int(a.get("bins") or 20), fit=a.get("fit"),
    group_col=a.get("group_col")))


_forecast = _analysis_tool(lambda a, c, r: advanced.forecast(
    c, r, a.get("time_col"), a.get("value_col"), periods=int(a.get("periods") or 6),
    method=a.get("method") or "auto",
    season_length=int(a["season_length"]) if a.get("season_length") else None,
    exog=a.get("exog")))


_timeseries_analysis = _analysis_tool(lambda a, c, r: advanced.timeseries(
    c, r, a.get("time_col"), a.get("value_col"), window=int(a.get("window") or 3),
    season_length=int(a["season_length"]) if a.get("season_length") else None))


_bootstrap_estimate = _analysis_tool(lambda a, c, r: advanced.bootstrap(
    c, r, a.get("target"), statistic=a.get("statistic") or "mean",
    trials=int(a.get("trials") or 5000), group_col=a.get("group_col")))


# k は "auto" も受けるので、ここで数値に変換しない
_clustering = _analysis_tool(lambda a, c, r: advanced.clustering(
    c, r, a.get("features") or [], k=a.get("k") or 3,
    label_col=a.get("label_col"), categorical=a.get("categorical")))


_abc_analysis = _analysis_tool(lambda a, c, r: advanced.abc_analysis(
    c, r, a.get("label_col"), a.get("value_col"), thresholds=a.get("thresholds")))


def _monte_carlo_simulation(args: dict, scope: list[dict]) -> dict:
    columns = rows = None
    if args.get("sql"):
        try:
            columns, rows, _ = _select_for(args, scope)
        except advanced.AnalysisError as e:
            return _err(str(e))
        except Exception as e:
            return _err(f"実データの取得に失敗しました: {e}")
    try:
        res = advanced.monte_carlo(
            args.get("formula", ""), args.get("variables") or {},
            trials=int(args.get("trials") or 10000), columns=columns, rows=rows,
            targets=args.get("targets"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"シミュレーションに失敗しました: {e}")
    if args.get("title"):
        res["title"] = args["title"]
    return _report_result(res)


def _scenario_analysis(args: dict, scope: list[dict]) -> dict:
    try:
        res = advanced.scenario(args.get("formula", ""), args.get("scenarios") or {},
                                base=args.get("base"))
    except advanced.AnalysisError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"シナリオ分析に失敗しました: {e}")
    if args.get("title"):
        res["title"] = args["title"]
    return _report_result(res)

HANDLERS_stats = {
    "hypothesis_test": _hypothesis_test,
    "regression": _regression,
    "distribution_analysis": _distribution_analysis,
    "forecast": _forecast,
    "timeseries_analysis": _timeseries_analysis,
    "monte_carlo_simulation": _monte_carlo_simulation,
    "scenario_analysis": _scenario_analysis,
    "bootstrap_estimate": _bootstrap_estimate,
    "clustering": _clustering,
    "abc_analysis": _abc_analysis,
}

# scenario_analysis と monte_carlo_simulation は SQL が任意なので含めない
SQL_TOOLS_stats = {"hypothesis_test", "regression", "distribution_analysis", "forecast",
             "timeseries_analysis", "bootstrap_estimate", "clustering", "abc_analysis"}


# ==========================================================================
# ===== 元 tools/usage.py
# このアプリ自身の使われ方を調べるツール（読むだけ）。
#
# usage.py の集計をLLMに公開する。材料はチャット履歴と取り込みの記録なので、
# 分析対象のDBを選んでいなくても答えられる（SQLでは答えられない話でもある）。
#
# 他人の質問文まで見えるため、管理者にだけ渡す。tools/files.py と同じ扱い。
# ==========================================================================
import usage


def _analyze_usage(args: dict, scope: list[dict]) -> dict:
    days = args.get("days")
    try:
        res = usage.analyze(str(args.get("method") or "summary"),
                            days=int(days) if days else None,
                            user=(args.get("user") or "").strip() or None)
    except ValueError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"利用状況の集計に失敗しました: {e}")

    if args.get("title"):
        res["title"] = args["title"]
    # scope を渡して表に result_id を付ける。グラフ化やExcel出力にそのまま繋げられる。
    return _report_result(res, scope=scope)


HANDLERS_usage = {"analyze_usage": _analyze_usage}

# SQLは受け取らない（材料はDBではなく履歴ファイル）
SQL_TOOLS_usage: set = set()

# 他の利用者の質問・失敗まで見えるので管理者だけに渡す。
ADMIN_TOOLS_usage = {"analyze_usage"}


# ==========================================================================
# ===== 元 tools/__init__.py
# LLM(OpenAI互換 function calling)に渡すツール定義と実行ロジック。
#
# 組み込みツールの系統:
#   調べる    run_sql_query / describe_table
#   集計する  pivot_table / analyze_stats
#   描く      plot_comparison / plot_trend / plot_composition / plot_distribution /
#             plot_relationship / plot_kpi / plot_dual_axis
#   統計      hypothesis_test / regression / distribution_analysis
#   時系列    forecast / timeseries_analysis / detect_anomalies
#   試算      monte_carlo_simulation / scenario_analysis / bootstrap_estimate
#   分ける    clustering / abc_analysis
#   業務分析  compare_periods / funnel_analysis / cohort_analysis / market_basket /
#             survival_analysis / data_quality
#   ファイル  explore_import_files（取り込み元フォルダを読むだけ。管理者のみ）
#   自己分析  analyze_usage（このアプリ自身の使われ方。管理者のみ）
#   出す      export_excel / export_csv / export_text / export_pptx / export_docx /
#             build_report
#   送る      find_mail_recipients / compose_email
#             ※ 実際の送信はユーザーが画面のボタンを押したときだけ。
#               誤送信は取り消せないので、LLMには下書きまでしかさせない。
#
# データを取るツールは sql の代わりに result_id を受け取れる。前のツールが返した
# データをそのまま使えるので、同じSQLを何度も流さずに済む（results.py 参照）。
#
# これに加えて、ユーザーが画面から定義したSQLテンプレート型ツール
# （各DBの .meta.yaml の tools:）を実行時に合成する。custom_tools.py 参照。
# 組み込みツールは .meta.yaml の builtin_tools: で無効化・説明の上書きができる。
#
# dispatch(name, arguments_json, scope, entries, admin) の戻り値:
#   {
#     "ok": bool,
#     "llm_content": str,        # LLMへ返すテキスト(JSON)。トークン節約のため要約。
#     "render": dict | None,     # UI描画用アイテム(app側が kind を見て解釈)
#   }
#
# scope は質問ごとに自動判定された対象DB群（web/chat_bp.py の _auto_scope）:
#   [{"path": str, "alias": str, "tables": [...]}, ...]
#
# ファイルの分かれ方:
#   schemas.py  … LLMに見せるツールの宣言（JSON Schema）。処理は書かない
#   common.py   … 実処理が共通で使う小道具（データの取り出しもここ）
#   results.py  … 取ったデータの置き場（result_id で使い回す）
#   query.py    … 調べる・集計する・描く・出す
#   stats.py    … 統計と試算
#   business.py … 期間比較・ファネル・コホート・併売・異常検知・生存時間・品質
#   reports.py  … PowerPoint / Word / 画面用レポート
#   mail.py     … 宛先探しと下書き
#   files.py    … 取り込み元フォルダの調査（管理者のみ）
#   usage.py    … このアプリ自身の利用状況（管理者のみ）
# ツールを1つ足すときは、宣言(schemas)と実処理(各モジュール)の2箇所を触る。
# 実処理を置いたモジュールの HANDLERS に名前を登録すれば dispatch から引ける。
# 管理者だけに渡したいものは、そのモジュールの ADMIN_TOOLS にも名前を入れる。
# ==========================================================================
import json

import config
import custom_tools
import db
import excel
import exports
import verify

# (実処理, SQLを受け取るもの, 管理者専用) をモジュールごとに並べる。
# 並び順は統合前の (query, stats, reports, mail, business, files, usage) のまま。
# 名前が重なったときにどちらが残るかを変えないため、順序は動かさないこと。
_MODULES = (
    (HANDLERS_query, SQL_TOOLS_query, ()),
    (HANDLERS_stats, SQL_TOOLS_stats, ()),
    (HANDLERS_reports, SQL_TOOLS_reports, ()),
    (HANDLERS_mail, SQL_TOOLS_mail, ()),
    (HANDLERS_business, SQL_TOOLS_business, ()),
    (HANDLERS_files, SQL_TOOLS_files, ADMIN_TOOLS_files),
    (HANDLERS_usage, SQL_TOOLS_usage, ADMIN_TOOLS_usage),
)

# ツール名 -> 実処理。各モジュールが自分のぶんを申告する。
_HANDLERS = {name: fn for m in _MODULES for name, fn in m[0].items()}

# SQLを受け取る組み込みツール（実行前プレビュー表示の対象）
SQL_TOOLS = {name for m in _MODULES for name in m[1]}

# 管理者にだけ渡すツール。画面側で管理者専用になっているものは、
# AI経由でも同じ制限にしないと抜け道になる（取り込み元フォルダの中身など）。
ADMIN_TOOLS = {name for m in _MODULES for name in m[2]}


def render_sql(tool: dict) -> str:
    """UIのプレビュー用。実行時は :name のままバインドするので置換はしない。"""
    return str(tool.get("sql") or "").strip()


def _run_custom(tool: dict, args: dict, scope: list[dict]) -> dict:
    sql = render_sql(tool)
    try:
        params = custom_tools.coerce_params(tool, args)
    except ValueError as e:
        return _err(str(e))
    # ツールは作るときにDBを意識させないので、SQLが選択外のDBに入ることがある。
    # 必要なぶんは繋いでから実行する（結果を預ける先も同じ範囲にする）。
    scope = db.widen_scope(sql, scope)
    # ファイルに出すツールは全行（Excelはシート上限で丸める）。画面用は2,000行。
    kind_ = tool.get("render") or "table"
    cap = (min(config.EXPORT_MAX_ROWS, 1_048_575) if kind_ in ("excel", "csv")
           else None)
    try:
        columns, rows, truncated = db.run_select(sql, scope, params=params,
                                                 max_rows=cap)
    except Exception as e:
        return _err(f"ツール '{tool.get('name')}' のSQL実行エラー: {e}")

    kind = tool.get("render") or "table"
    chart = tool.get("chart") or {}
    sample = rows[: config.SAMPLE_ROWS_FOR_LLM]

    # 取った表を預けて result_id を返す。組み込みツールは前からこうしているのに
    # ユーザー定義ツールだけ返しておらず、「このツールの結果をグラフにして」と
    # 言われてもAIには渡す手段が無かった（SQLはAIに見せていないので取り直せない）。
    # これがあれば、表で作ったツールでも後からグラフ・Excel・統計に回せる。
    keep = rows[: config.MAX_RESULT_ROWS]
    rid = _results.put(scope, columns, keep, truncated or len(rows) > len(keep),
                       sql=sql, label=f"{tool.get('name')}（ユーザー定義ツール）")
    llm_content = _json({
        "tool": tool.get("name"),
        "columns": columns,
        "row_count": len(rows),
        "truncated": truncated,
        "rows": [list(r) for r in sample],
        "result_id": rid,
        "note": ((f"全{len(rows)}行中 先頭{len(sample)}行を表示。"
                  if len(rows) > len(sample) else "")
                 + f"この結果は result_id '{rid}' で他のツールに渡せます"
                   "（グラフを描く・集計する・統計をかける・"
                   "Excel/CSV/PowerPoint/Word にする、など）。"),
    })

    if kind == "none":
        return {"ok": True, "llm_content": llm_content, "render": None}
    if kind in ("excel", "csv"):
        sheet = {"name": chart.get("title") or tool.get("name") or "Sheet1",
                 "columns": columns, "rows": rows,
                 "note": f"{config.MAX_RESULT_ROWS}行で切り詰め" if truncated else ""}
        base = chart.get("filename") or tool.get("name")
        try:
            if kind == "excel":
                data = excel.build_excel([sheet], title=tool.get("description"))
                filename, mime = exports.safe_filename(base, "xlsx"), exports.XLSX_MIME
            else:
                enc = chart.get("encoding") or exports.DEFAULT_ENCODING
                data = exports.build_csv(columns, rows, enc)
                filename, mime = exports.safe_filename(base, "csv"), exports.CSV_MIME
        except Exception as e:
            return _err(f"ファイルの作成に失敗しました: {e}")
        return {"ok": True, "llm_content": _json({
            "status": "file_ready", "tool": tool.get("name"),
            "filename": filename, "columns": columns, "row_count": len(rows),
            "note": "ユーザーの画面に保存済み。",
        }), "render": {
            "role": "assistant", "kind": "file", "filename": filename,
            "mime": mime, "data": data, "sheets": [sheet],
        }}
    if kind == "chart":
        missing = [c for c in (chart.get("x"), chart.get("y")) if c and c not in columns]
        if missing:
            return _err(f"グラフ用の列が結果にありません: {missing} / 利用可能: {columns}")
        return {"ok": True, "llm_content": llm_content, "render": {
            "role": "assistant", "kind": "chart", "columns": columns, "rows": rows,
            "chart_type": chart.get("chart_type", "bar"),
            "x": chart.get("x"), "y": chart.get("y"), "color": chart.get("color"),
            "barmode": chart.get("barmode"), "title": chart.get("title", "") or tool.get("name"),
        }}
    if kind == "chart_dual":
        bar_y = chart.get("bar_y") or []
        line_y = chart.get("line_y") or []
        needed = [chart.get("x")] + list(bar_y) + list(line_y)
        missing = [c for c in needed if c and c not in columns]
        if missing:
            return _err(f"2軸グラフ用の列が結果にありません: {missing} / 利用可能: {columns}")
        return {"ok": True, "llm_content": llm_content, "render": {
            "role": "assistant", "kind": "chart_dual", "columns": columns, "rows": rows,
            "x": chart.get("x"), "bar_y": bar_y, "line_y": line_y,
            "left_title": chart.get("left_title"), "right_title": chart.get("right_title"),
            "title": chart.get("title", "") or tool.get("name"),
        }}
    return {"ok": True, "llm_content": llm_content, "render": {
        "role": "assistant", "kind": "table",
        "columns": columns, "rows": rows, "truncated": truncated,
    }}


def build_tools(entries: list[dict], admin: bool = False) -> list[dict]:
    """組み込み（無効化・説明上書きを反映）＋ユーザー定義 のツール定義一覧。

    admin=False のときは管理者専用のツールを渡さない。渡さなければ
    AIはその存在を知らないので、呼ばれること自体が起きない。
    """
    ov = custom_tools.builtin_overrides(entries)
    out = []
    for t in BUILTIN_TOOLS:
        name = t["function"]["name"]
        if name in ADMIN_TOOLS and not admin:
            continue
        o = ov.get(name) or {}
        if o.get("enabled") is False:
            continue
        if o.get("description"):
            t = {**t, "function": {**t["function"], "description": o["description"]}}
        out.append(t)
    for tool in custom_tools.collect_everywhere(entries):
        if not custom_tools.validate_custom_tool(tool, set()):     # 壊れた定義はAIに渡さない
            out.append(custom_tools.to_schema(tool))
    return out


def _required_params() -> dict:
    """スキーマで必須になっている引数。{ツール名: (引数名, ...)}"""
    return {t["function"]["name"]: tuple((t["function"].get("parameters") or {})
                                         .get("required") or ())
            for t in BUILTIN_TOOLS}


_REQUIRED = _required_params()


#: スキーマ上は必須でも、実処理が既定値を持っている引数。
#: ここまで止めると、これまで通っていた呼び出しが弾かれてしまう
#: （表題やファイル名は無ければツール側が付ける）。
_HAS_DEFAULT = {"title", "filename", "chart_type", "purpose"}


def _missing_required(name: str, args: dict) -> list[str]:
    """必須なのに渡ってこなかった引数（既定値を持つものは除く）。

    LLMは required を落とすことがある。そのまま実処理へ渡すと、
    pandas の "'[None] not in index'" のような内部エラーになって返る。
    これでは何を直せばよいか分からず、同じ呼び出しを繰り返して打ち切られる。
    ここで止めて、足りない引数の名前をそのまま返す。
    """
    if not isinstance(args, dict):
        return []
    out = []
    for k in _REQUIRED.get(name) or ():
        if k in _HAS_DEFAULT:
            continue
        v = args.get(k)
        # 0 や False は正しい値なので、空とみなすのは None と空の入れ物だけ
        if v is None or (isinstance(v, (str, list, dict, tuple)) and len(v) == 0):
            out.append(k)
    return out


def _string_list_params() -> dict:
    """スキーマ上「文字列の配列」になっている引数。{ツール名: {引数名, ...}}"""
    out: dict = {}
    for t in BUILTIN_TOOLS:
        fn = t["function"]
        props = (fn.get("parameters") or {}).get("properties", {})
        names = {k for k, v in props.items()
                 if v.get("type") == "array"
                 and (v.get("items") or {}).get("type") == "string"}
        if names:
            out[fn["name"]] = names
    return out


_LIST_PARAMS = _string_list_params()


def _coerce_lists(name: str, args: dict) -> dict:
    """配列で受ける引数に文字列が1つ来たら、要素1つの配列として扱う。

    LLMは列名が1つのとき index="地域" のように素の文字列で渡してくることがある。
    そのまま渡すと文字列が1文字ずつに散り、「'地' という列がありません」という
    人には意味の分からないエラーになる（日本語の列名だと必ずこうなる）。
    ここで直せば、13個ある同じ形の引数すべてに効く。
    """
    wanted = _LIST_PARAMS.get(name)
    if not wanted or not isinstance(args, dict):
        return args
    for k in wanted:
        v = args.get(k)
        if isinstance(v, str):
            args[k] = [v.strip()] if v.strip() else []
    return args


def _gather_sqls(node, scope: list[dict], acc: list) -> None:
    """呼び出しの引数から、実行されるSQLを全部拾う。

    レポートの節・Excelのシートのように入れ子の中にも sql がある。
    result_id で前の結果を使い回している場合は、その元のSQLを引く。
    """
    if isinstance(node, dict):
        for k, v in node.items():
            if k == "sql" and isinstance(v, str) and v.strip():
                acc.append(v)
            elif k == "result_id" and isinstance(v, str) and v.strip():
                entry = _results.get(scope, v)
                if entry and entry.get("sql"):
                    acc.append(entry["sql"])
            else:
                _gather_sqls(v, scope, acc)
    elif isinstance(node, list):
        for v in node:
            _gather_sqls(v, scope, acc)


def _attach_verification(res: dict, sqls: list[str], scope: list[dict]) -> dict:
    """実行後の相互検証。触れたテーブルに関係する検算を突き合わせる。

    不一致があれば res["verify_alerts"] に積む。画面とLLMへの出し方は
    呼び出し側（chat側）が決める（同じ警告を会話の中で繰り返さないため）。
    検証自体の失敗で回答を止めない。
    """
    if not res.get("ok") or not sqls:
        return res
    try:
        alerts = verify.alerts_for(sqls, scope)
    except Exception as e:
        print(f"[verify] 検算でエラー（回答は続行）: {e}")
        return res
    if alerts:
        res["verify_alerts"] = alerts
    return res


def dispatch(name: str, arguments_json: str | None, scope: list[dict],
             entries: list[dict] | None = None, admin: bool = False) -> dict:
    try:
        args = json.loads(arguments_json) if arguments_json else {}
    except json.JSONDecodeError as e:
        return _err(f"ツール引数のJSON解析に失敗しました: {e}")

    # 渡していないツールを名指しで呼ばれても実行しない（守りは2箇所で持つ）
    if name in ADMIN_TOOLS and not admin:
        return _err(f"'{name}' は管理者だけが使えます。")

    args = _coerce_lists(name, args)
    missing = _missing_required(name, args)
    if missing:
        return _err(f"'{name}' の必須の引数が指定されていません: {'、'.join(missing)}。"
                    f"（{name} の必須引数は {'、'.join(_REQUIRED[name])}）"
                    "この引数を入れて呼び直してください。列名が分からないときは、"
                    "先に describe_table か run_sql_query で列を確認すること。")

    sqls: list = []
    _gather_sqls(args, scope, sqls)

    handler = _HANDLERS.get(name)
    if handler:
        try:
            return _attach_verification(handler(args, scope), sqls, scope)
        except Exception as e:  # ツールの例外でアプリを落とさない
            return _err(f"ツール '{name}' の実行でエラー: {e}")

    tool = next((t for t in custom_tools.collect_everywhere(entries or []) if t.get("name") == name), None)
    if tool is None:
        return {"ok": False, "llm_content": _json({"error": f"未知のツール: {name}"}), "render": None}
    try:
        sqls.append(render_sql(tool))
        return _attach_verification(_run_custom(tool, args, scope), sqls, scope)
    except Exception as e:
        return _err(f"ツール '{name}' の実行でエラー: {e}")


# ==========================================================================
# ===== 元 llm.py
# OpenAI / OpenAI互換API クライアント / system prompt 生成 / AI下書き。
# ==========================================================================
import json
import re
from datetime import datetime
from pathlib import Path

from openai import OpenAI

import catalog
import config
import custom_tools
import tools

_client: OpenAI | None = None


def is_configured() -> bool:
    return bool(config.OPENAI_BASE_URL and config.OPENAI_API_KEY)


def client() -> OpenAI:
    global _client
    if _client is None:
        _client = OpenAI(
            base_url=config.OPENAI_BASE_URL or None,
            api_key=config.OPENAI_API_KEY or "not-set",
        )
    return _client


# --- モデルごとの作法の違いを吸収する ---------------------------------------------
#
# 同じ OpenAI互換API でも、モデルによって受け付ける引数が違う。実測した例:
#   gpt-5.6-sol : ツールを使うなら reasoning_effort='none' が必須。temperature も既定値のみ
#   gpt-4o 系   : reasoning_effort を送ると「Unrecognized request argument」で拒否
# モデル名の一覧を持って場合分けすると、ゲートウェイや新モデルのたびに保守が要る。
# そこで「1回投げて、断られた理由を読んで直して、覚える」方式にする。
# 400 は推論前に弾かれるので、やり直しても費用はかからない。
#
# 覚えた内容はプロセスの寿命だけ持つ。再起動後の最初の1回だけ余計に往復する。

#: モデル名 -> 学習した調整 {"set": {引数: 値}, "drop": {引数, ...}}
_QUIRKS: dict[str, dict] = {}
#: 1回の呼び出しで引数を直しにいく上限。無限に投げ続けないための歯止め。
_MAX_FIX = 4


def _apply_quirks(kwargs: dict) -> dict:
    q = _QUIRKS.get(str(kwargs.get("model") or ""))
    if not q:
        return kwargs
    out = {k: v for k, v in kwargs.items() if k not in q.get("drop", set())}
    out.update(q.get("set", {}))
    return out


def _learn(model: str, *, set_: dict | None = None, drop: str | None = None) -> None:
    q = _QUIRKS.setdefault(model, {"set": {}, "drop": set()})
    if set_:
        q["set"].update(set_)
    if drop:
        q["drop"].add(drop)
        q["set"].pop(drop, None)
    print(f"[llm] {model} の呼び出し方を調整しました: "
          f"set={q['set']} drop={sorted(q['drop'])}")


def _fix_for(message: str, kwargs: dict) -> tuple | None:
    """エラー文から、次に試す直し方を決める。戻り値: (set_, drop) か None。"""
    low = message.lower()

    # ツールと推論モードを同時に使えない → 推論を切れば chat/completions で通る
    if "reasoning_effort" in low and "function tools" in low:
        return ({"reasoning_effort": "none"}, None)
    # 値が受け付けられない（'minimal' など）→ 'none' に寄せる
    if "reasoning_effort" in low and "does not support" in low:
        return ({"reasoning_effort": "none"}, None)
    # そもそもこの引数を知らないモデル → 落とす
    if "reasoning_effort" in low and "unrecognized" in low:
        return (None, "reasoning_effort")
    # temperature / top_p が既定値しか許されないモデル → 落として既定に任せる
    for name in ("temperature", "top_p"):
        if f"'{name}'" in low and ("does not support" in low
                                   or "only the default" in low):
            return (None, name)
    # 新しいモデルは max_tokens ではなく max_completion_tokens を使う
    if "max_tokens" in low and "max_completion_tokens" in low:
        v = kwargs.get("max_tokens")
        if v is not None:
            return ({"max_completion_tokens": v}, "max_tokens")
    return None


def _create(**kwargs):
    """chat.completions.create の呼び出し口。

    モデルが受け付けない引数を、エラーの内容を見て直しながら投げ直す。
    直し方が分からないエラーはそのまま投げる（画面にそのまま出す）。
    """
    model = str(kwargs.get("model") or "")
    attempt = _apply_quirks(kwargs)
    for _ in range(_MAX_FIX):
        try:
            return client().chat.completions.create(**attempt)
        except Exception as e:
            fix = _fix_for(str(e), attempt)
            if fix is None:
                raise
            set_, drop = fix
            if set_ and all(attempt.get(k) == v for k, v in set_.items()):
                raise                      # 同じ直しを繰り返している
            if drop and drop not in attempt:
                raise
            _learn(model, set_=set_, drop=drop)
            attempt = _apply_quirks(kwargs)
    return client().chat.completions.create(**attempt)


# --- DBルーター（質問に関係するDBだけを選ぶ） --------------------------------------

_ROUTE_SYSTEM = """あなたはデータ分析アプリの振り分け係です。
利用者の質問に答えるために必要なDBを、下のDB一覧から選んでください。

出力はJSONの配列だけ（説明文は書かない）:
  ["demo_sales.db", "demo_master.db"]

守ること:
- 質問の言葉とDBの説明・テーブル・用語を突き合わせて選ぶ。
- 結合キーの一覧に他のDBの名前が出ていたら、そのDBも一緒に選ぶ
  （例: 売上の質問で顧客名が要るなら、顧客マスタのあるDBも含める）。
- 迷ったら含める。外しすぎて答えられないより、多めの方がよい。
- 全DBが要る・判断できないときは ["*"] と書く。
- 直前の会話の続き（「それを」「さっきの」）なら、その話題のDBを選ぶ。"""


def route_dbs(question: str, history: list[str] | None = None) -> list[str] | None:
    """質問に関係するDBファイル名を選ぶ。判断できなければ None（=全DB）。

    本番の回答に無関係なDBのカタログを入れると、列やコード値の情報が薄まって
    精度が下がる。そこで前段に小さな1回を挟み、要約だけを見せて選ばせる。
    要約は全DBぶんでも詳細版の1/4ほどで、内容が変わらない限りキャッシュに乗る。

    失敗したら None を返して全DBで進める。ルーターの不調で答えられなくなるのが
    いちばん悪いので、この関数は例外を外に出さない。
    """
    import db                       # 循環importを避けるため、使うときに読む

    files = db.list_db_files()
    if len(files) <= 1:
        return None
    known = {f.name for f in files}
    summaries = "\n".join(
        catalog.db_text_cached(db.alias_for(f), f, None, full=False) for f in files)

    ask = []
    for h in (history or [])[-3:]:
        ask.append(f"（直前の質問: {h}）")
    ask.append(f"今回の質問: {question}")
    try:
        resp = _create(
            model=config.OPENAI_MODEL,
            messages=[{"role": "system", "content": _ROUTE_SYSTEM},
                      {"role": "user", "content": f"{summaries}\n\n{chr(10).join(ask)}"}],
            temperature=0, max_tokens=120,
        )
        m = re.search(r"\[.*?\]", resp.choices[0].message.content or "", re.DOTALL)
        picked = json.loads(m.group(0)) if m else []
    except Exception as e:
        print(f"[router] 振り分けに失敗したため全DBで進めます: {e}")
        return None
    if "*" in picked:
        return None
    names = [str(p) for p in picked if str(p) in known]
    return names or None


# --- system prompt -----------------------------------------------------------

def build_system_prompt(scope: list[dict], admin: bool = False,
                        model: str | None = None) -> str:
    """選択スコープのデータカタログを埋め込んだ system prompt を組み立てる。

    admin は「管理者だけに渡すツール」を一覧に載せるかどうか。
    渡していないツールを説明に書くと、AIが呼ぼうとして失敗するだけになる。
    model を渡すと、カタログをインラインするかの判定を「そのモデルが読める量」で行う
    （渡さなければ管理者設定/envの上限）。
    """
    inline_cap = None
    if model:
        import models as models_mod        # 循環importを避ける
        inline_cap = models_mod.inline_limit_for(model)
    aliases = [s["alias"] for s in scope]

    # 設定どおりに更新できていないテーブルがあれば、AIに教えておく。
    # そのテーブルを使った回答に「データが古い可能性」を添えさせるため。
    import jobs as jobs_mod
    stale_lines = []
    in_scope = {s["name"] for s in scope}
    for (db_file, table), ps in jobs_mod.problems_by_table().items():
        if db_file in in_scope:
            alias = next((s["alias"] for s in scope if s["name"] == db_file), db_file)
            since = (ps[0].get("since") or "")[:16].replace("T", " ")
            stale_lines.append(f"- {alias}.{table}: {ps[0]['message']}"
                               + (f"（{since} 以降）" if since else ""))
    stale_note = ""
    if stale_lines:
        # 「古い」だけではない。取り込めたが数値列が文字に落ちた場合（degraded）は
        # データは新しく、集計の方が信用できない。どちらかを断定せず両方を伝える。
        stale_note = ("\n# 状態に問題があるデータ（重要）\n"
                      "次のテーブルは定期取り込みが設定どおりに動いていない。"
                      "中身が古いか、値の型が想定と違う（数値の列が文字で入っている）可能性がある。\n"
                      + "\n".join(stale_lines) + "\n"
                      "これらのテーブルを使って答えるときは、回答の冒頭に"
                      "「※ このデータは○月○日以降、正しく更新できていない可能性があります（理由）」"
                      "と必ず一言添える。数値の列が文字になっている場合は、合計・平均がずれうることも書く。"
                      "使わない質問では触れなくてよい。\n")
    if len(aliases) > 1:
        naming = (f"複数のDBが対象です（{', '.join(aliases)}）。"
                  "テーブル名は必ず『エイリアス.テーブル名』で修飾すること"
                  f"（例: {aliases[0]}.xxx）。DBをまたぐ JOIN も可能。")
    elif aliases:
        naming = (f"対象のDBは {aliases[0]} の1つ。テーブル名はそのまま書いてよい"
                  f"（{aliases[0]}.テーブル名 と修飾しても可）。")
    else:
        naming = "対象にできるDBがありません。「データ取り込み」でDBを作るよう案内すること。"

    # 実際に渡すツール一覧（無効化・説明の上書き・ユーザー定義ツールを反映）
    lines = []
    for t in tools.build_tools(scope, admin=admin):
        fn = t["function"]
        args = ", ".join((fn.get("parameters") or {}).get("properties", {}).keys())
        lines.append(f"- {fn['name']}({args}) : {fn['description']}")
    custom = custom_tools.collect_everywhere(scope)
    if custom:
        lines.append("※ 上記のうち次はこの環境専用に用意されたツールです。"
                     "目的が合致するときは自分でSQLを書かずにこちらを優先して使ってください: "
                     + ", ".join(t["name"] for t in custom))
    tool_list = "\n".join(lines)

    return f"""あなたはSQLiteデータベースの分析アシスタントです。
読み取り専用(SELECTのみ)でデータベースにアクセスできます。

# 振る舞い
- ユーザーの質問に答えるため、必要に応じてツールを呼び出し、必ず実データに基づいて回答する。
- 推測で数値を答えてはいけない。データが必要なら run_sql_query を使う。
- 列の意味や値の実体が不確かなら、SQLを書く前に describe_table で確認する。
- ツールを使うか・どのSQLを書くかはあなたが判断する(挨拶や一般的な雑談ならツール不要)。
- 回答は日本語。まず結論、次に根拠(表やグラフの要点)を簡潔に述べる。
- 後述の「業務用語」に載っている言葉が質問に出たら、必ずその定義に従う。自分の常識で解釈し直さない。
- カタログを育てる提案は、ユーザーが頼まなくても、あなたから聞く。
  ただし押しつけない: 聞くのは1回の回答につき最大1つ、同じものを何度も聞かない。
  - 用語: ユーザーが業務用語の意味を教えてくれたら（「有効な受注とはキャンセル以外のこと」等）、
    回答の最後に「この定義を用語集に登録しますか？」と一言添える。
  - 例文: run_sql_query の結果に example_registered: false が入っていて、
    回答が問題なく出せたら、回答の最後に
    「この質問と答え方を例文として登録しますか？（似た質問に強くなります）」と一言添える。
    example_registered: true なら聞かない（すでに登録済み）。
    エラーで言い直した回答や、雑談・確認だけのやり取りでは聞かない。
  - 前向きな返事（「はい」「お願い」「登録して」等）が来たら、用語は propose_glossary_term、
    例文は propose_example でカードを出す。登録するかはユーザーがカードのボタンで決める。
    例文の sql は実行して正しかったものをそのまま使う。
  - 「SQL式:」が書かれている用語は、その式をそのまま WHERE や SELECT に埋め込む。
  - SQL式が無く説明文だけの用語は、その説明と列情報から自分でSQLを組み立てる。
    どの列をどう使ったかを回答の中で一言添える（人が誤りに気づけるようにするため）。
- ツールの結果に verification_warnings（検算の不一致）が入っていることがある。
  これは「同じ数字を別の経路で数えたら食い違った」という自動検算の結果で、
  あなたのSQLの誤りとは限らない。入っていたら、回答の末尾で
  「どの数字を使ったか」と「別の経路では値が異なること」を必ず1〜2文で注記する。
  差異の原因は、検算結果に示された内訳の範囲でだけ述べ、推測で断定しない。

# 可視化の方針（チャットにグラフを描く）
- ユーザーが「グラフ」「可視化」「チャート」「推移」「トレンド」「割合」「内訳」「分布」等を求めたら、必ず plot_* のツールを呼んでグラフを描く。
- 明示が無くても、結果が次に該当するなら積極的にグラフにする。目的でツールを選ぶ：
  - 項目どうしの比較・順位 → plot_comparison（bar / hbar / stacked_bar / pareto / radar など）
  - 時系列・推移 → plot_trend（line / step / area / calendar など）
  - 構成比・内訳・増減の要因 → plot_composition（pie / donut / treemap / funnel / waterfall / sankey）
  - ばらつき・分布 → plot_distribution（histogram / box / violin など）
  - 2つ以上の項目の関係 → plot_relationship（scatter / bubble / heatmap など）
  - 1つの数字を大きく・目標との対比 → plot_kpi（indicator / gauge / bullet）
- sql は集計済み(GROUP BY)にし、x/y にする列を AS で明示する。色分けは color に列名を渡す。
- 棒グラフの積み方は種別で指定する：「積み上げ」なら chart_type="stacked_bar"、「横並び」「比較」なら "bar"。
- 「2軸」「二軸」「棒と折れ線」「件数と比率を一緒に」など、単位の異なる2指標を重ねたい時は plot_dual_axis を使う。
  bar_y(左軸=棒, 件数など) と line_y(右軸=折れ線, 比率など) に列名を渡す。
- 必要なら run_sql_query で数値を確認しつつ、可視化は plot_* で別途描く（両方呼んでよい）。

# SQLで書けないこと（必ず専用ツールを使う）
このSQLiteには STDDEV / VARIANCE / MEDIAN / CORR / PERCENTILE / SQRT / POWER が無く、
PIVOT構文も無い。次はSQLで計算しようとせず、必ずツールを使うこと。
- 「クロス集計」「行に○○・列に△△」「マトリクスで」 → pivot_table
  （sql では集計せず、必要な列を返すだけにする。集計はツールが行う）
- 「相関」「中央値」「ばらつき」「標準偏差」「四分位」「分布の要約」 → analyze_stats
- 「外れ値」「異常値」「突出しているもの」 → analyze_stats の method="outliers"
  （sql は集計せず明細を返す。1行1件の状態にしてから渡すこと）

# ファイル出力
- 「エクセル」「Excel」「xlsx」→ export_excel。観点が複数なら sheets に複数の SELECT を渡し、
  1ブックに複数シートでまとめる。
- 「CSV」「csvにして」「取り込み用」→ export_csv。複数指定するとZIPにまとめて渡される。
  文字コードは既定の utf-8-sig でよい（Excelで文字化けしない）。Shift_JIS を求められたときだけ cp932。
- 「テキストで」「レポートにして」「議事録」「まとめを文書で」→ export_text。
  body に自分で文章を書き、集計表を入れたい箇所に {{見出し}} と書いて sections に SELECT を指定する。
- ファイル出力ツールを呼んだ後は、画面に保存済み。中身の全件を文章で繰り返さず、
  何を入れたかだけ簡潔に伝える。

# 利用可能なツール
{tool_list}

{stale_note}# SQLルール
- SQLite方言。SELECT(または WITH ... SELECT)のみ。INSERT/UPDATE/DELETE/DDL/PRAGMA等は禁止(実行されません)。
- {naming}
- 1回の呼び出しで1ステートメント。末尾セミコロン不要。
- 集計は GROUP BY を使い、列に AS で日本語の別名を付けると表示が分かりやすい。
- 日付は文字列で保存されていることが多い。date() / strftime() を活用する（例: strftime('%Y-%m', 列) で月別）。
- 行数が多くなりそうなら LIMIT や集計で絞る。
- 「テーブルを見せて」「中身を全部見たい」のようにデータそのものを見たいと言われたら、SELECT * を打って先頭数行を貼るのではなく open_table を使う（全行を辿れる画面へのリンクが出る）。出したあとは何のテーブルかを1〜2文添えるだけでよい。

# 選択中のデータカタログ
{catalog.prompt_for_scope(scope, limit=inline_cap)}

現在時刻: {datetime.now().isoformat(timespec="seconds")}
"""


# --- 文脈の使用量の見積もり ------------------------------------------------------
#
# 「カタログを増やしてよいか」を判断するには、モデルが一度に読める量に対して
# いまどれだけ使っているかが要る。正確なトークン数はAPIに投げないと分からないが、
# それでは画面を開くたびに課金が発生する。そこで実測から係数を出して概算する。
#
# 実測（gpt-4o-mini・11DB選択）:
#   要約版 system 28,851字 + ツール定義 44,377字 → 入力 29,265 トークン
#   全文版 system 71,879字 + ツール定義 44,377字 → 入力 52,395 トークン
#   差分から、日本語主体の本文は 43,028字 → 23,130トークン = 0.537 トークン/字
# 下の係数で上の2例を計算すると、実測に対して +1.2% / +1.7% に収まる（やや多めに出る）。
# 見積もりは「余裕がある」と言いすぎない方が安全なので、多めに出るぶんには構わない。

#: 日本語主体の文章（カタログ・指示）の 1文字あたりトークン数
TOKENS_PER_CHAR_TEXT = 0.55
#: ツール定義のJSON（英字と記号が多い）の 1文字あたりトークン数
TOKENS_PER_CHAR_JSON = 0.31


def tokens_for(chars: int, kind: str = "text") -> int:
    """文字数からトークン数の概算を出す。"""
    ratio = TOKENS_PER_CHAR_JSON if kind == "json" else TOKENS_PER_CHAR_TEXT
    return int(max(0, chars) * ratio)


def budget(scope: list[dict], model: str | None = None, admin: bool = False) -> dict:
    """このスコープ・このモデルで、文脈をどれだけ使うかの概算。

    「いま」と「上限までカタログが育ったとき」の両方を返す。
    上限を決める画面で、変えた結果どうなるかを見せるため。
    """
    import models as models_mod

    model = model or config.OPENAI_MODEL
    context, known = models_mod.context_window(model)
    limit = models_mod.inline_limit_for(model)      # そのモデルの自動上限

    # 推定せず、実際に組み立てたものを測る（カタログはキャッシュ済みなので速い）
    system = build_system_prompt(scope, admin=admin, model=model)
    used_catalog = len(catalog.prompt_for_scope(scope))
    catalog_chars = catalog.inline_length(scope)       # 全文にした場合の長さ
    tool_chars = len(json.dumps(tools.build_tools(scope, admin=admin), ensure_ascii=False))

    # カタログ以外（SQLルール・ツールの使い分け・ツール名の一覧など）
    fixed_chars = max(0, len(system) - used_catalog)

    tool_tokens = tokens_for(tool_chars, "json")
    now = tokens_for(len(system)) + tool_tokens
    at_limit = tokens_for(fixed_chars + limit) + tool_tokens

    def pct(n: int) -> float:
        return round(n / context * 100, 1) if context else 0.0

    return {
        "model": model, "context": context, "context_known": known,
        "limit_chars": limit,
        "catalog_chars": catalog_chars,
        "catalog_inlined": catalog_chars <= limit,
        "tool_tokens": tool_tokens,
        # カタログ以外の固定ぶん。画面で上限を動かしたときに、
        # サーバと同じ式で計算し直せるように渡す。
        "base_tokens": tool_tokens + tokens_for(fixed_chars),
        "tokens_per_char": TOKENS_PER_CHAR_TEXT,
        "now_tokens": now, "now_pct": pct(now), "headroom_pct": round(100 - pct(now), 1),
        "at_limit_tokens": at_limit, "at_limit_pct": pct(at_limit),
        # 上限をここまで上げても文脈の半分に収まる、という目安
        "suggest_max_chars": max(0, int((context * 0.5 - tool_tokens
                                         - tokens_for(fixed_chars))
                                        / TOKENS_PER_CHAR_TEXT)),
    }


# --- 画像つきのメッセージ --------------------------------------------------------

# 受け付ける画像。ここに無い形式は送らない（APIが解釈できないため）
IMAGE_MIMES = {"image/png", "image/jpeg", "image/gif", "image/webp"}


def user_message(text: str, images: list[dict] | None = None) -> dict:
    """ユーザー発言を作る。画像があればマルチモーダル形式にする。

    images: [{"mime": "image/png", "b64": "..."}] （b64はデータ本体のみ）
    画像が無いときは、これまで通り content が文字列のメッセージを返す
    （画像非対応のモデルに配列を渡すと弾かれることがあるため）。
    """
    if not images:
        return {"role": "user", "content": text}
    parts: list[dict] = []
    if text:
        parts.append({"type": "text", "text": text})
    for img in images:
        mime = img.get("mime") or "image/png"
        if mime not in IMAGE_MIMES:
            continue
        parts.append({"type": "image_url",
                      "image_url": {"url": f"data:{mime};base64,{img['b64']}",
                                    "detail": img.get("detail") or "auto"}})
    return {"role": "user", "content": parts or text}


# --- チャット補完 --------------------------------------------------------------

def chat(messages: list[dict], tool_defs: list[dict] | None = None,
         model: str | None = None):
    """messages を渡して1回の補完を取得。tools 付き。message オブジェクトを返す。

    tool_defs を省略した場合は組み込みツールのみ。通常はチャット側で
    tools.build_tools(entries) を渡し、ユーザー定義ツールも含める。
    model は画面で選ばれたモデル。省略時は env の既定。
    """
    kwargs = dict(
        model=model or config.OPENAI_MODEL,
        messages=messages,
        tools=tool_defs if tool_defs is not None else tools.BUILTIN_TOOLS,
        tool_choice="auto",
        temperature=config.OPENAI_TEMPERATURE,
    )
    if config.OPENAI_TOP_P is not None:
        kwargs["top_p"] = config.OPENAI_TOP_P
    if config.OPENAI_MAX_TOKENS is not None:
        kwargs["max_tokens"] = config.OPENAI_MAX_TOKENS
    resp = _create(**kwargs)
    return resp.choices[0].message


class StreamedMessage:
    """ストリーミングで組み立てた1回ぶんの応答。

    chat() が返す message オブジェクトと同じ形（content / tool_calls）に
    見えるようにしておく。呼び出し側はどちらでも同じ扱いができる。
    """

    class _Fn:
        def __init__(self, name="", arguments=""):
            self.name, self.arguments = name, arguments

    class _Call:
        def __init__(self, id="", name="", arguments=""):
            self.id, self.type = id, "function"
            self.function = StreamedMessage._Fn(name, arguments)

    def __init__(self):
        self.content = ""
        self.tool_calls = None
        self._parts: dict[int, dict] = {}

    def _finish(self):
        if not self._parts:
            self.tool_calls = None
            return
        self.tool_calls = [
            StreamedMessage._Call(p["id"], p["name"], p["arguments"])
            for _, p in sorted(self._parts.items())
        ]


def chat_stream(messages: list[dict], tool_defs: list[dict] | None = None,
                model: str | None = None):
    """1回の補完をストリーミングで受け取る。

    文字が届くたびに ("text", 差分) を yield し、
    最後に ("done", StreamedMessage) を1回だけ yield する。
    ツール呼び出しは途中経過を出さない（引数のJSONは途中では読めないため）。
    """
    kwargs = dict(
        model=model or config.OPENAI_MODEL,
        messages=messages,
        tools=tool_defs if tool_defs is not None else tools.BUILTIN_TOOLS,
        tool_choice="auto",
        temperature=config.OPENAI_TEMPERATURE,
        stream=True,
    )
    if config.OPENAI_TOP_P is not None:
        kwargs["top_p"] = config.OPENAI_TOP_P
    if config.OPENAI_MAX_TOKENS is not None:
        kwargs["max_tokens"] = config.OPENAI_MAX_TOKENS

    out = StreamedMessage()
    for chunk in _create(**kwargs):
        if not chunk.choices:
            continue
        delta = chunk.choices[0].delta
        if getattr(delta, "content", None):
            out.content += delta.content
            yield ("text", delta.content)
        for tc in (getattr(delta, "tool_calls", None) or []):
            # 同じツール呼び出しが複数のチャンクに分かれて届くので、index で束ねる
            part = out._parts.setdefault(tc.index, {"id": "", "name": "", "arguments": ""})
            if tc.id:
                part["id"] = tc.id
            fn = getattr(tc, "function", None)
            if fn is not None:
                if getattr(fn, "name", None):
                    part["name"] += fn.name
                if getattr(fn, "arguments", None):
                    part["arguments"] += fn.arguments
    out._finish()
    yield ("done", out)


# --- AI下書き（データカタログ用） ----------------------------------------------

_DRAFT_SYSTEM = """あなたはデータカタログ作成の専門家です。
与えられたテーブルのプロファイル（列名・型・実値の分布・サンプル行）から、
テーブルと各列の業務的な説明文を日本語で推測し、JSONだけを出力してください。

出力形式（JSON以外の文字を含めないこと）:
{
  "description": "テーブルの説明。1行 = 何のレコードかを必ず含める。",
  "columns": {
    "列名": {
      "description": "列の説明",
      "values": {"コード値": "意味"}   // 値がコード(区分値)と思われる列のみ。それ以外は省略
    }
  }
}

注意:
- 確信が持てない場合は「〜と思われる」と書く。
- values は実値一覧にある値だけを対象にする。
- すべての列に説明を付ける。"""


_GLOSSARY_SYSTEM = """あなたはSQLiteに詳しいデータカタログ作成の専門家です。
与えられたテーブル定義（列・型・実値の分布・サンプル行）をもとに、
業務用語の「自然言語の説明」をSQLの式に翻訳してください。

出力形式（JSON以外の文字を含めないこと）:
{"用語": "SQL式", "用語2": "SQL式"}

守ること:
- WHERE にそのまま入る条件式（例: status != '9' AND amount >= 1000000）か、
  SELECT にそのまま入る計算式（例: SUM(amount) * 1.0 / COUNT(*)）だけを書く。
- SELECT や FROM で始まる文全体は書かない。末尾にセミコロンを付けない。
- 列名は与えられたテーブルに実在するものだけを使う。値は実値一覧にあるものを使う。
- SQLiteに無い関数(STDDEV, MEDIAN, PERCENTILE_CONT, SQRT, POWER など)は使わない。
- 説明があいまいで確信が持てない用語は、キーごと省略する（推測で書かない）。"""


def draft_glossary_sql(db_path, table_name: str | None, terms: list[dict]) -> dict:
    """業務用語の説明文からSQL式の下書きを作る。

    terms: [{"term": 用語, "description": 自然言語の説明}, ...]
    戻り値: {用語: SQL式}（翻訳できなかった用語は含まれない）
    """
    if not terms:
        return {}
    profile = catalog.profile_db(Path(db_path))
    meta = catalog.load_meta(Path(db_path))
    if table_name:
        context = catalog.table_text("db", table_name, profile, meta, full=True)
    else:   # テーブルをまたぐ用語。DB全体を見せる
        context = catalog.db_text("db", Path(db_path), None, full=True)
    asked = "\n".join(f"- {t['term']}: {t.get('description') or ''}" for t in terms)

    resp = _create(
        model=config.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": _GLOSSARY_SYSTEM},
            {"role": "user", "content": f"{context}\n\n翻訳したい業務用語:\n{asked}"},
        ],
        temperature=0,
    )
    content = resp.choices[0].message.content or ""
    m = re.search(r"\{.*\}", content, re.DOTALL)
    if not m:
        raise ValueError(f"AIの応答をJSONとして解析できませんでした: {content[:200]}")
    data = json.loads(m.group(0))
    if not isinstance(data, dict):
        raise ValueError("AIの応答が想定した形式ではありません。")
    wanted = {t["term"] for t in terms}
    return {k: str(v).strip() for k, v in data.items() if k in wanted and str(v).strip()}


def draft_table_meta(db_path, table_name: str) -> dict:
    """テーブルのプロファイルからメタ情報の下書きを生成する。"""
    profile = catalog.profile_db(Path(db_path))
    meta = catalog.load_meta(Path(db_path))
    text = catalog.table_text("db", table_name, profile, meta, full=True)
    resp = _create(
        model=config.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": _DRAFT_SYSTEM},
            {"role": "user", "content": f"ファイル名: {Path(db_path).name}\n\n{text}"},
        ],
        temperature=0,
    )
    content = resp.choices[0].message.content or ""
    # ```json ... ``` フェンスを剥がしてから解析
    m = re.search(r"\{.*\}", content, re.DOTALL)
    if not m:
        raise ValueError(f"AI下書きの応答をJSONとして解析できませんでした: {content[:200]}")
    data = json.loads(m.group(0))
    if not isinstance(data, dict):
        raise ValueError("AI下書きの応答が想定した形式ではありません。")
    return data


# --- ユーザー定義ツールの下書き ---------------------------------------------------
#
# SQLを書けない人でもツールを作れるようにするための入口。
# 「何をするツールか」を日本語で書いてもらい、SQLとパラメータはAIに起こさせる。
# 起こしたSQLは呼び出し側で必ず実データに当てて確かめる（推測のまま保存させない）。

_TOOL_SYSTEM = """あなたはSQLiteに詳しいデータ分析アプリの設定担当です。
利用者が日本語で書いた「やりたいこと」を、AIが呼び出せるツールの定義に変換してください。

出力形式（JSON以外の文字を含めないこと）:
{
  "name": "英小文字と_のみの短い名前（例: monthly_sales）",
  "description": "このツールが何を返すかの説明。AIがこれを読んで使うかどうかを決める",
  "sql": "SELECT ...（1文だけ。末尾のセミコロンは不要）",
  "parameters": [
    {"name": "year", "type": "string", "description": "対象年 YYYY",
     "required": true, "example": "2026"}
  ],
  "chart": {"chart_type": "line", "x": "月", "y": "売上", "title": "月別売上"}
}

守ること:
- SQLは SELECT（または WITH ... SELECT）だけ。書き込み・DDLは書かない。
- 列名・テーブル名は、与えられたカタログに実在するものだけを使う。推測で作らない。
- 「毎回変えたい値」は、利用者の日本語から自分で見極めて parameters にする。
  聞かれ方が変わるたびに差し替える値（「指定した年の」「ある部署の」「任意の期間で」など）は
  パラメータにし、SQLでは :名前 の形で参照する。
  一方「部署ごと」「月別」のような集計の切り口は、パラメータではなく GROUP BY で表す。
  迷ったらパラメータにしない。引数が増えるほどAIは呼びにくくなる。
- parameters の type は string / integer / number / boolean のいずれか。
- parameters には description（日本語）と example を必ず書く。
  example は「カタログの実値・期間に実在し、実際に行が返る値」にする。
  この値で試し実行して見せるので、0行になる値を書かないこと。
- 複数のDBにまたがるときは「DB名.テーブル名」で修飾する。
- 列には日本語の別名を AS で付ける（画面にそのまま出るため）。
- SQLiteに無い関数(STDDEV, MEDIAN, PERCENTILE_CONT, SQRT, POWER)やPIVOT構文は使わない。
- 見せ方が「グラフ」のときだけ chart を書く。x と y には SELECT の別名をそのまま使う。
  グラフでないときは chart を省略する。"""


def draft_tool(db_path, purpose: str, params_wanted: list[str] | None = None,
               render: str = "table", previous: dict | None = None,
               error: str | None = None) -> dict:
    """日本語の「やりたいこと」から、ユーザー定義ツールの下書きを起こす。

    db_path        … None なら全DBのカタログを見せる（作る人にDBを選ばせない）。
                     特定のDBに限りたいときだけパスを渡す。
    purpose        … 何をするツールか（日本語）。毎回変えたい値もこの文から読み取らせる
                     ので、呼び出し側が指定を組み立てる必要はない。
    params_wanted  … 毎回変えたい項目を明示したいときだけ渡す（例: ["対象年", "部署"]）。
                     省略すれば purpose の書き方からAIが判断する。
    render         … 結果の見せ方（table / chart / chart_dual / excel / csv / none）
    previous/error … 前回の下書きが実データで失敗したときの、SQLとエラー文。
                     渡すと「どこが間違っていたか」を踏まえて書き直す。
    """
    import db                       # 循環importを避けるため、使うときに読む

    # db_path が None なら全DBを見せる。どのDBに書くかは、やりたいことを読んだAIが
    # 決める（作る人にDBを選ばせない）。量が上限を超えるときは要約に落ちる。
    if db_path is None:
        paths = db.list_db_files()
    else:
        paths = [Path(db_path)]
    context = catalog.prompt_for_scope(
        [{"path": str(p), "alias": db.alias_for(p), "tables": None} for p in paths])
    ask = [f"やりたいこと: {purpose}"]
    if params_wanted:
        ask.append("毎回変えたい項目: " + "、".join(params_wanted))
    ask.append(f"結果の見せ方: {render}")
    if previous and error:
        ask.append("\n前回の下書きは実際のデータで失敗しました。原因を直して書き直してください。")
        ask.append(f"前回のSQL:\n{previous.get('sql', '')}")
        ask.append(f"エラー: {error}")

    resp = _create(
        model=config.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": _TOOL_SYSTEM},
            {"role": "user", "content": f"{context}\n\n{chr(10).join(ask)}"},
        ],
        temperature=0,
    )
    content = resp.choices[0].message.content or ""
    m = re.search(r"\{.*\}", content, re.DOTALL)
    if not m:
        raise ValueError(f"AIの応答をJSONとして解析できませんでした: {content[:200]}")
    data = json.loads(m.group(0))
    if not isinstance(data, dict):
        raise ValueError("AIの応答が想定した形式ではありません。")

    out = {
        "name": str(data.get("name") or "").strip(),
        "description": str(data.get("description") or purpose).strip(),
        "sql": str(data.get("sql") or "").strip().rstrip(";"),
        "parameters": [],
        "render": render,
        "enabled": True,
    }
    for p in (data.get("parameters") or []):
        if not isinstance(p, dict) or not str(p.get("name") or "").strip():
            continue
        t = str(p.get("type") or "string")
        item = {
            "name": str(p["name"]).strip(),
            "type": t if t in custom_tools.PARAM_TYPES else "string",
            "description": str(p.get("description") or "").strip(),
            "required": p.get("required", True) is not False,
        }
        # 試し実行に使う値。空のまま流すと0行になり、動くかどうか確かめられない。
        if p.get("example") not in (None, ""):
            item["example"] = p["example"]
        out["parameters"].append(item)
    if render in ("chart", "chart_dual") and isinstance(data.get("chart"), dict):
        out["chart"] = data["chart"]
    return out
